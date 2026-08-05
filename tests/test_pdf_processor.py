from __future__ import annotations

import tempfile
import unittest
import warnings
from pathlib import Path
from unittest import mock

from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from pypdf.constants import UserAccessPermissions
from pypdf.errors import LimitReachedError
from pypdf.generic import ArrayObject, RectangleObject, TextStringObject
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, letter
from reportlab.pdfgen import canvas
from reportlab.platypus import Table, TableStyle

from docuforge.processors import pdf as pdf_processor


def _make_text_pdf(path: Path, labels: tuple[str, ...]) -> Path:
    document = canvas.Canvas(str(path), pagesize=letter)
    for label in labels:
        document.setFont("Helvetica", 16)
        document.drawString(72, 700, label)
        document.showPage()
    document.save()
    return path


def _make_blank_pdf(path: Path) -> Path:
    document = canvas.Canvas(str(path), pagesize=letter)
    document.showPage()
    document.save()
    return path


def _make_watermark_geometry_pdf(path: Path) -> Path:
    base = path.with_name(f"{path.stem}_base.pdf")
    geometries = (
        ((0, 0, 612, 792), (50, 100, 450, 600), 0),
        ((0, 0, 612, 792), (100, 50, 500, 550), 0),
        ((0, 0, 612, 792), (50, 100, 450, 600), 180),
        ((0, 0, 612, 792), (50, 100, 450, 600), 90),
        ((0, 0, 612, 792), (50, 100, 450, 600), 270),
        ((0, 0, 792, 612), (0, 0, 792, 612), 0),
    )
    document = canvas.Canvas(str(base), pagesize=(792, 792))
    for _geometry in geometries:
        document.showPage()
    document.save()

    reader = PdfReader(base)
    writer = PdfWriter()
    try:
        for page, (media_box, crop_box, rotation) in zip(
            reader.pages, geometries, strict=True
        ):
            page.mediabox = RectangleObject(media_box)
            page.cropbox = RectangleObject(crop_box)
            if rotation:
                page.rotate(rotation)
            writer.add_page(page)
        with path.open("wb") as stream:
            writer.write(stream)
    finally:
        writer.close()
        reader.close()
    return path


def _ink_component_count(pixmap: object, *, red_only: bool = False) -> int:
    width = int(getattr(pixmap, "width"))
    height = int(getattr(pixmap, "height"))
    channels = int(getattr(pixmap, "n"))
    samples = bytes(getattr(pixmap, "samples"))
    mask = bytearray(width * height)
    for index in range(width * height):
        offset = index * channels
        red, green, blue = samples[offset : offset + 3]
        if (red > 200 and green < 50 and blue < 50) if red_only else red < 245:
            mask[index] = 1

    seen = bytearray(width * height)
    components = 0
    for start in range(width * height):
        if not mask[start] or seen[start]:
            continue
        components += 1
        seen[start] = 1
        pending = [start]
        for position in pending:
            x = position % width
            y = position // width
            neighbours = []
            if x:
                neighbours.append(position - 1)
            if x + 1 < width:
                neighbours.append(position + 1)
            if y:
                neighbours.append(position - width)
            if y + 1 < height:
                neighbours.append(position + width)
            for neighbour in neighbours:
                if mask[neighbour] and not seen[neighbour]:
                    seen[neighbour] = 1
                    pending.append(neighbour)
    return components


def _make_table_pdf(path: Path) -> Path:
    document = canvas.Canvas(str(path), pagesize=A4)
    data = [["Name", "Value"], ["Alice", "10"], ["Bob", "20"]]
    table = Table(data, colWidths=(120, 90), rowHeights=(28, 28, 28))
    table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 11),
            ]
        )
    )
    table.wrapOn(document, 300, 300)
    table.drawOn(document, 72, 650)
    document.showPage()
    document.save()
    return path


def _page_texts(path: Path, password: str | None = None) -> list[str]:
    reader = PdfReader(path)
    try:
        if reader.is_encrypted:
            result = reader.decrypt(password or "")
            if int(result) == 0:
                raise AssertionError("test helper could not decrypt PDF")
        return [(page.extract_text() or "").strip() for page in reader.pages]
    finally:
        reader.close()


def _fake_convert_from_path(_source: str, **kwargs: object) -> list[str]:
    output_folder = Path(str(kwargs["output_folder"]))
    first_page = int(kwargs["first_page"])
    last_page = int(kwargs["last_page"])
    image_format = str(kwargs["fmt"])
    suffix = ".jpg" if image_format == "jpeg" else ".png"
    outputs: list[str] = []
    for page_number in range(first_page, last_page + 1):
        path = output_folder / f"rendered-{page_number}{suffix}"
        image = Image.new("RGB", (200 + page_number, 300), (page_number * 30, 80, 120))
        image.save(path, format="JPEG" if suffix == ".jpg" else "PNG")
        outputs.append(str(path))
    return outputs


class PDFProcessorTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name) / "中文路径"
        self.root.mkdir()

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def test_structural_page_operations(self) -> None:
        first = _make_text_pdf(self.root / "甲.pdf", ("A1", "A2"))
        second = _make_text_pdf(self.root / "乙.pdf", ("B1",))

        merged = self.root / "合并.pdf"
        self.assertEqual(
            pdf_processor.merge_pdfs([first, second], merged),
            [merged.resolve()],
        )
        self.assertEqual(_page_texts(merged), ["A1", "A2", "B1"])

        split_outputs = pdf_processor.split_pdf(
            merged, [(1, 2), (3, 3)], self.root / "拆分", prefix="部分"
        )
        self.assertEqual(len(split_outputs), 2)
        self.assertEqual([len(_page_texts(path)) for path in split_outputs], [2, 1])

        extracted = self.root / "提取.pdf"
        pdf_processor.extract_pages(merged, "3,1", extracted)
        self.assertEqual(_page_texts(extracted), ["B1", "A1"])

        deleted = self.root / "删除.pdf"
        pdf_processor.delete_pages(merged, [2], deleted)
        self.assertEqual(_page_texts(deleted), ["A1", "B1"])

        inserted = self.root / "插入.pdf"
        pdf_processor.insert_pages(first, second, 2, inserted)
        self.assertEqual(_page_texts(inserted), ["A1", "B1", "A2"])

        rotated = self.root / "旋转.pdf"
        pdf_processor.rotate_pages(merged, "2", 90, rotated)
        reader = PdfReader(rotated)
        try:
            self.assertEqual(reader.pages[0].rotation, 0)
            self.assertEqual(reader.pages[1].rotation, 90)
        finally:
            reader.close()

        with self.assertRaises(FileExistsError):
            pdf_processor.extract_pages(merged, 1, extracted)

    def test_atomic_overwrite_protection_keeps_existing_file(self) -> None:
        source = _make_text_pdf(self.root / "输入.pdf", ("one",))
        target = self.root / "已存在.pdf"
        target.write_bytes(b"do-not-replace")
        with self.assertRaises(FileExistsError):
            pdf_processor.merge_pdfs([source], target)
        self.assertEqual(target.read_bytes(), b"do-not-replace")

    def test_compress_encrypt_and_decrypt(self) -> None:
        source = _make_text_pdf(self.root / "原文.pdf", ("Secret", "Second"))
        compressed = self.root / "压缩.pdf"
        self.assertEqual(
            pdf_processor.compress_pdf(source, compressed, level=9),
            [compressed.resolve()],
        )
        self.assertEqual(_page_texts(compressed), ["Secret", "Second"])

        encrypted = self.root / "加密.pdf"
        pdf_processor.encrypt_pdf(compressed, encrypted, "open-password")
        encrypted_reader = PdfReader(encrypted)
        try:
            self.assertTrue(encrypted_reader.is_encrypted)
        finally:
            encrypted_reader.close()

        with self.assertRaises(pdf_processor.PDFPasswordError):
            pdf_processor.decrypt_pdf(encrypted, self.root / "错误密码.pdf", "wrong")

        decrypted = self.root / "解密.pdf"
        pdf_processor.decrypt_pdf(encrypted, decrypted, "open-password")
        decrypted_reader = PdfReader(decrypted)
        try:
            self.assertFalse(decrypted_reader.is_encrypted)
        finally:
            decrypted_reader.close()
        self.assertEqual(_page_texts(decrypted), ["Secret", "Second"])

    def test_lossless_compression_skips_oversized_decoded_streams(self) -> None:
        source = _make_text_pdf(self.root / "超大内容流.pdf", ("Page one", "Page two"))
        output = self.root / "无损优化.pdf"

        with mock.patch(
            "pypdf._page.PageObject.compress_content_streams",
            side_effect=LimitReachedError("Limit reached while decompressing"),
        ), warnings.catch_warnings(record=True) as caught:
            warnings.simplefilter("always")
            self.assertEqual(
                pdf_processor.compress_pdf(source, output, level=9),
                [output.resolve()],
            )

        self.assertEqual(_page_texts(output), ["Page one", "Page two"])
        self.assertTrue(any("超过安全阈值" in str(item.message) for item in caught))

    def test_lossless_compression_rebuilds_after_deduplication_limit(self) -> None:
        source = _make_text_pdf(self.root / "去重超限.pdf", ("Safe content",))
        output = self.root / "安全回退.pdf"

        with mock.patch.object(
            PdfWriter,
            "compress_identical_objects",
            side_effect=LimitReachedError("oversized decoded object"),
        ), warnings.catch_warnings(record=True) as caught:
            warnings.simplefilter("always")
            self.assertEqual(
                pdf_processor.compress_pdf(source, output, level=6),
                [output.resolve()],
            )

        self.assertEqual(_page_texts(output), ["Safe content"])
        self.assertTrue(any("跳过重复对象去重" in str(item.message) for item in caught))

    def test_high_fidelity_lossy_compression_rebuilds_pages_in_jpeg(self) -> None:
        source = _make_text_pdf(self.root / "大文件.pdf", ("Page one", "Page two"))
        output = self.root / "高精度有损.pdf"

        with mock.patch.object(
            pdf_processor, "convert_from_path", side_effect=_fake_convert_from_path
        ) as renderer:
            with warnings.catch_warnings():
                warnings.simplefilter("ignore", pdf_processor.PDFProcessingWarning)
                self.assertEqual(
                    pdf_processor.compress_pdf_lossy(
                        source,
                        output,
                        strategy="raster",
                        dpi=220,
                        jpeg_quality=88,
                        color_mode="grayscale",
                    ),
                    [output.resolve()],
                )

        reader = PdfReader(output)
        try:
            self.assertEqual(len(reader.pages), 2)
            self.assertAlmostEqual(float(reader.pages[0].mediabox.width), 612, places=1)
            self.assertAlmostEqual(
                float(reader.pages[0].mediabox.height), 792, places=1
            )
        finally:
            reader.close()
        kwargs = renderer.call_args.kwargs
        self.assertEqual(kwargs["fmt"], "jpeg")
        self.assertEqual(kwargs["dpi"], 220)
        self.assertEqual(kwargs["jpegopt"]["quality"], 88)
        self.assertTrue(kwargs["grayscale"])

    def test_smart_lossy_compression_preserves_pdf_structure_and_line_art(self) -> None:
        pymupdf = __import__("pymupdf")
        source = self.root / "结构保留输入.pdf"
        photo_path = self.root / "高分辨率照片.jpg"
        line_art_path = self.root / "黑白线稿.png"

        Image.effect_noise((1800, 1200), 70).convert("RGB").save(
            photo_path, format="JPEG", quality=98
        )
        line_art = Image.new("1", (1200, 600), 1)
        for x in range(40, 1160, 80):
            for y in range(20, 580):
                line_art.putpixel((x, y), 0)
        line_art.save(line_art_path, format="PNG")

        document = pymupdf.open()
        page = document.new_page(width=612, height=792)
        page.insert_text((72, 72), "Selectable structure text", fontsize=14)
        page.draw_line((72, 95), (540, 95), color=(0, 0, 0), width=1)
        page.insert_image(pymupdf.Rect(72, 120, 288, 264), filename=str(photo_path))
        page.insert_image(pymupdf.Rect(72, 300, 288, 408), filename=str(line_art_path))
        page.insert_link(
            {
                "kind": pymupdf.LINK_URI,
                "from": pymupdf.Rect(72, 50, 260, 78),
                "uri": "https://example.com/structure",
            }
        )
        document.set_toc([[1, "Structure start", 1]])
        document.save(str(source), garbage=4, deflate=True)
        document.close()

        original = pymupdf.open(str(source))
        try:
            original_images = original[0].get_images(full=True)
            original_bitonal = next(item for item in original_images if item[4] == 1)
            original_bitonal_size = (original_bitonal[2], original_bitonal[3])
        finally:
            original.close()

        output = self.root / "结构保留智能压缩.pdf"
        self.assertEqual(
            pdf_processor.compress_pdf_lossy(
                source,
                output,
                strategy="smart",
                dpi=150,
                jpeg_quality=72,
            ),
            [output.resolve()],
        )
        self.assertLess(output.stat().st_size, source.stat().st_size)

        compressed = pymupdf.open(str(output))
        try:
            self.assertEqual(compressed.page_count, 1)
            compressed_page = compressed[0]
            self.assertIn("Selectable structure text", compressed_page.get_text())
            self.assertTrue(compressed_page.get_drawings())
            self.assertEqual(compressed.get_toc(), [[1, "Structure start", 1]])
            links = compressed_page.get_links()
            self.assertEqual(links[0]["uri"], "https://example.com/structure")
            compressed_images = compressed_page.get_images(full=True)
            compressed_bitonal = next(
                item for item in compressed_images if item[4] == 1
            )
            self.assertEqual(
                (compressed_bitonal[2], compressed_bitonal[3]),
                original_bitonal_size,
            )
        finally:
            compressed.close()

    def test_smart_lossy_compression_rejects_unknown_strategy(self) -> None:
        source = _make_text_pdf(self.root / "策略.pdf", ("content",))
        with self.assertRaisesRegex(pdf_processor.ValidationError, "smart 或 raster"):
            pdf_processor.compress_pdf_lossy(
                source,
                self.root / "输出.pdf",
                strategy="unknown",
            )

    def test_aes_permissions_and_distinct_owner_password(self) -> None:
        source = _make_text_pdf(self.root / "permissions.pdf", ("Secret",))

        for owner_password in (None, "reader-password"):
            with self.subTest(owner_password=owner_password):
                target = self.root / f"invalid-{owner_password or 'missing'}.pdf"
                with self.assertRaisesRegex(
                    pdf_processor.ValidationError, "所有者密码"
                ):
                    pdf_processor.encrypt_pdf(
                        source,
                        target,
                        "reader-password",
                        owner_password=owner_password,
                        allow_copy=False,
                    )
                self.assertFalse(target.exists())

        encrypted = self.root / "restricted-aes.pdf"
        pdf_processor.encrypt_pdf(
            source,
            encrypted,
            "reader-password",
            owner_password="owner-password",
            algorithm="AES-256-R5",
            allow_print=True,
            allow_modify=False,
            allow_copy=False,
            allow_annotate=False,
            allow_fill_forms=False,
            allow_assemble=False,
        )

        reader = PdfReader(encrypted)
        try:
            encryption = reader.trailer["/Encrypt"].get_object()
            self.assertEqual(int(encryption["/V"]), 5)
            self.assertEqual(int(encryption["/R"]), 5)
            self.assertEqual(str(encryption["/CF"]["/StdCF"]["/CFM"]), "/AESV3")
            self.assertEqual(int(reader.decrypt("reader-password")), 1)
            self.assertTrue(reader.are_permissions_valid)
            permissions = reader.user_access_permissions
            self.assertIsNotNone(permissions)
            assert permissions is not None
            self.assertTrue(permissions & UserAccessPermissions.PRINT)
            self.assertTrue(permissions & UserAccessPermissions.PRINT_TO_REPRESENTATION)
            for denied in (
                UserAccessPermissions.MODIFY,
                UserAccessPermissions.EXTRACT,
                UserAccessPermissions.EXTRACT_TEXT_AND_GRAPHICS,
                UserAccessPermissions.ADD_OR_MODIFY,
                UserAccessPermissions.FILL_FORM_FIELDS,
                UserAccessPermissions.ASSEMBLE_DOC,
            ):
                with self.subTest(permission=denied.name):
                    self.assertFalse(permissions & denied)
        finally:
            reader.close()

        owner_reader = PdfReader(encrypted)
        try:
            self.assertEqual(int(owner_reader.decrypt("owner-password")), 2)
        finally:
            owner_reader.close()

    def test_aes128_accepts_reportlab_text_string_file_identifier(self) -> None:
        original = _make_text_pdf(self.root / "reportlab-original.pdf", ("AES-128",))
        source = self.root / "reportlab-text-id.pdf"
        original_reader = PdfReader(original)
        writer = PdfWriter()
        try:
            writer.clone_document_from_reader(original_reader)
            writer._ID = ArrayObject(
                (
                    TextStringObject("deterministic-text-identifier"),
                    TextStringObject("deterministic-text-identifier"),
                )
            )
            with source.open("wb") as stream:
                writer.write(stream)
        finally:
            writer.close()
            original_reader.close()
        source_reader = PdfReader(source)
        try:
            identifiers = source_reader.trailer.get("/ID") or []
            self.assertTrue(identifiers)
            self.assertTrue(any(isinstance(item, str) for item in identifiers))
        finally:
            source_reader.close()

        encrypted = self.root / "reportlab-id-encrypted.pdf"
        pdf_processor.encrypt_pdf(
            source,
            encrypted,
            "reader-password",
            owner_password="owner-password",
            algorithm="AES-128",
        )

        reader = PdfReader(encrypted)
        try:
            self.assertTrue(reader.is_encrypted)
            self.assertEqual(int(reader.decrypt("reader-password")), 1)
            self.assertEqual(len(reader.pages), 1)
        finally:
            reader.close()

    def test_rotate_preserves_acroform_and_metadata(self) -> None:
        source = self.root / "form.pdf"
        document = canvas.Canvas(str(source), pagesize=(300, 200))
        document.setTitle("Customer form")
        document.setAuthor("LayoutLoom tests")
        document.drawString(20, 170, "Customer")
        document.acroForm.textfield(
            name="customer_name", x=20, y=100, width=160, height=24
        )
        document.showPage()
        document.save()

        rotated = self.root / "form-rotated.pdf"
        pdf_processor.rotate_pages(source, "1", 90, rotated)

        reader = PdfReader(rotated)
        try:
            self.assertEqual(reader.pages[0].rotation, 90)
            self.assertEqual(reader.metadata.title, "Customer form")
            self.assertEqual(reader.metadata.author, "LayoutLoom tests")
            self.assertIn("/AcroForm", reader.trailer["/Root"])
            fields = reader.get_fields() or {}
            self.assertIn("customer_name", fields)
            self.assertEqual(str(fields["customer_name"]["/FT"]), "/Tx")
        finally:
            reader.close()

    def test_compress_never_increases_size_and_warns_on_fallback(self) -> None:
        source = self.root / "already-compact.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=100, height=100)
        writer.add_metadata({"/Title": "Compact"})
        with source.open("wb") as stream:
            writer.write(stream)
        writer.close()

        original_write_pdf = pdf_processor._write_pdf

        def write_oversized_pdf(
            pending_writer: PdfWriter, target: Path, overwrite: bool
        ) -> list[Path]:
            outputs = original_write_pdf(pending_writer, target, overwrite)
            with target.open("ab") as stream:
                stream.write(b"\n% force compression fallback in this test\n")
            return outputs

        compressed = self.root / "compact-output.pdf"
        with mock.patch.object(
            pdf_processor, "_write_pdf", side_effect=write_oversized_pdf
        ):
            with warnings.catch_warnings(record=True) as caught:
                warnings.simplefilter("always")
                pdf_processor.compress_pdf(source, compressed)

        self.assertLessEqual(compressed.stat().st_size, source.stat().st_size)
        self.assertEqual(compressed.read_bytes(), source.read_bytes())
        self.assertTrue(
            any(
                issubclass(item.category, pdf_processor.PDFProcessingWarning)
                for item in caught
            )
        )

    def test_text_watermark_header_footer_and_text_export(self) -> None:
        source = _make_text_pdf(self.root / "正文.pdf", ("PageOne", "PageTwo"))
        watermarked = self.root / "水印.pdf"
        pdf_processor.add_watermark(
            source,
            watermarked,
            text="CONFIDENTIAL",
            pages=1,
            opacity=0.3,
            angle=30,
        )
        watermarked_text = _page_texts(watermarked)
        self.assertIn("CONFIDENTIAL", watermarked_text[0])
        self.assertNotIn("CONFIDENTIAL", watermarked_text[1])

        decorated = self.root / "页眉页脚.pdf"
        pdf_processor.add_header_footer(
            watermarked,
            decorated,
            header="Header {page}",
            footer="Footer",
            add_page_numbers=True,
            page_number_format="{page}/{total}",
        )
        decorated_text = _page_texts(decorated)
        self.assertIn("Header 1", decorated_text[0])
        self.assertIn("Footer", decorated_text[0])
        self.assertIn("1/2", decorated_text[0])

        text_output = self.root / "正文.txt"
        self.assertEqual(
            pdf_processor.pdf_to_text(source, text_output),
            [text_output.resolve()],
        )
        content = text_output.read_text(encoding="utf-8")
        self.assertIn("PageOne", content)
        self.assertIn("PageTwo", content)

    def test_image_watermark(self) -> None:
        source = _make_text_pdf(self.root / "图片水印源.pdf", ("Body",))
        watermark = self.root / "标志.png"
        Image.new("RGBA", (80, 40), (255, 0, 0, 128)).save(watermark)
        output = self.root / "图片水印.pdf"
        pdf_processor.add_watermark(
            source,
            output,
            image_path=watermark,
            opacity=0.5,
            angle=0,
        )
        self.assertEqual(len(_page_texts(output)), 1)
        self.assertGreater(output.stat().st_size, 0)

    def test_multiple_watermarks_are_evenly_tiled_and_count_is_validated(self) -> None:
        source = _make_text_pdf(self.root / "多水印源.pdf", ("Body",))
        for count in (1, 2, 3, 5, 6, 99, 100):
            with self.subTest(valid_count=count):
                output = self.root / f"多水印_{count}.pdf"
                pdf_processor.add_watermark(
                    source,
                    output,
                    text="INTERNAL",
                    opacity=0.25,
                    angle=30,
                    font_size=18,
                    count=count,
                )
                page_text = _page_texts(output)[0]
                self.assertEqual(page_text.count("INTERNAL"), count)

        with self.assertRaisesRegex(pdf_processor.ValidationError, "1 到 100"):
            pdf_processor.add_watermark(
                source,
                self.root / "非法数量.pdf",
                text="INTERNAL",
                count=0,
            )

        for invalid_count in (True, 1.5, float("nan"), float("inf"), 101):
            with self.subTest(count=invalid_count), self.assertRaisesRegex(
                pdf_processor.ValidationError, "1 到 100"
            ):
                pdf_processor.add_watermark(
                    source,
                    self.root / f"非法数量_{invalid_count!s}.pdf",
                    text="INTERNAL",
                    count=invalid_count,
                )

    def test_watermark_uses_visible_cropbox_and_reuses_static_overlay(self) -> None:
        source = _make_text_pdf(self.root / "裁切页源.pdf", ("one", "two"))
        reader = PdfReader(source)
        cropped = self.root / "裁切页.pdf"
        writer = PdfWriter()
        try:
            for page in reader.pages:
                page.mediabox = RectangleObject((0, 0, 612, 792))
                page.cropbox = RectangleObject((50, 100, 450, 600))
                writer.add_page(page)
            with cropped.open("wb") as stream:
                writer.write(stream)
        finally:
            writer.close()
            reader.close()

        output = self.root / "裁切页水印.pdf"
        with mock.patch.object(
            pdf_processor,
            "_overlay_page",
            wraps=pdf_processor._overlay_page,
        ) as make_overlay:
            pdf_processor.add_watermark(
                cropped,
                output,
                text="CENTER",
                opacity=0.25,
                angle=0,
                font_size=24,
            )
        self.assertEqual(make_overlay.call_count, 1)

        pymupdf = __import__("pymupdf")
        document = pymupdf.open(str(output))
        try:
            self.assertEqual(document.page_count, 2)
            for page in document:
                words = [word for word in page.get_text("words") if word[4] == "CENTER"]
                self.assertEqual(len(words), 1)
                left, top, right, bottom = words[0][:4]
                self.assertAlmostEqual((left + right) / 2, page.rect.width / 2, delta=2)
                self.assertAlmostEqual(
                    (top + bottom) / 2, page.rect.height / 2, delta=4
                )
        finally:
            document.close()

    def test_watermark_grid_handles_crop_rotation_landscape_and_cache(self) -> None:
        source = _make_watermark_geometry_pdf(self.root / "watermark_geometry.pdf")
        pymupdf = __import__("pymupdf")

        for count in (1, 2, 3, 5, 99, 100):
            with self.subTest(count=count):
                output = self.root / f"watermark_geometry_{count}.pdf"
                with mock.patch.object(
                    pdf_processor,
                    "_overlay_page",
                    wraps=pdf_processor._overlay_page,
                ) as make_overlay:
                    pdf_processor.add_watermark(
                        source,
                        output,
                        text="I",
                        opacity=1,
                        angle=30,
                        font_size=18,
                        count=count,
                    )
                self.assertEqual(make_overlay.call_count, 3)

                document = pymupdf.open(str(output))
                try:
                    self.assertEqual(document.page_count, 6)
                    for page in document:
                        words = [
                            word for word in page.get_text("words") if word[4] == "I"
                        ]
                        self.assertEqual(len(words), count)
                        if page.rotation == 0:
                            for word in words:
                                bounds = pymupdf.Rect(word[:4])
                                self.assertTrue(
                                    page.rect.contains(bounds),
                                    f"watermark outside visible page: {bounds} not in {page.rect}",
                                )
                        pixmap = page.get_pixmap(alpha=False)
                        self.assertEqual(_ink_component_count(pixmap), count)
                finally:
                    document.close()

    def test_image_watermark_counts_are_not_dropped_or_clipped(self) -> None:
        source = _make_watermark_geometry_pdf(
            self.root / "image_watermark_geometry.pdf"
        )
        watermark = self.root / "solid_red_watermark.png"
        Image.new("RGB", (80, 40), (255, 0, 0)).save(watermark)
        pymupdf = __import__("pymupdf")

        for count in (1, 2, 3, 5, 99, 100):
            with self.subTest(count=count):
                output = self.root / f"image_watermark_geometry_{count}.pdf"
                with mock.patch.object(
                    pdf_processor,
                    "_overlay_page",
                    wraps=pdf_processor._overlay_page,
                ) as make_overlay:
                    pdf_processor.add_watermark(
                        source,
                        output,
                        image_path=watermark,
                        opacity=1,
                        angle=0,
                        count=count,
                    )
                self.assertEqual(make_overlay.call_count, 3)

                reader = PdfReader(output)
                try:
                    self.assertEqual(len(reader.pages), 6)
                    for page in reader.pages:
                        content = page.get_contents()
                        self.assertIsNotNone(content)
                        image_draws = sum(
                            operator == b"Do"
                            for _operands, operator in content.operations
                        )
                        self.assertEqual(image_draws, count)
                finally:
                    reader.close()

                if count in (1, 100):
                    document = pymupdf.open(str(output))
                    try:
                        for page in document:
                            pixmap = page.get_pixmap(alpha=False)
                            self.assertEqual(
                                _ink_component_count(pixmap, red_only=True), count
                            )
                    finally:
                        document.close()

    def test_watermark_preserves_rotated_page_links_and_geometry(self) -> None:
        base = self.root / "rotated_link_base.pdf"
        document = canvas.Canvas(str(base), pagesize=letter)
        document.setFont("Helvetica", 20)
        document.drawString(72, 700, "LINK")
        document.linkURL(
            "https://example.com",
            (70, 680, 200, 730),
            relative=0,
        )
        document.showPage()
        document.save()

        source = self.root / "rotated_link.pdf"
        reader = PdfReader(base)
        writer = PdfWriter()
        try:
            page = reader.pages[0]
            page.rotate(90)
            writer.add_page(page)
            with source.open("wb") as stream:
                writer.write(stream)
        finally:
            writer.close()
            reader.close()

        source_reader = PdfReader(source)
        try:
            source_page = source_reader.pages[0]
            source_rotation = source_page.rotation
            source_annotation = source_page["/Annots"][0].get_object()
            source_annotation_rect = tuple(
                float(value) for value in source_annotation["/Rect"]
            )
        finally:
            source_reader.close()

        output = self.root / "rotated_link_watermark.pdf"
        pdf_processor.add_watermark(
            source,
            output,
            text="WM",
            opacity=0.5,
            angle=0,
            font_size=24,
            count=3,
        )

        output_reader = PdfReader(output)
        try:
            output_page = output_reader.pages[0]
            self.assertEqual(output_page.rotation, source_rotation)
            output_annotation = output_page["/Annots"][0].get_object()
            self.assertEqual(
                tuple(float(value) for value in output_annotation["/Rect"]),
                source_annotation_rect,
            )
            self.assertEqual(output_annotation["/A"]["/URI"], "https://example.com")
        finally:
            output_reader.close()

        pymupdf = __import__("pymupdf")
        source_document = pymupdf.open(str(source))
        output_document = pymupdf.open(str(output))
        try:
            source_page = source_document[0]
            output_page = output_document[0]
            self.assertEqual(output_page.rotation, source_page.rotation)
            self.assertEqual(output_page.rect, source_page.rect)

            source_link = source_page.get_links()[0]["from"]
            output_link = output_page.get_links()[0]["from"]
            for actual, expected in zip(output_link, source_link, strict=True):
                self.assertAlmostEqual(actual, expected, places=4)

            source_word = next(
                word for word in source_page.get_text("words") if word[4] == "LINK"
            )
            output_word = next(
                word for word in output_page.get_text("words") if word[4] == "LINK"
            )
            for actual, expected in zip(output_word[:4], source_word[:4], strict=True):
                self.assertAlmostEqual(actual, expected, places=4)

            watermark_words = [
                word for word in output_page.get_text("words") if word[4] == "WM"
            ]
            self.assertEqual(len(watermark_words), 3)
            for word in watermark_words:
                self.assertTrue(output_page.rect.contains(pymupdf.Rect(word[:4])))
        finally:
            output_document.close()
            source_document.close()

    def test_watermark_rejects_non_finite_numeric_parameters(self) -> None:
        source = _make_text_pdf(self.root / "参数源.pdf", ("Body",))
        for parameter, value in (
            ("opacity", True),
            ("opacity", float("nan")),
            ("angle", float("inf")),
            ("scale", float("-inf")),
            ("font_size", float("nan")),
        ):
            with self.subTest(parameter=parameter), self.assertRaisesRegex(
                pdf_processor.ValidationError, "数字"
            ):
                pdf_processor.add_watermark(
                    source,
                    self.root / f"非法参数_{parameter}.pdf",
                    text="INTERNAL",
                    **{parameter: value},
                )

    def test_blank_pdf_requires_ocr_for_text(self) -> None:
        blank = _make_blank_pdf(self.root / "扫描件.pdf")
        with self.assertRaisesRegex(pdf_processor.PDFTextExtractionError, "OCR"):
            pdf_processor.pdf_to_text(blank, self.root / "扫描件.txt")
        self.assertFalse((self.root / "扫描件.txt").exists())

    def test_pdf_table_to_excel(self) -> None:
        source = _make_table_pdf(self.root / "表格.pdf")
        output = self.root / "表格.xlsx"
        self.assertEqual(
            pdf_processor.pdf_to_excel(source, output),
            [output.resolve()],
        )
        workbook = load_workbook(output, read_only=True)
        try:
            worksheet = workbook[workbook.sheetnames[0]]
            self.assertEqual(worksheet.cell(1, 1).value, "Name")
            self.assertEqual(worksheet.cell(2, 1).value, "Alice")
            self.assertEqual(worksheet.cell(3, 2).value, "20")
        finally:
            workbook.close()

    def test_no_table_gives_actionable_error(self) -> None:
        source = _make_text_pdf(self.root / "无表格.pdf", ("Plain paragraph",))
        with self.assertRaisesRegex(pdf_processor.PDFTableExtractionError, "OCR"):
            pdf_processor.pdf_to_excel(source, self.root / "无表格.xlsx")

    def test_pdf_to_images_and_visual_ppt(self) -> None:
        source = _make_text_pdf(self.root / "演示.pdf", ("Slide1", "Slide2"))
        with mock.patch.object(
            pdf_processor, "convert_from_path", side_effect=_fake_convert_from_path
        ) as renderer, mock.patch.object(
            pdf_processor, "optimal_worker_count", return_value=3
        ) as worker_count:
            images = pdf_processor.pdf_to_images(
                source,
                self.root / "页面图",
                image_format="jpg",
                dpi=144,
                prefix="页",
            )
            self.assertEqual(
                [path.name for path in images], ["页_001.jpg", "页_002.jpg"]
            )
            self.assertTrue(all(path.is_file() for path in images))

            presentation_path = self.root / "视觉保真.pptx"
            self.assertEqual(
                pdf_processor.pdf_to_ppt(source, presentation_path, dpi=144),
                [presentation_path.resolve()],
            )

        self.assertEqual(
            worker_count.call_args_list, [mock.call(2, cap=4), mock.call(2, cap=4)]
        )
        self.assertEqual(
            [call.kwargs["thread_count"] for call in renderer.call_args_list], [3, 3]
        )

        presentation = Presentation(presentation_path)
        self.assertEqual(len(presentation.slides), 2)
        self.assertEqual(len(presentation.slides[0].shapes), 1)

    def test_missing_renderer_error_is_clear(self) -> None:
        source = _make_text_pdf(self.root / "待渲染.pdf", ("one",))
        error = pdf_processor.PDFInfoNotInstalledError("missing")
        with mock.patch.object(pdf_processor, "convert_from_path", side_effect=error):
            with self.assertRaisesRegex(
                pdf_processor.PDFRendererUnavailableError, "Poppler"
            ):
                pdf_processor.pdf_to_images(source, self.root / "渲染输出")

    def test_images_to_pdf(self) -> None:
        first = self.root / "红色.png"
        second = self.root / "透明.png"
        Image.new("RGB", (160, 100), (255, 0, 0)).save(first)
        Image.new("RGBA", (80, 160), (0, 0, 255, 100)).save(second)
        output = self.root / "图片合集.pdf"
        self.assertEqual(
            pdf_processor.images_to_pdf([first, second], output),
            [output.resolve()],
        )
        reader = PdfReader(output)
        try:
            self.assertEqual(len(reader.pages), 2)
        finally:
            reader.close()


if __name__ == "__main__":
    unittest.main()
