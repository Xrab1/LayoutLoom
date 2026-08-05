from __future__ import annotations

import json
import tempfile
import types
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch
from xml.etree import ElementTree as ET

from docuforge import engines
from docuforge.models import MissingEngineError, ValidationError
from docuforge.processors import office

try:
    import openpyxl
    from docx import Document
    from docx.shared import Inches
    from openpyxl.drawing.image import Image as ExcelImage
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches as PptInches

    HAS_OFFICE_LIBRARIES = True
except ImportError:
    HAS_OFFICE_LIBRARIES = False


def _make_png(path: Path, size: tuple[int, int] = (96, 64)) -> None:
    image = Image.new("RGBA", size, (220, 30, 50, 180))
    image.save(path)


def _rewrite_zip(
    source: Path, target: Path, replacements: dict[str, bytes], additions=None
) -> None:
    additions = additions or {}
    with zipfile.ZipFile(source) as reader, zipfile.ZipFile(target, "w") as writer:
        for info in reader.infolist():
            writer.writestr(info, replacements.get(info.filename, reader.read(info)))
        for name, payload in additions.items():
            writer.writestr(name, payload)


@unittest.skipUnless(
    HAS_OFFICE_LIBRARIES, "Office OOXML test dependencies are not installed"
)
class OfficeProcessorTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory(prefix="docuforge-office-")
        self.root = Path(self.temporary.name)
        self.output = self.root / "中文输出"
        self.output.mkdir()
        self.image = self.root / "示例图片.png"
        _make_png(self.image)

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def _word_document(self) -> Path:
        path = self.root / "合同模板.docx"
        document = Document()
        document.add_heading("标题", level=1)
        paragraph = document.add_paragraph()
        first = paragraph.add_run("客户：{{na")
        first.bold = True
        paragraph.add_run("me}}，编号 {{id}}")
        document.add_paragraph("   ")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "项目"
        table.cell(0, 1).text = "金额"
        table.cell(1, 0).text = "服务"
        table.cell(1, 1).text = "100"
        document.add_picture(str(self.image), width=Inches(1))
        document.sections[0].header.paragraphs[0].text = "旧页眉"
        document.save(path)
        return path

    def _excel_document(self) -> Path:
        path = self.root / "报表.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "数据"
        sheet.append(["姓名", "分数", "组合", "公式", None])
        sheet.append(["Bob", 2, "A-B", "=B2*2", None])
        sheet.append([None, None, None, None, None])
        sheet.append(["Alice", 3, "C-D", "=B4*2", None])
        sheet.append(["Bob", 2, "A-B", "=B5*2", None])
        image = ExcelImage(str(self.image))
        sheet.add_image(image, "G2")
        workbook.create_sheet("附表")["A1"] = "内容"
        workbook.save(path)
        return path

    def _presentation(self) -> Path:
        path = self.root / "演示文稿.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(
            PptInches(1), PptInches(1), PptInches(4), PptInches(1)
        )
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "测试文字"
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0, 0, 0)
        slide.shapes.add_picture(str(self.image), PptInches(1), PptInches(2))
        presentation.save(path)
        return path

    def _simple_excel_document(self) -> Path:
        path = self.root / "简单数据.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "数据"
        sheet.append(["姓名", "分数", "组合", None])
        sheet.append(["Bob", 2, "A-B", None])
        sheet.append([None, None, None, None])
        sheet.append(["Alice", 3, "C-D", None])
        sheet.append(["Bob", 2, "A-B", None])
        workbook.save(path)
        return path

    def test_word_semantic_conversions_are_valid_and_non_overwriting(self) -> None:
        source = self._word_document()
        txt = office.word_to_txt(source, self.output)[0]
        markdown = office.word_to_markdown(source, self.output)[0]
        html_path = office.word_to_html(source, self.output)[0]
        epub = office.word_to_epub(source, self.output, title="测试书")[0]

        self.assertIn("客户：{{name}}", txt.read_text(encoding="utf-8"))
        self.assertIn("| 项目 | 金额 |", markdown.read_text(encoding="utf-8"))
        self.assertIn("data:image/png;base64,", html_path.read_text(encoding="utf-8"))
        with zipfile.ZipFile(epub) as archive:
            self.assertEqual(archive.infolist()[0].filename, "mimetype")
            self.assertEqual(archive.read("mimetype"), b"application/epub+zip")
            ET.fromstring(archive.read("OEBPS/package.opf"))
            self.assertIn("OEBPS/images/image_1.png", archive.namelist())

        second_txt = office.word_to_txt(source, self.output)[0]
        self.assertNotEqual(txt, second_txt)
        self.assertTrue(second_txt.name.endswith("_1.txt"))

    def test_word_edits_preserve_cross_run_replacement_and_media_safety(self) -> None:
        source = self._word_document()
        replaced = office.word_replace_text(
            source, self.output, {"{{name}}": "张三", "{{id}}": "A-01", "张": "张张"}
        )[0]
        replaced_doc = Document(replaced)
        # ``张 -> 张张`` confirms expanding replacements terminate rather than
        # repeatedly matching their own newly inserted text.
        self.assertIn(
            "客户：张张三", "\n".join(p.text for p in replaced_doc.paragraphs)
        )

        no_blanks = office.word_remove_blank_paragraphs(replaced, self.output)[0]
        self.assertFalse(any(p.text == "   " for p in Document(no_blanks).paragraphs))

        formatted = office.word_set_typography(
            no_blanks,
            self.output,
            font_name="微软雅黑",
            font_size_pt=12,
            line_spacing=1.5,
        )[0]
        with_header = office.word_set_headers_footers(
            formatted, self.output, header_text="新页眉", footer_text="第 1 页"
        )[0]
        checked = Document(with_header)
        self.assertEqual(checked.sections[0].header.paragraphs[0].text, "新页眉")
        self.assertEqual(checked.sections[0].footer.paragraphs[0].text, "第 1 页")

        extracted = office.word_extract_images(source, self.output)
        self.assertEqual(len(extracted), 1)
        self.assertGreater(extracted[0].stat().st_size, 0)
        no_images = office.word_remove_images(source, self.output)[0]
        with zipfile.ZipFile(no_images) as archive:
            self.assertFalse(
                any(name.startswith("word/media/") for name in archive.namelist())
            )

    def test_ooxml_media_extraction_is_naturally_sorted_and_byte_exact(self) -> None:
        payloads: dict[str, bytes] = {}
        for name, size in (
            ("image1.png", (31, 17)),
            ("image2.png", (37, 19)),
            ("image10.png", (41, 23)),
        ):
            image_path = self.root / f"fixture-{name}"
            _make_png(image_path, size)
            payloads[name] = image_path.read_bytes()

        cases = (
            ("word", ".docx", "word/media/", office.word_extract_images),
            ("excel", ".xlsx", "xl/media/", office.excel_extract_images),
            ("powerpoint", ".pptx", "ppt/media/", office.ppt_extract_media),
        )
        expected_names = ["image1.png", "image2.png", "image10.png"]

        for family, extension, media_prefix, extractor in cases:
            with self.subTest(family=family):
                source = self.root / f"multi-media-{family}{extension}"
                with zipfile.ZipFile(source, "w") as archive:
                    archive.writestr(media_prefix, b"")
                    for name in ("image10.png", "image2.png", "image1.png"):
                        archive.writestr(f"{media_prefix}{name}", payloads[name])
                    archive.writestr("docProps/ignored.png", b"not package media")

                output = self.root / f"extracted-{family}"
                extracted = extractor(source, output)

                self.assertEqual(len(extracted), len(expected_names))
                self.assertEqual(
                    [path.name for path in extracted],
                    [
                        f"{source.stem}_{index:03d}_{name}"
                        for index, name in enumerate(expected_names, 1)
                    ],
                )
                self.assertEqual(
                    [path.read_bytes() for path in extracted],
                    [payloads[name] for name in expected_names],
                )

    def test_word_mail_merge_and_revision_accept_reject(self) -> None:
        source = self._word_document()
        merged = office.word_mail_merge(
            source,
            self.output,
            [{"name": "甲公司", "id": "001"}, {"name": "乙公司", "id": "002"}],
            filename_template="{name}_{id}",
        )
        self.assertEqual(
            [path.name for path in merged], ["甲公司_001.docx", "乙公司_002.docx"]
        )
        self.assertIn(
            "甲公司", "\n".join(p.text for p in Document(merged[0]).paragraphs)
        )

        base = self.root / "revision-base.docx"
        document = Document()
        document.add_paragraph("placeholder")
        document.save(base)
        with zipfile.ZipFile(base) as archive:
            root = ET.fromstring(archive.read("word/document.xml"))
        namespace = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        paragraph = root.find(f".//{{{namespace}}}p")
        self.assertIsNotNone(paragraph)
        paragraph.clear()
        run = ET.SubElement(paragraph, f"{{{namespace}}}r")
        ET.SubElement(run, f"{{{namespace}}}t").text = "Base "
        deleted = ET.SubElement(paragraph, f"{{{namespace}}}del")
        deleted_run = ET.SubElement(deleted, f"{{{namespace}}}r")
        ET.SubElement(deleted_run, f"{{{namespace}}}delText").text = "old"
        inserted = ET.SubElement(paragraph, f"{{{namespace}}}ins")
        inserted_run = ET.SubElement(inserted, f"{{{namespace}}}r")
        ET.SubElement(inserted_run, f"{{{namespace}}}t").text = "new"
        revised = self.root / "带修订.docx"
        _rewrite_zip(
            base,
            revised,
            {
                "word/document.xml": ET.tostring(
                    root, encoding="utf-8", xml_declaration=True
                )
            },
        )

        accepted = office.word_accept_revisions(revised, self.output)[0]
        rejected = office.word_reject_revisions(revised, self.output)[0]
        self.assertEqual(Document(accepted).paragraphs[0].text, "Base new")
        self.assertEqual(Document(rejected).paragraphs[0].text, "Base old")

    def test_word_comment_and_hyperlink_cleanup(self) -> None:
        base = self.root / "comment-base.docx"
        document = Document()
        document.add_paragraph("可见文字")
        document.save(base)
        with zipfile.ZipFile(base) as archive:
            doc_root = ET.fromstring(archive.read("word/document.xml"))
            rel_root = ET.fromstring(archive.read("word/_rels/document.xml.rels"))
            type_root = ET.fromstring(archive.read("[Content_Types].xml"))
        w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        paragraph = doc_root.find(f".//{{{w}}}p")
        start = ET.Element(f"{{{w}}}commentRangeStart", {f"{{{w}}}id": "0"})
        paragraph.insert(0, start)
        end = ET.SubElement(paragraph, f"{{{w}}}commentRangeEnd", {f"{{{w}}}id": "0"})
        reference_run = ET.SubElement(paragraph, f"{{{w}}}r")
        ET.SubElement(reference_run, f"{{{w}}}commentReference", {f"{{{w}}}id": "0"})
        hyperlink = ET.SubElement(
            paragraph, f"{{{w}}}hyperlink", {f"{{{r}}}id": "rLink"}
        )
        hyperlink_run = ET.SubElement(hyperlink, f"{{{w}}}r")
        ET.SubElement(hyperlink_run, f"{{{w}}}t").text = "链接文字"
        ET.SubElement(
            rel_root,
            f"{{{rel_ns}}}Relationship",
            {"Id": "rComment", "Type": f"{r}/comments", "Target": "comments.xml"},
        )
        ET.SubElement(
            rel_root,
            f"{{{rel_ns}}}Relationship",
            {
                "Id": "rLink",
                "Type": f"{r}/hyperlink",
                "Target": "https://example.com",
                "TargetMode": "External",
            },
        )
        ET.SubElement(
            type_root,
            f"{{{ct_ns}}}Override",
            {
                "PartName": "/word/comments.xml",
                "ContentType": "application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml",
            },
        )
        comments = (
            f'<?xml version="1.0" encoding="UTF-8"?><w:comments xmlns:w="{w}">'
            '<w:comment w:id="0"><w:p><w:r><w:t>批注</w:t></w:r></w:p></w:comment></w:comments>'
        ).encode()
        commented = self.root / "批注链接.docx"
        _rewrite_zip(
            base,
            commented,
            {
                "word/document.xml": ET.tostring(
                    doc_root, encoding="utf-8", xml_declaration=True
                ),
                "word/_rels/document.xml.rels": ET.tostring(
                    rel_root, encoding="utf-8", xml_declaration=True
                ),
                "[Content_Types].xml": ET.tostring(
                    type_root, encoding="utf-8", xml_declaration=True
                ),
            },
            {"word/comments.xml": comments},
        )
        clean = office.word_clean_comments_and_revisions(
            commented, self.output, revision_mode="keep"
        )[0]
        unlinked = office.word_remove_hyperlinks(clean, self.output)[0]
        with zipfile.ZipFile(unlinked) as archive:
            self.assertNotIn("word/comments.xml", archive.namelist())
            document_xml = archive.read("word/document.xml")
            self.assertNotIn(b"commentReference", document_xml)
            self.assertNotIn(b"hyperlink", document_xml)
        self.assertIn("链接文字", Document(unlinked).paragraphs[0].text)

    def test_excel_exports_and_formula_cache_refusal(self) -> None:
        source = self._excel_document()
        csv_outputs = office.excel_to_csv(source, self.output, data_only=False)
        txt_outputs = office.excel_to_txt(source, self.output, data_only=False)
        json_path = office.excel_to_json(source, self.output, data_only=False)[0]
        xml_path = office.excel_to_xml(source, self.output, data_only=False)[0]
        self.assertEqual({path.suffix for path in csv_outputs}, {".csv"})
        self.assertEqual({path.suffix for path in txt_outputs}, {".txt"})
        payload = json.loads(json_path.read_text(encoding="utf-8"))
        self.assertEqual(payload["数据"][0]["姓名"], "Bob")
        self.assertEqual(payload["数据"][0]["公式"], "=B2*2")
        ET.parse(xml_path)
        with self.assertRaisesRegex(ValidationError, "没有可用的计算缓存"):
            office.excel_formulas_to_values(source, self.output)

    def test_excel_data_operations(self) -> None:
        source = self._excel_document()
        sorted_path = office.excel_sort_rows(
            source, self.output, column="B", sheet_name="数据"
        )[0]
        workbook = openpyxl.load_workbook(sorted_path, data_only=False)
        self.assertEqual(workbook["数据"]["A2"].value, "Bob")
        workbook.close()

        filtered = office.excel_filter_rows(
            source, self.output, column="A", value="Bob", sheet_name="数据"
        )[0]
        workbook = openpyxl.load_workbook(filtered)
        self.assertTrue(workbook["数据"].row_dimensions[4].hidden)
        self.assertFalse(workbook["数据"].row_dimensions[2].hidden)
        workbook.close()

        simple_source = self._simple_excel_document()
        deduped = office.excel_remove_duplicates(
            simple_source, self.output, columns=["A", "B", "C"], sheet_name="数据"
        )[0]
        workbook = openpyxl.load_workbook(deduped)
        values = [workbook["数据"].cell(row=row, column=1).value for row in range(2, 6)]
        self.assertEqual(values.count("Bob"), 1)
        workbook.close()

        cleaned = office.excel_remove_blank_rows_columns(
            simple_source, self.output, sheet_name="数据"
        )[0]
        replaced = office.excel_replace_text(
            cleaned, self.output, {"Bob": "鲍勃"}, sheet_names=["数据"], exact=True
        )[0]
        split = office.excel_split_column(
            replaced, self.output, column="C", delimiter="-", sheet_name="数据"
        )[0]
        merged = office.excel_merge_cells(
            split, self.output, "H1:I1", sheet_name="数据"
        )[0]
        unmerged = office.excel_unmerge_cells(
            merged, self.output, "H1:I1", sheet_name="数据", fill=True
        )[0]
        conditional = office.excel_apply_conditional_format(
            unmerged, self.output, cell_range="B2:B10", threshold=2, sheet_name="数据"
        )[0]
        workbook = openpyxl.load_workbook(conditional)
        self.assertEqual(workbook["数据"]["A2"].value, "鲍勃")
        self.assertEqual(workbook["数据"]["C2"].value, "A")
        self.assertEqual(workbook["数据"]["D2"].value, "B")
        self.assertTrue(bool(workbook["数据"].conditional_formatting))
        workbook.close()

        managed = office.excel_manage_sheets(
            source,
            self.output,
            rename={"附表": "说明"},
            copy_sheets={"说明": "说明副本"},
            order=["说明副本", "数据", "说明"],
        )[0]
        workbook = openpyxl.load_workbook(managed)
        self.assertEqual(workbook.sheetnames, ["说明副本", "数据", "说明"])
        workbook.close()

        images = office.excel_extract_images(source, self.output)
        self.assertEqual(len(images), 1)

    def test_excel_sort_moves_row_layout_with_its_data(self) -> None:
        source = self.root / "带行属性.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "数据"
        sheet.append(["名称", "排序值"])
        sheet.append(["高", 3])
        sheet.append(["低", 1])
        sheet.append(["中", 2])
        sheet.row_dimensions[2].height = 42
        sheet.row_dimensions[2].hidden = True
        sheet.row_dimensions[2].outlineLevel = 2
        sheet.row_dimensions[3].height = 18
        workbook.save(source)
        workbook.close()

        sorted_path = office.excel_sort_rows(
            source, self.output, column="B", sheet_name="数据"
        )[0]
        sorted_workbook = openpyxl.load_workbook(sorted_path)
        try:
            sorted_sheet = sorted_workbook["数据"]
            self.assertEqual(
                [sorted_sheet.cell(row=row, column=1).value for row in range(2, 5)],
                ["低", "中", "高"],
            )
            self.assertEqual(sorted_sheet.row_dimensions[2].height, 18)
            self.assertEqual(sorted_sheet.row_dimensions[4].height, 42)
            self.assertTrue(sorted_sheet.row_dimensions[4].hidden)
            self.assertEqual(sorted_sheet.row_dimensions[4].outlineLevel, 2)
        finally:
            sorted_workbook.close()

    def test_powerpoint_edits_extract_and_compress(self) -> None:
        source = self._presentation()
        replaced = office.ppt_replace_fonts(source, self.output, {"Arial": "微软雅黑"})[
            0
        ]
        presentation = Presentation(replaced)
        run = presentation.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "微软雅黑")

        watermarked = office.ppt_add_watermark(
            replaced, self.output, "机密", opacity=0.25
        )[0]
        presentation = Presentation(watermarked)
        self.assertIn(
            "机密",
            [
                shape.text
                for shape in presentation.slides[0].shapes
                if shape.has_text_frame
            ],
        )

        media = office.ppt_extract_media(source, self.output)
        self.assertEqual(len(media), 1)
        compressed = office.ppt_compress_images(
            source, self.output, quality=60, max_dimension=48
        )[0]
        self.assertTrue(compressed.is_file())
        Presentation(compressed)  # package remains readable

    def test_engine_probe_and_validation(self) -> None:
        statuses = office.detect_office_engines()
        self.assertEqual(
            set(statuses),
            {
                "microsoft_office",
                "microsoft_word",
                "microsoft_excel",
                "microsoft_powerpoint",
                "libreoffice",
            },
        )
        self.assertIsInstance(statuses["microsoft_office"], office.OfficeEngineStatus)
        with self.assertRaises(ValidationError):
            office.convert_with_office(
                self._word_document(), self.output, engine="imaginary"
            )


class OfficeEngineSelectionTests(unittest.TestCase):
    @staticmethod
    def _statuses(
        *, word: bool, excel: bool, powerpoint: bool, libreoffice: bool
    ) -> dict[str, office.OfficeEngineStatus]:
        components = {
            "microsoft_word": office.OfficeEngineStatus(word, reason="word status"),
            "microsoft_excel": office.OfficeEngineStatus(excel, reason="excel status"),
            "microsoft_powerpoint": office.OfficeEngineStatus(
                powerpoint, reason="powerpoint status"
            ),
        }
        return {
            "microsoft_office": office.OfficeEngineStatus(
                any((word, excel, powerpoint)), reason="aggregate"
            ),
            **components,
            "libreoffice": office.OfficeEngineStatus(
                libreoffice,
                executable=Path("soffice") if libreoffice else None,
                reason="libreoffice status",
            ),
        }

    @staticmethod
    def _wps_statuses(
        *, writer: bool = False, spreadsheets: bool = False, presentation: bool = False
    ) -> dict[str, SimpleNamespace]:
        return {
            "writer": SimpleNamespace(available=writer, reason="writer status"),
            "spreadsheets": SimpleNamespace(
                available=spreadsheets, reason="spreadsheets status"
            ),
            "presentation": SimpleNamespace(
                available=presentation, reason="presentation status"
            ),
        }

    def test_component_probe_reports_word_excel_and_powerpoint_independently(
        self,
    ) -> None:
        def registered(prog_id: str) -> bool:
            return prog_id == "Word.Application"

        with patch.object(office.sys, "platform", "win32"), patch.object(
            office, "_registered_office_application", side_effect=registered
        ), patch.object(office, "_pywin32_available", return_value=True), patch.object(
            office,
            "_detect_libreoffice",
            return_value=office.OfficeEngineStatus(False, reason="no libreoffice"),
        ):
            statuses = office.detect_office_engines()
        self.assertTrue(statuses["microsoft_word"].available)
        self.assertFalse(statuses["microsoft_excel"].available)
        self.assertFalse(statuses["microsoft_powerpoint"].available)
        self.assertTrue(statuses["microsoft_office"].available)
        self.assertIn("Word", statuses["microsoft_office"].reason)

    def test_auto_selection_is_family_specific_and_prefers_wps_before_libreoffice(
        self,
    ) -> None:
        word_only = self._statuses(
            word=True, excel=False, powerpoint=False, libreoffice=True
        )
        no_wps = self._wps_statuses()
        self.assertEqual(
            office._select_conversion_engine(
                Path("报表.xlsx"), "auto", word_only, no_wps
            ),
            "libreoffice",
        )
        with self.assertRaisesRegex(MissingEngineError, "excel status"):
            office._select_conversion_engine(
                Path("报表.xlsx"), "microsoft_office", word_only, no_wps
            )

        wps_excel = self._wps_statuses(spreadsheets=True)
        self.assertEqual(
            office._select_conversion_engine(
                Path("报表.xlsx"), "auto", word_only, wps_excel
            ),
            "wps",
        )
        self.assertEqual(
            office._select_conversion_engine(
                Path("合同.docx"), "auto", word_only, self._wps_statuses(writer=True)
            ),
            "wps",
        )

    def test_office_timeout_must_be_finite_and_positive(self) -> None:
        with tempfile.TemporaryDirectory() as folder:
            source = Path(folder) / "source.docx"
            source.write_bytes(b"source")
            for invalid in (0, -1, float("nan"), float("inf")):
                with self.subTest(timeout=invalid), self.assertRaisesRegex(
                    ValidationError, "有限且大于 0"
                ):
                    office.convert_with_office(source, folder, timeout=invalid)

    def test_convert_with_office_delegates_auto_excel_to_wps_when_only_word_exists(
        self,
    ) -> None:
        statuses = self._statuses(
            word=True, excel=False, powerpoint=False, libreoffice=True
        )
        wps_statuses = self._wps_statuses(spreadsheets=True)
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "报表.xlsx"
            source.write_bytes(b"placeholder")
            expected = root / "out" / "报表.pdf"

            def convert_wps(*_args: object, **_kwargs: object) -> list[Path]:
                expected.parent.mkdir(parents=True, exist_ok=True)
                expected.write_bytes(b"converted")
                return [expected]

            with patch.object(
                office, "detect_office_engines", return_value=statuses
            ), patch(
                "docuforge.processors.wps.detect_wps_engines",
                return_value=wps_statuses,
            ), patch(
                "docuforge.processors.wps.convert_with_wps",
                side_effect=convert_wps,
            ) as convert_wps:
                result = office.convert_with_office(
                    source, root / "out", "pdf", engine="auto"
                )
        self.assertEqual(result, [expected])
        convert_wps.assert_called_once_with(
            source.resolve(),
            root / "out",
            "pdf",
            overwrite=False,
            excel_pdf_layout="smart",
            excel_pdf_paper="auto",
            excel_pdf_orientation="auto",
            excel_pdf_margin="auto",
        )

    def test_auto_falls_back_to_microsoft_after_wps_fails_quickly(self) -> None:
        statuses = self._statuses(
            word=True, excel=True, powerpoint=True, libreoffice=False
        )
        wps_statuses = self._wps_statuses(spreadsheets=True)
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "报表.xlsx"
            source.write_bytes(b"placeholder")
            expected = root / "out" / "报表.pdf"
            with patch.object(
                office, "detect_office_engines", return_value=statuses
            ), patch(
                "docuforge.processors.wps.detect_wps_engines",
                return_value=wps_statuses,
            ), patch.object(
                office,
                "_convert_with_selected_engine",
                side_effect=[MissingEngineError("WPS 启动失败"), [expected]],
            ) as convert_selected:
                result = office.convert_with_office(
                    source, root / "out", "pdf", engine="auto"
                )

        self.assertEqual(result, [expected])
        self.assertEqual(
            [item.kwargs["selected"] for item in convert_selected.call_args_list],
            ["wps", "microsoft_office"],
        )

    def test_auto_fallback_order_is_wps_then_microsoft_then_libreoffice(
        self,
    ) -> None:
        statuses = self._statuses(
            word=True, excel=True, powerpoint=True, libreoffice=True
        )
        wps_statuses = self._wps_statuses(spreadsheets=True)
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "报表.xlsx"
            source.write_bytes(b"placeholder")
            expected = root / "out" / "报表.pdf"
            with patch.object(
                office, "detect_office_engines", return_value=statuses
            ), patch(
                "docuforge.processors.wps.detect_wps_engines",
                return_value=wps_statuses,
            ), patch.object(
                office,
                "_convert_with_selected_engine",
                side_effect=[
                    MissingEngineError("WPS 导出失败"),
                    MissingEngineError("Excel 导出失败"),
                    [expected],
                ],
            ) as convert_selected:
                result = office.convert_with_office(
                    source, root / "out", "pdf", engine="auto"
                )

        self.assertEqual(result, [expected])
        self.assertEqual(
            [item.kwargs["selected"] for item in convert_selected.call_args_list],
            ["wps", "microsoft_office", "libreoffice"],
        )

    def test_explicit_microsoft_failure_does_not_fall_back(self) -> None:
        statuses = self._statuses(
            word=True, excel=True, powerpoint=True, libreoffice=False
        )
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "报表.xlsx"
            source.write_bytes(b"placeholder")
            with patch.object(
                office, "detect_office_engines", return_value=statuses
            ), patch(
                "docuforge.processors.wps.detect_wps_engines"
            ) as detect_wps, patch.object(
                office,
                "_convert_with_selected_engine",
                side_effect=MissingEngineError("Excel 启动失败"),
            ) as convert_selected:
                with self.assertRaisesRegex(MissingEngineError, "Excel 启动失败"):
                    office.convert_with_office(
                        source,
                        root / "out",
                        "pdf",
                        engine="microsoft_office",
                    )

        detect_wps.assert_not_called()
        self.assertEqual(convert_selected.call_count, 1)
        self.assertEqual(
            convert_selected.call_args.kwargs["selected"], "microsoft_office"
        )

    def test_supervised_com_conversion_approves_only_new_exact_process(self) -> None:
        parent_messages = [
            {"type": "office_process", "pid": 321, "prog_id": "Excel.Application"},
            {"type": "result", "ok": True},
        ]
        ownership_responses: list[dict[str, object]] = []

        class ParentConnection:
            def poll(self, _timeout: float) -> bool:
                return bool(parent_messages)

            def recv(self) -> dict[str, object]:
                return parent_messages.pop(0)

            def send(self, message: dict[str, object]) -> None:
                ownership_responses.append(message)

            def close(self) -> None:
                return None

        class ChildConnection:
            def close(self) -> None:
                return None

        class Process:
            def __init__(self, args: tuple[object, ...]) -> None:
                self.args = args
                self.alive = False

            def start(self) -> None:
                self.alive = True
                Path(str(self.args[2])).write_bytes(b"converted")

            def is_alive(self) -> bool:
                return self.alive

            def join(self, _timeout: float) -> None:
                self.alive = False

        class Context:
            def __init__(self) -> None:
                self.process: Process | None = None

            def Pipe(self, *, duplex: bool) -> tuple[ParentConnection, ChildConnection]:
                self.assert_duplex = duplex
                return ParentConnection(), ChildConnection()

            def Process(self, **kwargs: object) -> Process:
                self.process = Process(kwargs["args"])  # type: ignore[arg-type]
                return self.process

        identity = office._OfficeProcessIdentity(
            321, Path("C:/Program Files/Microsoft Office/EXCEL.EXE"), "created"
        )
        context = Context()
        with tempfile.TemporaryDirectory() as folder, patch(
            "multiprocessing.get_context", return_value=context
        ), patch.object(
            office, "_windows_process_snapshot", return_value={}
        ), patch.object(
            office, "_windows_process_identity", side_effect=[identity, None]
        ):
            root = Path(folder)
            source = root / "source.xlsx"
            target = root / "target.pdf"
            source.write_bytes(b"source")
            office._convert_com_supervised(source, target, "pdf", timeout=5)

        self.assertTrue(context.assert_duplex)
        self.assertEqual(ownership_responses, [{"type": "ownership", "approved": True}])

    def test_supervised_com_timeout_stops_worker(self) -> None:
        class ParentConnection:
            def poll(self, _timeout: float) -> bool:
                return False

            def close(self) -> None:
                return None

        class ChildConnection:
            def close(self) -> None:
                return None

        class Process:
            def start(self) -> None:
                return None

            def is_alive(self) -> bool:
                return True

        class Context:
            def Pipe(self, *, duplex: bool) -> tuple[ParentConnection, ChildConnection]:
                return ParentConnection(), ChildConnection()

            def Process(self, **_kwargs: object) -> Process:
                return Process()

        process = Process()
        context = Context()
        with patch("multiprocessing.get_context", return_value=context), patch.object(
            office, "_windows_process_snapshot", return_value={}
        ), patch.object(office.time, "monotonic", side_effect=[0.0, 1.0]), patch.object(
            office, "_stop_com_worker"
        ) as stop_worker:
            with self.assertRaisesRegex(MissingEngineError, "转换超时"):
                office._convert_com_supervised(
                    Path("source.xlsx"), Path("target.pdf"), "pdf", timeout=0.5
                )

        stop_worker.assert_called_once()

    def test_microsoft_excel_registration_rejects_wps_and_missing_servers(
        self,
    ) -> None:
        class RegistryKey:
            def __init__(self, path: str) -> None:
                self.path = path

            def __enter__(self) -> RegistryKey:
                return self

            def __exit__(self, *_args: object) -> None:
                return None

        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            wps_server = root / "et.exe"
            wps_server.touch()
            missing_excel = root / "missing" / "EXCEL.EXE"
            values = {
                "Excel.Application\\CLSID": "{FAKE-EXCEL-CLSID}",
                "CLSID\\{FAKE-EXCEL-CLSID}\\LocalServer32": "",
            }
            fake_winreg = types.ModuleType("winreg")
            fake_winreg.HKEY_CLASSES_ROOT = object()  # type: ignore[attr-defined]
            fake_winreg.KEY_READ = 0x20019  # type: ignore[attr-defined]
            fake_winreg.KEY_WOW64_64KEY = 0x0100  # type: ignore[attr-defined]
            fake_winreg.KEY_WOW64_32KEY = 0x0200  # type: ignore[attr-defined]
            fake_winreg.OpenKey = (  # type: ignore[attr-defined]
                lambda _root, path, *_args: RegistryKey(path)
            )
            fake_winreg.QueryValueEx = (  # type: ignore[attr-defined]
                lambda key, _name: (values[key.path], 1)
            )

            for server in (wps_server, missing_excel):
                with self.subTest(server=server.name):
                    values["CLSID\\{FAKE-EXCEL-CLSID}\\LocalServer32"] = (
                        f'"{server}" /automation'
                    )
                    with patch.object(office.sys, "platform", "win32"), patch.dict(
                        "sys.modules", {"winreg": fake_winreg}
                    ):
                        self.assertFalse(
                            office._registered_office_application("Excel.Application")
                        )

            genuine_excel = root / "EXCEL.EXE"
            genuine_excel.write_bytes(b"binary")
            genuine_clsid = office._MICROSOFT_COM_CLSIDS["Excel.Application"]
            values["Excel.Application\\CLSID"] = genuine_clsid
            values[f"CLSID\\{genuine_clsid}\\LocalServer32"] = (
                f'"{genuine_excel}" /automation'
            )
            with patch.object(office.sys, "platform", "win32"), patch.dict(
                "sys.modules", {"winreg": fake_winreg}
            ):
                self.assertTrue(
                    office._registered_office_application("Excel.Application")
                )

    def test_microsoft_com_identity_requires_matching_family_and_executable(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as folder:
            install_dir = Path(folder) / "Microsoft Office" / "root" / "Office16"
            install_dir.mkdir(parents=True)
            (install_dir / "winword.exe").write_bytes(b"binary")
            with patch.object(
                office, "_registered_office_application", return_value=True
            ):
                office._validate_microsoft_com_application(
                    SimpleNamespace(Name="Microsoft Word", Path=str(install_dir)),
                    "Word.Application",
                    "Word",
                )
                with self.assertRaisesRegex(MissingEngineError, "不是 Microsoft Word"):
                    office._validate_microsoft_com_application(
                        SimpleNamespace(Name="WPS Writer", Path=str(install_dir)),
                        "Word.Application",
                        "Word",
                    )
                with self.assertRaisesRegex(MissingEngineError, "程序路径无效"):
                    office._validate_microsoft_com_application(
                        SimpleNamespace(
                            Name="Microsoft Word", Path=str(install_dir / "bad")
                        ),
                        "Word.Application",
                        "Word",
                    )

            with patch.object(
                office, "_registered_office_application", return_value=False
            ), self.assertRaisesRegex(MissingEngineError, "正式 COM 注册信息"):
                office._validate_microsoft_com_application(
                    SimpleNamespace(Name="Microsoft Word", Path=str(install_dir)),
                    "Word.Application",
                    "Word",
                )

    def test_word_pdf_uses_export_as_fixed_format(self) -> None:
        events: list[object] = []

        class Document:
            def ExportAsFixedFormat(self, path: str, export_format: int) -> None:
                events.append(("export", export_format))
                Path(path).write_bytes(b"pdf")

            def SaveAs2(self, *_args: object, **_kwargs: object) -> None:
                raise AssertionError("PDF must not use SaveAs2")

            def Close(self, _save: bool) -> None:
                return None

        class Documents:
            def Open(self, _path: str, **_kwargs: object) -> Document:
                return Document()

        class Application:
            def __init__(self) -> None:
                self.Documents = Documents()

            def Quit(self) -> None:
                return None

        pythoncom = types.ModuleType("pythoncom")
        pythoncom.CoInitialize = lambda: None  # type: ignore[attr-defined]
        pythoncom.CoUninitialize = lambda: None  # type: ignore[attr-defined]
        client = types.ModuleType("win32com.client")
        client.DispatchEx = lambda _prog_id: Application()  # type: ignore[attr-defined]
        win32com = types.ModuleType("win32com")
        win32com.client = client  # type: ignore[attr-defined]

        with tempfile.TemporaryDirectory() as folder, patch.dict(
            "sys.modules",
            {"pythoncom": pythoncom, "win32com": win32com, "win32com.client": client},
        ), patch.object(office, "_validate_microsoft_com_application"), patch.object(
            office, "_disable_automation_macros"
        ):
            root = Path(folder)
            source = root / "合同.docx"
            target = root / "合同.pdf"
            source.write_bytes(b"source")
            office._convert_com(source, target, "pdf")

        self.assertEqual(events, [("export", 17)])

    def test_powerpoint_conversion_closes_presentation_without_arguments(self) -> None:
        events: list[object] = []

        class Presentation:
            def SaveAs(self, path: str, file_format: int) -> None:
                events.append(("save", file_format))
                Path(path).write_bytes(b"pdf")

            def Close(self) -> None:
                events.append("close")

        class Presentations:
            def Open(self, _path: str, **_kwargs: object) -> Presentation:
                events.append("open")
                return Presentation()

        class Application:
            def __init__(self) -> None:
                self.Presentations = Presentations()

            def Quit(self) -> None:
                events.append("quit")

        pythoncom = types.ModuleType("pythoncom")
        pythoncom.CoInitialize = lambda: events.append("initialize")  # type: ignore[attr-defined]
        pythoncom.CoUninitialize = lambda: events.append("uninitialize")  # type: ignore[attr-defined]
        client = types.ModuleType("win32com.client")
        client.DispatchEx = lambda _prog_id: Application()  # type: ignore[attr-defined]
        win32com = types.ModuleType("win32com")
        win32com.client = client  # type: ignore[attr-defined]

        with tempfile.TemporaryDirectory() as folder, patch.dict(
            "sys.modules",
            {"pythoncom": pythoncom, "win32com": win32com, "win32com.client": client},
        ), patch.object(office, "_validate_microsoft_com_application"), patch.object(
            office, "_disable_automation_macros"
        ):
            root = Path(folder)
            source = root / "slides.pptx"
            target = root / "slides.pdf"
            source.write_bytes(b"source")
            office._convert_com(source, target, "pdf")

        self.assertIn(("save", 32), events)
        self.assertLess(events.index("close"), events.index("quit"))
        self.assertEqual(events[-1], "uninitialize")

    def test_com_worker_reports_and_requires_process_ownership(self) -> None:
        messages: list[dict[str, object]] = []

        class Connection:
            def send(self, message: dict[str, object]) -> None:
                messages.append(message)

            def poll(self, _timeout: float) -> bool:
                return True

            def recv(self) -> dict[str, object]:
                return {"type": "ownership", "approved": True}

            def close(self) -> None:
                return None

        def convert(
            _source: Path,
            target: Path,
            _target_format: str,
            **kwargs: object,
        ) -> None:
            guard = kwargs["ownership_guard"]
            self.assertTrue(callable(guard))
            self.assertTrue(guard(SimpleNamespace(), "Word.Application"))
            target.write_bytes(b"converted")

        with tempfile.TemporaryDirectory() as folder, patch.object(
            office, "_office_application_pid", return_value=321
        ), patch.object(office, "_convert_com", side_effect=convert):
            root = Path(folder)
            office._convert_com_worker_entry(
                Connection(),
                str(root / "source.docx"),
                str(root / "target.pdf"),
                "pdf",
                {},
            )

        self.assertEqual(messages[0]["type"], "office_process")
        self.assertEqual(messages[0]["pid"], 321)
        self.assertEqual(messages[-1], {"type": "result", "ok": True})

    def test_unowned_com_instance_is_not_quit(self) -> None:
        events: list[str] = []

        class Application:
            def Quit(self) -> None:
                events.append("quit")

        pythoncom = types.ModuleType("pythoncom")
        pythoncom.CoInitialize = lambda: events.append("initialize")  # type: ignore[attr-defined]
        pythoncom.CoUninitialize = lambda: events.append("uninitialize")  # type: ignore[attr-defined]
        client = types.ModuleType("win32com.client")
        client.DispatchEx = lambda _prog_id: Application()  # type: ignore[attr-defined]
        win32com = types.ModuleType("win32com")
        win32com.client = client  # type: ignore[attr-defined]

        with tempfile.TemporaryDirectory() as folder, patch.dict(
            "sys.modules",
            {"pythoncom": pythoncom, "win32com": win32com, "win32com.client": client},
        ), patch.object(office, "_validate_microsoft_com_application"):
            root = Path(folder)
            source = root / "source.docx"
            source.write_bytes(b"source")
            with self.assertRaisesRegex(MissingEngineError, "保护用户已打开的文档"):
                office._convert_com(
                    source,
                    root / "target.pdf",
                    "pdf",
                    ownership_guard=lambda _application, _prog_id: False,
                )

        self.assertEqual(events, ["initialize", "uninitialize"])

    def test_powerpoint_png_jpg_single_file_conversion_is_rejected(self) -> None:
        for target_format in ("png", "jpg"):
            with self.subTest(target_format=target_format), self.assertRaisesRegex(
                ValidationError, "PPT 转图片序列"
            ):
                office._convert_com(
                    Path("slides.pptx"),
                    Path(f"slides.{target_format}"),
                    target_format,
                )

    def test_wps_returning_a_missing_output_is_rejected(self) -> None:
        statuses = self._statuses(
            word=True, excel=True, powerpoint=True, libreoffice=False
        )
        wps_statuses = self._wps_statuses(writer=True)
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "合同.docx"
            source.write_bytes(b"source")
            missing = root / "out" / "合同.pdf"
            with patch.object(
                office, "detect_office_engines", return_value=statuses
            ), patch(
                "docuforge.processors.wps.detect_wps_engines",
                return_value=wps_statuses,
            ), patch(
                "docuforge.processors.wps.convert_with_wps", return_value=[missing]
            ):
                with self.assertRaisesRegex(MissingEngineError, "未生成有效输出文件"):
                    office.convert_with_office(
                        source, root / "out", "pdf", engine="wps"
                    )

    def test_microsoft_macro_security_is_set_before_open(self) -> None:
        events: list[object] = []

        class Document:
            def SaveAs2(self, path: str, FileFormat: int) -> None:
                Path(path).write_bytes(b"converted")

            def Close(self, _save: bool) -> None:
                events.append("close")

        class Documents:
            def Open(self, _path: str, **_kwargs: object) -> Document:
                events.append("open")
                return Document()

        class Application:
            def __init__(self) -> None:
                object.__setattr__(self, "Documents", Documents())

            def __setattr__(self, name: str, value: object) -> None:
                if name == "AutomationSecurity":
                    events.append(("security", value))
                object.__setattr__(self, name, value)

            def Quit(self) -> None:
                events.append("quit")

        pythoncom = types.ModuleType("pythoncom")
        pythoncom.CoInitialize = lambda: events.append("initialize")  # type: ignore[attr-defined]
        pythoncom.CoUninitialize = lambda: events.append("uninitialize")  # type: ignore[attr-defined]
        client = types.ModuleType("win32com.client")
        client.DispatchEx = lambda _prog_id: Application()  # type: ignore[attr-defined]
        win32com = types.ModuleType("win32com")
        win32com.client = client  # type: ignore[attr-defined]

        with tempfile.TemporaryDirectory() as folder, patch.dict(
            "sys.modules",
            {"pythoncom": pythoncom, "win32com": win32com, "win32com.client": client},
        ), patch.object(office, "_validate_microsoft_com_application"):
            root = Path(folder)
            source = root / "合同.docx"
            target = root / "合同-copy.docx"
            source.write_bytes(b"placeholder")
            office._convert_com(source, target, "docx")

        self.assertLess(events.index(("security", 3)), events.index("open"))
        self.assertEqual(events[0], "initialize")
        self.assertEqual(events[-1], "uninitialize")

    def test_powerpoint_capability_does_not_use_word_component(self) -> None:
        statuses = self._statuses(
            word=True, excel=False, powerpoint=False, libreoffice=False
        )
        with patch(
            "docuforge.processors.office.detect_office_engines",
            return_value=statuses,
        ), patch(
            "docuforge.processors.wps.detect_wps_engines",
            return_value=self._wps_statuses(),
        ), patch.object(
            engines, "find_executable", return_value=None
        ), patch.object(
            engines, "_known_soffice_paths", return_value=[]
        ):
            engines.office_render_capability.cache_clear()
            capability = engines.office_render_capability("powerpoint")
            engines.microsoft_powerpoint_capability.cache_clear()
            native = engines.microsoft_powerpoint_capability()
        engines.office_render_capability.cache_clear()
        engines.microsoft_powerpoint_capability.cache_clear()
        self.assertFalse(capability.runnable)
        self.assertFalse(native.runnable)


if __name__ == "__main__":
    unittest.main()
