from __future__ import annotations

import hashlib
import io
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Mapping, Sequence

try:
    import cv2  # type: ignore[import-not-found]
    import numpy as np  # type: ignore[import-not-found]
except Exception:  # pragma: no cover - covered by capability checks
    cv2 = None
    np = None

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models import MissingEngineError, ValidationError
from ..runner import check_cancelled, report_progress
from ..utils import atomic_output, unique_path

PathLike = str | Path

PLAN_SCHEMA = "docuforge.video-ppt-repair.v1"
_VIDEO_EXTENSIONS = {
    ".mp4",
    ".mkv",
    ".mov",
    ".avi",
    ".webm",
    ".wmv",
    ".m4v",
    ".flv",
    ".mpeg",
    ".mpg",
    ".ts",
    ".mts",
    ".m2ts",
    ".3gp",
}


@dataclass(frozen=True)
class StableFrame:
    timestamp: float
    image: "np.ndarray"
    support: int


@dataclass(frozen=True)
class _DeckPage:
    blob: bytes
    image: "np.ndarray"
    source_page: int | None
    modified_region: tuple[float, float, float, float] | None = None


def _require_dependencies() -> None:
    if cv2 is None or np is None:
        raise MissingEngineError(
            "视频 PPT 人工补修需要 OpenCV 与 NumPy；请重新运行安装脚本"
        )


def _decode_image(payload: bytes, label: str) -> "np.ndarray":
    _require_dependencies()
    array = np.frombuffer(payload, dtype=np.uint8)
    image = cv2.imdecode(array, cv2.IMREAD_COLOR)
    if image is None or image.size == 0:
        raise ValidationError(f"无法读取{label}中的幻灯片图片")
    return image


def _encode_png(image: "np.ndarray") -> bytes:
    encoded, payload = cv2.imencode(
        ".png", image, [cv2.IMWRITE_PNG_COMPRESSION, 2]
    )
    if not encoded:
        raise ValidationError("无法编码补修后的幻灯片图片")
    return payload.tobytes()


def _largest_picture(slide: object) -> object:
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if not pictures:
        raise ValidationError("PPT 中存在没有图片的页面；本工具仅用于一页一图的提取 PPT")
    return max(pictures, key=lambda shape: int(shape.width) * int(shape.height))


def ppt_slide_count(path: PathLike) -> int:
    source = Path(path)
    if not source.is_file():
        raise ValidationError(f"PPT 不存在：{source}")
    presentation = Presentation(source)
    if not presentation.slides:
        raise ValidationError("PPT 中没有页面")
    return len(presentation.slides)


def read_ppt_slide_image(path: PathLike, page: int) -> "np.ndarray":
    source = Path(path)
    presentation = Presentation(source)
    if page < 1 or page > len(presentation.slides):
        raise ValidationError(f"PPT 页码超出范围：{page}")
    picture = _largest_picture(presentation.slides[page - 1])
    return _decode_image(picture.image.blob, f"第 {page} 页")


def _parse_report_timestamp(value: object) -> float | None:
    text = str(value or "").strip()
    pieces = text.split(":")
    if len(pieces) != 3:
        return None
    try:
        hours = int(pieces[0])
        minutes = int(pieces[1])
        seconds = float(pieces[2])
    except (TypeError, ValueError):
        return None
    result = hours * 3600.0 + minutes * 60.0 + seconds
    return result if math.isfinite(result) and result >= 0 else None


def companion_report_timestamp(path: PathLike, page: int) -> float | None:
    """Read the original selected-frame time from the adjacent extraction report."""

    source = Path(path)
    if page < 1:
        return None
    stem = source.stem
    if stem.endswith("_高清幻灯片"):
        stem = stem[: -len("_高清幻灯片")]
    candidates = [source.with_name(f"{stem}_提取报告.docx")]
    candidates.extend(sorted(source.parent.glob("*提取报告.docx")))
    seen: set[Path] = set()
    for report in candidates:
        if report in seen or not report.is_file():
            continue
        seen.add(report)
        try:
            from docx import Document

            document = Document(report)
            for table in document.tables:
                if not table.rows:
                    continue
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                if "PPT 页码" not in headers or "实际采用帧" not in headers:
                    continue
                page_column = headers.index("PPT 页码")
                time_column = headers.index("实际采用帧")
                for row in table.rows[1:]:
                    if len(row.cells) <= max(page_column, time_column):
                        continue
                    try:
                        row_page = int(row.cells[page_column].text.strip())
                    except ValueError:
                        continue
                    if row_page == page:
                        return _parse_report_timestamp(
                            row.cells[time_column].text.strip()
                        )
        except (OSError, ValueError, KeyError):
            continue
    return None


def _trim_uniform_dark_bars(frame: "np.ndarray") -> tuple[int, int, int, int]:
    """Remove only contiguous, genuinely dark outer video bars."""

    height, width = frame.shape[:2]
    gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
    row_mean = gray.mean(axis=1)
    row_std = gray.std(axis=1)
    column_mean = gray.mean(axis=0)
    column_std = gray.std(axis=0)

    def run(means: "np.ndarray", deviations: "np.ndarray", reverse: bool) -> int:
        indices = range(len(means) - 1, -1, -1) if reverse else range(len(means))
        count = 0
        for index in indices:
            if float(means[index]) <= 9.0 and float(deviations[index]) <= 5.0:
                count += 1
            else:
                break
        return count

    top = run(row_mean, row_std, False)
    bottom = run(row_mean, row_std, True)
    left = run(column_mean, column_std, False)
    right = run(column_mean, column_std, True)
    if top < height * 0.012:
        top = 0
    if bottom < height * 0.012:
        bottom = 0
    if left < width * 0.012:
        left = 0
    if right < width * 0.012:
        right = 0
    if width - left - right < width * 0.60 or height - top - bottom < height * 0.60:
        return 0, 0, width, height
    return left, top, width - right, height - bottom


def align_video_frame(
    frame: "np.ndarray", target_width: int, target_height: int
) -> "np.ndarray":
    """Align a video frame to the aspect and pixel size of an extracted slide."""

    if target_width <= 0 or target_height <= 0:
        raise ValidationError("PPT 页面图片尺寸无效")
    left, top, right, bottom = _trim_uniform_dark_bars(frame)
    cropped = frame[top:bottom, left:right]
    height, width = cropped.shape[:2]
    target_ratio = target_width / target_height
    current_ratio = width / max(1, height)
    if current_ratio > target_ratio + 1e-6:
        keep_width = max(1, int(round(height * target_ratio)))
        excess = width - keep_width
        # The extraction engine's automatic lecturer handling removes a
        # sustained lower-right presenter strip.  A large aspect mismatch is
        # therefore right-anchored; a small letterbox mismatch stays centred.
        trim_left = 0 if excess >= width * 0.08 else excess // 2
        cropped = cropped[:, trim_left : trim_left + keep_width]
    elif current_ratio < target_ratio - 1e-6:
        keep_height = max(1, int(round(width / target_ratio)))
        excess = height - keep_height
        trim_top = excess // 2
        cropped = cropped[trim_top : trim_top + keep_height, :]
    if cropped.shape[1] != target_width or cropped.shape[0] != target_height:
        interpolation = (
            cv2.INTER_AREA
            if cropped.shape[1] >= target_width and cropped.shape[0] >= target_height
            else cv2.INTER_LANCZOS4
        )
        cropped = cv2.resize(
            cropped, (target_width, target_height), interpolation=interpolation
        )
    return cropped


def _open_video(path: PathLike) -> tuple[object, float, float]:
    _require_dependencies()
    source = Path(path)
    if not source.is_file() or source.suffix.casefold() not in _VIDEO_EXTENSIONS:
        raise ValidationError(f"原视频无效：{source}")
    capture = cv2.VideoCapture(str(source))
    if not capture.isOpened():
        capture.release()
        raise MissingEngineError(f"OpenCV 无法读取原视频：{source.name}")
    fps = float(capture.get(cv2.CAP_PROP_FPS) or 25.0)
    frame_count = max(1, int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 1))
    duration = max(0.0, (frame_count - 1) / max(0.001, fps))
    return capture, fps, duration


def read_aligned_video_frame(
    video_path: PathLike,
    timestamp: float,
    target_width: int,
    target_height: int,
) -> "np.ndarray":
    capture, _fps, duration = _open_video(video_path)
    try:
        moment = min(max(0.0, float(timestamp)), duration)
        capture.set(cv2.CAP_PROP_POS_MSEC, moment * 1000.0)
        ok, frame = capture.read()
        if not ok:
            raise ValidationError(f"无法读取视频 {moment:.3f} 秒处画面")
        return align_video_frame(frame, target_width, target_height)
    finally:
        capture.release()


def _frame_features(image: "np.ndarray") -> tuple["np.ndarray", "np.ndarray", float]:
    width = 320
    height = max(90, int(round(image.shape[0] * width / image.shape[1])))
    sample = cv2.resize(image, (width, height), interpolation=cv2.INTER_AREA)
    gray = cv2.cvtColor(sample, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 55, 150) > 0
    sharpness = float(cv2.Laplacian(gray, cv2.CV_32F).var())
    return gray, edges, sharpness


def _frame_distance(
    first: tuple["np.ndarray", "np.ndarray", float],
    second: tuple["np.ndarray", "np.ndarray", float],
) -> tuple[float, float]:
    first_gray, first_edges, _ = first
    second_gray, second_edges, _ = second
    mean = float(
        np.mean(np.abs(first_gray.astype(np.int16) - second_gray.astype(np.int16)))
        / 255.0
    )
    first_dilated = cv2.dilate(first_edges.astype(np.uint8), np.ones((3, 3), np.uint8))
    second_dilated = cv2.dilate(second_edges.astype(np.uint8), np.ones((3, 3), np.uint8))
    containment = min(
        float(np.count_nonzero(first_edges & (second_dilated > 0)))
        / max(1, int(np.count_nonzero(first_edges))),
        float(np.count_nonzero(second_edges & (first_dilated > 0)))
        / max(1, int(np.count_nonzero(second_edges))),
    )
    return mean, containment


def find_stable_aligned_frame(
    video_path: PathLike,
    timestamp: float,
    target_width: int,
    target_height: int,
    *,
    radius_seconds: float = 1.2,
) -> StableFrame:
    """Choose a fully settled frame near the user's timeline position."""

    capture, fps, duration = _open_video(video_path)
    try:
        centre = min(max(0.0, float(timestamp)), duration)
        step = max(0.08, min(0.18, 3.0 / max(fps, 1.0)))
        start = max(0.0, centre - max(0.2, radius_seconds))
        end = min(duration, centre + max(0.2, radius_seconds))
        moments = np.arange(start, end + step * 0.25, step).tolist()
        candidates: list[tuple[float, np.ndarray, tuple[np.ndarray, np.ndarray, float]]] = []
        for moment in moments:
            capture.set(cv2.CAP_PROP_POS_MSEC, float(moment) * 1000.0)
            ok, frame = capture.read()
            if not ok:
                continue
            aligned = align_video_frame(frame, target_width, target_height)
            candidates.append((float(moment), aligned, _frame_features(aligned)))
        if not candidates:
            raise ValidationError("附近没有可读取的视频帧")
        scores: list[tuple[int, float, float, int]] = []
        for index, (moment, _image, features) in enumerate(candidates):
            support = 1
            for neighbour in range(max(0, index - 3), min(len(candidates), index + 4)):
                if neighbour == index:
                    continue
                mean, containment = _frame_distance(features, candidates[neighbour][2])
                if mean <= 0.024 and containment >= 0.88:
                    support += 1
            _gray, edges, sharpness = features
            completeness = float(np.count_nonzero(edges))
            # Support dominates.  Later frames win ties so a half-finished
            # transition or progressively revealed slide is not selected.
            scores.append((support, completeness + 0.08 * sharpness, moment, index))
        support, _quality, _moment, selected = max(scores)
        moment, image, _features = candidates[selected]
        return StableFrame(moment, image, support)
    finally:
        capture.release()


def _parse_region(value: object) -> tuple[float, float, float, float]:
    if not isinstance(value, Sequence) or isinstance(value, (str, bytes)) or len(value) != 4:
        raise ValidationError("补修区域必须包含 x、y、宽、高四个百分比")
    try:
        x, y, width, height = (float(item) for item in value)
    except (TypeError, ValueError) as exc:
        raise ValidationError("补修区域坐标无效") from exc
    if not all(math.isfinite(item) for item in (x, y, width, height)):
        raise ValidationError("补修区域坐标必须是有限数字")
    if x < 0 or y < 0 or width <= 0 or height <= 0 or x + width > 100.0001 or y + height > 100.0001:
        raise ValidationError("补修区域必须位于页面 0–100% 范围内")
    return x, y, width, height


def _pixel_region(
    region: tuple[float, float, float, float], width: int, height: int
) -> tuple[int, int, int, int]:
    x, y, box_width, box_height = region
    left = max(0, min(width - 1, int(round(x * width / 100.0))))
    top = max(0, min(height - 1, int(round(y * height / 100.0))))
    right = max(left + 1, min(width, int(round((x + box_width) * width / 100.0))))
    bottom = max(top + 1, min(height, int(round((y + box_height) * height / 100.0))))
    return left, top, right, bottom


def _same_page_frames(
    video_path: PathLike,
    timestamp: float,
    reference: "np.ndarray",
    *,
    window_seconds: float = 5.0,
) -> tuple[list["np.ndarray"], int]:
    capture, _fps, duration = _open_video(video_path)
    try:
        target_height, target_width = reference.shape[:2]
        start = max(0.0, timestamp - window_seconds)
        end = min(duration, timestamp + window_seconds)
        moments = np.arange(start, end + 0.001, 0.35).tolist()
        reference_features = _frame_features(reference)
        frames: list[np.ndarray] = []
        selected_index = 0
        selected_distance = float("inf")
        for moment in moments:
            check_cancelled("已取消视频 PPT 人工补修")
            capture.set(cv2.CAP_PROP_POS_MSEC, float(moment) * 1000.0)
            ok, frame = capture.read()
            if not ok:
                continue
            aligned = align_video_frame(frame, target_width, target_height)
            mean, containment = _frame_distance(reference_features, _frame_features(aligned))
            if mean <= 0.095 and containment >= 0.60:
                distance = abs(float(moment) - timestamp)
                if distance < selected_distance:
                    selected_index = len(frames)
                    selected_distance = distance
                frames.append(aligned)
        if len(frames) < 3:
            stable = find_stable_aligned_frame(
                video_path, timestamp, target_width, target_height, radius_seconds=1.5
            )
            return [reference.copy(), stable.image], 1
        return frames, selected_index
    finally:
        capture.release()


def _local_temporal_repair(
    base: "np.ndarray",
    frames: Sequence["np.ndarray"],
    region: tuple[float, float, float, float],
) -> "np.ndarray":
    """Repair only low-support temporal overlays inside a user box."""

    output = base.copy()
    height, width = output.shape[:2]
    left, top, right, bottom = _pixel_region(region, width, height)
    if len(frames) < 2:
        return output
    stack = np.stack([frame[top:bottom, left:right] for frame in frames], axis=0)
    target = output[top:bottom, left:right]
    modal, modal_support = _modal_medoid(stack)
    variation = np.max(
        stack.max(axis=0).astype(np.int16) - stack.min(axis=0).astype(np.int16),
        axis=2,
    )
    similarity = np.max(
        np.abs(stack.astype(np.int16) - target[None, ...].astype(np.int16)), axis=3
    ) <= 13
    target_support = np.count_nonzero(similarity, axis=0)
    modal_difference = np.max(
        np.abs(target.astype(np.int16) - modal.astype(np.int16)), axis=2
    )
    support_limit = max(2, int(math.ceil(len(frames) * 0.38)))
    modal_minimum = max(2, int(math.ceil(len(frames) * 0.25)))
    mask = (
        (variation >= 13)
        & (modal_difference >= 10)
        & (target_support <= support_limit)
        & (modal_support >= modal_minimum)
    ).astype(np.uint8)
    if np.any(mask):
        mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, np.ones((3, 3), np.uint8))
        mask = cv2.dilate(mask, np.ones((3, 3), np.uint8), iterations=1)
        target[mask > 0] = modal[mask > 0]
    return output


def _modal_medoid(stack: "np.ndarray") -> tuple["np.ndarray", "np.ndarray"]:
    if stack.shape[0] == 1:
        return stack[0].copy(), np.ones(stack.shape[1:3], dtype=np.uint8)
    codes = (
        ((stack[:, :, :, 0].astype(np.uint16) >> 4) << 8)
        | ((stack[:, :, :, 1].astype(np.uint16) >> 4) << 4)
        | (stack[:, :, :, 2].astype(np.uint16) >> 4)
    )
    ordered = np.sort(codes, axis=0)
    current = ordered[0].copy()
    current_count = np.ones(current.shape, dtype=np.uint8)
    best_code = current.copy()
    best_count = current_count.copy()
    for index in range(1, ordered.shape[0]):
        same = ordered[index] == current
        current_count = np.where(same, current_count + 1, 1).astype(np.uint8)
        current = ordered[index]
        better = current_count > best_count
        best_count[better] = current_count[better]
        best_code[better] = current[better]
    centre = np.stack(
        (
            ((best_code >> 8) & 15).astype(np.int16) * 16 + 8,
            ((best_code >> 4) & 15).astype(np.int16) * 16 + 8,
            (best_code & 15).astype(np.int16) * 16 + 8,
        ),
        axis=2,
    )
    distance = np.zeros(codes.shape, dtype=np.int16)
    for channel in range(3):
        distance += np.abs(stack[:, :, :, channel].astype(np.int16) - centre[None, :, :, channel])
    distance[codes != best_code[None, ...]] = np.iinfo(np.int16).max
    chosen = np.argmin(distance, axis=0)
    rows, columns = np.indices(chosen.shape)
    return stack[chosen, rows, columns], best_count


def _fill_region(
    base: "np.ndarray",
    region: tuple[float, float, float, float],
    method: str,
    colour: str,
) -> "np.ndarray":
    output = base.copy()
    height, width = output.shape[:2]
    left, top, right, bottom = _pixel_region(region, width, height)
    mask = np.zeros((height, width), dtype=np.uint8)
    mask[top:bottom, left:right] = 255
    if method == "background":
        radius = max(3, min(17, min(right - left, bottom - top) // 6))
        return cv2.inpaint(output, mask, radius, cv2.INPAINT_TELEA)
    if method != "color":
        raise ValidationError("人工补修填充方式无效")
    text = str(colour or "#FFFFFF").strip().lstrip("#")
    if len(text) != 6:
        raise ValidationError("指定填充色必须是 #RRGGBB")
    try:
        red, green, blue = (int(text[index : index + 2], 16) for index in (0, 2, 4))
    except ValueError as exc:
        raise ValidationError("指定填充色必须是 #RRGGBB") from exc
    output[top:bottom, left:right] = (blue, green, red)
    return output


def repair_region_image(
    base: "np.ndarray",
    video_path: PathLike,
    timestamp: float,
    region: tuple[float, float, float, float],
    *,
    method: str = "temporal",
    colour: str = "#FFFFFF",
) -> "np.ndarray":
    if method == "temporal":
        stable = find_stable_aligned_frame(
            video_path, timestamp, base.shape[1], base.shape[0]
        )
        frames, _selected = _same_page_frames(
            video_path, stable.timestamp, stable.image
        )
        result = _local_temporal_repair(base, frames, region)
    elif method in {"background", "color"}:
        result = _fill_region(base, region, method, colour)
    else:
        raise ValidationError("人工补修方式无效")
    _validate_outside_region_unchanged(base, result, region)
    return result


def repair_regions_on_selected_frame(
    base: "np.ndarray",
    video_path: PathLike,
    timestamp: float,
    regions: Sequence[tuple[float, float, float, float]],
    *,
    method: str = "temporal",
    colour: str = "#FFFFFF",
) -> "np.ndarray":
    """Clean several declared watermark boxes without changing the chosen frame."""

    output = base.copy()
    if not regions:
        return output
    if method == "temporal":
        frames, _selected = _same_page_frames(video_path, timestamp, base)
        for region in regions:
            output = _local_temporal_repair(output, frames, region)
    elif method in {"background", "color"}:
        for region in regions:
            output = _fill_region(output, region, method, colour)
    else:
        raise ValidationError("人工补修方式无效")
    _validate_outside_regions_unchanged(base, output, regions)
    return output


def _validate_outside_region_unchanged(
    before: "np.ndarray",
    after: "np.ndarray",
    region: tuple[float, float, float, float],
) -> None:
    if before.shape != after.shape:
        raise ValidationError("局部补修二重检查失败：页面尺寸发生变化")
    height, width = before.shape[:2]
    left, top, right, bottom = _pixel_region(region, width, height)
    outside = np.ones((height, width), dtype=bool)
    outside[top:bottom, left:right] = False
    if not np.array_equal(before[outside], after[outside]):
        raise ValidationError("局部补修二重检查失败：框选区域外像素被修改")


def _validate_outside_regions_unchanged(
    before: "np.ndarray",
    after: "np.ndarray",
    regions: Sequence[tuple[float, float, float, float]],
) -> None:
    if before.shape != after.shape:
        raise ValidationError("多区域补修二重检查失败：页面尺寸发生变化")
    height, width = before.shape[:2]
    outside = np.ones((height, width), dtype=bool)
    for region in regions:
        left, top, right, bottom = _pixel_region(region, width, height)
        outside[top:bottom, left:right] = False
    if not np.array_equal(before[outside], after[outside]):
        raise ValidationError("多区域补修二重检查失败：所有框选区域外像素被修改")


def preview_action(
    pptx_path: PathLike,
    video_path: PathLike,
    action: Mapping[str, Any],
    *,
    base_page: "np.ndarray | None" = None,
) -> tuple["np.ndarray", "np.ndarray"]:
    kind = str(action.get("kind", ""))
    timestamp = float(action.get("timestamp", 0.0))
    method = str(action.get("method", "temporal"))
    colour = str(action.get("colour", "#FFFFFF"))
    if kind == "repair_region":
        page = int(action.get("page", 0))
        before = (
            base_page.copy()
            if base_page is not None
            else read_ppt_slide_image(pptx_path, page)
        )
        region = _parse_region(action.get("region"))
        after = repair_region_image(
            before, video_path, timestamp, region, method=method, colour=colour
        )
        return before, after
    if kind == "insert_page":
        presentation = Presentation(Path(pptx_path))
        if not presentation.slides:
            raise ValidationError("PPT 中没有可用于确定页面尺寸的幻灯片")
        reference = _decode_image(
            _largest_picture(presentation.slides[0]).image.blob, "PPT 首页面"
        )
        stable = find_stable_aligned_frame(
            video_path, timestamp, reference.shape[1], reference.shape[0]
        )
        before = stable.image
        raw_region = action.get("region")
        if raw_region in (None, "", []):
            return before, before.copy()
        region = _parse_region(raw_region)
        after = repair_region_image(
            before, video_path, stable.timestamp, region, method=method, colour=colour
        )
        return before, after
    if kind == "replace_page_frame":
        page = int(action.get("page", 0))
        reference = (
            base_page
            if base_page is not None
            else read_ppt_slide_image(pptx_path, page)
        )
        before = read_aligned_video_frame(
            video_path,
            timestamp,
            reference.shape[1],
            reference.shape[0],
        )
        raw_regions = action.get("regions")
        if raw_regions in (None, "", []):
            return before, before.copy()
        if not isinstance(raw_regions, Sequence) or isinstance(raw_regions, (str, bytes)):
            raise ValidationError("多水印区域格式无效")
        regions = [_parse_region(region) for region in raw_regions]
        after = repair_regions_on_selected_frame(
            before,
            video_path,
            timestamp,
            regions,
            method=method,
            colour=colour,
        )
        return before, after
    raise ValidationError("人工补修动作类型无效")


def render_page_after_actions(
    pptx_path: PathLike,
    video_path: PathLike,
    page: int,
    actions: Sequence[Mapping[str, Any]],
) -> "np.ndarray":
    """Render one original PPT page after all confirmed page actions.

    Insertions deliberately do not participate here because their positions are
    defined against the original deck.  Keeping this preview in the same action
    order as :func:`repair_video_ppt` prevents the workbench from showing an old
    source page after the user has already confirmed a replacement or repair.
    """

    current = read_ppt_slide_image(pptx_path, page)
    for action in actions:
        kind = str(action.get("kind", ""))
        if kind not in {"repair_region", "replace_page_frame"}:
            continue
        if int(action.get("page", 0)) != page:
            continue
        _before, current = preview_action(
            pptx_path,
            video_path,
            action,
            base_page=current,
        )
    return current


def normalize_plan(value: object, slide_count: int) -> dict[str, Any]:
    if isinstance(value, str):
        text = value.strip()
        if not text:
            raise ValidationError("请先打开人工补修工作台并保存补修方案")
        try:
            payload = json.loads(text)
        except json.JSONDecodeError as exc:
            raise ValidationError("人工补修方案已损坏，请重新打开工作台") from exc
    elif isinstance(value, Mapping):
        payload = dict(value)
    else:
        raise ValidationError("人工补修方案格式无效")
    if payload.get("schema") != PLAN_SCHEMA:
        raise ValidationError("人工补修方案版本无效，请重新打开工作台")
    actions = payload.get("actions")
    if not isinstance(actions, list) or not actions:
        raise ValidationError("人工补修方案中没有任何页面操作")
    normalized: list[dict[str, Any]] = []
    for raw in actions:
        if not isinstance(raw, Mapping):
            raise ValidationError("人工补修动作格式无效")
        kind = str(raw.get("kind", ""))
        timestamp = float(raw.get("timestamp", -1.0))
        if not math.isfinite(timestamp) or timestamp < 0:
            raise ValidationError("人工补修视频时间无效")
        method = str(raw.get("method", "temporal"))
        if method not in {"temporal", "background", "color"}:
            raise ValidationError("人工补修方式无效")
        item: dict[str, Any] = {
            "kind": kind,
            "timestamp": timestamp,
            "method": method,
            "colour": str(raw.get("colour", "#FFFFFF")),
        }
        if kind == "repair_region":
            page = int(raw.get("page", 0))
            if page < 1 or page > slide_count:
                raise ValidationError(f"待修复页码超出范围：{page}")
            item["page"] = page
            item["region"] = _parse_region(raw.get("region"))
        elif kind == "insert_page":
            position = int(raw.get("position", 0))
            if position < 1 or position > slide_count + 1:
                raise ValidationError(f"漏页插入位置超出范围：{position}")
            item["position"] = position
            if raw.get("region") not in (None, "", []):
                item["region"] = _parse_region(raw.get("region"))
        elif kind == "replace_page_frame":
            page = int(raw.get("page", 0))
            if page < 1 or page > slide_count:
                raise ValidationError(f"待替换页码超出范围：{page}")
            item["page"] = page
            raw_regions = raw.get("regions", [])
            if not isinstance(raw_regions, Sequence) or isinstance(
                raw_regions, (str, bytes)
            ):
                raise ValidationError("多水印区域格式无效")
            item["regions"] = [_parse_region(region) for region in raw_regions]
        else:
            raise ValidationError("人工补修动作类型无效")
        normalized.append(item)
    return {"schema": PLAN_SCHEMA, "actions": normalized}


def make_plan(actions: Sequence[Mapping[str, Any]]) -> str:
    return json.dumps(
        {"schema": PLAN_SCHEMA, "actions": [dict(item) for item in actions]},
        ensure_ascii=False,
        separators=(",", ":"),
    )


def _load_deck(path: Path) -> tuple[object, list[_DeckPage], str]:
    presentation = Presentation(path)
    if not presentation.slides:
        raise ValidationError("PPT 中没有页面")
    pages: list[_DeckPage] = []
    for index, slide in enumerate(presentation.slides, start=1):
        picture = _largest_picture(slide)
        if (
            picture.left != 0
            or picture.top != 0
            or picture.width != presentation.slide_width
            or picture.height != presentation.slide_height
        ):
            raise ValidationError(
                f"第 {index} 页不是完整铺满的一页一图结构，无法保证局部替换精度"
            )
        blob = bytes(picture.image.blob)
        pages.append(_DeckPage(blob, _decode_image(blob, f"第 {index} 页"), index))
    source_digest = hashlib.sha256(path.read_bytes()).hexdigest()
    return presentation, pages, source_digest


def _write_deck(
    source_presentation: object, pages: Sequence[_DeckPage], target: Path
) -> None:
    output = Presentation()
    output.slide_width = source_presentation.slide_width
    output.slide_height = source_presentation.slide_height
    blank = output.slide_layouts[6]
    for page in pages:
        slide = output.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(page.blob),
            0,
            0,
            width=output.slide_width,
            height=output.slide_height,
        )
    if len(output.slides) > len(pages):
        slide_id = output.slides._sldIdLst[0]
        output.part.drop_rel(slide_id.rId)
        del output.slides._sldIdLst[0]
    with atomic_output(target) as temporary:
        output.save(temporary)


def _validate_output(
    source: Path,
    target: Path,
    pages: Sequence[_DeckPage],
    source_digest: str,
) -> None:
    if source.resolve() == target.resolve():
        raise ValidationError("人工补修不得覆盖原 PPT")
    if hashlib.sha256(source.read_bytes()).hexdigest() != source_digest:
        raise ValidationError("二重检查失败：原 PPT 在处理过程中发生变化")
    if not target.is_file() or target.stat().st_size < 1000:
        raise ValidationError("人工补修 PPT 输出为空或不完整")
    presentation = Presentation(target)
    if len(presentation.slides) != len(pages):
        raise ValidationError("二重检查失败：输出 PPT 页数不正确")
    for index, (slide, expected) in enumerate(zip(presentation.slides, pages), start=1):
        pictures = [
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        if len(pictures) != 1:
            raise ValidationError(f"二重检查失败：第 {index} 页图片结构异常")
        picture = pictures[0]
        if (
            picture.left != 0
            or picture.top != 0
            or picture.width != presentation.slide_width
            or picture.height != presentation.slide_height
        ):
            raise ValidationError(f"二重检查失败：第 {index} 页未完整铺满")
        embedded = bytes(picture.image.blob)
        if hashlib.sha256(embedded).digest() != hashlib.sha256(expected.blob).digest():
            raise ValidationError(f"二重检查失败：第 {index} 页图片写入不一致")


def repair_video_ppt(
    pptx_path: PathLike,
    video_path: PathLike,
    output_dir: PathLike,
    repair_plan: object,
) -> Path:
    """Apply user-confirmed local repairs and missing-page insertions."""

    _require_dependencies()
    source = Path(pptx_path)
    video = Path(video_path)
    if not source.is_file() or source.suffix.casefold() != ".pptx":
        raise ValidationError("请选择由视频提取功能生成的 PPTX 文件")
    if not video.is_file() or video.suffix.casefold() not in _VIDEO_EXTENSIONS:
        raise ValidationError("请选择对应的原视频")
    source_presentation, pages, source_digest = _load_deck(source)
    plan = normalize_plan(repair_plan, len(pages))
    actions = plan["actions"]
    report_progress(0.08, "读取原 PPT 并核对一页一图结构")

    page_actions = [
        item
        for item in actions
        if item["kind"] in {"repair_region", "replace_page_frame"}
    ]
    for action_index, action in enumerate(page_actions, start=1):
        check_cancelled("已取消视频 PPT 人工补修")
        page_index = int(action["page"]) - 1
        page = pages[page_index]
        if action["kind"] == "replace_page_frame":
            before = read_aligned_video_frame(
                video,
                float(action["timestamp"]),
                page.image.shape[1],
                page.image.shape[0],
            )
            after = repair_regions_on_selected_frame(
                before,
                video,
                float(action["timestamp"]),
                action.get("regions", []),
                method=str(action["method"]),
                colour=str(action["colour"]),
            )
            changed_region = None
        else:
            before = page.image
            after = repair_region_image(
                before,
                video,
                float(action["timestamp"]),
                action["region"],
                method=str(action["method"]),
                colour=str(action["colour"]),
            )
            changed_region = action["region"]
        pages[page_index] = _DeckPage(
            _encode_png(after), after, page.source_page, changed_region
        )
        report_progress(
            0.08 + 0.48 * action_index / max(1, len(page_actions)),
            (
                f"使用手选最佳帧重建原 PPT 第 {page_index + 1} 页"
                if action["kind"] == "replace_page_frame"
                else f"局部补修原 PPT 第 {page_index + 1} 页"
            ),
        )

    insert_actions = [item for item in actions if item["kind"] == "insert_page"]
    insert_actions = sorted(
        enumerate(insert_actions), key=lambda item: (int(item[1]["position"]), item[0])
    )
    offset = 0
    reference_height, reference_width = pages[0].image.shape[:2]
    for insert_index, (_order, action) in enumerate(insert_actions, start=1):
        check_cancelled("已取消视频 PPT 人工补修")
        stable = find_stable_aligned_frame(
            video,
            float(action["timestamp"]),
            reference_width,
            reference_height,
            radius_seconds=1.5,
        )
        image = stable.image
        if "region" in action:
            image = repair_region_image(
                image,
                video,
                stable.timestamp,
                action["region"],
                method=str(action["method"]),
                colour=str(action["colour"]),
            )
        insertion = int(action["position"]) - 1 + offset
        pages.insert(
            insertion,
            _DeckPage(
                _encode_png(image),
                image,
                None,
                action.get("region"),
            ),
        )
        offset += 1
        report_progress(
            0.56 + 0.25 * insert_index / max(1, len(insert_actions)),
            f"插入漏页到第 {insertion + 1} 页",
        )

    target_root = Path(output_dir)
    target_root.mkdir(parents=True, exist_ok=True)
    target = unique_path(target_root / f"{source.stem}_人工补修.pptx")
    report_progress(0.84, "写入新的 PPT，不覆盖原文件")
    _write_deck(source_presentation, pages, target)
    report_progress(0.94, "二重检查页数、顺序、图片结构与原文件完整性")
    _validate_output(source, target, pages, source_digest)
    report_progress(1.0, f"完成：输出 {len(pages)} 页人工补修 PPT")
    return target
