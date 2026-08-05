"""Offline PDF processing primitives.

The public functions in this module deliberately use ordinary ``pathlib``
friendly signatures and always return a list of generated :class:`Path`
objects.  Page numbers exposed to callers are one-based.

Structural operations are performed with :mod:`pypdf`; text and table
extraction use :mod:`pdfplumber`; generated overlays use ReportLab.  Raster
conversion is delegated to :mod:`pdf2image`, which additionally requires a
working Poppler installation (``pdftoppm``/``pdfinfo``).
"""

from __future__ import annotations

import io
import math
import os
import re
import shutil
import tempfile
import warnings
from contextlib import ExitStack, contextmanager
from hashlib import sha1
from pathlib import Path
from typing import Any, Callable, Iterator, Mapping, Sequence

import pdfplumber
from pdf2image import convert_from_path
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError,
)
from pypdf import PdfReader, PdfWriter
from pypdf.constants import UserAccessPermissions
from pypdf.errors import DependencyError, LimitReachedError, PdfReadError
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from ..models import DocuForgeError, MissingEngineError, ValidationError
from ..utils import optimal_worker_count

try:  # Optional at import time; the feature raises a clear error if missing.
    from openpyxl import Workbook
    from openpyxl.utils import get_column_letter
except ImportError:  # pragma: no cover - exercised only in reduced installs
    Workbook = None  # type: ignore[assignment]
    get_column_letter = None  # type: ignore[assignment]

try:
    from PIL import Image, ImageOps
except ImportError:  # pragma: no cover - exercised only in reduced installs
    Image = None  # type: ignore[assignment]
    ImageOps = None  # type: ignore[assignment]

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:  # pragma: no cover - exercised only in reduced installs
    Presentation = None  # type: ignore[assignment]
    Emu = None  # type: ignore[assignment]


PathLike = str | os.PathLike[str]
PageSpec = int | str | Sequence[int]
RangeSpec = str | tuple[int, int] | Sequence[tuple[int, int] | int | str]


class PDFProcessorError(DocuForgeError):
    """Base exception for an unsuccessful PDF operation."""


class PDFPasswordError(PDFProcessorError):
    """Raised when an encrypted PDF cannot be opened with the supplied password."""


class PDFRendererUnavailableError(MissingEngineError):
    """Raised when Poppler is unavailable for a raster-based conversion."""


class PDFTextExtractionError(PDFProcessorError):
    """Raised when a PDF contains no extractable text layer."""


class PDFTableExtractionError(PDFProcessorError):
    """Raised when no table can be recovered from a PDF."""


class PDFProcessingWarning(UserWarning):
    """Non-fatal quality warning produced by a PDF operation."""


__all__ = [
    "PDFPasswordError",
    "PDFProcessingWarning",
    "PDFProcessorError",
    "PDFRendererUnavailableError",
    "PDFTableExtractionError",
    "PDFTextExtractionError",
    "add_header_footer",
    "add_watermark",
    "compress_pdf",
    "compress_pdf_lossy",
    "decrypt_pdf",
    "delete_pages",
    "encrypt_pdf",
    "extract_pages",
    "images_to_pdf",
    "insert_pages",
    "merge_pdfs",
    "pdf_to_excel",
    "pdf_to_images",
    "pdf_to_ppt",
    "pdf_to_text",
    "rotate_pages",
    "split_pdf",
]


_INVALID_FILENAME = re.compile(r'[<>:"/\\|?*\x00-\x1f]')
_ILLEGAL_EXCEL_CHARS = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def _input_file(value: PathLike, label: str = "输入文件") -> Path:
    path = Path(value).expanduser().resolve()
    if not path.is_file():
        raise ValidationError(f"{label}不存在或不是文件：{path}")
    return path


def _output_file(value: PathLike) -> Path:
    path = Path(value).expanduser()
    if not path.is_absolute():
        path = (Path.cwd() / path).resolve()
    else:
        path = path.resolve()
    if path.exists() and path.is_dir():
        raise ValidationError(f"输出位置是文件夹而不是文件：{path}")
    return path


def _output_directory(value: PathLike) -> Path:
    path = Path(value).expanduser()
    if not path.is_absolute():
        path = (Path.cwd() / path).resolve()
    else:
        path = path.resolve()
    if path.exists() and not path.is_dir():
        raise ValidationError(f"输出位置不是文件夹：{path}")
    path.mkdir(parents=True, exist_ok=True)
    return path


def _same_path(left: Path, right: Path) -> bool:
    try:
        return os.path.samefile(left, right)
    except (FileNotFoundError, OSError):
        return os.path.normcase(str(left.resolve())) == os.path.normcase(
            str(right.resolve())
        )


def _protect_inputs(target: Path, inputs: Sequence[Path]) -> None:
    if any(_same_path(target, source) for source in inputs):
        raise ValidationError("输出文件不能覆盖输入文件；请选择新的输出路径")


def _ensure_targets_available(targets: Sequence[Path], overwrite: bool) -> None:
    normalized = [os.path.normcase(str(path.resolve())) for path in targets]
    if len(set(normalized)) != len(normalized):
        raise ValidationError("多个输出将使用同一个文件名，请调整页码范围或前缀")
    if overwrite:
        return
    for target in targets:
        if target.exists():
            raise FileExistsError(
                f"输出文件已存在：{target}；如需替换请设置 overwrite=True"
            )


@contextmanager
def _atomic_output(target: Path, overwrite: bool) -> Iterator[Path]:
    """Yield a same-volume temporary file and commit it only on success."""

    target.parent.mkdir(parents=True, exist_ok=True)
    _ensure_targets_available([target], overwrite)
    descriptor, temp_name = tempfile.mkstemp(
        prefix=f".{target.stem}.", suffix=f".tmp{target.suffix}", dir=target.parent
    )
    os.close(descriptor)
    temporary = Path(temp_name)
    try:
        yield temporary
        if not temporary.is_file() or temporary.stat().st_size == 0:
            raise PDFProcessorError(f"处理器没有生成有效的输出文件：{target.name}")
        # Recheck to reduce the chance of an accidental race-time overwrite.
        if target.exists() and not overwrite:
            raise FileExistsError(
                f"输出文件已存在：{target}；如需替换请设置 overwrite=True"
            )
        os.replace(temporary, target)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise


@contextmanager
def _open_reader(path: Path, password: str | None = None) -> Iterator[PdfReader]:
    try:
        reader = PdfReader(path, strict=False)
    except (PdfReadError, OSError, ValueError) as exc:
        raise PDFProcessorError(f"无法读取 PDF：{path.name}（{exc}）") from exc

    try:
        if reader.is_encrypted:
            if password is None:
                raise PDFPasswordError(f"PDF 已加密，需要提供密码：{path.name}")
            try:
                result = reader.decrypt(password)
            except Exception as exc:
                raise PDFPasswordError(
                    f"PDF 密码无效或加密方式不受支持：{path.name}"
                ) from exc
            if int(result) == 0:
                raise PDFPasswordError(f"PDF 密码错误：{path.name}")
        yield reader
    finally:
        reader.close()


def _copy_metadata(reader: PdfReader, writer: PdfWriter) -> None:
    metadata = reader.metadata
    if not metadata:
        return
    cleaned = {
        str(key): str(value) for key, value in metadata.items() if value is not None
    }
    if cleaned:
        writer.add_metadata(cleaned)


def _has_applied_signature(reader: PdfReader) -> bool:
    fields = reader.get_fields() or {}
    for field in fields.values():
        if str(field.get("/FT", "")) != "/Sig":
            continue
        value = field.get("/V")
        value_object = value.get_object() if hasattr(value, "get_object") else value
        if (
            value_object
            and hasattr(value_object, "get")
            and (
                value_object.get("/ByteRange") is not None
                or value_object.get("/Contents") is not None
            )
        ):
            return True
    root = reader.trailer.get("/Root")
    root_object = root.get_object() if hasattr(root, "get_object") else root
    return bool(
        root_object
        and hasattr(root_object, "get")
        and root_object.get("/Perms") is not None
    )


def _reject_applied_signature(reader: PdfReader) -> None:
    if _has_applied_signature(reader):
        raise PDFProcessorError(
            "该 PDF 包含数字签名；重新加密或解密会使现有签名失效，因此已停止。"
        )


def _write_pdf(writer: PdfWriter, target: Path, overwrite: bool) -> list[Path]:
    try:
        with _atomic_output(target, overwrite) as temporary:
            with temporary.open("wb") as stream:
                writer.write(stream)
    finally:
        writer.close()
    return [target]


def _parse_page_spec(pages: PageSpec | None, page_count: int) -> list[int]:
    """Return zero-based page indexes while accepting one-based public values."""

    if page_count < 1:
        raise ValidationError("PDF 不包含任何页面")
    if pages is None:
        return list(range(page_count))

    raw_pages: list[int] = []
    if isinstance(pages, bool):
        raise ValidationError("页码必须是从 1 开始的整数")
    if isinstance(pages, int):
        raw_pages = [pages]
    elif isinstance(pages, str):
        text = (
            pages.strip()
            .replace("，", ",")
            .replace("；", ",")
            .replace("—", "-")
            .replace("–", "-")
            .replace("－", "-")
        )
        if not text or text.lower() in {"all", "全部"}:
            return list(range(page_count))
        for token in text.split(","):
            token = token.strip()
            if not token:
                continue
            if "-" in token:
                if token.count("-") != 1:
                    raise ValidationError(f"无效的页码范围：{token}")
                left, right = (part.strip() for part in token.split("-", 1))
                try:
                    start = int(left) if left else 1
                    end = int(right) if right else page_count
                except ValueError as exc:
                    raise ValidationError(f"无效的页码范围：{token}") from exc
                if start > end:
                    raise ValidationError(f"页码范围起点不能大于终点：{token}")
                raw_pages.extend(range(start, end + 1))
            else:
                try:
                    raw_pages.append(int(token))
                except ValueError as exc:
                    raise ValidationError(f"无效的页码：{token}") from exc
    else:
        try:
            raw_pages = [int(value) for value in pages]
        except (TypeError, ValueError) as exc:
            raise ValidationError("页码必须是整数、页码字符串或整数序列") from exc

    if not raw_pages:
        raise ValidationError("至少需要选择一页")
    invalid = next((page for page in raw_pages if page < 1 or page > page_count), None)
    if invalid is not None:
        raise ValidationError(f"页码 {invalid} 超出范围；文档共 {page_count} 页")

    result: list[int] = []
    seen: set[int] = set()
    for page in raw_pages:
        index = page - 1
        if index not in seen:
            seen.add(index)
            result.append(index)
    return result


def _parse_ranges(ranges: RangeSpec, page_count: int) -> list[tuple[int, int]]:
    """Normalize inclusive, one-based split ranges."""

    items: list[tuple[int, int] | int | str]
    if isinstance(ranges, str):
        items = [
            part.strip()
            for part in ranges.replace("，", ",").split(",")
            if part.strip()
        ]
    elif (
        isinstance(ranges, tuple)
        and len(ranges) == 2
        and all(
            isinstance(value, int) and not isinstance(value, bool) for value in ranges
        )
    ):
        items = [ranges]
    else:
        items = list(ranges)

    normalized: list[tuple[int, int]] = []
    for item in items:
        if isinstance(item, bool):
            raise ValidationError("拆分范围必须使用从 1 开始的页码")
        if isinstance(item, int):
            start = end = item
        elif isinstance(item, str):
            token = item.strip().replace("—", "-").replace("–", "-").replace("－", "-")
            if "-" in token:
                left, right = (part.strip() for part in token.split("-", 1))
                try:
                    start = int(left) if left else 1
                    end = int(right) if right else page_count
                except ValueError as exc:
                    raise ValidationError(f"无效的拆分范围：{item}") from exc
            else:
                try:
                    start = end = int(token)
                except ValueError as exc:
                    raise ValidationError(f"无效的拆分范围：{item}") from exc
        else:
            try:
                start, end = item
                start, end = int(start), int(end)
            except (TypeError, ValueError) as exc:
                raise ValidationError(f"无效的拆分范围：{item}") from exc
        if start < 1 or end < 1 or start > end or end > page_count:
            raise ValidationError(
                f"拆分范围 {start}-{end} 无效；文档共 {page_count} 页"
            )
        normalized.append((start, end))

    if not normalized:
        raise ValidationError("至少需要提供一个拆分范围")
    return normalized


def _safe_stem(value: str, fallback: str) -> str:
    cleaned = _INVALID_FILENAME.sub("_", value).strip(" .")
    return cleaned or fallback


def merge_pdfs(
    input_paths: Sequence[PathLike],
    output_path: PathLike,
    *,
    passwords: Sequence[str | None] | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Merge PDFs in order without rasterizing their pages."""

    sources = [_input_file(path, "PDF 输入文件") for path in input_paths]
    if not sources:
        raise ValidationError("至少需要选择一个 PDF 文件")
    target = _output_file(output_path)
    _protect_inputs(target, sources)
    _ensure_targets_available([target], overwrite)

    if passwords is None:
        password_values: list[str | None] = [None] * len(sources)
    elif isinstance(passwords, str):
        if len(sources) != 1:
            raise ValidationError("多个 PDF 的密码必须按输入顺序逐个提供")
        password_values = [passwords]
    else:
        password_values = list(passwords)
        if len(password_values) != len(sources):
            raise ValidationError("passwords 数量必须与输入 PDF 数量一致")

    writer = PdfWriter()
    with ExitStack() as stack:
        readers = [
            stack.enter_context(_open_reader(source, password))
            for source, password in zip(sources, password_values)
        ]
        for reader in readers:
            for page in reader.pages:
                writer.add_page(page)
        if readers:
            _copy_metadata(readers[0], writer)
        return _write_pdf(writer, target, overwrite)


def split_pdf(
    input_pdf: PathLike,
    ranges: RangeSpec,
    output_dir: PathLike,
    *,
    prefix: str | None = None,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Split a PDF into inclusive one-based page ranges."""

    source = _input_file(input_pdf, "PDF 输入文件")
    directory = _output_directory(output_dir)
    with _open_reader(source, password) as reader:
        normalized = _parse_ranges(ranges, len(reader.pages))
        base = _safe_stem(prefix or source.stem, "split")
        targets = [directory / f"{base}_{start}-{end}.pdf" for start, end in normalized]
        _ensure_targets_available(targets, overwrite)
        outputs: list[Path] = []
        for (start, end), target in zip(normalized, targets):
            writer = PdfWriter()
            for index in range(start - 1, end):
                writer.add_page(reader.pages[index])
            _copy_metadata(reader, writer)
            outputs.extend(_write_pdf(writer, target, overwrite))
        return outputs


def extract_pages(
    input_pdf: PathLike,
    pages: PageSpec,
    output_path: PathLike,
    *,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Extract selected one-based pages into a new PDF."""

    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password) as reader:
        indexes = _parse_page_spec(pages, len(reader.pages))
        writer = PdfWriter()
        for index in indexes:
            writer.add_page(reader.pages[index])
        _copy_metadata(reader, writer)
        return _write_pdf(writer, target, overwrite)


def delete_pages(
    input_pdf: PathLike,
    pages: PageSpec,
    output_path: PathLike,
    *,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Create a copy of a PDF with selected one-based pages removed."""

    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password) as reader:
        removed = set(_parse_page_spec(pages, len(reader.pages)))
        if len(removed) == len(reader.pages):
            raise ValidationError("不能删除 PDF 的全部页面")
        writer = PdfWriter()
        for index, page in enumerate(reader.pages):
            if index not in removed:
                writer.add_page(page)
        _copy_metadata(reader, writer)
        return _write_pdf(writer, target, overwrite)


def insert_pages(
    input_pdf: PathLike,
    insert_pdf: PathLike,
    position: int,
    output_path: PathLike,
    *,
    pages: PageSpec | None = None,
    password: str | None = None,
    insert_password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Insert pages before the one-based ``position``; ``page_count + 1`` appends."""

    source = _input_file(input_pdf, "主 PDF")
    inserted_source = _input_file(insert_pdf, "待插入 PDF")
    target = _output_file(output_path)
    _protect_inputs(target, [source, inserted_source])
    with ExitStack() as stack:
        reader = stack.enter_context(_open_reader(source, password))
        inserted_reader = stack.enter_context(
            _open_reader(inserted_source, insert_password)
        )
        if (
            isinstance(position, bool)
            or position < 1
            or position > len(reader.pages) + 1
        ):
            raise ValidationError(f"插入位置必须在 1 到 {len(reader.pages) + 1} 之间")
        inserted_indexes = _parse_page_spec(pages, len(inserted_reader.pages))
        writer = PdfWriter()
        insertion_index = position - 1
        for index in range(len(reader.pages) + 1):
            if index == insertion_index:
                for inserted_index in inserted_indexes:
                    writer.add_page(inserted_reader.pages[inserted_index])
            if index < len(reader.pages):
                writer.add_page(reader.pages[index])
        _copy_metadata(reader, writer)
        return _write_pdf(writer, target, overwrite)


def rotate_pages(
    input_pdf: PathLike,
    pages: PageSpec,
    angle: int,
    output_path: PathLike,
    *,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Rotate selected pages clockwise by a multiple of 90 degrees."""

    if isinstance(angle, bool) or int(angle) % 90 != 0:
        raise ValidationError("旋转角度必须是 90 的整数倍")
    normalized_angle = int(angle) % 360
    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password) as reader:
        selected = set(_parse_page_spec(pages, len(reader.pages)))
        writer = PdfWriter()
        writer.clone_document_from_reader(reader)
        for index, page in enumerate(writer.pages):
            if index in selected and normalized_angle:
                page.rotate(normalized_angle)
        return _write_pdf(writer, target, overwrite)


def compress_pdf(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    password: str | None = None,
    level: int = 9,
    overwrite: bool = False,
) -> list[Path]:
    """Apply lossless stream compression and remove duplicate/unreferenced objects."""

    if isinstance(level, bool) or not isinstance(level, int) or not 0 <= level <= 9:
        raise ValidationError("压缩级别必须是 0 到 9 的整数")
    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password) as reader:
        _reject_applied_signature(reader)

        def build_writer() -> tuple[PdfWriter, list[int]]:
            current_writer = PdfWriter()
            current_writer.clone_document_from_reader(reader)
            current_skipped_pages: list[int] = []
            for page_number, page in enumerate(current_writer.pages, start=1):
                try:
                    page.compress_content_streams(level=level)
                except LimitReachedError:
                    # A single decoded page content stream can legitimately
                    # exceed pypdf's decompression safety limit. Keeping that
                    # stream encoded is lossless and lets other pages continue.
                    current_skipped_pages.append(page_number)
            return current_writer, current_skipped_pages

        writer, skipped_pages = build_writer()
        deduplication_skipped = False
        try:
            writer.compress_identical_objects(
                remove_duplicates=True, remove_unreferenced=True
            )
        except LimitReachedError:
            # compress_identical_objects() can remove objects before it reaches
            # a later oversized decoded stream. Rebuild from the reader rather
            # than risking a partially-mutated writer.
            writer.close()
            writer, skipped_pages = build_writer()
            deduplication_skipped = True
        outputs = _write_pdf(writer, target, overwrite)
    if skipped_pages:
        page_list = "、".join(str(value) for value in skipped_pages[:12])
        suffix = "等" if len(skipped_pages) > 12 else ""
        warnings.warn(
            f"第 {page_list} 页{suffix}的解压后内容流超过安全阈值，已保留这些页面的原始编码，"
            "其余页面仍完成无损结构优化。",
            PDFProcessingWarning,
            stacklevel=2,
        )
    if deduplication_skipped:
        warnings.warn(
            "部分对象的解压后数据超过安全阈值，已自动跳过重复对象去重；"
            "内容流压缩和安全写回仍已完成。",
            PDFProcessingWarning,
            stacklevel=2,
        )
    if target.stat().st_size >= source.stat().st_size:
        with _atomic_output(target, overwrite=True) as temporary:
            shutil.copy2(source, temporary)
        warnings.warn(
            "该 PDF 已充分压缩；结构优化未能减小体积，输出已保留原始字节以避免文件变大。"
            "如仍需继续缩小，请改用“PDF 高精度有损压缩”。",
            PDFProcessingWarning,
            stacklevel=2,
        )
    return outputs


def compress_pdf_lossy(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    strategy: str = "smart",
    dpi: int = 220,
    jpeg_quality: int = 88,
    color_mode: str = "color",
    password: str | None = None,
    poppler_path: PathLike | None = None,
    thread_count: int | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Compress costly images while preserving PDF structure by default.

    ``strategy="smart"`` keeps text, embedded fonts, vector graphics, links,
    bookmarks, forms and annotations intact. Only high-resolution raster images
    are downsampled or recompressed. Existing JPEG-like images may be encoded as
    JPEG again; lossless color/gray images stay losslessly encoded, and bitonal
    line art is left untouched.

    ``strategy="raster"`` is the compatibility path for scans or malformed
    PDFs. It renders pages in bounded batches and rebuilds them as JPEG pages,
    which intentionally flattens all interactive and selectable content.

    Neither strategy applies an input file-size limit.
    """

    normalized_strategy = str(strategy).strip().lower()
    if normalized_strategy not in {"smart", "raster"}:
        raise ValidationError("压缩策略仅支持 smart 或 raster")
    if isinstance(dpi, bool) or not isinstance(dpi, int) or not 72 <= dpi <= 600:
        raise ValidationError("有损压缩清晰度必须是 72 到 600 之间的整数 DPI")
    if (
        isinstance(jpeg_quality, bool)
        or not isinstance(jpeg_quality, int)
        or not 40 <= jpeg_quality <= 100
    ):
        raise ValidationError("JPEG 质量必须是 40 到 100 之间的整数")
    normalized_color = str(color_mode).strip().lower()
    if normalized_color not in {"color", "grayscale"}:
        raise ValidationError("色彩模式仅支持 color 或 grayscale")
    if thread_count is not None and (
        isinstance(thread_count, bool)
        or not isinstance(thread_count, int)
        or thread_count < 1
    ):
        raise ValidationError("thread_count 必须至少为 1")

    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])

    with _open_reader(source, password) as reader:
        _reject_applied_signature(reader)
        if not reader.pages:
            raise ValidationError("PDF 不包含任何页面")

    if normalized_strategy == "smart":
        return _compress_pdf_lossy_smart(
            source,
            target,
            dpi=dpi,
            jpeg_quality=jpeg_quality,
            color_mode=normalized_color,
            password=password,
            overwrite=overwrite,
        )
    return _compress_pdf_lossy_raster(
        source,
        target,
        dpi=dpi,
        jpeg_quality=jpeg_quality,
        color_mode=normalized_color,
        password=password,
        poppler_path=poppler_path,
        thread_count=thread_count,
        overwrite=overwrite,
    )


def _compress_pdf_lossy_smart(
    source: Path,
    target: Path,
    *,
    dpi: int,
    jpeg_quality: int,
    color_mode: str,
    password: str | None,
    overwrite: bool,
) -> list[Path]:
    """Rewrite expensive images without flattening the PDF page structure."""

    _ensure_targets_available([target], overwrite)
    target.parent.mkdir(parents=True, exist_ok=True)

    try:
        import pymupdf
        from pymupdf import mupdf
    except ImportError as exc:  # pragma: no cover - dependency is required in builds
        raise MissingEngineError(
            "结构保留智能压缩需要 PyMuPDF；请重新运行安装程序补齐依赖"
        ) from exc

    try:
        document = pymupdf.open(str(source))
    except Exception as exc:
        raise PDFProcessorError(
            f"无法使用智能压缩引擎打开 PDF：{source.name}（{exc}）"
        ) from exc

    try:
        if document.needs_pass and not bool(
            password and document.authenticate(password)
        ):
            raise PDFPasswordError(f"PDF 密码错误：{source.name}")
        if not getattr(document, "is_pdf", False):
            raise ValidationError(f"输入文件不是有效 PDF：{source.name}")
        if document.page_count < 1:
            raise ValidationError("PDF 不包含任何页面")
        if not hasattr(document, "rewrite_images"):
            raise MissingEngineError(
                "当前 PyMuPDF 版本不支持结构保留图片重写；请升级后重试，"
                "或选择“整页栅格兼容压缩”"
            )

        # Only images materially above the requested target are downsampled.
        # Keeping a margin prevents repeated compression from needlessly
        # resampling images that are already close to the chosen resolution.
        dpi_threshold = max(dpi + 1, int(round(dpi * 1.15)))
        options = mupdf.PdfImageRewriterOptions()
        options.recompress_when = mupdf.FZ_RECOMPRESS_WHEN_SMALLER
        jpeg_lossless_images = _lossless_images_are_photographic(
            document,
            dpi_threshold=dpi_threshold,
        )

        for prefix in ("color", "gray"):
            # Existing photographic encodings can safely use JPEG again.
            setattr(
                options,
                f"{prefix}_lossy_image_subsample_method",
                mupdf.FZ_SUBSAMPLE_BICUBIC,
            )
            setattr(
                options,
                f"{prefix}_lossy_image_subsample_threshold",
                dpi_threshold,
            )
            setattr(options, f"{prefix}_lossy_image_subsample_to", dpi)
            setattr(
                options,
                f"{prefix}_lossy_image_recompress_method",
                mupdf.FZ_RECOMPRESS_JPEG,
            )
            setattr(
                options,
                f"{prefix}_lossy_image_recompress_quality",
                str(jpeg_quality),
            )

            # PNG/Flate-style images can include diagrams, screenshots or
            # transparency. Downsample only when they are genuinely expensive,
            # then retain lossless encoding instead of introducing JPEG noise.
            setattr(
                options,
                f"{prefix}_lossless_image_subsample_method",
                mupdf.FZ_SUBSAMPLE_BICUBIC,
            )
            setattr(
                options,
                f"{prefix}_lossless_image_subsample_threshold",
                dpi_threshold,
            )
            setattr(options, f"{prefix}_lossless_image_subsample_to", dpi)
            setattr(
                options,
                f"{prefix}_lossless_image_recompress_method",
                (
                    mupdf.FZ_RECOMPRESS_JPEG
                    if jpeg_lossless_images
                    else mupdf.FZ_RECOMPRESS_LOSSLESS
                ),
            )
            if jpeg_lossless_images:
                setattr(
                    options,
                    f"{prefix}_lossless_image_recompress_quality",
                    str(jpeg_quality),
                )

        # Do not set bitonal rewrite options. Fax/JBIG2/1-bit images are often
        # text or line art, where lossy JPEG conversion is visibly destructive.
        source_page_count = document.page_count
        source_text = [page.get_text("text") for page in document]
        document.rewrite_images(
            options=options,
            set_to_gray=color_mode == "grayscale",
        )

        with _atomic_output(target, overwrite) as temporary:
            document.save(
                str(temporary),
                garbage=4,
                deflate=True,
                deflate_images=True,
                deflate_fonts=True,
                use_objstms=True,
                compression_effort=100,
                encryption=pymupdf.PDF_ENCRYPT_NONE,
            )
            try:
                verification = pymupdf.open(str(temporary))
            except Exception as exc:
                raise PDFProcessorError(f"智能压缩结果无法重新打开：{exc}") from exc
            try:
                if verification.page_count != source_page_count:
                    raise PDFProcessorError(
                        "智能压缩后的页数发生变化，已停止写出以保护原文档"
                    )
                output_text = [page.get_text("text") for page in verification]
                if output_text != source_text:
                    raise PDFProcessorError(
                        "智能压缩后的可提取文本发生变化，已停止写出以保护原文档"
                    )
            finally:
                verification.close()
    except (PDFProcessorError, PDFPasswordError, MissingEngineError, ValidationError):
        raise
    except Exception as exc:
        raise PDFProcessorError(
            "结构保留智能压缩失败；原文件未被修改。可改用“整页栅格兼容压缩”"
            f"处理扫描件或兼容性异常 PDF（{exc}）"
        ) from exc
    finally:
        document.close()

    if target.stat().st_size >= source.stat().st_size:
        warnings.warn(
            "结构保留智能压缩已完成，但结果没有小于原文件；源 PDF 可能已充分优化，"
            "可降低目标 DPI / 图片质量，或对纯扫描件改用“整页栅格兼容压缩”。",
            PDFProcessingWarning,
            stacklevel=2,
        )
    return [target]


def _lossless_images_are_photographic(
    document: Any,
    *,
    dpi_threshold: int,
) -> bool:
    """Return true only when every costly lossless image looks photographic.

    PyMuPDF's image rewriter applies one policy to all lossless color/gray
    images. A deliberately conservative all-or-nothing decision prevents one
    diagram, screenshot or line-art image from being JPEG encoded merely
    because the same document also contains photographs.
    """

    if Image is None:
        return False

    seen: set[int] = set()
    candidate_count = 0
    has_costly_payload = False
    try:
        for page in document:
            for item in page.get_images(full=True):
                xref, _smask, width, height, bits_per_component = item[:5]
                if int(bits_per_component or 0) <= 1:
                    continue
                rectangles = page.get_image_rects(xref)
                is_high_dpi = any(
                    rectangle.width > 0
                    and rectangle.height > 0
                    and max(
                        float(width) * 72.0 / rectangle.width,
                        float(height) * 72.0 / rectangle.height,
                    )
                    >= dpi_threshold
                    for rectangle in rectangles
                )
                if not is_high_dpi or int(xref) in seen:
                    continue
                seen.add(int(xref))
                extracted = document.extract_image(xref)
                extension = str(extracted.get("ext", "")).lower()
                if extension in {"jpeg", "jpg", "jpx", "jp2"}:
                    continue
                payload = extracted.get("image")
                if not isinstance(payload, bytes):
                    return False
                candidate_count += 1
                has_costly_payload = has_costly_payload or len(payload) >= 128 * 1024

                with Image.open(io.BytesIO(payload)) as source_image:
                    if source_image.mode in {"1", "P"}:
                        return False
                    sample = source_image.convert("RGB")
                    sample.thumbnail((256, 256), Image.Resampling.LANCZOS)
                    pixel_count = sample.width * sample.height
                    if pixel_count < 1024:
                        return False

                    colors = sample.getcolors(maxcolors=4096)
                    if colors is not None:
                        if len(colors) < 512:
                            return False
                        dominant_ratio = (
                            max(count for count, _color in colors) / pixel_count
                        )
                        if dominant_ratio > 0.25:
                            return False

                    if float(sample.entropy()) < 5.0:
                        return False

                    luminance = sample.convert("L")
                    histogram = luminance.histogram()
                    near_white_ratio = sum(histogram[248:]) / pixel_count
                    if near_white_ratio > 0.55:
                        return False
    except Exception:
        return False

    return candidate_count > 0 and has_costly_payload


def _compress_pdf_lossy_raster(
    source: Path,
    target: Path,
    *,
    dpi: int,
    jpeg_quality: int,
    color_mode: str,
    password: str | None,
    poppler_path: PathLike | None,
    thread_count: int | None,
    overwrite: bool,
) -> list[Path]:
    """Rasterize pages in bounded batches for maximum compatibility."""

    _ensure_targets_available([target], overwrite)
    target.parent.mkdir(parents=True, exist_ok=True)

    with _open_reader(source, password) as reader:
        page_sizes: list[tuple[float, float]] = []
        for page in reader.pages:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            if int(page.rotation or 0) % 180:
                width, height = height, width
            page_sizes.append((width, height))
    if not page_sizes:
        raise ValidationError("PDF 不包含任何页面")

    workers = (
        optimal_worker_count(len(page_sizes), cap=4)
        if thread_count is None
        else thread_count
    )
    batch_size = max(4, workers * 2)
    with tempfile.TemporaryDirectory(
        prefix=".docuforge-lossy-", dir=target.parent
    ) as temp_name, _atomic_output(target, overwrite) as temporary:
        output_canvas = canvas.Canvas(str(temporary), pageCompression=1)
        for first_page in range(1, len(page_sizes) + 1, batch_size):
            last_page = min(len(page_sizes), first_page + batch_size - 1)
            rendered = _render_pdf_pages(
                source,
                Path(temp_name),
                image_format="jpeg",
                dpi=dpi,
                first_page=first_page,
                last_page=last_page,
                password=password,
                poppler_path=poppler_path,
                thread_count=workers,
                jpeg_quality=jpeg_quality,
                grayscale=color_mode == "grayscale",
            )
            expected = last_page - first_page + 1
            if len(rendered) != expected:
                raise PDFProcessorError(
                    f"有损压缩渲染页数不一致：预期 {expected} 页，实际 {len(rendered)} 页"
                )
            for offset, rendered_path in enumerate(rendered):
                width, height = page_sizes[first_page - 1 + offset]
                output_canvas.setPageSize((width, height))
                output_canvas.drawImage(
                    ImageReader(str(rendered_path)),
                    0,
                    0,
                    width=width,
                    height=height,
                    preserveAspectRatio=False,
                    mask="auto",
                )
                output_canvas.showPage()
                rendered_path.unlink(missing_ok=True)
        output_canvas.save()

    if target.stat().st_size >= source.stat().st_size:
        warnings.warn(
            "整页栅格压缩结果没有小于原文件；如以体积为优先，请选择较低 DPI 或 JPEG 质量。",
            PDFProcessingWarning,
            stacklevel=2,
        )
    return [target]


def encrypt_pdf(
    input_pdf: PathLike,
    output_path: PathLike,
    user_password: str,
    *,
    owner_password: str | None = None,
    algorithm: str = "AES-256-R5",
    allow_print: bool = True,
    allow_modify: bool = True,
    allow_copy: bool = True,
    allow_annotate: bool = True,
    allow_fill_forms: bool = True,
    allow_assemble: bool = True,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Encrypt a PDF and optionally restrict permissions for the user password."""

    if not user_password:
        raise ValidationError("打开密码不能为空")
    restrictions_requested = not all(
        (
            allow_print,
            allow_modify,
            allow_copy,
            allow_annotate,
            allow_fill_forms,
            allow_assemble,
        )
    )
    if restrictions_requested and (
        not owner_password or owner_password == user_password
    ):
        raise ValidationError("限制 PDF 权限时必须设置与打开密码不同的所有者密码")
    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password) as reader:
        _reject_applied_signature(reader)
        writer = PdfWriter()
        writer.clone_document_from_reader(reader)
        # Some valid PDFs (notably ReportLab output) expose trailer /ID values
        # as TextStringObject even though encryption revision 4 requires raw
        # bytes for the key hash.  Regenerate binary identifiers in the writer
        # instead of passing those decoded strings into pypdf's AES-128 path.
        writer._ID = None
        writer.generate_file_identifiers()
        permissions = UserAccessPermissions(0xFFFFF0C0)
        if allow_print:
            permissions |= (
                UserAccessPermissions.PRINT
                | UserAccessPermissions.PRINT_TO_REPRESENTATION
            )
        if allow_modify:
            permissions |= UserAccessPermissions.MODIFY
        if allow_copy:
            permissions |= (
                UserAccessPermissions.EXTRACT
                | UserAccessPermissions.EXTRACT_TEXT_AND_GRAPHICS
            )
        if allow_annotate:
            permissions |= UserAccessPermissions.ADD_OR_MODIFY
        if allow_fill_forms:
            permissions |= UserAccessPermissions.FILL_FORM_FIELDS
        if allow_assemble:
            permissions |= UserAccessPermissions.ASSEMBLE_DOC
        try:
            writer.encrypt(
                user_password=user_password,
                owner_password=owner_password or user_password,
                permissions_flag=permissions,
                algorithm=algorithm,
            )
        except (DependencyError, ValueError) as exc:
            writer.close()
            raise PDFProcessorError(f"无法使用加密算法 {algorithm}：{exc}") from exc
        return _write_pdf(writer, target, overwrite)


def decrypt_pdf(
    input_pdf: PathLike,
    output_path: PathLike,
    password: str,
    *,
    overwrite: bool = False,
) -> list[Path]:
    """Remove PDF encryption when the opening/owner password is known."""

    if password is None:
        raise ValidationError("解密必须提供已知密码")
    source = _input_file(input_pdf, "加密 PDF")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password) as reader:
        _reject_applied_signature(reader)
        if not reader.is_encrypted:
            warnings.warn(
                "输入 PDF 未加密，将生成普通副本", PDFProcessingWarning, stacklevel=2
            )
        writer = PdfWriter()
        writer.clone_document_from_reader(reader)
        return _write_pdf(writer, target, overwrite)


def _auto_font_path() -> Path | None:
    windows = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    for name in ("msyh.ttc", "msyh.ttf", "simsun.ttc", "simhei.ttf", "arial.ttf"):
        candidate = windows / name
        if candidate.is_file():
            return candidate
    return None


def _reportlab_font(font_path: PathLike | None, text: str) -> str:
    candidate: Path | None
    if font_path is not None:
        candidate = _input_file(font_path, "字体文件")
    elif any(ord(character) > 255 for character in text):
        candidate = _auto_font_path()
        if candidate is None:
            raise ValidationError(
                "文字包含中文或其他 Unicode 字符，请通过 font_path 提供 TrueType 字体"
            )
    else:
        return "Helvetica"

    digest = sha1(str(candidate).encode("utf-8")).hexdigest()[:12]
    font_name = f"LayoutLoom_{digest}"
    try:
        pdfmetrics.getFont(font_name)
    except KeyError:
        try:
            pdfmetrics.registerFont(TTFont(font_name, str(candidate)))
        except Exception as exc:
            raise ValidationError(f"无法加载字体文件：{candidate}（{exc}）") from exc
    return font_name


def _prepare_page_for_overlay(
    page: Any,
) -> tuple[float, float, tuple[float, float, float, float, float, float]]:
    # Watermarks and other decorations must be positioned inside the visible
    # page, not merely inside the usually larger MediaBox.  This matters for
    # scanned/cropped PDFs whose CropBox has a non-zero origin.
    visible_box = page.cropbox
    left = float(visible_box.left)
    bottom = float(visible_box.bottom)
    width = float(visible_box.width)
    height = float(visible_box.height)
    if not all(math.isfinite(value) for value in (left, bottom, width, height)):
        raise PDFProcessorError("PDF 页面可见区域包含无效坐标")
    if width <= 0 or height <= 0:
        raise PDFProcessorError("PDF 页面可见区域尺寸无效")

    # Keep /Rotate and annotation/widget rectangles untouched.  Converting a
    # rotated page's content in-place can move links and form fields outside
    # the visible page because pypdf does not transform annotation rectangles.
    # Instead, author the overlay in the user's visible coordinate system and
    # map it back into the original page coordinates only while merging.
    rotation = int(page.rotation or 0) % 360
    if rotation == 0:
        return width, height, (1.0, 0.0, 0.0, 1.0, left, bottom)
    if rotation == 90:
        return height, width, (0.0, 1.0, -1.0, 0.0, left + width, bottom)
    if rotation == 180:
        return width, height, (-1.0, 0.0, 0.0, -1.0, left + width, bottom + height)
    if rotation == 270:
        return height, width, (0.0, -1.0, 1.0, 0.0, left, bottom + height)
    raise PDFProcessorError("PDF 页面旋转角度必须是 0、90、180 或 270 度")


def _overlay_page(
    width: float,
    height: float,
    draw: Callable[[canvas.Canvas], None],
) -> tuple[PdfReader, io.BytesIO]:
    buffer = io.BytesIO()
    overlay_canvas = canvas.Canvas(buffer, pagesize=(width, height), pageCompression=1)
    draw(overlay_canvas)
    overlay_canvas.showPage()
    overlay_canvas.save()
    buffer.seek(0)
    return PdfReader(buffer), buffer


def _apply_overlay(
    source: Path,
    target: Path,
    pages: PageSpec | None,
    password: str | None,
    overwrite: bool,
    draw_page: Callable[[canvas.Canvas, float, float, int, int], None],
    *,
    static_by_page_size: bool = False,
) -> list[Path]:
    _protect_inputs(target, [source])
    handles: list[tuple[PdfReader, io.BytesIO]] = []
    try:
        with _open_reader(source, password) as reader:
            writer = PdfWriter()
            writer.clone_document_from_reader(reader)
            selected = set(_parse_page_spec(pages, len(writer.pages)))
            total = len(writer.pages)
            static_overlays: dict[tuple[float, float], tuple[PdfReader, io.BytesIO]] = (
                {}
            )
            for index, page in enumerate(writer.pages):
                if index not in selected:
                    continue
                width, height, transformation = _prepare_page_for_overlay(page)

                def draw(
                    current_canvas: canvas.Canvas, page_number: int = index + 1
                ) -> None:
                    draw_page(current_canvas, width, height, page_number, total)

                cache_key = (round(width, 6), round(height, 6))
                cached = static_overlays.get(cache_key) if static_by_page_size else None
                if cached is None:
                    overlay_reader, buffer = _overlay_page(width, height, draw)
                    handles.append((overlay_reader, buffer))
                    if static_by_page_size:
                        static_overlays[cache_key] = (overlay_reader, buffer)
                else:
                    overlay_reader, buffer = cached
                page.merge_transformed_page(
                    overlay_reader.pages[0],
                    transformation,
                    over=True,
                    expand=False,
                )
            return _write_pdf(writer, target, overwrite)
    finally:
        for overlay_reader, buffer in handles:
            overlay_reader.close()
            if not buffer.closed:
                buffer.close()


def add_watermark(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    text: str | None = None,
    image_path: PathLike | None = None,
    pages: PageSpec | None = None,
    opacity: float = 0.2,
    angle: float = 45.0,
    scale: float = 0.35,
    font_size: float = 48.0,
    count: int = 1,
    font_path: PathLike | None = None,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Add one or more evenly tiled text or image watermarks to selected pages."""

    has_text = bool(text)
    has_image = image_path is not None
    if has_text == has_image:
        raise ValidationError("文字水印和图片水印必须且只能选择一种")
    if any(isinstance(value, bool) for value in (opacity, angle, scale, font_size)):
        raise ValidationError("水印参数必须是有效数字")
    try:
        opacity_value = float(opacity)
        angle_value = float(angle)
        scale_value = float(scale)
        font_size_value = float(font_size)
    except (TypeError, ValueError, OverflowError) as exc:
        raise ValidationError("水印参数必须是有效数字") from exc
    if not all(
        math.isfinite(value)
        for value in (opacity_value, angle_value, scale_value, font_size_value)
    ):
        raise ValidationError("水印参数必须是有限数字")
    if not 0 <= opacity_value <= 1:
        raise ValidationError("水印透明度必须在 0 到 1 之间")
    if scale_value <= 0:
        raise ValidationError("图片水印比例必须大于 0")
    if font_size_value <= 0:
        raise ValidationError("水印字号必须大于 0")
    if isinstance(count, bool):
        raise ValidationError("每页水印数量必须是 1 到 100 的整数")
    try:
        watermark_count = int(count)
    except (TypeError, ValueError, OverflowError) as exc:
        raise ValidationError("每页水印数量必须是 1 到 100 的整数") from exc
    try:
        is_exact_integer = float(count) == watermark_count
    except (TypeError, ValueError, OverflowError):
        is_exact_integer = False
    if not is_exact_integer or not 1 <= watermark_count <= 100:
        raise ValidationError("每页水印数量必须是 1 到 100 的整数")

    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    font_name = _reportlab_font(font_path, text or "") if has_text else "Helvetica"

    image_reader: ImageReader | None = None
    image_buffer: io.BytesIO | None = None
    image_width = image_height = 0
    if has_image:
        if Image is None or ImageOps is None:
            raise MissingEngineError("图片水印需要安装 Pillow")
        watermark_image = _input_file(image_path, "水印图片")  # type: ignore[arg-type]
        try:
            with Image.open(watermark_image) as opened:
                normalized = ImageOps.exif_transpose(opened)
                normalized.load()
                image_width, image_height = normalized.size
                image_buffer = io.BytesIO()
                normalized.save(image_buffer, format="PNG")
                image_buffer.seek(0)
            image_reader = ImageReader(image_buffer)
        except Exception as exc:
            if image_buffer is not None:
                image_buffer.close()
            raise PDFProcessorError(
                f"无法读取水印图片：{watermark_image.name}（{exc}）"
            ) from exc

    def draw_watermark(
        current_canvas: canvas.Canvas,
        width: float,
        height: float,
        _page: int,
        _total: int,
    ) -> None:
        lines = (text or "").splitlines() or [text or ""]
        leading = font_size_value * 1.2
        if has_text:
            mark_width = max(
                1.0,
                *(
                    pdfmetrics.stringWidth(line, font_name, font_size_value)
                    for line in lines
                ),
            )
            mark_height = max(
                font_size_value,
                font_size_value + (len(lines) - 1) * leading,
            )
        else:
            assert image_width and image_height
            image_factor = min(
                width * scale_value / image_width,
                height * scale_value / image_height,
            )
            mark_width = image_width * image_factor
            mark_height = image_height * image_factor

        radians = math.radians(angle_value % 180.0)
        cosine = abs(math.cos(radians))
        sine = abs(math.sin(radians))
        rotated_width = max(1.0, mark_width * cosine + mark_height * sine)
        rotated_height = max(1.0, mark_width * sine + mark_height * cosine)

        if watermark_count == 1:
            positions = [(width / 2.0, height / 2.0)]
            fit_scale = min(
                1.0,
                width * 0.86 / rotated_width,
                height * 0.86 / rotated_height,
            )
        else:
            usable_width = width * 0.92
            usable_height = height * 0.92
            best_layout: tuple[float, int, int, float, float] | None = None
            for rows in range(1, watermark_count + 1):
                columns = math.ceil(watermark_count / rows)
                cell_width = usable_width / columns
                cell_height = usable_height / rows
                fit = min(
                    cell_width / rotated_width,
                    cell_height / rotated_height,
                )
                empty_cells = rows * columns - watermark_count
                # Prefer the grid that gives each watermark the most room while
                # mildly penalising unused cells in the final row.
                score = fit / (1.0 + empty_cells * 0.08)
                candidate = (score, -empty_cells, -rows, cell_width, cell_height)
                if best_layout is None or candidate[:3] > best_layout[:3]:
                    best_layout = candidate

            assert best_layout is not None
            _score, _empty_cells, negative_rows, cell_width, cell_height = best_layout
            rows = -negative_rows
            columns = math.ceil(watermark_count / rows)
            fit_scale = min(
                1.0,
                cell_width * 0.78 / rotated_width,
                cell_height * 0.78 / rotated_height,
            )
            positions: list[tuple[float, float]] = []
            remaining = watermark_count
            for row in range(rows):
                items_in_row = min(columns, remaining)
                if items_in_row <= 0:
                    break
                y = height / 2.0 + (rows - 1 - 2 * row) * cell_height / 2.0
                for column in range(items_in_row):
                    x = width / 2.0 + (items_in_row - 1 - 2 * column) * cell_width / 2.0
                    positions.append((x, y))
                remaining -= items_in_row

        current_canvas.saveState()
        if hasattr(current_canvas, "setFillAlpha"):
            current_canvas.setFillAlpha(opacity_value)
        current_canvas.setFillColorRGB(0.45, 0.45, 0.45)
        for x, y in positions:
            current_canvas.saveState()
            current_canvas.translate(x, y)
            current_canvas.rotate(angle_value)
            if has_text:
                effective_font_size = font_size_value * fit_scale
                effective_leading = leading * fit_scale
                current_canvas.setFont(font_name, effective_font_size)
                ascent, descent = pdfmetrics.getAscentDescent(
                    font_name, effective_font_size
                )
                first_y = ((len(lines) - 1) * effective_leading - ascent - descent) / 2
                for line_number, line in enumerate(lines):
                    current_canvas.drawCentredString(
                        0,
                        first_y - line_number * effective_leading,
                        line,
                    )
            else:
                assert image_reader is not None
                draw_width = mark_width * fit_scale
                draw_height = mark_height * fit_scale
                current_canvas.drawImage(
                    image_reader,
                    -draw_width / 2,
                    -draw_height / 2,
                    width=draw_width,
                    height=draw_height,
                    preserveAspectRatio=True,
                    mask="auto",
                )
            current_canvas.restoreState()
        current_canvas.restoreState()

    try:
        return _apply_overlay(
            source,
            target,
            pages,
            password,
            overwrite,
            draw_watermark,
            static_by_page_size=True,
        )
    finally:
        if image_buffer is not None:
            image_buffer.close()


def _format_page_text(template: str, page: int, total: int, label: str) -> str:
    try:
        return template.format(page=page, total=total)
    except (KeyError, IndexError, ValueError) as exc:
        raise ValidationError(
            f"{label}格式无效；仅支持 {{page}} 和 {{total}} 占位符"
        ) from exc


def add_header_footer(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    header: str | None = None,
    footer: str | None = None,
    add_page_numbers: bool = False,
    page_number_format: str = "{page}/{total}",
    pages: PageSpec | None = None,
    font_size: float = 10.0,
    margin: float = 24.0,
    font_path: PathLike | None = None,
    password: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Add centered headers/footers and optional page numbers."""

    if not header and not footer and not add_page_numbers:
        raise ValidationError("请至少设置页眉、页脚或页码中的一项")
    if float(font_size) <= 0:
        raise ValidationError("页眉页脚字号必须大于 0")
    if float(margin) < 0:
        raise ValidationError("页边距不能小于 0")
    all_text = "".join(value or "" for value in (header, footer, page_number_format))
    font_name = _reportlab_font(font_path, all_text)
    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)

    def draw_header_footer(
        current_canvas: canvas.Canvas,
        width: float,
        height: float,
        page: int,
        total: int,
    ) -> None:
        current_canvas.saveState()
        current_canvas.setFont(font_name, float(font_size))
        current_canvas.setFillColorRGB(0.15, 0.15, 0.15)
        if header:
            current_canvas.drawCentredString(
                width / 2,
                max(float(margin), height - float(margin) - float(font_size)),
                _format_page_text(header, page, total, "页眉"),
            )
        if footer:
            current_canvas.drawCentredString(
                width / 2,
                float(margin),
                _format_page_text(footer, page, total, "页脚"),
            )
        if add_page_numbers:
            page_text = _format_page_text(page_number_format, page, total, "页码")
            if footer:
                current_canvas.drawRightString(
                    width - float(margin), float(margin), page_text
                )
            else:
                current_canvas.drawCentredString(width / 2, float(margin), page_text)
        current_canvas.restoreState()

    return _apply_overlay(
        source, target, pages, password, overwrite, draw_header_footer
    )


def pdf_to_text(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    password: str | None = None,
    layout: bool = False,
    page_separator: str = "\n\f\n",
    overwrite: bool = False,
) -> list[Path]:
    """Extract a PDF text layer as UTF-8 text.

    A document with no extractable text raises :class:`PDFTextExtractionError`
    instead of silently creating an empty file.  Empty pages in a mixed
    document produce a warning that recommends OCR.
    """

    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    # Validate encryption first so pdfminer errors become an actionable message.
    with _open_reader(source, password):
        pass

    try:
        with pdfplumber.open(source, password=password) as document:
            texts: list[str] = []
            empty_pages: list[int] = []
            for page_number, page in enumerate(document.pages, start=1):
                value = page.extract_text(layout=layout) or ""
                texts.append(value)
                if not value.strip():
                    empty_pages.append(page_number)
    except PDFProcessorError:
        raise
    except Exception as exc:
        raise PDFProcessorError(f"提取 PDF 文本失败：{source.name}（{exc}）") from exc

    if not any(value.strip() for value in texts):
        raise PDFTextExtractionError(
            "PDF 没有可提取的文本层；文件可能是扫描件，请先执行 OCR 识别"
        )
    if empty_pages:
        page_list = "、".join(str(value) for value in empty_pages[:10])
        suffix = "等" if len(empty_pages) > 10 else ""
        warnings.warn(
            f"第 {page_list} 页{suffix}没有提取到文本；若为扫描页请先执行 OCR",
            PDFProcessingWarning,
            stacklevel=2,
        )
    with _atomic_output(target, overwrite) as temporary:
        temporary.write_text(page_separator.join(texts), encoding="utf-8")
    return [target]


def _render_pdf_pages(
    source: Path,
    output_folder: Path,
    *,
    image_format: str,
    dpi: int,
    first_page: int,
    last_page: int,
    password: str | None,
    poppler_path: PathLike | None,
    thread_count: int,
    jpeg_quality: int | None = None,
    grayscale: bool = False,
) -> list[Path]:
    format_name = "jpeg" if image_format in {"jpg", "jpeg"} else "png"
    try:
        options: dict[str, Any] = {}
        if format_name == "jpeg" and jpeg_quality is not None:
            options["jpegopt"] = {
                "quality": int(jpeg_quality),
                "progressive": True,
                "optimize": True,
            }
        rendered = convert_from_path(
            str(source),
            dpi=int(dpi),
            first_page=first_page,
            last_page=last_page,
            fmt=format_name,
            output_folder=str(output_folder),
            paths_only=True,
            userpw=password,
            poppler_path=str(Path(poppler_path)) if poppler_path else None,
            thread_count=int(thread_count),
            grayscale=bool(grayscale),
            **options,
        )
    except PDFInfoNotInstalledError as exc:
        raise PDFRendererUnavailableError(
            "PDF 转图片/PPT 需要 Poppler；请安装 pdfinfo 与 pdftoppm，或配置 poppler_path"
        ) from exc
    except PDFPageCountError as exc:
        message = str(exc)
        if "password" in message.lower() or "encrypted" in message.lower():
            raise PDFPasswordError("渲染器无法打开 PDF，请检查密码") from exc
        raise PDFProcessorError(f"无法读取 PDF 页数：{message}") from exc
    except PDFSyntaxError as exc:
        raise PDFProcessorError(f"PDF 语法错误，无法渲染：{exc}") from exc
    except (FileNotFoundError, OSError) as exc:
        raise PDFRendererUnavailableError(
            "未找到可用的 PDF 渲染器；请安装 Poppler 或配置 poppler_path"
        ) from exc
    paths = [Path(value) for value in rendered]
    if not paths or any(not path.is_file() for path in paths):
        raise PDFProcessorError("PDF 渲染器没有生成有效图片")
    return paths


def pdf_to_images(
    input_pdf: PathLike,
    output_dir: PathLike,
    *,
    image_format: str = "png",
    dpi: int = 200,
    prefix: str | None = None,
    first_page: int | None = None,
    last_page: int | None = None,
    password: str | None = None,
    poppler_path: PathLike | None = None,
    thread_count: int | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Render PDF pages to a numbered PNG or JPG sequence."""

    normalized_format = image_format.lower().lstrip(".")
    if normalized_format not in {"png", "jpg", "jpeg"}:
        raise ValidationError("PDF 图片输出格式仅支持 PNG 或 JPG")
    suffix = ".jpg" if normalized_format in {"jpg", "jpeg"} else ".png"
    if isinstance(dpi, bool) or int(dpi) < 36 or int(dpi) > 1200:
        raise ValidationError("DPI 必须是 36 到 1200 之间的整数")
    if thread_count is not None and (
        isinstance(thread_count, bool) or int(thread_count) < 1
    ):
        raise ValidationError("thread_count 必须至少为 1")

    source = _input_file(input_pdf, "PDF 输入文件")
    directory = _output_directory(output_dir)
    with _open_reader(source, password) as reader:
        total = len(reader.pages)
    start = 1 if first_page is None else int(first_page)
    end = total if last_page is None else int(last_page)
    if start < 1 or end < start or end > total:
        raise ValidationError(f"渲染页码范围必须在 1 到 {total} 之间")
    workers = (
        optimal_worker_count(end - start + 1, cap=4)
        if thread_count is None
        else int(thread_count)
    )

    base = _safe_stem(prefix or source.stem, "page")
    digits = max(3, len(str(total)))
    targets = [
        directory / f"{base}_{page:0{digits}d}{suffix}"
        for page in range(start, end + 1)
    ]
    _ensure_targets_available(targets, overwrite)

    with tempfile.TemporaryDirectory(
        prefix=".docuforge-render-", dir=directory
    ) as temp_name:
        rendered = _render_pdf_pages(
            source,
            Path(temp_name),
            image_format=normalized_format,
            dpi=int(dpi),
            first_page=start,
            last_page=end,
            password=password,
            poppler_path=poppler_path,
            thread_count=workers,
        )
        if len(rendered) != len(targets):
            raise PDFProcessorError(
                f"渲染页数不一致：预期 {len(targets)} 页，实际生成 {len(rendered)} 页"
            )
        outputs: list[Path] = []
        for rendered_path, target in zip(rendered, targets):
            with _atomic_output(target, overwrite) as temporary:
                shutil.copyfile(rendered_path, temporary)
            outputs.append(target)
    return outputs


def _excel_value(value: Any) -> Any:
    if isinstance(value, str):
        return _ILLEGAL_EXCEL_CHARS.sub("", value)
    return value


def pdf_to_excel(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    password: str | None = None,
    table_settings: Mapping[str, Any] | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Extract detected tables to an XLSX workbook, with one sheet per PDF page."""

    if Workbook is None or get_column_letter is None:
        raise MissingEngineError("PDF 转 Excel 需要安装 openpyxl")
    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    with _open_reader(source, password):
        pass

    workbook = Workbook()
    workbook.remove(workbook.active)
    table_count = 0
    pages_without_tables: list[int] = []
    try:
        with pdfplumber.open(source, password=password) as document:
            for page_number, page in enumerate(document.pages, start=1):
                try:
                    tables = page.extract_tables(
                        table_settings=dict(table_settings or {})
                    )
                except Exception as exc:
                    raise PDFTableExtractionError(
                        f"第 {page_number} 页表格识别失败：{exc}"
                    ) from exc
                usable = [table for table in tables if table and any(table)]
                if not usable:
                    pages_without_tables.append(page_number)
                    continue
                worksheet = workbook.create_sheet(title=f"第{page_number}页")
                current_row = 1
                for table_number, table in enumerate(usable, start=1):
                    if table_number > 1:
                        current_row += 1
                    for row in table:
                        values = list(row or [])
                        for column_number, value in enumerate(values, start=1):
                            worksheet.cell(
                                row=current_row,
                                column=column_number,
                                value=_excel_value(value),
                            )
                        current_row += 1
                    table_count += 1

                for column_number in range(1, worksheet.max_column + 1):
                    maximum = max(
                        (
                            len(
                                str(
                                    worksheet.cell(row=row, column=column_number).value
                                    or ""
                                )
                            )
                            for row in range(1, worksheet.max_row + 1)
                        ),
                        default=0,
                    )
                    worksheet.column_dimensions[
                        get_column_letter(column_number)
                    ].width = min(max(maximum + 2, 8), 60)
    except PDFTableExtractionError:
        workbook.close()
        raise
    except Exception as exc:
        workbook.close()
        raise PDFProcessorError(f"提取 PDF 表格失败：{source.name}（{exc}）") from exc

    if table_count == 0:
        workbook.close()
        raise PDFTableExtractionError(
            "PDF 中没有检测到可提取表格；扫描件需要先执行表格 OCR"
        )
    if pages_without_tables:
        page_list = "、".join(str(value) for value in pages_without_tables[:10])
        suffix = "等" if len(pages_without_tables) > 10 else ""
        warnings.warn(
            f"第 {page_list} 页{suffix}未检测到表格，已跳过",
            PDFProcessingWarning,
            stacklevel=2,
        )
    try:
        with _atomic_output(target, overwrite) as temporary:
            workbook.save(temporary)
    finally:
        workbook.close()
    return [target]


def pdf_to_ppt(
    input_pdf: PathLike,
    output_path: PathLike,
    *,
    dpi: int = 150,
    password: str | None = None,
    poppler_path: PathLike | None = None,
    thread_count: int | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Create a visual-fidelity PPTX by placing each rendered page on one slide."""

    if Presentation is None or Emu is None:
        raise MissingEngineError("PDF 转 PPT 需要安装 python-pptx")
    if Image is None:
        raise MissingEngineError("PDF 转 PPT 需要安装 Pillow")
    if isinstance(dpi, bool) or int(dpi) < 36 or int(dpi) > 1200:
        raise ValidationError("DPI 必须是 36 到 1200 之间的整数")
    if thread_count is not None and (
        isinstance(thread_count, bool) or int(thread_count) < 1
    ):
        raise ValidationError("thread_count 必须至少为 1")

    source = _input_file(input_pdf, "PDF 输入文件")
    target = _output_file(output_path)
    _protect_inputs(target, [source])
    _ensure_targets_available([target], overwrite)
    target.parent.mkdir(parents=True, exist_ok=True)

    with _open_reader(source, password) as reader:
        total = len(reader.pages)
        if total < 1:
            raise ValidationError("PDF 不包含任何页面")
        first = reader.pages[0]
        width_points = float(first.mediabox.width)
        height_points = float(first.mediabox.height)
        if int(first.rotation or 0) % 180:
            width_points, height_points = height_points, width_points
    workers = (
        optimal_worker_count(total, cap=4)
        if thread_count is None
        else int(thread_count)
    )

    with tempfile.TemporaryDirectory(
        prefix=".docuforge-ppt-", dir=target.parent
    ) as temp_name:
        rendered = _render_pdf_pages(
            source,
            Path(temp_name),
            image_format="png",
            dpi=int(dpi),
            first_page=1,
            last_page=total,
            password=password,
            poppler_path=poppler_path,
            thread_count=workers,
        )
        if len(rendered) != total:
            raise PDFProcessorError(
                f"渲染页数不一致：预期 {total} 页，实际生成 {len(rendered)} 页"
            )

        presentation = Presentation()
        presentation.slide_width = Emu(round(width_points * 12700))
        presentation.slide_height = Emu(round(height_points * 12700))
        blank_layout = (
            presentation.slide_layouts[6]
            if len(presentation.slide_layouts) > 6
            else presentation.slide_layouts[-1]
        )
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)

        for image_path in rendered:
            with Image.open(image_path) as rendered_image:
                image_width, image_height = rendered_image.size
            scale = min(slide_width / image_width, slide_height / image_height)
            draw_width = round(image_width * scale)
            draw_height = round(image_height * scale)
            left = round((slide_width - draw_width) / 2)
            top = round((slide_height - draw_height) / 2)
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path), left, top, width=draw_width, height=draw_height
            )

        with _atomic_output(target, overwrite) as temporary:
            presentation.save(temporary)
    return [target]


def images_to_pdf(
    image_paths: Sequence[PathLike],
    output_path: PathLike,
    *,
    dpi: float = 96.0,
    page_size: tuple[float, float] | None = None,
    background: tuple[int, int, int] = (255, 255, 255),
    overwrite: bool = False,
) -> list[Path]:
    """Combine images into a PDF, preserving aspect ratio and input order.

    ``page_size`` is expressed in PDF points.  When omitted, each PDF page is
    sized to the image at the requested DPI.
    """

    if Image is None or ImageOps is None:
        raise MissingEngineError("图片转 PDF 需要安装 Pillow")
    if float(dpi) <= 0:
        raise ValidationError("DPI 必须大于 0")
    if page_size is not None and (
        len(page_size) != 2 or float(page_size[0]) <= 0 or float(page_size[1]) <= 0
    ):
        raise ValidationError("page_size 必须是两个大于 0 的 PDF 点数")
    if len(background) != 3 or any(value < 0 or value > 255 for value in background):
        raise ValidationError("background 必须是 0 到 255 的 RGB 三元组")

    sources = [_input_file(path, "输入图片") for path in image_paths]
    if not sources:
        raise ValidationError("至少需要选择一张图片")
    target = _output_file(output_path)
    _protect_inputs(target, sources)

    image_buffers: list[io.BytesIO] = []
    try:
        with _atomic_output(target, overwrite) as temporary:
            pdf_canvas = canvas.Canvas(str(temporary), pageCompression=1)
            for source in sources:
                try:
                    with Image.open(source) as opened:
                        if getattr(opened, "n_frames", 1) > 1:
                            warnings.warn(
                                f"{source.name} 包含多个动画/图像帧，仅使用第一帧",
                                PDFProcessingWarning,
                                stacklevel=2,
                            )
                        opened.seek(0)
                        normalized = ImageOps.exif_transpose(opened)
                        normalized.load()
                        if normalized.mode in {"RGBA", "LA"} or (
                            normalized.mode == "P" and "transparency" in normalized.info
                        ):
                            rgba = normalized.convert("RGBA")
                            flattened = Image.new("RGB", rgba.size, tuple(background))
                            flattened.paste(rgba, mask=rgba.getchannel("A"))
                            normalized = flattened
                        elif normalized.mode != "RGB":
                            normalized = normalized.convert("RGB")
                        pixel_width, pixel_height = normalized.size
                        buffer = io.BytesIO()
                        normalized.save(buffer, format="PNG", optimize=True)
                        buffer.seek(0)
                        image_buffers.append(buffer)
                except Exception as exc:
                    raise PDFProcessorError(
                        f"无法读取图片：{source.name}（{exc}）"
                    ) from exc

                natural_width = pixel_width * 72.0 / float(dpi)
                natural_height = pixel_height * 72.0 / float(dpi)
                if page_size is None:
                    page_width, page_height = natural_width, natural_height
                    draw_width, draw_height = page_width, page_height
                    left = bottom = 0.0
                else:
                    page_width, page_height = float(page_size[0]), float(page_size[1])
                    factor = min(
                        page_width / natural_width, page_height / natural_height
                    )
                    draw_width = natural_width * factor
                    draw_height = natural_height * factor
                    left = (page_width - draw_width) / 2
                    bottom = (page_height - draw_height) / 2
                pdf_canvas.setPageSize((page_width, page_height))
                pdf_canvas.drawImage(
                    ImageReader(image_buffers[-1]),
                    left,
                    bottom,
                    width=draw_width,
                    height=draw_height,
                    preserveAspectRatio=True,
                    mask="auto",
                )
                pdf_canvas.showPage()
            pdf_canvas.save()
    finally:
        for buffer in image_buffers:
            buffer.close()
    return [target]
