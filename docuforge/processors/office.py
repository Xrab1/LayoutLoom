"""Reliable, offline Office document processors.

The functions in this module deliberately separate *structural* conversions from
high-fidelity rendering.  ``python-docx``, ``openpyxl`` and ``python-pptx`` are
excellent at editing OOXML packages, but they are not layout/rendering engines.
When visual fidelity matters, use :func:`convert_with_office`, which only runs
when WPS Office automation, Microsoft Office automation, or LibreOffice is
actually available.

Every processing function accepts ``pathlib``-compatible paths, writes through
a same-directory temporary file, avoids overwriting by default, and returns a
list of generated :class:`~pathlib.Path` objects.
"""

from __future__ import annotations

import base64
import csv
import datetime as _datetime
import html
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import uuid
import warnings
import zipfile
from contextlib import contextmanager
from copy import copy, deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Iterable, Mapping, Sequence
from xml.etree import ElementTree as ET

from docuforge.models import MissingEngineError, ValidationError
from docuforge.utils import atomic_output, ensure_output_dir, safe_filename, unique_path

PathLike = str | os.PathLike[str]

_WORD_EXTENSIONS = {".docx", ".docm", ".dotx", ".dotm"}
_EXCEL_EXTENSIONS = {".xlsx", ".xlsm", ".xltx", ".xltm"}
_PPT_EXTENSIONS = {".pptx", ".pptm", ".potx", ".potm", ".ppsx", ".ppsm"}

_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_V = "urn:schemas-microsoft-com:vml"
_PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_CONTENT_TYPES = "http://schemas.openxmlformats.org/package/2006/content-types"

for _prefix, _uri in {
    "w": _W,
    "r": _R,
    "a": _A,
    "v": _V,
}.items():
    ET.register_namespace(_prefix, _uri)


def _require_source(source: PathLike, extensions: set[str], family: str) -> Path:
    from docuforge.runner import check_cancelled

    check_cancelled("任务已取消；已完成的文件会保留")
    path = Path(source).expanduser().resolve()
    if not path.is_file():
        raise ValidationError(f"文件不存在：{path}")
    if path.suffix.lower() not in extensions:
        legacy = {"Word": ".doc", "Excel": ".xls", "PowerPoint": ".ppt"}
        if path.suffix.lower() == legacy.get(family):
            raise ValidationError(
                f"旧版 {path.suffix} 需要 WPS、Microsoft Office 或 LibreOffice 转换，"
                "不能由 OOXML 离线编辑器直接处理"
            )
        raise ValidationError(
            f"不支持的 {family} 文件格式：{path.suffix or '无扩展名'}"
        )
    return path


def _output_path(
    source: Path,
    output_dir: PathLike,
    *,
    suffix: str | None = None,
    tag: str | None = None,
    overwrite: bool = False,
) -> Path:
    directory = ensure_output_dir(output_dir)
    extension = suffix if suffix is not None else source.suffix
    if extension and not extension.startswith("."):
        extension = f".{extension}"
    stem = source.stem if not tag else f"{source.stem}_{safe_filename(tag)}"
    target = directory / f"{stem}{extension}"
    # Never replace the input in place, even when overwrite=True.
    if target.resolve() == source.resolve():
        target = directory / f"{source.stem}_output{extension}"
        overwrite = False
    return unique_path(target, overwrite=overwrite)


def _save_document(document: Any, target: Path) -> list[Path]:
    with atomic_output(target) as temporary:
        document.save(temporary)
    return [target]


def _load_word(source: PathLike) -> tuple[Path, Any]:
    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - depends on installation
        raise MissingEngineError("缺少 python-docx，无法处理 Word OOXML 文件") from exc
    try:
        return source_path, Document(source_path)
    except Exception as exc:
        raise ValidationError(
            f"无法读取 Word 文件：{source_path.name}（{exc}）"
        ) from exc


def _load_workbook(source: PathLike, *, data_only: bool = False) -> tuple[Path, Any]:
    source_path = _require_source(source, _EXCEL_EXTENSIONS, "Excel")
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover - depends on installation
        raise MissingEngineError("缺少 openpyxl，无法处理 Excel OOXML 文件") from exc
    try:
        workbook = openpyxl.load_workbook(
            source_path,
            data_only=data_only,
            keep_vba=source_path.suffix.lower() in {".xlsm", ".xltm"},
        )
        return source_path, workbook
    except Exception as exc:
        raise ValidationError(
            f"无法读取 Excel 文件：{source_path.name}（{exc}）"
        ) from exc


def _load_presentation(source: PathLike) -> tuple[Path, Any]:
    source_path = _require_source(source, _PPT_EXTENSIONS, "PowerPoint")
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - depends on installation
        raise MissingEngineError(
            "缺少 python-pptx，无法处理 PowerPoint OOXML 文件"
        ) from exc
    try:
        return source_path, Presentation(source_path)
    except Exception as exc:
        raise ValidationError(
            f"无法读取 PowerPoint 文件：{source_path.name}（{exc}）"
        ) from exc


def _iter_table_paragraphs(table: Any, seen: set[int]) -> Iterable[Any]:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                marker = id(paragraph._p)
                if marker not in seen:
                    seen.add(marker)
                    yield paragraph
            for nested in cell.tables:
                yield from _iter_table_paragraphs(nested, seen)


def _iter_word_paragraphs(
    document: Any, *, include_headers: bool = True
) -> Iterable[Any]:
    """Yield body, table, header and footer paragraphs without merged-cell duplicates."""

    seen: set[int] = set()
    for paragraph in document.paragraphs:
        marker = id(paragraph._p)
        if marker not in seen:
            seen.add(marker)
            yield paragraph
    for table in document.tables:
        yield from _iter_table_paragraphs(table, seen)
    if include_headers:
        for section in document.sections:
            for part in (
                section.header,
                section.first_page_header,
                section.even_page_header,
                section.footer,
                section.first_page_footer,
                section.even_page_footer,
            ):
                for paragraph in part.paragraphs:
                    marker = id(paragraph._p)
                    if marker not in seen:
                        seen.add(marker)
                        yield paragraph
                for table in part.tables:
                    yield from _iter_table_paragraphs(table, seen)


def _paragraph_runs(paragraph: Any) -> list[Any]:
    """Return direct and hyperlink-contained runs in visual document order."""

    iterator = getattr(paragraph, "iter_inner_content", None)
    if callable(iterator):
        runs: list[Any] = []
        for item in iterator():
            if hasattr(item, "_r"):
                runs.append(item)
            else:
                runs.extend(getattr(item, "runs", ()))
        return runs
    # python-docx 1.1 fallback.
    from docx.text.run import Run

    runs = []
    for child in paragraph._p.iterchildren():
        if child.tag == f"{{{_W}}}r":
            runs.append(Run(child, paragraph))
        elif child.tag == f"{{{_W}}}hyperlink":
            runs.extend(
                Run(node, paragraph) for node in child if node.tag == f"{{{_W}}}r"
            )
    return runs


def _paragraph_has_nontext_content(paragraph: Any) -> bool:
    tags = {
        f"{{{_W}}}drawing",
        f"{{{_W}}}pict",
        f"{{{_W}}}object",
        f"{{{_W}}}br",
        f"{{{_W}}}lastRenderedPageBreak",
        f"{{{_W}}}fldChar",
        f"{{{_W}}}instrText",
        f"{{{_W}}}footnoteReference",
        f"{{{_W}}}endnoteReference",
        f"{{{_W}}}commentReference",
        f"{{{_W}}}bookmarkStart",
        f"{{{_W}}}bookmarkEnd",
    }
    return any(element.tag in tags for element in paragraph._p.iter())


def _replace_in_paragraph(
    paragraph: Any,
    replacements: Mapping[str, Any],
    *,
    case_sensitive: bool,
) -> int:
    """Replace across run boundaries while retaining unaffected run formatting."""

    changed = 0
    if not _paragraph_runs(paragraph):
        return changed
    for needle, raw_replacement in replacements.items():
        needle = str(needle)
        replacement = str(raw_replacement)
        if not needle:
            continue
        original_text = "".join(run.text or "" for run in _paragraph_runs(paragraph))
        matches = list(
            re.finditer(
                re.escape(needle),
                original_text,
                flags=0 if case_sensitive else re.IGNORECASE,
            )
        )
        # Work backwards so earlier offsets stay valid.  Looking up matches once
        # also prevents infinite replacement for mappings such as ``a -> aa``.
        for match in reversed(matches):
            position, end_position = match.span()
            runs = _paragraph_runs(paragraph)
            texts = [run.text or "" for run in runs]
            offsets: list[tuple[int, int]] = []
            cursor = 0
            for text in texts:
                offsets.append((cursor, cursor + len(text)))
                cursor += len(text)
            start_index = next(
                (index for index, (_, end) in enumerate(offsets) if position < end),
                len(runs) - 1,
            )
            end_index = next(
                (
                    index
                    for index, (_, end) in enumerate(offsets)
                    if end_position <= end
                ),
                len(runs) - 1,
            )
            start_offset = position - offsets[start_index][0]
            end_offset = end_position - offsets[end_index][0]
            prefix = runs[start_index].text[:start_offset]
            suffix = runs[end_index].text[end_offset:]
            if start_index == end_index:
                runs[start_index].text = prefix + replacement + suffix
            else:
                runs[start_index].text = prefix + replacement
                for index in range(start_index + 1, end_index):
                    runs[index].text = ""
                runs[end_index].text = suffix
            changed += 1
    return changed


def _paragraph_markdown(paragraph: Any) -> str:
    chunks: list[str] = []
    for run in _paragraph_runs(paragraph):
        value = run.text.replace("\\", "\\\\").replace("*", "\\*").replace("_", "\\_")
        if not value:
            continue
        if run.bold and run.italic:
            value = f"***{value}***"
        elif run.bold:
            value = f"**{value}**"
        elif run.italic:
            value = f"*{value}*"
        if run.font.strike:
            value = f"~~{value}~~"
        chunks.append(value)
    text = "".join(chunks) if chunks else paragraph.text
    style_name = (getattr(paragraph.style, "name", "") or "").lower()
    heading = re.match(r"heading\s*([1-6])", style_name)
    if heading:
        return f"{'#' * int(heading.group(1))} {text}".rstrip()
    if "list bullet" in style_name:
        return f"- {text}".rstrip()
    if "list number" in style_name:
        return f"1. {text}".rstrip()
    return text


def _table_markdown(table: Any) -> str:
    rows = [
        [cell.text.replace("|", "\\|").replace("\n", "<br>") for cell in row.cells]
        for row in table.rows
    ]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    lines = [
        "| " + " | ".join(normalized[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in normalized[1:])
    return "\n".join(lines)


def _image_from_run(run: Any) -> list[tuple[str, bytes]]:
    images: list[tuple[str, bytes]] = []
    for blip in run._r.iter(f"{{{_A}}}blip"):
        relationship_id = blip.attrib.get(f"{{{_R}}}embed")
        if not relationship_id:
            continue
        try:
            part = run.part.related_parts[relationship_id]
            images.append((getattr(part, "content_type", "image/png"), part.blob))
        except (KeyError, AttributeError):
            continue
    return images


def _paragraph_html(paragraph: Any, image_resolver: Callable[[str, bytes], str]) -> str:
    chunks: list[str] = []
    for run in _paragraph_runs(paragraph):
        value = html.escape(run.text).replace("\n", "<br>")
        if value:
            if run.bold:
                value = f"<strong>{value}</strong>"
            if run.italic:
                value = f"<em>{value}</em>"
            if run.underline:
                value = f"<u>{value}</u>"
            if run.font.strike:
                value = f"<s>{value}</s>"
            chunks.append(value)
        for content_type, blob in _image_from_run(run):
            source = html.escape(image_resolver(content_type, blob), quote=True)
            chunks.append(f'<img src="{source}" alt="embedded image">')
    content = "".join(chunks) or html.escape(paragraph.text)
    style_name = (getattr(paragraph.style, "name", "") or "").lower()
    heading = re.match(r"heading\s*([1-6])", style_name)
    if heading:
        level = heading.group(1)
        return f"<h{level}>{content}</h{level}>"
    if "list bullet" in style_name or "list number" in style_name:
        return f"<li>{content}</li>"
    return f"<p>{content}</p>"


def _table_html(table: Any, image_resolver: Callable[[str, bytes], str]) -> str:
    rows: list[str] = []
    for row_index, row in enumerate(table.rows):
        cell_tag = "th" if row_index == 0 else "td"
        cells = []
        for cell in row.cells:
            body = "".join(
                _paragraph_html(paragraph, image_resolver)
                for paragraph in cell.paragraphs
            )
            cells.append(f"<{cell_tag}>{body}</{cell_tag}>")
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return f"<table>{''.join(rows)}</table>"


def _document_html_fragment(
    document: Any, image_resolver: Callable[[str, bytes], str]
) -> str:
    # document.element.body preserves paragraph/table order.
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    fragments: list[str] = []
    body = document.element.body
    for child in body.iterchildren():
        if child.tag == f"{{{_W}}}p":
            fragments.append(
                _paragraph_html(Paragraph(child, document), image_resolver)
            )
        elif child.tag == f"{{{_W}}}tbl":
            fragments.append(_table_html(Table(child, document), image_resolver))
    return "\n".join(fragments)


def word_to_txt(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
    include_tables: bool = True,
    encoding: str = "utf-8",
) -> list[Path]:
    """Extract Word's semantic text.  This is not a visual layout conversion."""

    source_path, document = _load_word(source)
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    lines: list[str] = []
    for child in document.element.body.iterchildren():
        if child.tag == f"{{{_W}}}p":
            lines.append(Paragraph(child, document).text)
        elif include_tables and child.tag == f"{{{_W}}}tbl":
            table = Table(child, document)
            lines.extend(
                "\t".join(cell.text for cell in row.cells) for row in table.rows
            )
    target = _output_path(source_path, output_dir, suffix=".txt", overwrite=overwrite)
    with atomic_output(target) as temporary:
        temporary.write_text("\n".join(lines), encoding=encoding)
    return [target]


def word_to_markdown(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
) -> list[Path]:
    """Convert Word paragraphs and tables to semantic Markdown."""

    source_path, document = _load_word(source)
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if child.tag == f"{{{_W}}}p":
            blocks.append(_paragraph_markdown(Paragraph(child, document)))
        elif child.tag == f"{{{_W}}}tbl":
            blocks.append(_table_markdown(Table(child, document)))
    target = _output_path(source_path, output_dir, suffix=".md", overwrite=overwrite)
    with atomic_output(target) as temporary:
        temporary.write_text("\n\n".join(blocks).rstrip() + "\n", encoding="utf-8")
    return [target]


def word_to_md(
    source: PathLike, output_dir: PathLike, *, overwrite: bool = False
) -> list[Path]:
    """Alias for :func:`word_to_markdown`."""

    return word_to_markdown(source, output_dir, overwrite=overwrite)


def word_to_html(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
    title: str | None = None,
) -> list[Path]:
    """Create a standalone semantic HTML document with embedded data-URI images."""

    source_path, document = _load_word(source)

    def data_uri(content_type: str, blob: bytes) -> str:
        return f"data:{content_type};base64,{base64.b64encode(blob).decode('ascii')}"

    fragment = _document_html_fragment(document, data_uri)
    page_title = title or source_path.stem
    payload = (
        '<!doctype html>\n<html lang="zh-CN"><head><meta charset="utf-8">'
        f"<title>{html.escape(page_title)}</title>"
        "<style>body{max-width:960px;margin:2rem auto;font-family:sans-serif;line-height:1.6}"
        "table{border-collapse:collapse}th,td{border:1px solid #999;padding:.35rem}"
        "img{max-width:100%;height:auto}</style></head><body>"
        f"{fragment}</body></html>"
    )
    target = _output_path(source_path, output_dir, suffix=".html", overwrite=overwrite)
    with atomic_output(target) as temporary:
        temporary.write_text(payload, encoding="utf-8")
    return [target]


def word_to_epub(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
    title: str | None = None,
    language: str = "zh-CN",
    author: str = "",
) -> list[Path]:
    """Create a valid EPUB 3 containing semantic text, tables and embedded images."""

    source_path, document = _load_word(source)
    image_entries: list[tuple[str, str, bytes]] = []
    image_by_blob: dict[bytes, str] = {}

    def epub_image(content_type: str, blob: bytes) -> str:
        if blob in image_by_blob:
            return image_by_blob[blob]
        extension = {
            "image/jpeg": ".jpg",
            "image/png": ".png",
            "image/gif": ".gif",
            "image/svg+xml": ".svg",
            "image/tiff": ".tiff",
            "image/webp": ".webp",
        }.get(content_type, ".bin")
        name = f"images/image_{len(image_entries) + 1}{extension}"
        image_by_blob[blob] = name
        image_entries.append((name, content_type, blob))
        return name

    fragment = _document_html_fragment(document, epub_image)
    book_title = title or source_path.stem
    identifier = f"urn:uuid:{uuid.uuid4()}"
    manifest_images = "".join(
        f'<item id="img{index}" href="{html.escape(name, quote=True)}" media-type="{media}"/>'
        for index, (name, media, _blob) in enumerate(image_entries, 1)
    )
    content_xhtml = (
        '<?xml version="1.0" encoding="utf-8"?>'
        f'<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="{html.escape(language)}">'
        f'<head><title>{html.escape(book_title)}</title><link rel="stylesheet" href="style.css"/>'
        f"</head><body>{fragment}</body></html>"
    ).encode("utf-8")
    nav_xhtml = (
        '<?xml version="1.0" encoding="utf-8"?>'
        '<html xmlns="http://www.w3.org/1999/xhtml" '
        'xmlns:epub="http://www.idpf.org/2007/ops"><head><title>目录</title></head>'
        f'<body><nav epub:type="toc"><ol><li><a href="content.xhtml">'
        f"{html.escape(book_title)}</a></li></ol></nav></body></html>"
    ).encode("utf-8")
    package_opf = (
        '<?xml version="1.0" encoding="utf-8"?>'
        '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="book-id">'
        f'<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
        f'<dc:identifier id="book-id">{identifier}</dc:identifier>'
        f"<dc:title>{html.escape(book_title)}</dc:title><dc:language>{html.escape(language)}</dc:language>"
        f"<dc:creator>{html.escape(author)}</dc:creator>"
        f'<meta property="dcterms:modified">{_datetime.datetime.now(_datetime.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")}</meta>'
        '</metadata><manifest><item id="content" href="content.xhtml" media-type="application/xhtml+xml"/>'
        '<item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>'
        '<item id="css" href="style.css" media-type="text/css"/>'
        f'{manifest_images}</manifest><spine><itemref idref="content"/></spine></package>'
    ).encode("utf-8")
    container_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
        '<rootfiles><rootfile full-path="OEBPS/package.opf" '
        'media-type="application/oebps-package+xml"/></rootfiles></container>'
    ).encode("utf-8")
    css = b"body{font-family:serif;line-height:1.5}table{border-collapse:collapse}th,td{border:1px solid #777;padding:.3em}img{max-width:100%;height:auto}"
    target = _output_path(source_path, output_dir, suffix=".epub", overwrite=overwrite)
    with atomic_output(target) as temporary:
        with zipfile.ZipFile(temporary, "w") as archive:
            archive.writestr(
                "mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED
            )
            archive.writestr(
                "META-INF/container.xml",
                container_xml,
                compress_type=zipfile.ZIP_DEFLATED,
            )
            archive.writestr(
                "OEBPS/content.xhtml", content_xhtml, compress_type=zipfile.ZIP_DEFLATED
            )
            archive.writestr(
                "OEBPS/nav.xhtml", nav_xhtml, compress_type=zipfile.ZIP_DEFLATED
            )
            archive.writestr(
                "OEBPS/package.opf", package_opf, compress_type=zipfile.ZIP_DEFLATED
            )
            archive.writestr("OEBPS/style.css", css, compress_type=zipfile.ZIP_DEFLATED)
            for name, _media, blob in image_entries:
                archive.writestr(
                    f"OEBPS/{name}", blob, compress_type=zipfile.ZIP_DEFLATED
                )
    return [target]


def word_replace_text(
    source: PathLike,
    output_dir: PathLike,
    replacements: Mapping[str, Any],
    *,
    case_sensitive: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, document = _load_word(source)
    if not replacements:
        raise ValidationError("替换映射不能为空")
    for paragraph in _iter_word_paragraphs(document):
        _replace_in_paragraph(paragraph, replacements, case_sensitive=case_sensitive)
    target = _output_path(source_path, output_dir, tag="替换", overwrite=overwrite)
    return _save_document(document, target)


def word_remove_blank_paragraphs(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
) -> list[Path]:
    source_path, document = _load_word(source)
    for paragraph in list(_iter_word_paragraphs(document)):
        if paragraph.text.strip() or _paragraph_has_nontext_content(paragraph):
            continue
        p_pr = paragraph._p.pPr
        if p_pr is not None and getattr(p_pr, "sectPr", None) is not None:
            continue
        parent = paragraph._p.getparent()
        if parent is not None:
            parent.remove(paragraph._p)
    target = _output_path(source_path, output_dir, tag="无空行", overwrite=overwrite)
    return _save_document(document, target)


def _rewrite_zip(
    source: Path,
    target: Path,
    transform: Callable[[str, bytes], bytes | None],
) -> list[Path]:
    try:
        with zipfile.ZipFile(source, "r") as reader, atomic_output(target) as temporary:
            with zipfile.ZipFile(temporary, "w") as writer:
                for info in reader.infolist():
                    payload = transform(info.filename, reader.read(info))
                    if payload is None:
                        continue
                    writer.writestr(info, payload)
    except zipfile.BadZipFile as exc:
        raise ValidationError(f"不是有效的 OOXML 文件：{source.name}") from exc
    return [target]


def _remove_elements(root: ET.Element, predicate: Callable[[ET.Element], bool]) -> None:
    changed = True
    while changed:
        changed = False
        for parent in root.iter():
            for child in list(parent):
                if predicate(child):
                    parent.remove(child)
                    changed = True


def word_remove_images(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
) -> list[Path]:
    """Remove raster/vector pictures while leaving charts and other drawings intact."""

    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    target = _output_path(source_path, output_dir, tag="无图片", overwrite=overwrite)

    def transform(name: str, payload: bytes) -> bytes | None:
        if name.startswith("word/media/"):
            return None
        if name.startswith("word/") and name.endswith(".xml"):
            root = ET.fromstring(payload)

            def is_picture(element: ET.Element) -> bool:
                if element.tag == f"{{{_W}}}drawing":
                    return any(node.tag == f"{{{_A}}}blip" for node in element.iter())
                if element.tag == f"{{{_W}}}pict":
                    return any(
                        node.tag == f"{{{_V}}}imagedata" for node in element.iter()
                    )
                return False

            _remove_elements(root, is_picture)
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
        if name.startswith("word/") and name.endswith(".rels"):
            root = ET.fromstring(payload)
            for relationship in list(root):
                if relationship.attrib.get("Type", "").endswith("/image"):
                    root.remove(relationship)
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
        return payload

    return _rewrite_zip(source_path, target, transform)


_REVISION_WRAPPERS_ACCEPT = {f"{{{_W}}}ins", f"{{{_W}}}moveTo"}
_REVISION_WRAPPERS_REJECT = {f"{{{_W}}}del", f"{{{_W}}}moveFrom"}
_REVISION_PROPERTY_CHANGES = {
    f"{{{_W}}}pPrChange",
    f"{{{_W}}}rPrChange",
    f"{{{_W}}}tblPrChange",
    f"{{{_W}}}trPrChange",
    f"{{{_W}}}tcPrChange",
    f"{{{_W}}}sectPrChange",
    f"{{{_W}}}tblGridChange",
}
_REVISION_MARKERS = {
    f"{{{_W}}}moveFromRangeStart",
    f"{{{_W}}}moveFromRangeEnd",
    f"{{{_W}}}moveToRangeStart",
    f"{{{_W}}}moveToRangeEnd",
    f"{{{_W}}}customXmlInsRangeStart",
    f"{{{_W}}}customXmlInsRangeEnd",
    f"{{{_W}}}customXmlDelRangeStart",
    f"{{{_W}}}customXmlDelRangeEnd",
    f"{{{_W}}}customXmlMoveFromRangeStart",
    f"{{{_W}}}customXmlMoveFromRangeEnd",
    f"{{{_W}}}customXmlMoveToRangeStart",
    f"{{{_W}}}customXmlMoveToRangeEnd",
}
_UNSUPPORTED_REVISIONS = {
    f"{{{_W}}}cellIns",
    f"{{{_W}}}cellDel",
    f"{{{_W}}}cellMerge",
    f"{{{_W}}}numberingChange",
}


def _assert_supported_revisions(root: ET.Element) -> None:
    unsupported = {element.tag for element in root.iter()} & _UNSUPPORTED_REVISIONS
    if unsupported:
        names = ", ".join(sorted(tag.rsplit("}", 1)[-1] for tag in unsupported))
        raise ValidationError(
            f"文档含无法安全自动处理的表格/编号修订（{names}）；请用 Microsoft Word 审阅"
        )


def _unwrap(parent: ET.Element, child: ET.Element) -> None:
    index = list(parent).index(child)
    parent.remove(child)
    for nested in list(child):
        parent.insert(index, nested)
        index += 1


def _convert_deleted_text(element: ET.Element) -> None:
    for node in element.iter():
        if node.tag == f"{{{_W}}}delText":
            node.tag = f"{{{_W}}}t"
        elif node.tag == f"{{{_W}}}delInstrText":
            node.tag = f"{{{_W}}}instrText"


def _apply_revisions(root: ET.Element, mode: str) -> None:
    if mode not in {"accept", "reject"}:
        raise ValidationError("revision_mode 必须是 accept、reject 或 keep")
    _assert_supported_revisions(root)

    if mode == "reject":
        # Property change records contain the previous property set.
        for parent in list(root.iter()):
            for change in list(parent):
                if change.tag not in _REVISION_PROPERTY_CHANGES:
                    continue
                snapshot = next(
                    (item for item in list(change) if item.tag == parent.tag), None
                )
                if snapshot is None:
                    raise ValidationError(
                        f"修订记录 {change.tag.rsplit('}', 1)[-1]} 缺少原始属性，无法安全拒绝"
                    )
                parent.attrib.clear()
                parent.attrib.update(snapshot.attrib)
                for existing in list(parent):
                    parent.remove(existing)
                for previous in list(snapshot):
                    parent.append(deepcopy(previous))
    else:
        _remove_elements(
            root, lambda element: element.tag in _REVISION_PROPERTY_CHANGES
        )

    changed = True
    while changed:
        changed = False
        for parent in root.iter():
            for child in list(parent):
                if child.tag in _REVISION_MARKERS:
                    parent.remove(child)
                    changed = True
                elif mode == "accept" and child.tag in _REVISION_WRAPPERS_ACCEPT:
                    _unwrap(parent, child)
                    changed = True
                elif mode == "accept" and child.tag in _REVISION_WRAPPERS_REJECT:
                    parent.remove(child)
                    changed = True
                elif mode == "reject" and child.tag in _REVISION_WRAPPERS_ACCEPT:
                    parent.remove(child)
                    changed = True
                elif mode == "reject" and child.tag in _REVISION_WRAPPERS_REJECT:
                    _convert_deleted_text(child)
                    _unwrap(parent, child)
                    changed = True


def _revision_transform(
    name: str, payload: bytes, mode: str, remove_comments: bool
) -> bytes | None:
    if remove_comments and name.startswith("word/comments") and name.endswith(".xml"):
        return None
    if name == "[Content_Types].xml" and remove_comments:
        root = ET.fromstring(payload)
        for item in list(root):
            if "/word/comments" in item.attrib.get("PartName", ""):
                root.remove(item)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    if name.startswith("word/") and name.endswith(".rels") and remove_comments:
        root = ET.fromstring(payload)
        for relationship in list(root):
            rel_type = relationship.attrib.get("Type", "")
            target = relationship.attrib.get("Target", "")
            if "/comments" in rel_type or target.startswith("comments"):
                root.remove(relationship)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    if not (name.startswith("word/") and name.endswith(".xml")):
        return payload
    root = ET.fromstring(payload)
    if remove_comments:
        comment_tags = {
            f"{{{_W}}}commentRangeStart",
            f"{{{_W}}}commentRangeEnd",
            f"{{{_W}}}commentReference",
        }
        _remove_elements(root, lambda element: element.tag in comment_tags)
    if mode in {"accept", "reject"}:
        _apply_revisions(root, mode)
    if name == "word/settings.xml" and mode in {"accept", "reject"}:
        _remove_elements(root, lambda element: element.tag == f"{{{_W}}}trackRevisions")
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def word_accept_revisions(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
) -> list[Path]:
    """Accept common run/paragraph/property revisions directly in OOXML.

    Complex table-cell and numbering revisions are refused instead of guessed.
    """

    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    target = _output_path(source_path, output_dir, tag="接受修订", overwrite=overwrite)
    return _rewrite_zip(
        source_path,
        target,
        lambda name, payload: _revision_transform(name, payload, "accept", False),
    )


def word_reject_revisions(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
) -> list[Path]:
    """Reject common run/paragraph/property revisions directly in OOXML."""

    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    target = _output_path(source_path, output_dir, tag="拒绝修订", overwrite=overwrite)
    return _rewrite_zip(
        source_path,
        target,
        lambda name, payload: _revision_transform(name, payload, "reject", False),
    )


def word_clean_comments_and_revisions(
    source: PathLike,
    output_dir: PathLike,
    *,
    revision_mode: str = "accept",
    overwrite: bool = False,
) -> list[Path]:
    """Remove comments and either accept, reject, or keep tracked revisions."""

    mode = revision_mode.lower().strip()
    if mode not in {"accept", "reject", "keep"}:
        raise ValidationError("revision_mode 必须是 accept、reject 或 keep")
    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    target = _output_path(
        source_path, output_dir, tag="清理批注修订", overwrite=overwrite
    )
    return _rewrite_zip(
        source_path,
        target,
        lambda name, payload: _revision_transform(name, payload, mode, True),
    )


def word_set_typography(
    source: PathLike,
    output_dir: PathLike,
    *,
    font_name: str | None = None,
    font_size_pt: float | None = None,
    line_spacing: float | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path, document = _load_word(source)
    if font_size_pt is not None and font_size_pt <= 0:
        raise ValidationError("字号必须大于 0")
    if line_spacing is not None and line_spacing <= 0:
        raise ValidationError("行距必须大于 0")
    try:
        from docx.oxml.ns import qn
        from docx.shared import Pt
    except ImportError as exc:  # pragma: no cover
        raise MissingEngineError("缺少 python-docx") from exc
    for paragraph in _iter_word_paragraphs(document):
        if line_spacing is not None:
            paragraph.paragraph_format.line_spacing = line_spacing
        for run in _paragraph_runs(paragraph):
            if font_name:
                run.font.name = font_name
                run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), font_name)
            if font_size_pt is not None:
                run.font.size = Pt(font_size_pt)
    normal = document.styles["Normal"] if "Normal" in document.styles else None
    if normal is not None:
        if font_name:
            normal.font.name = font_name
            normal._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), font_name)
        if font_size_pt is not None:
            normal.font.size = Pt(font_size_pt)
    target = _output_path(source_path, output_dir, tag="统一格式", overwrite=overwrite)
    return _save_document(document, target)


def _set_story_text(story: Any, text: str, replace: bool) -> None:
    if not story.paragraphs:
        story.add_paragraph(text)
        return
    if replace:
        first = story.paragraphs[0]
        first.text = text
        for paragraph in list(story.paragraphs[1:]):
            parent = paragraph._p.getparent()
            if parent is not None:
                parent.remove(paragraph._p)
    else:
        story.add_paragraph(text)


def word_set_headers_footers(
    source: PathLike,
    output_dir: PathLike,
    *,
    header_text: str | None = None,
    footer_text: str | None = None,
    replace: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, document = _load_word(source)
    if header_text is None and footer_text is None:
        raise ValidationError("header_text 和 footer_text 至少填写一项")
    seen: set[int] = set()
    for section in document.sections:
        if header_text is not None:
            for story in (
                section.header,
                section.first_page_header,
                section.even_page_header,
            ):
                if id(story._element) not in seen:
                    seen.add(id(story._element))
                    _set_story_text(story, str(header_text), replace)
        if footer_text is not None:
            for story in (
                section.footer,
                section.first_page_footer,
                section.even_page_footer,
            ):
                if id(story._element) not in seen:
                    seen.add(id(story._element))
                    _set_story_text(story, str(footer_text), replace)
    target = _output_path(source_path, output_dir, tag="页眉页脚", overwrite=overwrite)
    return _save_document(document, target)


def _extract_zip_members(
    source: Path,
    output_dir: PathLike,
    prefix_path: str,
    *,
    prefix: str | None,
    overwrite: bool,
    member_filter: Callable[[str], bool] | None = None,
) -> list[Path]:
    directory = ensure_output_dir(output_dir)
    outputs: list[Path] = []
    try:
        with zipfile.ZipFile(source) as archive:
            members = sorted(
                (
                    info
                    for info in archive.infolist()
                    if not info.is_dir()
                    and info.filename.startswith(prefix_path)
                    and (member_filter is None or member_filter(info.filename))
                ),
                key=lambda info: tuple(
                    (1, int(part)) if part.isdigit() else (0, part.casefold())
                    for part in re.split(r"(\d+)", info.filename)
                ),
            )
            for index, info in enumerate(members, 1):
                original_name = safe_filename(
                    Path(info.filename).name, f"media_{index}"
                )
                name_prefix = safe_filename(prefix or source.stem)
                target = unique_path(
                    directory / f"{name_prefix}_{index:03d}_{original_name}", overwrite
                )
                with atomic_output(target) as temporary:
                    temporary.write_bytes(archive.read(info))
                outputs.append(target)
    except zipfile.BadZipFile as exc:
        raise ValidationError(f"不是有效的 OOXML 文件：{source.name}") from exc
    return outputs


def word_extract_images(
    source: PathLike,
    output_dir: PathLike,
    *,
    prefix: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    return _extract_zip_members(
        source_path, output_dir, "word/media/", prefix=prefix, overwrite=overwrite
    )


def word_mail_merge(
    source: PathLike,
    output_dir: PathLike,
    records: Sequence[Mapping[str, Any]] | Mapping[str, Any],
    *,
    filename_template: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    """Generate documents by replacing ``{{name}}`` placeholders.

    This is deterministic placeholder merging, not Word's database-backed MERGEFIELD
    engine.  Placeholders may span multiple runs.
    """

    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    normalized = [records] if isinstance(records, Mapping) else list(records)
    if not normalized:
        raise ValidationError("邮件合并记录不能为空")
    directory = ensure_output_dir(output_dir)
    outputs: list[Path] = []
    for index, record in enumerate(normalized, 1):
        _path, document = _load_word(source_path)
        replacements = {f"{{{{{key}}}}}": value for key, value in record.items()}
        for paragraph in _iter_word_paragraphs(document):
            _replace_in_paragraph(paragraph, replacements, case_sensitive=True)
        if filename_template:
            try:
                requested = filename_template.format_map(_SafeFormatDict(record))
            except (ValueError, KeyError) as exc:
                raise ValidationError(f"文件名模板无效：{exc}") from exc
            stem = safe_filename(requested, f"{source_path.stem}_{index:03d}")
        else:
            stem = f"{source_path.stem}_{index:03d}"
        if stem.lower().endswith(source_path.suffix.lower()):
            stem = stem[: -len(source_path.suffix)]
        target = unique_path(directory / f"{stem}{source_path.suffix}", overwrite)
        _save_document(document, target)
        outputs.append(target)
    return outputs


class _SafeFormatDict(dict[str, Any]):
    def __missing__(self, key: str) -> str:
        raise KeyError(key)


def word_remove_hyperlinks(
    source: PathLike,
    output_dir: PathLike,
    *,
    overwrite: bool = False,
) -> list[Path]:
    """Remove explicit OOXML hyperlinks while preserving their visible text."""

    source_path = _require_source(source, _WORD_EXTENSIONS, "Word")
    target = _output_path(source_path, output_dir, tag="无超链接", overwrite=overwrite)

    def transform(name: str, payload: bytes) -> bytes | None:
        if name.startswith("word/") and name.endswith(".xml"):
            root = ET.fromstring(payload)
            changed = True
            while changed:
                changed = False
                for parent in root.iter():
                    for child in list(parent):
                        if child.tag == f"{{{_W}}}hyperlink":
                            _unwrap(parent, child)
                            changed = True
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
        if name.startswith("word/") and name.endswith(".rels"):
            root = ET.fromstring(payload)
            for relationship in list(root):
                if relationship.attrib.get("Type", "").endswith("/hyperlink"):
                    root.remove(relationship)
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
        return payload

    return _rewrite_zip(source_path, target, transform)


# ---------------------------------------------------------------------------
# Excel


def _selected_sheets(workbook: Any, names: Sequence[str] | None) -> list[Any]:
    if names is None:
        return list(workbook.worksheets)
    if isinstance(names, str):
        names = [names]
    missing = [name for name in names if name not in workbook.sheetnames]
    if missing:
        raise ValidationError(f"工作表不存在：{missing[0]}")
    return [workbook[name] for name in names]


def _excel_serializable(value: Any) -> Any:
    if isinstance(value, (_datetime.datetime, _datetime.date, _datetime.time)):
        return value.isoformat()
    if isinstance(value, _datetime.timedelta):
        return value.total_seconds()
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return str(value)


def _sheet_file_target(
    source: Path,
    output_dir: PathLike,
    sheet_name: str,
    suffix: str,
    overwrite: bool,
) -> Path:
    directory = ensure_output_dir(output_dir)
    stem = safe_filename(f"{source.stem}_{sheet_name}")
    return unique_path(directory / f"{stem}{suffix}", overwrite)


def excel_to_csv(
    source: PathLike,
    output_dir: PathLike,
    *,
    sheet_names: Sequence[str] | None = None,
    delimiter: str = ",",
    encoding: str = "utf-8-sig",
    data_only: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source, data_only=data_only)
    if len(delimiter) != 1:
        raise ValidationError("CSV delimiter 必须是单个字符")
    outputs: list[Path] = []
    for worksheet in _selected_sheets(workbook, sheet_names):
        target = _sheet_file_target(
            source_path, output_dir, worksheet.title, ".csv", overwrite
        )
        with atomic_output(target) as temporary:
            with temporary.open("w", encoding=encoding, newline="") as handle:
                writer = csv.writer(handle, delimiter=delimiter)
                for row in worksheet.iter_rows(values_only=True):
                    writer.writerow([_excel_serializable(value) for value in row])
        outputs.append(target)
    workbook.close()
    return outputs


def _sheet_records(worksheet: Any, header: bool) -> Any:
    rows = [
        [_excel_serializable(value) for value in row]
        for row in worksheet.iter_rows(values_only=True)
    ]
    if not header:
        return rows
    if not rows:
        return []
    raw_headers = rows[0]
    seen: dict[str, int] = {}
    headers: list[str] = []
    for index, value in enumerate(raw_headers, 1):
        base = str(value) if value not in (None, "") else f"column_{index}"
        count = seen.get(base, 0) + 1
        seen[base] = count
        headers.append(base if count == 1 else f"{base}_{count}")
    return [
        dict(zip(headers, row + [None] * (len(headers) - len(row)))) for row in rows[1:]
    ]


def excel_to_json(
    source: PathLike,
    output_dir: PathLike,
    *,
    sheet_names: Sequence[str] | None = None,
    header: bool = True,
    data_only: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source, data_only=data_only)
    payload = {
        sheet.title: _sheet_records(sheet, header)
        for sheet in _selected_sheets(workbook, sheet_names)
    }
    target = _output_path(source_path, output_dir, suffix=".json", overwrite=overwrite)
    with atomic_output(target) as temporary:
        temporary.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    workbook.close()
    return [target]


def excel_to_xml(
    source: PathLike,
    output_dir: PathLike,
    *,
    sheet_names: Sequence[str] | None = None,
    data_only: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source, data_only=data_only)
    try:
        from openpyxl.utils import get_column_letter
    except ImportError as exc:  # pragma: no cover
        raise MissingEngineError("缺少 openpyxl") from exc
    root = ET.Element("workbook", {"source": source_path.name})
    for worksheet in _selected_sheets(workbook, sheet_names):
        sheet_element = ET.SubElement(root, "worksheet", {"name": worksheet.title})
        for row_index, row in enumerate(worksheet.iter_rows(values_only=True), 1):
            row_element = ET.SubElement(sheet_element, "row", {"index": str(row_index)})
            for column_index, value in enumerate(row, 1):
                cell = ET.SubElement(
                    row_element, "cell", {"column": get_column_letter(column_index)}
                )
                if value is not None:
                    cell.text = str(_excel_serializable(value))
    target = _output_path(source_path, output_dir, suffix=".xml", overwrite=overwrite)
    with atomic_output(target) as temporary:
        ET.ElementTree(root).write(temporary, encoding="utf-8", xml_declaration=True)
    workbook.close()
    return [target]


def excel_to_txt(
    source: PathLike,
    output_dir: PathLike,
    *,
    sheet_names: Sequence[str] | None = None,
    delimiter: str = "\t",
    encoding: str = "utf-8",
    data_only: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source, data_only=data_only)
    if len(delimiter) != 1:
        raise ValidationError("TXT delimiter 必须是单个字符")
    outputs: list[Path] = []
    for worksheet in _selected_sheets(workbook, sheet_names):
        target = _sheet_file_target(
            source_path, output_dir, worksheet.title, ".txt", overwrite
        )
        with atomic_output(target) as temporary:
            with temporary.open("w", encoding=encoding, newline="") as handle:
                writer = csv.writer(handle, delimiter=delimiter, lineterminator="\n")
                for row in worksheet.iter_rows(values_only=True):
                    writer.writerow([_excel_serializable(value) for value in row])
        outputs.append(target)
    workbook.close()
    return outputs


def _resolve_column(column: int | str) -> int:
    if isinstance(column, int):
        if column < 1:
            raise ValidationError("列号从 1 开始")
        return column
    text = str(column).strip()
    if text.isdigit():
        return _resolve_column(int(text))
    try:
        from openpyxl.utils import column_index_from_string

        return column_index_from_string(text.upper())
    except (ImportError, ValueError) as exc:
        raise ValidationError(f"无效列：{column}") from exc


def _choose_sheet(workbook: Any, sheet_name: str | None) -> Any:
    if sheet_name is None:
        return workbook.active
    if sheet_name not in workbook.sheetnames:
        raise ValidationError(f"工作表不存在：{sheet_name}")
    return workbook[sheet_name]


def _save_workbook(
    source: Path, workbook: Any, output_dir: PathLike, tag: str, overwrite: bool
) -> list[Path]:
    target = _output_path(source, output_dir, tag=tag, overwrite=overwrite)
    with atomic_output(target) as temporary:
        workbook.save(temporary)
    workbook.close()
    return [target]


def _cell_snapshot(cell: Any) -> dict[str, Any]:
    return {
        "coordinate": cell.coordinate,
        "value": cell.value,
        "style": copy(cell._style),
        "number_format": cell.number_format,
        "font": copy(cell.font),
        "fill": copy(cell.fill),
        "border": copy(cell.border),
        "alignment": copy(cell.alignment),
        "protection": copy(cell.protection),
        "hyperlink": copy(cell.hyperlink),
        "comment": copy(cell.comment),
    }


def _restore_cell(cell: Any, snapshot: Mapping[str, Any]) -> None:
    value = snapshot["value"]
    if (
        isinstance(value, str)
        and value.startswith("=")
        and snapshot["coordinate"] != cell.coordinate
    ):
        try:
            from openpyxl.formula.translate import Translator

            value = Translator(value, origin=snapshot["coordinate"]).translate_formula(
                cell.coordinate
            )
        except Exception:
            warnings.warn(
                f"公式 {snapshot['coordinate']} 移动到 {cell.coordinate} 时无法重写引用，已保留原公式",
                RuntimeWarning,
                stacklevel=2,
            )
    cell.value = value
    cell._style = copy(snapshot["style"])
    cell.number_format = snapshot["number_format"]
    cell.font = copy(snapshot["font"])
    cell.fill = copy(snapshot["fill"])
    cell.border = copy(snapshot["border"])
    cell.alignment = copy(snapshot["alignment"])
    cell.protection = copy(snapshot["protection"])
    cell.hyperlink = copy(snapshot["hyperlink"])
    cell.comment = copy(snapshot["comment"])


def _sort_key(value: Any) -> tuple[int, Any]:
    if value is None:
        return (3, "")
    if isinstance(value, bool):
        return (0, int(value))
    if isinstance(value, (int, float)):
        return (0, value)
    if isinstance(value, (_datetime.date, _datetime.datetime, _datetime.time)):
        return (1, value.isoformat())
    return (2, str(value).casefold())


def excel_sort_rows(
    source: PathLike,
    output_dir: PathLike,
    *,
    column: int | str,
    sheet_name: str | None = None,
    header: bool = True,
    reverse: bool = False,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    column_index = _resolve_column(column)
    if column_index > worksheet.max_column:
        workbook.close()
        raise ValidationError(f"排序列超出数据范围：{column}")
    start_row = 2 if header else 1
    if any(
        cell_range.min_row <= worksheet.max_row and cell_range.max_row >= start_row
        for cell_range in worksheet.merged_cells.ranges
    ):
        workbook.close()
        raise ValidationError("排序区域含合并单元格，无法在不破坏结构的情况下可靠排序")
    rows = [
        (
            [
                _cell_snapshot(worksheet.cell(row=row_index, column=column_index_inner))
                for column_index_inner in range(1, worksheet.max_column + 1)
            ],
            copy(worksheet.row_dimensions.get(row_index)),
        )
        for row_index in range(start_row, worksheet.max_row + 1)
    ]
    rows.sort(
        key=lambda item: _sort_key(item[0][column_index - 1]["value"]),
        reverse=reverse,
    )
    for row_offset, (row, row_dimension) in enumerate(rows, start_row):
        for column_offset, snapshot in enumerate(row, 1):
            _restore_cell(
                worksheet.cell(row=row_offset, column=column_offset), snapshot
            )
        if row_dimension is None:
            worksheet.row_dimensions.pop(row_offset, None)
        else:
            row_dimension.index = row_offset
            worksheet.row_dimensions[row_offset] = row_dimension
    return _save_workbook(source_path, workbook, output_dir, "排序", overwrite)


def _filter_match(
    cell_value: Any, wanted: Any, operator: str, case_sensitive: bool
) -> bool:
    if operator in {"equals", "not_equals"}:
        left, right = cell_value, wanted
        if isinstance(left, str) and isinstance(right, str) and not case_sensitive:
            left, right = left.casefold(), right.casefold()
        result = left == right
        return result if operator == "equals" else not result
    if operator in {"contains", "not_contains", "starts_with", "ends_with"}:
        left = "" if cell_value is None else str(cell_value)
        right = "" if wanted is None else str(wanted)
        if not case_sensitive:
            left, right = left.casefold(), right.casefold()
        if operator in {"contains", "not_contains"}:
            result = right in left
            return result if operator == "contains" else not result
        return (
            left.startswith(right)
            if operator == "starts_with"
            else left.endswith(right)
        )
    if operator in {"greater_than", "greater_equal", "less_than", "less_equal"}:
        try:
            if operator == "greater_than":
                return cell_value > wanted
            if operator == "greater_equal":
                return cell_value >= wanted
            if operator == "less_than":
                return cell_value < wanted
            return cell_value <= wanted
        except TypeError:
            return False
    if operator == "is_blank":
        return cell_value in (None, "")
    if operator == "not_blank":
        return cell_value not in (None, "")
    raise ValidationError(f"不支持的筛选运算符：{operator}")


def _worksheet_deletion_hazards(worksheet: Any) -> list[str]:
    """Return structures openpyxl cannot safely retarget after row/column deletion."""

    hazards: list[str] = []
    if any(
        cell.data_type == "f"
        or (isinstance(cell.value, str) and cell.value.startswith("="))
        for row in worksheet.iter_rows()
        for cell in row
    ):
        hazards.append("公式引用")
    if worksheet.tables:
        hazards.append("Excel 表")
    if worksheet.merged_cells.ranges:
        hazards.append("合并单元格")
    if getattr(worksheet, "_images", None):
        hazards.append("图片锚点")
    if getattr(worksheet, "_charts", None):
        hazards.append("图表引用")
    if len(worksheet.conditional_formatting):
        hazards.append("条件格式范围")
    validations = getattr(
        getattr(worksheet, "data_validations", None), "dataValidation", ()
    )
    if validations:
        hazards.append("数据验证范围")
    return hazards


def _refuse_unsafe_deletion(worksheet: Any, action: str) -> None:
    hazards = _worksheet_deletion_hazards(worksheet)
    if hazards:
        raise ValidationError(
            f"工作表“{worksheet.title}”含{'、'.join(hazards)}；openpyxl 无法在{action}后"
            "可靠重写所有依赖。请先简化工作表，或使用 WPS/Microsoft Office/LibreOffice 高保真引擎"
        )


def excel_filter_rows(
    source: PathLike,
    output_dir: PathLike,
    *,
    column: int | str,
    value: Any = None,
    operator: str = "equals",
    sheet_name: str | None = None,
    header: bool = True,
    case_sensitive: bool = False,
    overwrite: bool = False,
) -> list[Path]:
    """Hide non-matching rows; data is preserved and can be unhidden later."""

    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    column_index = _resolve_column(column)
    start_row = 2 if header else 1
    for row_index in range(start_row, worksheet.max_row + 1):
        worksheet.row_dimensions[row_index].hidden = not _filter_match(
            worksheet.cell(row=row_index, column=column_index).value,
            value,
            operator,
            case_sensitive,
        )
    if worksheet.max_row and worksheet.max_column:
        from openpyxl.utils import get_column_letter

        worksheet.auto_filter.ref = (
            f"A1:{get_column_letter(worksheet.max_column)}{worksheet.max_row}"
        )
    return _save_workbook(source_path, workbook, output_dir, "筛选", overwrite)


def excel_remove_duplicates(
    source: PathLike,
    output_dir: PathLike,
    *,
    columns: Sequence[int | str] | None = None,
    sheet_name: str | None = None,
    header: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    column_indexes = (
        [_resolve_column(column) for column in columns]
        if columns
        else list(range(1, worksheet.max_column + 1))
    )
    seen: set[tuple[Any, ...]] = set()
    duplicate_rows: list[int] = []
    for row_index in range(2 if header else 1, worksheet.max_row + 1):
        key = tuple(
            worksheet.cell(row=row_index, column=column).value
            for column in column_indexes
        )
        try:
            already_seen = key in seen
            if not already_seen:
                seen.add(key)
        except TypeError:
            normalized = tuple(repr(value) for value in key)
            already_seen = normalized in seen
            if not already_seen:
                seen.add(normalized)
        if already_seen:
            duplicate_rows.append(row_index)
    if duplicate_rows:
        try:
            _refuse_unsafe_deletion(worksheet, "删除重复行")
        except ValidationError:
            workbook.close()
            raise
    for row_index in reversed(duplicate_rows):
        worksheet.delete_rows(row_index)
    return _save_workbook(source_path, workbook, output_dir, "去重", overwrite)


def excel_remove_blank_rows_columns(
    source: PathLike,
    output_dir: PathLike,
    *,
    sheet_name: str | None = None,
    remove_rows: bool = True,
    remove_columns: bool = True,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    worksheets = (
        [_choose_sheet(workbook, sheet_name)] if sheet_name else workbook.worksheets
    )
    for worksheet in worksheets:
        blank_rows = (
            [
                row_index
                for row_index in range(1, worksheet.max_row + 1)
                if all(
                    worksheet.cell(row=row_index, column=column).value in (None, "")
                    for column in range(1, worksheet.max_column + 1)
                )
            ]
            if remove_rows
            else []
        )
        blank_columns = (
            [
                column_index
                for column_index in range(1, worksheet.max_column + 1)
                if all(
                    worksheet.cell(row=row, column=column_index).value in (None, "")
                    for row in range(1, worksheet.max_row + 1)
                )
            ]
            if remove_columns
            else []
        )
        if blank_rows or blank_columns:
            try:
                _refuse_unsafe_deletion(worksheet, "删除空白行列")
            except ValidationError:
                workbook.close()
                raise
        if remove_rows:
            for row_index in reversed(blank_rows):
                worksheet.delete_rows(row_index)
        if remove_columns:
            for column_index in reversed(blank_columns):
                worksheet.delete_cols(column_index)
    return _save_workbook(source_path, workbook, output_dir, "无空行列", overwrite)


def excel_replace_text(
    source: PathLike,
    output_dir: PathLike,
    replacements: Mapping[str, Any],
    *,
    sheet_names: Sequence[str] | None = None,
    case_sensitive: bool = True,
    exact: bool = False,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    if not replacements:
        workbook.close()
        raise ValidationError("替换映射不能为空")
    for worksheet in _selected_sheets(workbook, sheet_names):
        for row in worksheet.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                value = cell.value
                for old, new in replacements.items():
                    old, new = str(old), str(new)
                    if exact:
                        matches = (
                            value == old
                            if case_sensitive
                            else value.casefold() == old.casefold()
                        )
                        if matches:
                            value = new
                    elif case_sensitive:
                        value = value.replace(old, new)
                    else:
                        value = re.sub(
                            re.escape(old),
                            lambda _match, replacement=new: replacement,
                            value,
                            flags=re.IGNORECASE,
                        )
                cell.value = value
    return _save_workbook(source_path, workbook, output_dir, "替换", overwrite)


def excel_formulas_to_values(
    source: PathLike,
    output_dir: PathLike,
    *,
    sheet_names: Sequence[str] | None = None,
    missing_cache: str = "error",
    overwrite: bool = False,
) -> list[Path]:
    """Replace formulas with workbook cached results; formulas are never calculated here.

    ``openpyxl`` does not calculate formulas.  When WPS/Excel/LibreOffice has not stored a
    cached result, ``missing_cache='error'`` refuses the conversion.  ``'warn'`` keeps
    such formulas unchanged and emits :class:`RuntimeWarning`.
    """

    policy = missing_cache.lower().strip()
    if policy not in {"error", "warn"}:
        raise ValidationError("missing_cache 必须是 error 或 warn")
    source_path, formulas = _load_workbook(source, data_only=False)
    _cached_path, cached = _load_workbook(source, data_only=True)
    names = [sheet.title for sheet in _selected_sheets(formulas, sheet_names)]
    missing: list[str] = []
    for name in names:
        formula_sheet = formulas[name]
        cached_sheet = cached[name]
        for row in formula_sheet.iter_rows():
            for cell in row:
                if cell.data_type == "f" or (
                    isinstance(cell.value, str) and cell.value.startswith("=")
                ):
                    cached_value = cached_sheet[cell.coordinate].value
                    if cached_value is None:
                        missing.append(f"{name}!{cell.coordinate}")
                    else:
                        cell.value = cached_value
    cached.close()
    if missing and policy == "error":
        formulas.close()
        preview = "、".join(missing[:5])
        raise ValidationError(
            f"{len(missing)} 个公式没有可用的计算缓存（如 {preview}）；"
            "请先用 WPS/Excel/LibreOffice 重新计算并保存，或使用 missing_cache='warn' 保留这些公式"
        )
    if missing:
        warnings.warn(
            f"{len(missing)} 个公式没有计算缓存，已原样保留：{'、'.join(missing[:5])}",
            RuntimeWarning,
            stacklevel=2,
        )
    return _save_workbook(source_path, formulas, output_dir, "公式转值", overwrite)


def excel_split_column(
    source: PathLike,
    output_dir: PathLike,
    *,
    column: int | str,
    delimiter: str,
    sheet_name: str | None = None,
    header: bool = True,
    maxsplit: int = -1,
    overwrite: bool = False,
) -> list[Path]:
    if delimiter == "":
        raise ValidationError("分隔符不能为空")
    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    column_index = _resolve_column(column)
    hazards = _worksheet_deletion_hazards(worksheet)
    if hazards:
        workbook.close()
        raise ValidationError(
            f"工作表“{worksheet.title}”含{'、'.join(hazards)}；插入分列会使引用或锚点失效，"
            "请使用 WPS/Microsoft Office/LibreOffice 高保真引擎"
        )
    start_row = 2 if header else 1
    split_rows: dict[int, list[str]] = {}
    width = 1
    for row_index in range(start_row, worksheet.max_row + 1):
        value = worksheet.cell(row=row_index, column=column_index).value
        parts = str(value).split(delimiter, maxsplit) if value is not None else [""]
        split_rows[row_index] = parts
        width = max(width, len(parts))
    if width > 1:
        worksheet.insert_cols(column_index + 1, width - 1)
    for row_index, parts in split_rows.items():
        for offset, part in enumerate(parts):
            worksheet.cell(row=row_index, column=column_index + offset).value = part
    return _save_workbook(source_path, workbook, output_dir, "分列", overwrite)


def excel_merge_cells(
    source: PathLike,
    output_dir: PathLike,
    ranges: Sequence[str] | str,
    *,
    sheet_name: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    normalized = [ranges] if isinstance(ranges, str) else list(ranges)
    if not normalized:
        workbook.close()
        raise ValidationError("合并范围不能为空")
    from openpyxl.utils.cell import range_boundaries

    for cell_range in normalized:
        min_col, min_row, max_col, max_row = range_boundaries(str(cell_range))
        occupied = [
            worksheet.cell(row=row, column=column).coordinate
            for row in range(min_row, max_row + 1)
            for column in range(min_col, max_col + 1)
            if (row, column) != (min_row, min_col)
            and worksheet.cell(row=row, column=column).value not in (None, "")
        ]
        if occupied:
            workbook.close()
            raise ValidationError(
                f"合并 {cell_range} 会丢弃非左上角单元格内容（如 {occupied[0]}），已拒绝操作"
            )
        worksheet.merge_cells(str(cell_range))
    return _save_workbook(source_path, workbook, output_dir, "合并单元格", overwrite)


def excel_unmerge_cells(
    source: PathLike,
    output_dir: PathLike,
    ranges: Sequence[str] | str | None = None,
    *,
    sheet_name: str | None = None,
    fill: bool = False,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    if ranges is None:
        normalized = [str(cell_range) for cell_range in worksheet.merged_cells.ranges]
    elif isinstance(ranges, str):
        normalized = [ranges]
    else:
        normalized = list(ranges)
    existing = {str(cell_range) for cell_range in worksheet.merged_cells.ranges}
    for cell_range in normalized:
        if str(cell_range) not in existing:
            workbook.close()
            raise ValidationError(f"不是现有合并范围：{cell_range}")
        merged = next(
            item
            for item in worksheet.merged_cells.ranges
            if str(item) == str(cell_range)
        )
        value = worksheet.cell(merged.min_row, merged.min_col).value
        bounds = (merged.min_row, merged.max_row, merged.min_col, merged.max_col)
        worksheet.unmerge_cells(str(cell_range))
        if fill:
            min_row, max_row, min_col, max_col = bounds
            for row in range(min_row, max_row + 1):
                for column in range(min_col, max_col + 1):
                    worksheet.cell(row, column).value = value
    return _save_workbook(source_path, workbook, output_dir, "拆分单元格", overwrite)


def excel_apply_conditional_format(
    source: PathLike,
    output_dir: PathLike,
    *,
    cell_range: str,
    rule: str = "cell",
    operator: str = "greaterThan",
    threshold: Any = None,
    formula: str | Sequence[str] | None = None,
    fill_color: str = "FFF2CC",
    sheet_name: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    worksheet = _choose_sheet(workbook, sheet_name)
    try:
        from openpyxl.formatting.rule import CellIsRule, ColorScaleRule, FormulaRule
        from openpyxl.styles import PatternFill
    except ImportError as exc:  # pragma: no cover
        workbook.close()
        raise MissingEngineError("缺少 openpyxl") from exc
    fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
    if rule == "cell":
        if threshold is None:
            workbook.close()
            raise ValidationError("cell 条件格式需要 threshold")
        worksheet.conditional_formatting.add(
            cell_range,
            CellIsRule(operator=operator, formula=[str(threshold)], fill=fill),
        )
    elif rule == "formula":
        if formula is None:
            workbook.close()
            raise ValidationError("formula 条件格式需要 formula")
        formulas = [formula] if isinstance(formula, str) else list(formula)
        worksheet.conditional_formatting.add(
            cell_range, FormulaRule(formula=formulas, fill=fill)
        )
    elif rule == "color_scale":
        worksheet.conditional_formatting.add(
            cell_range,
            ColorScaleRule(
                start_type="min",
                start_color="F8696B",
                mid_type="percentile",
                mid_value=50,
                mid_color="FFEB84",
                end_type="max",
                end_color="63BE7B",
            ),
        )
    else:
        workbook.close()
        raise ValidationError("rule 必须是 cell、formula 或 color_scale")
    return _save_workbook(source_path, workbook, output_dir, "条件格式", overwrite)


def excel_extract_images(
    source: PathLike,
    output_dir: PathLike,
    *,
    prefix: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path = _require_source(source, _EXCEL_EXTENSIONS, "Excel")
    return _extract_zip_members(
        source_path, output_dir, "xl/media/", prefix=prefix, overwrite=overwrite
    )


def excel_manage_sheets(
    source: PathLike,
    output_dir: PathLike,
    *,
    rename: Mapping[str, str] | None = None,
    delete: Sequence[str] | None = None,
    copy_sheets: Mapping[str, str] | None = None,
    order: Sequence[str] | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path, workbook = _load_workbook(source)
    for old, new in (rename or {}).items():
        if old not in workbook.sheetnames:
            workbook.close()
            raise ValidationError(f"工作表不存在：{old}")
        if new in workbook.sheetnames and new != old:
            workbook.close()
            raise ValidationError(f"工作表名称已存在：{new}")
        workbook[old].title = str(new)
    for name in delete or ():
        if name not in workbook.sheetnames:
            workbook.close()
            raise ValidationError(f"工作表不存在：{name}")
        if len(workbook.worksheets) == 1:
            workbook.close()
            raise ValidationError("不能删除工作簿中的最后一个工作表")
        workbook.remove(workbook[name])
    for original, new_name in (copy_sheets or {}).items():
        if original not in workbook.sheetnames:
            workbook.close()
            raise ValidationError(f"工作表不存在：{original}")
        if new_name in workbook.sheetnames:
            workbook.close()
            raise ValidationError(f"工作表名称已存在：{new_name}")
        workbook.copy_worksheet(workbook[original]).title = str(new_name)
    if order is not None:
        requested = list(order)
        if set(requested) != set(workbook.sheetnames) or len(requested) != len(
            workbook.sheetnames
        ):
            workbook.close()
            raise ValidationError("order 必须恰好包含当前所有工作表名称且不能重复")
        workbook._sheets = [workbook[name] for name in requested]
    return _save_workbook(source_path, workbook, output_dir, "工作表管理", overwrite)


# ---------------------------------------------------------------------------
# PowerPoint


def _iter_ppt_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == 6 and hasattr(
            shape, "shapes"
        ):  # MSO_GROUP
            yield from _iter_ppt_shapes(shape.shapes)


def _iter_shape_text_frames(shape: Any) -> Iterable[Any]:
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame


def ppt_replace_fonts(
    source: PathLike,
    output_dir: PathLike,
    replacements: Mapping[str, str] | None = None,
    *,
    default_font: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path, presentation = _load_presentation(source)
    mapping = {
        str(key).casefold(): str(value) for key, value in (replacements or {}).items()
    }
    if not mapping and not default_font:
        raise ValidationError("replacements 或 default_font 至少提供一项")
    shape_sets = [slide.shapes for slide in presentation.slides]
    shape_sets.extend(layout.shapes for layout in presentation.slide_layouts)
    shape_sets.extend(master.shapes for master in presentation.slide_masters)
    for shapes in shape_sets:
        for shape in _iter_ppt_shapes(shapes):
            for text_frame in _iter_shape_text_frames(shape):
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        current = run.font.name
                        replacement = (
                            mapping.get(current.casefold()) if current else None
                        )
                        if replacement or default_font:
                            run.font.name = replacement or default_font
    target = _output_path(source_path, output_dir, tag="替换字体", overwrite=overwrite)
    return _save_document(presentation, target)


def _set_color_alpha(color_element: Any, opacity: float) -> None:
    try:
        from pptx.oxml.xmlchemy import OxmlElement
    except ImportError:  # pragma: no cover
        return
    for child in list(color_element):
        if child.tag == f"{{{_A}}}alpha":
            color_element.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(round(opacity * 100000)))
    color_element.append(alpha)


def ppt_add_watermark(
    source: PathLike,
    output_dir: PathLike,
    text: str,
    *,
    font_size_pt: float = 32,
    color: str = "B7B7B7",
    rotation: float = -30,
    opacity: float = 0.35,
    overwrite: bool = False,
) -> list[Path]:
    if not text:
        raise ValidationError("水印文字不能为空")
    if not 0 <= opacity <= 1:
        raise ValidationError("opacity 必须在 0 到 1 之间")
    source_path, presentation = _load_presentation(source)
    try:
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Pt
    except ImportError as exc:  # pragma: no cover
        raise MissingEngineError("缺少 python-pptx") from exc
    try:
        rgb = RGBColor.from_string(color.lstrip("#"))
    except ValueError as exc:
        raise ValidationError(f"无效颜色：{color}") from exc
    width = int(presentation.slide_width * 0.8)
    height = int(presentation.slide_height * 0.18)
    left = int((presentation.slide_width - width) / 2)
    top = int((presentation.slide_height - height) / 2)
    for slide in presentation.slides:
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.rotation = rotation % 360
        shape.fill.background()
        shape.line.fill.background()
        frame = shape.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.bold = True
        run.font.color.rgb = rgb
        _set_color_alpha(run.font.color._color._srgbClr, opacity)
    target = _output_path(source_path, output_dir, tag="水印", overwrite=overwrite)
    return _save_document(presentation, target)


_IMAGE_SUFFIXES = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".bmp",
    ".tif",
    ".tiff",
    ".webp",
    ".svg",
    ".emf",
    ".wmf",
}
_AUDIO_SUFFIXES = {".mp3", ".wav", ".m4a", ".aac", ".wma", ".ogg", ".flac"}
_VIDEO_SUFFIXES = {".mp4", ".mov", ".avi", ".wmv", ".mkv", ".webm", ".mpeg", ".mpg"}


def ppt_extract_media(
    source: PathLike,
    output_dir: PathLike,
    *,
    include_images: bool = True,
    include_audio: bool = True,
    include_video: bool = True,
    prefix: str | None = None,
    overwrite: bool = False,
) -> list[Path]:
    source_path = _require_source(source, _PPT_EXTENSIONS, "PowerPoint")

    def allowed(name: str) -> bool:
        suffix = Path(name).suffix.lower()
        return (
            (include_images and suffix in _IMAGE_SUFFIXES)
            or (include_audio and suffix in _AUDIO_SUFFIXES)
            or (include_video and suffix in _VIDEO_SUFFIXES)
        )

    return _extract_zip_members(
        source_path,
        output_dir,
        "ppt/media/",
        prefix=prefix,
        overwrite=overwrite,
        member_filter=allowed,
    )


def ppt_compress_images(
    source: PathLike,
    output_dir: PathLike,
    *,
    quality: int = 82,
    max_dimension: int | None = 1920,
    overwrite: bool = False,
) -> list[Path]:
    """Re-encode JPEG/PNG images in-place inside a copied PPTX package.

    Other media and unsupported image formats are copied byte-for-byte.  The
    original file is never modified.
    """

    if not 1 <= quality <= 100:
        raise ValidationError("quality 必须在 1 到 100 之间")
    if max_dimension is not None and max_dimension < 1:
        raise ValidationError("max_dimension 必须大于 0 或为 None")
    source_path = _require_source(source, _PPT_EXTENSIONS, "PowerPoint")
    try:
        from PIL import Image
    except ImportError as exc:  # pragma: no cover
        raise MissingEngineError("缺少 Pillow，无法压缩 PPT 图片") from exc
    from io import BytesIO

    target = _output_path(source_path, output_dir, tag="压缩图片", overwrite=overwrite)

    def transform(name: str, payload: bytes) -> bytes | None:
        suffix = Path(name).suffix.lower()
        if not name.startswith("ppt/media/") or suffix not in {".jpg", ".jpeg", ".png"}:
            return payload
        try:
            with Image.open(BytesIO(payload)) as opened:
                opened.load()
                working = opened
                owns_working = False
                if max_dimension and max(working.size) > max_dimension:
                    ratio = max_dimension / max(working.size)
                    working = working.resize(
                        (
                            max(1, round(working.width * ratio)),
                            max(1, round(working.height * ratio)),
                        ),
                        Image.Resampling.LANCZOS,
                    )
                    owns_working = True
                try:
                    with BytesIO() as output:
                        if suffix in {".jpg", ".jpeg"}:
                            prepared = working
                            owns_prepared = False
                            try:
                                if working.mode not in {"RGB", "L", "CMYK"}:
                                    if "A" in working.getbands():
                                        rgba = working.convert("RGBA")
                                        try:
                                            prepared = Image.new(
                                                "RGB", working.size, "white"
                                            )
                                            prepared.paste(
                                                rgba, mask=rgba.getchannel("A")
                                            )
                                        finally:
                                            rgba.close()
                                    else:
                                        prepared = working.convert("RGB")
                                    owns_prepared = True
                                subsampling = (
                                    0 if quality >= 85 else 1 if quality >= 70 else 2
                                )
                                prepared.save(
                                    output,
                                    "JPEG",
                                    quality=quality,
                                    subsampling=subsampling,
                                    optimize=True,
                                    progressive=True,
                                )
                            finally:
                                if owns_prepared:
                                    prepared.close()
                        else:
                            working.save(
                                output, "PNG", optimize=False, compress_level=6
                            )
                        compressed = output.getvalue()
                    # Never make a presentation larger merely to claim it was compressed.
                    return compressed if len(compressed) < len(payload) else payload
                finally:
                    if owns_working:
                        working.close()
        except Exception:
            return payload

    return _rewrite_zip(source_path, target, transform)


# ---------------------------------------------------------------------------
# High-fidelity Office/LibreOffice adapter


@dataclass(frozen=True)
class OfficeEngineStatus:
    available: bool
    executable: Path | None = None
    reason: str = ""


@dataclass(frozen=True)
class _OfficeProcessIdentity:
    pid: int
    executable: Path
    created: str


_MICROSOFT_COMPONENTS = {
    "microsoft_word": ("Word.Application", "Word"),
    "microsoft_excel": ("Excel.Application", "Excel"),
    "microsoft_powerpoint": ("PowerPoint.Application", "PowerPoint"),
}

_MICROSOFT_COM_EXECUTABLES = {
    "Word.Application": "winword.exe",
    "Excel.Application": "excel.exe",
    "PowerPoint.Application": "powerpnt.exe",
}

_MICROSOFT_COM_CLSIDS = {
    "Word.Application": "{000209FF-0000-0000-C000-000000000046}",
    "Excel.Application": "{00024500-0000-0000-C000-000000000046}",
    "PowerPoint.Application": "{91493441-5A91-11CF-8700-00AA0060263B}",
}

_MICROSOFT_COM_NAMES = {
    "Word.Application": "word",
    "Excel.Application": "excel",
    "PowerPoint.Application": "powerpoint",
}

_COMPONENT_FOR_SUFFIX = {
    **{suffix: "microsoft_word" for suffix in _WORD_EXTENSIONS | {".doc", ".rtf"}},
    **{suffix: "microsoft_excel" for suffix in _EXCEL_EXTENSIONS | {".xls", ".csv"}},
    **{suffix: "microsoft_powerpoint" for suffix in _PPT_EXTENSIONS | {".ppt"}},
}

_WPS_KIND_FOR_COMPONENT = {
    "microsoft_word": "writer",
    "microsoft_excel": "spreadsheets",
    "microsoft_powerpoint": "presentation",
}


def _office_application_pid(application: Any) -> int | None:
    """Return the process id for a newly started Office COM application."""

    if sys.platform != "win32":
        return None
    window_handle = 0
    for attribute in ("Hwnd", "HWND"):
        try:
            raw_handle = getattr(application, attribute)
            if callable(raw_handle):
                raw_handle = raw_handle()
            window_handle = int(raw_handle)
        except (AttributeError, TypeError, ValueError, OSError):
            continue
        except Exception:
            continue
        if window_handle:
            break
    if not window_handle:
        return None
    try:
        import win32process

        _thread_id, process_id = win32process.GetWindowThreadProcessId(window_handle)
    except Exception:
        return None
    return int(process_id) if process_id else None


def _windows_process_identity(process_id: int) -> _OfficeProcessIdentity | None:
    """Read an exact Windows process identity without relying on its display name."""

    if sys.platform != "win32" or process_id <= 0:
        return None
    try:
        import win32api
        import win32process

        # PROCESS_QUERY_LIMITED_INFORMATION | PROCESS_VM_READ.  The latter is
        # needed by GetModuleFileNameEx on older supported Windows versions.
        handle = win32api.OpenProcess(0x1000 | 0x0010, False, int(process_id))
        try:
            executable = Path(
                str(win32process.GetModuleFileNameEx(handle, 0))
            ).resolve()
            created_value = win32process.GetProcessTimes(handle)["CreationTime"]
        finally:
            handle.Close()
    except Exception:
        return None
    created = (
        created_value.isoformat()
        if hasattr(created_value, "isoformat")
        else str(created_value)
    )
    return _OfficeProcessIdentity(int(process_id), executable, created)


def _windows_process_snapshot(
    executable_name: str,
) -> dict[int, _OfficeProcessIdentity]:
    """Snapshot running processes for one exact executable family."""

    if sys.platform != "win32":
        return {}
    try:
        import win32process

        process_ids = win32process.EnumProcesses()
    except Exception:
        return {}
    expected = executable_name.casefold()
    snapshot: dict[int, _OfficeProcessIdentity] = {}
    for process_id in process_ids:
        identity = _windows_process_identity(int(process_id))
        if identity is None or identity.executable.name.casefold() != expected:
            continue
        snapshot[identity.pid] = identity
    return snapshot


def _new_owned_office_process(
    before: Mapping[int, _OfficeProcessIdentity],
    *,
    expected_executable: Path,
    reported_pid: int | None = None,
) -> _OfficeProcessIdentity | None:
    """Return the one newly-created Office process that is safe to own.

    Hidden Word instances do not always expose ``Application.Hwnd``.  In that
    case the COM worker cannot report a PID, so compare exact process snapshots
    instead.  We only accept a unique new PID/path/creation tuple and never
    guess when Office started more than one matching process.
    """

    try:
        expected_path = expected_executable.expanduser().resolve()
    except OSError:
        expected_path = expected_executable.expanduser().absolute()
    expected_name = expected_path.name.casefold()

    if reported_pid is not None and reported_pid > 0:
        identity = _windows_process_identity(int(reported_pid))
        previous = before.get(int(reported_pid))
        if (
            identity is not None
            and identity != previous
            and identity.executable == expected_path
        ):
            return identity
        return None

    after = _windows_process_snapshot(expected_name)
    candidates = [
        identity
        for process_id, identity in after.items()
        if identity != before.get(process_id) and identity.executable == expected_path
    ]
    return candidates[0] if len(candidates) == 1 else None


def _terminate_owned_office_process(identity: _OfficeProcessIdentity) -> bool:
    """Terminate only the exact Office instance previously approved as ours."""

    current = _windows_process_identity(identity.pid)
    if current != identity:
        return False
    try:
        import win32api

        # PROCESS_TERMINATE | SYNCHRONIZE
        handle = win32api.OpenProcess(0x0001 | 0x00100000, False, identity.pid)
        try:
            win32api.TerminateProcess(handle, 1)
        finally:
            handle.Close()
    except Exception:
        return False
    return True


def _wait_for_owned_office_exit(
    identity: _OfficeProcessIdentity | None, timeout: float = 2.0
) -> None:
    if identity is None:
        return
    deadline = time.monotonic() + max(0.0, float(timeout))
    while time.monotonic() < deadline:
        if _windows_process_identity(identity.pid) != identity:
            return
        time.sleep(0.05)
    _terminate_owned_office_process(identity)


@contextmanager
def _word_pdf_reflow_prompt_suppressed(version: str = "16.0") -> Iterable[None]:
    """Temporarily suppress Word's informational PDF conversion prompt.

    ``Application.DisplayAlerts`` does not cover this modal dialog.  The
    documented UI checkbox stores a per-user option, so set it only while the
    supervised candidate is running and restore the exact prior registry value
    afterwards.  This does not alter Protected View or macro security.
    """

    if sys.platform != "win32":
        yield
        return
    try:
        import winreg
    except ImportError:
        yield
        return

    key_path = rf"Software\Microsoft\Office\{version}\Word\Options"
    value_name = "DisableConvertPrompt"
    key = None
    had_value = False
    previous_value: Any = None
    previous_type = int(getattr(winreg, "REG_DWORD", 4))
    try:
        key = winreg.CreateKeyEx(
            winreg.HKEY_CURRENT_USER,
            key_path,
            0,
            int(getattr(winreg, "KEY_READ", 0))
            | int(getattr(winreg, "KEY_SET_VALUE", 0)),
        )
        try:
            previous_value, previous_type = winreg.QueryValueEx(key, value_name)
            had_value = True
        except OSError:
            pass
        winreg.SetValueEx(
            key,
            value_name,
            0,
            int(getattr(winreg, "REG_DWORD", 4)),
            1,
        )
        yield
    finally:
        if key is not None:
            try:
                if had_value:
                    winreg.SetValueEx(
                        key, value_name, 0, int(previous_type), previous_value
                    )
                else:
                    try:
                        winreg.DeleteValue(key, value_name)
                    except OSError:
                        pass
            finally:
                try:
                    winreg.CloseKey(key)
                except OSError:
                    pass


def _dismiss_owned_pdf_reflow_prompt_once(process_id: int) -> bool:
    """Silently confirm the NUI dialog owned by our exact PDFREFLOW process.

    The NetUI prompt does not expose a normal Win32 button.  Sending a global
    Enter key used to work around that limitation, but it also stole focus,
    flashed the dialog and could play the Windows warning sound.  Hide the
    verified task-owned window first, then send Enter directly to its NetUI
    child and parent windows without foreground activation or global input.
    """

    if sys.platform != "win32" or process_id <= 0:
        return False
    try:
        import win32con
        import win32gui
        import win32process
    except ImportError:
        return False

    windows: list[int] = []
    try:
        win32gui.EnumWindows(lambda handle, _extra: windows.append(handle), None)
    except Exception:
        return False

    for window_handle in windows:
        try:
            _thread_id, owner_pid = win32process.GetWindowThreadProcessId(
                window_handle
            )
        except Exception:
            continue
        if int(owner_pid or 0) != int(process_id):
            continue
        try:
            class_name = str(win32gui.GetClassName(window_handle) or "")
            title = str(win32gui.GetWindowText(window_handle) or "").strip()
        except Exception:
            continue
        if class_name != "NUIDialog" or title.casefold() != "microsoft word":
            continue
        try:
            win32gui.ShowWindow(
                window_handle,
                int(getattr(win32con, "SW_HIDE", 0)),
            )
            win32gui.SendMessageTimeout(
                window_handle,
                int(getattr(win32con, "WM_COMMAND", 0x0111)),
                int(getattr(win32con, "IDOK", 1)),
                0,
                int(getattr(win32con, "SMTO_ABORTIFHUNG", 0x0002)),
                500,
            )
            child_windows: list[int] = []
            win32gui.EnumChildWindows(
                window_handle,
                lambda child, _extra: child_windows.append(child),
                None,
            )
            for target_handle in [*child_windows, window_handle]:
                for message, key_code, flags in (
                    (
                        int(getattr(win32con, "WM_KEYDOWN", 0x0100)),
                        int(getattr(win32con, "VK_RETURN", 0x0D)),
                        0,
                    ),
                    (
                        int(getattr(win32con, "WM_CHAR", 0x0102)),
                        13,
                        0,
                    ),
                    (
                        int(getattr(win32con, "WM_KEYUP", 0x0101)),
                        int(getattr(win32con, "VK_RETURN", 0x0D)),
                        0xC0000001,
                    ),
                ):
                    win32gui.SendMessageTimeout(
                        target_handle,
                        message,
                        key_code,
                        flags,
                        int(getattr(win32con, "SMTO_ABORTIFHUNG", 0x0002)),
                        500,
                    )
        except Exception:
            continue
        return True
    return False


def _registered_office_executable(prog_id: str) -> Path | None:
    """Return the verified Microsoft Office executable for *prog_id*.

    WPS can register compatibility ProgIDs such as ``Excel.Application``.  A
    ProgID/CLSID key alone therefore produces false positives and can make auto
    conversion wait for a COM server that will never start.  Validate both
    32-bit and 64-bit registry views and require both Microsoft's stable
    application CLSID and an existing Office executable with the expected
    filename.
    """

    if sys.platform != "win32":
        return None
    try:
        import winreg
    except ImportError:
        return None
    expected_name = _MICROSOFT_COM_EXECUTABLES.get(prog_id)
    expected_clsid = _MICROSOFT_COM_CLSIDS.get(prog_id)
    if expected_name is None or expected_clsid is None:
        return None

    access_modes: list[int] = []
    for attribute in ("KEY_WOW64_64KEY", "KEY_WOW64_32KEY"):
        flag = int(getattr(winreg, attribute, 0))
        if flag not in access_modes:
            access_modes.append(flag)
    if 0 not in access_modes:
        access_modes.append(0)

    def query_default(path: str, view: int) -> str:
        access = int(getattr(winreg, "KEY_READ", 0)) | view
        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, path, 0, access) as registry_key:
            value, _value_type = winreg.QueryValueEx(registry_key, None)
        return str(value).strip()

    for view in access_modes:
        try:
            clsid = query_default(f"{prog_id}\\CLSID", view)
            if clsid.casefold() != expected_clsid.casefold():
                continue
            command = query_default(f"CLSID\\{clsid}\\LocalServer32", view)
        except (OSError, ValueError, TypeError):
            continue
        expanded = os.path.expandvars(command).strip()
        if expanded.startswith('"'):
            closing_quote = expanded.find('"', 1)
            executable_text = expanded[1:closing_quote] if closing_quote > 1 else ""
        else:
            match = re.match(r"(?i)^(.+?\.exe)(?:\s|$)", expanded)
            executable_text = match.group(1).strip() if match else ""
        if not executable_text:
            continue
        executable = Path(executable_text)
        if executable.name.casefold() != expected_name or not executable.is_file():
            continue
        return executable.resolve()
    return None


def _registered_office_application(prog_id: str) -> bool:
    """Return whether *prog_id* resolves to the real Microsoft application."""

    return _registered_office_executable(prog_id) is not None


def _pywin32_available() -> bool:
    try:
        import win32com.client  # noqa: F401
    except ImportError:
        return False
    return True


def _detect_microsoft_components() -> dict[str, OfficeEngineStatus]:
    if sys.platform != "win32":
        return {
            key: OfficeEngineStatus(
                False, reason=f"Microsoft {label} COM 仅在 Windows 上可用"
            )
            for key, (_prog_id, label) in _MICROSOFT_COMPONENTS.items()
        }
    registered = {
        key: _registered_office_application(prog_id)
        for key, (prog_id, _label) in _MICROSOFT_COMPONENTS.items()
    }
    has_pywin32 = _pywin32_available()
    statuses: dict[str, OfficeEngineStatus] = {}
    for key, (_prog_id, label) in _MICROSOFT_COMPONENTS.items():
        if not registered[key]:
            statuses[key] = OfficeEngineStatus(
                False, reason=f"未检测到 Microsoft {label} COM 注册信息"
            )
        elif not has_pywin32:
            statuses[key] = OfficeEngineStatus(
                False, reason=f"检测到 Microsoft {label}，但缺少 pywin32 COM 运行库"
            )
        else:
            statuses[key] = OfficeEngineStatus(
                True, reason=f"已检测到 Microsoft {label} COM"
            )
    return statuses


def _aggregate_microsoft_status(
    components: Mapping[str, OfficeEngineStatus],
) -> OfficeEngineStatus:
    available_labels = [
        _MICROSOFT_COMPONENTS[key][1]
        for key in _MICROSOFT_COMPONENTS
        if components[key].available
    ]
    if available_labels:
        return OfficeEngineStatus(
            True,
            reason=f"已检测到 Microsoft Office 组件：{', '.join(available_labels)}",
        )
    reasons = "；".join(components[key].reason for key in _MICROSOFT_COMPONENTS)
    return OfficeEngineStatus(False, reason=reasons)


def _detect_com() -> OfficeEngineStatus:
    """Backward-compatible aggregate probe; component probes are authoritative."""

    return _aggregate_microsoft_status(_detect_microsoft_components())


def _detect_libreoffice() -> OfficeEngineStatus:
    candidates = [shutil.which("soffice"), shutil.which("libreoffice")]
    explicit = (
        os.environ.get("LAYOUTLOOM_LIBREOFFICE_PATH", "").strip()
        or os.environ.get("DOCUFORGE_LIBREOFFICE_PATH", "").strip()
    )
    if explicit:
        configured = Path(explicit).expanduser()
        candidates.append(
            configured / "program/soffice.exe" if configured.is_dir() else configured
        )
    bundle_root = getattr(sys, "_MEIPASS", None)
    if bundle_root:
        candidates.append(Path(bundle_root) / "libreoffice/program/soffice.exe")
    executable_root = Path(sys.executable).resolve().parent
    candidates.extend(
        [
            executable_root / "libreoffice/program/soffice.exe",
            executable_root / "_internal/libreoffice/program/soffice.exe",
        ]
    )
    project_root = Path(__file__).resolve().parents[2]
    candidates.append(
        project_root / "third_party/libreoffice/program/soffice.exe"
    )
    if sys.platform == "win32":
        candidates.extend(
            [
                Path(os.environ.get("ProgramFiles", r"C:\Program Files"))
                / "LibreOffice/program/soffice.exe",
                Path(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"))
                / "LibreOffice/program/soffice.exe",
            ]
        )
    for candidate in candidates:
        if candidate and Path(candidate).is_file():
            return OfficeEngineStatus(
                True, Path(candidate).resolve(), "已检测到 LibreOffice"
            )
    return OfficeEngineStatus(
        False, reason="未在 PATH 或标准安装目录找到 LibreOffice/soffice"
    )


def detect_office_engines() -> dict[str, OfficeEngineStatus]:
    """Probe conversion engines without launching Word, Excel or PowerPoint."""

    components = _detect_microsoft_components()
    return {
        "microsoft_office": _aggregate_microsoft_status(components),
        **components,
        "libreoffice": _detect_libreoffice(),
    }


def _component_for_source(source: Path) -> str:
    try:
        return _COMPONENT_FOR_SUFFIX[source.suffix.lower()]
    except KeyError as exc:
        raise ValidationError(
            f"不支持的 Office 输入格式：{source.suffix or '无扩展名'}"
        ) from exc


def _convert_libreoffice(
    source: Path,
    target: Path,
    target_format: str,
    executable: Path,
    timeout: float,
) -> None:
    from docuforge.runner import cancellation_callback, check_cancelled

    with tempfile.TemporaryDirectory(prefix="docuforge-lo-") as temporary_name:
        temporary_dir = Path(temporary_name)
        profile_dir = temporary_dir / "profile"
        profile_dir.mkdir()
        command = [
            str(executable),
            "--headless",
            "--nologo",
            "--nodefault",
            "--norestore",
            f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
            "--convert-to",
            target_format,
            "--outdir",
            str(temporary_dir),
            str(source),
        ]
        process = None
        try:
            process = subprocess.Popen(
                command,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                creationflags=(
                    subprocess.CREATE_NO_WINDOW if sys.platform == "win32" else 0
                ),
            )
            started = time.monotonic()
            with cancellation_callback(
                lambda: process.terminate() if process.poll() is None else None
            ):
                while True:
                    remaining = timeout - (time.monotonic() - started)
                    if remaining <= 0:
                        process.kill()
                        process.communicate()
                        raise subprocess.TimeoutExpired(command, timeout)
                    try:
                        stdout, stderr = process.communicate(timeout=min(0.15, remaining))
                        break
                    except subprocess.TimeoutExpired:
                        check_cancelled("任务已取消；正在终止 LibreOffice 转换")
            result = subprocess.CompletedProcess(
                command,
                process.returncode,
                stdout,
                stderr,
            )
        except subprocess.TimeoutExpired as exc:
            raise MissingEngineError(f"LibreOffice 转换超时（{timeout:g} 秒）") from exc
        except BaseException:
            if process is not None and process.poll() is None:
                try:
                    process.terminate()
                    process.wait(timeout=1.0)
                except (OSError, subprocess.TimeoutExpired):
                    try:
                        process.kill()
                        process.wait(timeout=1.0)
                    except (OSError, subprocess.TimeoutExpired):
                        pass
            raise
        candidates = list(temporary_dir.glob(f"{source.stem}.*"))
        generated = next(
            (
                path
                for path in candidates
                if path.suffix.lower() == f".{target_format.lower()}"
            ),
            None,
        )
        if (
            result.returncode != 0
            or generated is None
            or not generated.is_file()
            or generated.stat().st_size == 0
        ):
            diagnostic = (result.stderr or result.stdout or "未生成目标文件").strip()
            raise MissingEngineError(f"LibreOffice 转换失败：{diagnostic[:500]}")
        with atomic_output(target) as temporary:
            shutil.copy2(generated, temporary)
            if not temporary.is_file() or temporary.stat().st_size == 0:
                raise MissingEngineError("LibreOffice 未生成有效输出文件")


def _disable_automation_macros(application: Any, display_name: str) -> None:
    try:
        application.AutomationSecurity = 3  # msoAutomationSecurityForceDisable
    except Exception as exc:
        raise MissingEngineError(
            f"无法为 Microsoft {display_name} 强制禁用宏，已拒绝打开文件：{exc}"
        ) from exc


def _validate_microsoft_com_application(
    application: Any, prog_id: str, display_name: str
) -> None:
    """Reject compatibility ProgIDs that resolve to another Office family."""

    if not _registered_office_application(prog_id):
        raise MissingEngineError(
            f"{prog_id} 未指向 Microsoft {display_name} 的正式 COM 注册信息；"
            "可能被兼容软件接管，已拒绝继续"
        )
    expected_executable = _MICROSOFT_COM_EXECUTABLES[prog_id]
    expected_name = _MICROSOFT_COM_NAMES[prog_id]
    try:
        application_name = str(application.Name).strip()
        application_path = Path(str(application.Path).strip()).expanduser()
    except Exception as exc:
        raise MissingEngineError(
            f"无法验证 Microsoft {display_name} COM 组件身份：{exc}"
        ) from exc
    normalized_name = application_name.casefold().replace(" ", "")
    if (
        expected_name not in normalized_name
        or "wps" in normalized_name
        or "kingsoft" in normalized_name
    ):
        raise MissingEngineError(
            f"{prog_id} 实际启动的是“{application_name or '未知程序'}”，"
            f"不是 Microsoft {display_name}；已拒绝继续"
        )
    executable = application_path / expected_executable
    if not executable.is_file():
        raise MissingEngineError(
            f"Microsoft {display_name} COM 返回的程序路径无效：{executable}"
        )


def _convert_com(
    source: Path,
    target: Path,
    target_format: str,
    *,
    excel_pdf_layout: str = "smart",
    excel_pdf_paper: str = "auto",
    excel_pdf_orientation: str = "auto",
    excel_pdf_margin: str = "auto",
    ownership_guard: Callable[[Any, str], bool | int | None] | None = None,
) -> None:
    suffix = source.suffix.lower()
    if suffix in _PPT_EXTENSIONS | {".ppt"} and target_format in {"png", "jpg"}:
        raise ValidationError(
            "PowerPoint 的 PNG/JPG SaveAs 会生成图片目录而非单个目标文件；"
            "请使用“PPT 转图片序列”"
        )
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:  # pragma: no cover - platform-specific
        raise MissingEngineError("缺少 pywin32，无法调用 Microsoft Office COM") from exc
    pythoncom.CoInitialize()
    application = None
    document = None
    staged_source_directory: tempfile.TemporaryDirectory[str] | None = None
    application_owned = ownership_guard is None

    def claim_application(prog_id: str, display_name: str) -> None:
        nonlocal application_owned
        if ownership_guard is None:
            application_owned = True
            return
        # Until the parent confirms the exact PID/path/creation tuple, never
        # assume a DispatchEx result is safe to Quit or terminate.
        application_owned = False
        ownership = ownership_guard(application, prog_id)
        if not bool(ownership):
            raise MissingEngineError(
                f"Microsoft {display_name} 复用了无法确认归属的现有进程；"
                "为保护用户已打开的文档，已停止自动化"
            )
        application_owned = True

    try:
        if suffix == ".pdf":
            if target_format != "docx":
                raise ValidationError("Microsoft Word PDF Reflow 仅支持输出 DOCX")
            application = win32com.client.DispatchEx("Word.Application")
            _validate_microsoft_com_application(application, "Word.Application", "Word")
            claim_application("Word.Application", "Word")
            _disable_automation_macros(application, "Word")
            application.Visible = False
            application.DisplayAlerts = 0
            try:
                application.Options.ConfirmConversions = False
            except Exception:
                pass
            # A previous interrupted PDF Reflow can make Word remember the
            # original path as a failed document and show an unsuppressible
            # recovery prompt.  A byte-identical copy at a fresh private path
            # avoids inheriting that stale per-file recovery state.
            staged_source_directory = tempfile.TemporaryDirectory(
                prefix="layoutloom-word-reflow-"
            )
            staged_source = (
                Path(staged_source_directory.name) / f"input-{uuid.uuid4().hex}.pdf"
            )
            shutil.copyfile(source, staged_source)
            document = application.Documents.Open(
                str(staged_source),
                ConfirmConversions=False,
                ReadOnly=True,
                AddToRecentFiles=False,
                Revert=False,
                Visible=False,
                OpenAndRepair=False,
                NoEncodingDialog=True,
            )
            document.SaveAs2(str(target), FileFormat=16, AddToRecentFiles=False)
        elif suffix in _WORD_EXTENSIONS | {".doc", ".rtf"}:
            application = win32com.client.DispatchEx("Word.Application")
            _validate_microsoft_com_application(application, "Word.Application", "Word")
            claim_application("Word.Application", "Word")
            _disable_automation_macros(application, "Word")
            application.Visible = False
            application.DisplayAlerts = 0
            try:
                application.Options.ConfirmConversions = False
                application.Options.SaveNormalPrompt = False
            except Exception:
                pass
            document = application.Documents.Open(
                str(source),
                ConfirmConversions=False,
                ReadOnly=True,
                AddToRecentFiles=False,
                Revert=False,
                Visible=False,
                OpenAndRepair=False,
                NoEncodingDialog=True,
            )
            if target_format == "pdf":
                document.ExportAsFixedFormat(str(target), 17)
            else:
                formats = {"docx": 16, "txt": 2, "html": 10}
                if target_format not in formats:
                    raise ValidationError(f"Word COM 不支持此目标格式：{target_format}")
                document.SaveAs2(str(target), FileFormat=formats[target_format])
        elif suffix in _EXCEL_EXTENSIONS | {".xls", ".csv"}:
            application = win32com.client.DispatchEx("Excel.Application")
            _validate_microsoft_com_application(
                application, "Excel.Application", "Excel"
            )
            claim_application("Excel.Application", "Excel")
            _disable_automation_macros(application, "Excel")
            application.Visible = False
            application.DisplayAlerts = False
            for property_name, value in (
                ("ScreenUpdating", False),
                ("EnableEvents", False),
                ("AskToUpdateLinks", False),
            ):
                try:
                    setattr(application, property_name, value)
                except Exception:
                    pass
            document = application.Workbooks.Open(
                str(source),
                ReadOnly=True,
                UpdateLinks=0,
                AddToMru=False,
                IgnoreReadOnlyRecommended=True,
                Notify=False,
            )
            if target_format == "pdf":
                from .excel_pdf_layout import prepare_excel_workbook_for_pdf

                prepare_excel_workbook_for_pdf(
                    document,
                    application,
                    layout=excel_pdf_layout,
                    paper=excel_pdf_paper,
                    orientation=excel_pdf_orientation,
                    margin=excel_pdf_margin,
                )
                document.ExportAsFixedFormat(0, str(target))
            else:
                formats = {"xlsx": 51, "xls": 56, "csv": 62, "xml": 46, "txt": 42}
                if target_format not in formats:
                    raise ValidationError(
                        f"Excel COM 不支持此目标格式：{target_format}"
                    )
                document.SaveAs(str(target), FileFormat=formats[target_format])
        elif suffix in _PPT_EXTENSIONS | {".ppt"}:
            application = win32com.client.DispatchEx("PowerPoint.Application")
            _validate_microsoft_com_application(
                application, "PowerPoint.Application", "PowerPoint"
            )
            claim_application("PowerPoint.Application", "PowerPoint")
            _disable_automation_macros(application, "PowerPoint")
            try:
                application.Visible = False
                application.DisplayAlerts = 1  # ppAlertsNone
            except Exception:
                pass
            document = application.Presentations.Open(
                str(source), WithWindow=False, ReadOnly=True, Untitled=False
            )
            formats = {"pdf": 32, "pptx": 24}
            if target_format not in formats:
                raise ValidationError(
                    f"PowerPoint COM 不支持此目标格式：{target_format}"
                )
            document.SaveAs(str(target), formats[target_format])
        else:
            raise ValidationError(
                f"Microsoft Office COM 不支持此输入格式：{source.suffix}"
            )
        if not target.is_file() or target.stat().st_size == 0:
            raise MissingEngineError(
                f"Microsoft Office COM 未生成有效输出文件：{target.name}"
            )
    except (ValidationError, MissingEngineError):
        raise
    except Exception as exc:  # pragma: no cover - platform-specific
        raise MissingEngineError(f"Microsoft Office COM 转换失败：{exc}") from exc
    finally:  # pragma: no cover - platform-specific
        if document is not None:
            try:
                if suffix in _PPT_EXTENSIONS | {".ppt"}:
                    document.Close()
                else:
                    document.Close(False)
            except Exception:
                pass
        if application is not None and application_owned:
            try:
                application.Quit()
            except Exception:
                pass
        if staged_source_directory is not None:
            try:
                staged_source_directory.cleanup()
            except OSError:
                pass
        pythoncom.CoUninitialize()


def _send_com_worker_message(connection: Any, message: Mapping[str, Any]) -> None:
    try:
        connection.send(dict(message))
    except (BrokenPipeError, EOFError, OSError):
        pass


def _convert_com_worker_entry(
    connection: Any,
    source: str,
    target: str,
    target_format: str,
    options: Mapping[str, Any],
) -> None:
    """Spawn-safe worker for one isolated Microsoft Office COM conversion."""

    def confirm_ownership(application: Any, prog_id: str) -> bool | int:
        process_id = _office_application_pid(application)
        expected_name = _MICROSOFT_COM_EXECUTABLES.get(prog_id, "")
        expected_path = ""
        if expected_name:
            try:
                expected_path = str(
                    (Path(str(application.Path).strip()) / expected_name).resolve()
                )
            except Exception:
                expected_path = ""
        _send_com_worker_message(
            connection,
            {
                "type": "office_process",
                "pid": process_id,
                "prog_id": prog_id,
                "executable": expected_path,
            },
        )
        try:
            if not connection.poll(15.0):
                return False
            response = connection.recv()
        except (EOFError, OSError):
            return False
        if not (
            isinstance(response, Mapping)
            and response.get("type") == "ownership"
            and response.get("approved") is True
        ):
            return False
        try:
            approved_pid = int(response.get("pid") or 0)
        except (TypeError, ValueError):
            approved_pid = 0
        return approved_pid if approved_pid > 0 else True

    try:
        _convert_com(
            Path(source),
            Path(target),
            target_format,
            excel_pdf_layout=str(options.get("excel_pdf_layout", "smart")),
            excel_pdf_paper=str(options.get("excel_pdf_paper", "auto")),
            excel_pdf_orientation=str(options.get("excel_pdf_orientation", "auto")),
            excel_pdf_margin=str(options.get("excel_pdf_margin", "auto")),
            ownership_guard=confirm_ownership,
        )
    except ValidationError as exc:
        _send_com_worker_message(
            connection,
            {"type": "result", "ok": False, "kind": "validation", "error": str(exc)},
        )
    except FileExistsError as exc:
        _send_com_worker_message(
            connection,
            {"type": "result", "ok": False, "kind": "exists", "error": str(exc)},
        )
    except MissingEngineError as exc:
        _send_com_worker_message(
            connection,
            {"type": "result", "ok": False, "kind": "engine", "error": str(exc)},
        )
    except Exception as exc:  # pragma: no cover - defensive worker boundary
        _send_com_worker_message(
            connection,
            {
                "type": "result",
                "ok": False,
                "kind": "engine",
                "error": f"Microsoft Office COM 子进程异常：{exc}",
            },
        )
    else:
        _send_com_worker_message(connection, {"type": "result", "ok": True})
    finally:
        try:
            connection.close()
        except OSError:
            pass


def _stop_com_worker(
    process: Any,
    identity: _OfficeProcessIdentity | None,
    auxiliary_identities: Sequence[_OfficeProcessIdentity] = (),
) -> None:
    """Stop only the worker and Office instance whose ownership was confirmed."""

    for auxiliary in auxiliary_identities:
        _terminate_owned_office_process(auxiliary)
    if identity is not None:
        _terminate_owned_office_process(identity)
    try:
        process.join(1.0)
    except Exception:
        pass
    try:
        alive = bool(process.is_alive())
    except Exception:
        alive = False
    if alive:
        try:
            process.terminate()
        except Exception:
            pass
        try:
            process.join(1.0)
        except Exception:
            pass
    try:
        alive = bool(process.is_alive())
    except Exception:
        alive = False
    if alive and hasattr(process, "kill"):
        try:
            process.kill()
            process.join(1.0)
        except Exception:
            pass
    _wait_for_owned_office_exit(identity, timeout=0.5)
    for auxiliary in auxiliary_identities:
        _wait_for_owned_office_exit(auxiliary, timeout=0.5)


def _convert_com_supervised(
    source: Path,
    target: Path,
    target_format: str,
    *,
    timeout: float,
    excel_pdf_layout: str = "smart",
    excel_pdf_paper: str = "auto",
    excel_pdf_orientation: str = "auto",
    excel_pdf_margin: str = "auto",
    component_override: str | None = None,
    auxiliary_executables: Sequence[Path] = (),
) -> None:
    """Run COM in an isolated process so timeout and cancellation are enforceable."""

    import multiprocessing

    from docuforge.runner import check_cancelled

    component = component_override or _component_for_source(source)
    if component not in _MICROSOFT_COMPONENTS:
        raise ValidationError(f"未知 Microsoft Office 组件：{component}")
    prog_id, display_name = _MICROSOFT_COMPONENTS[component]
    expected_executable = _MICROSOFT_COM_EXECUTABLES[prog_id]
    before = _windows_process_snapshot(expected_executable)
    auxiliary_paths = tuple(path.expanduser().resolve() for path in auxiliary_executables)
    auxiliary_before = {
        path: _windows_process_snapshot(path.name) for path in auxiliary_paths
    }
    context = multiprocessing.get_context("spawn")
    parent_connection, child_connection = context.Pipe(duplex=True)
    process = context.Process(
        target=_convert_com_worker_entry,
        args=(
            child_connection,
            str(source),
            str(target),
            target_format,
            {
                "excel_pdf_layout": excel_pdf_layout,
                "excel_pdf_paper": excel_pdf_paper,
                "excel_pdf_orientation": excel_pdf_orientation,
                "excel_pdf_margin": excel_pdf_margin,
            },
        ),
        name=f"docuforge-{display_name.casefold()}-com",
        daemon=False,
    )
    try:
        process.start()
    except Exception as exc:
        parent_connection.close()
        child_connection.close()
        raise MissingEngineError(
            f"无法启动 Microsoft {display_name} 隔离转换进程：{exc}"
        ) from exc
    child_connection.close()

    deadline = time.monotonic() + timeout
    owned_identity: _OfficeProcessIdentity | None = None
    owned_auxiliary: dict[Path, _OfficeProcessIdentity] = {}
    dismissed_auxiliary_prompts: set[int] = set()
    result: Mapping[str, Any] | None = None
    try:
        while result is None:
            check_cancelled("任务已取消；正在终止 Microsoft Office 转换")
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                hint = ""
                if component == "microsoft_word" and source.suffix.lower() == ".pdf":
                    hint = (
                        "；Word PDF Reflow 可能被提示窗口、异常 PDF 文字层或残留进程阻塞，"
                        "请关闭 Word 后重试，或改用页织工坊内置 PDF 转 Word 模式"
                    )
                raise MissingEngineError(
                    f"Microsoft {display_name} COM 转换超时（{timeout:g} 秒）{hint}"
                )
            try:
                has_message = parent_connection.poll(min(0.1, remaining))
            except (EOFError, OSError):
                has_message = False
            if has_message:
                try:
                    message = parent_connection.recv()
                except (EOFError, OSError):
                    message = None
                if isinstance(message, Mapping):
                    message_type = message.get("type")
                    if message_type == "office_process":
                        try:
                            process_id = int(message.get("pid") or 0)
                        except (TypeError, ValueError):
                            process_id = 0
                        reported_prog_id = str(message.get("prog_id") or "")
                        reported_executable = Path(
                            str(message.get("executable") or "")
                        )
                        identity = None
                        if (
                            reported_prog_id == prog_id
                            and reported_executable.name.casefold()
                            == expected_executable.casefold()
                        ):
                            identity = _new_owned_office_process(
                                before,
                                expected_executable=reported_executable,
                                reported_pid=process_id or None,
                            )
                        approved = identity is not None
                        if approved:
                            owned_identity = identity
                        try:
                            parent_connection.send(
                                {
                                    "type": "ownership",
                                    "approved": approved,
                                    "pid": identity.pid if identity is not None else None,
                                }
                            )
                        except (BrokenPipeError, EOFError, OSError):
                            pass
                    elif message_type == "result":
                        result = message
            if owned_identity is not None:
                for executable_path in auxiliary_paths:
                    if executable_path in owned_auxiliary:
                        continue
                    auxiliary = _new_owned_office_process(
                        auxiliary_before[executable_path],
                        expected_executable=executable_path,
                    )
                    if auxiliary is not None:
                        owned_auxiliary[executable_path] = auxiliary
                for executable_path, auxiliary in owned_auxiliary.items():
                    if (
                        executable_path.name.casefold() == "pdfreflow.exe"
                        and auxiliary.pid not in dismissed_auxiliary_prompts
                        and _dismiss_owned_pdf_reflow_prompt_once(auxiliary.pid)
                    ):
                        dismissed_auxiliary_prompts.add(auxiliary.pid)
            if result is None and not process.is_alive():
                # Drain a final message that may have arrived immediately before exit.
                try:
                    if parent_connection.poll(0):
                        final_message = parent_connection.recv()
                        if (
                            isinstance(final_message, Mapping)
                            and final_message.get("type") == "result"
                        ):
                            result = final_message
                except (EOFError, OSError):
                    pass
                if result is None:
                    break
    except BaseException:
        # PDFREFLOW can appear shortly before a timeout, after the regular
        # polling pass last looked for auxiliaries.  Re-scan once while the
        # exact Word instance is still owned, then terminate only the unique
        # new executable/path/creation tuple that belongs to this operation.
        if owned_identity is not None:
            for executable_path in auxiliary_paths:
                if executable_path in owned_auxiliary:
                    continue
                auxiliary = _new_owned_office_process(
                    auxiliary_before[executable_path],
                    expected_executable=executable_path,
                )
                if auxiliary is not None:
                    owned_auxiliary[executable_path] = auxiliary
        if process.is_alive() or owned_identity is not None:
            _stop_com_worker(
                process, owned_identity, tuple(owned_auxiliary.values())
            )
        raise
    finally:
        try:
            parent_connection.close()
        except OSError:
            pass

    try:
        process.join(2.0)
    except Exception:
        pass
    if process.is_alive():
        _stop_com_worker(process, owned_identity, tuple(owned_auxiliary.values()))
        raise MissingEngineError(f"Microsoft {display_name} COM 转换子进程未正常退出")
    _wait_for_owned_office_exit(owned_identity)
    for auxiliary in owned_auxiliary.values():
        _wait_for_owned_office_exit(auxiliary)
    if result is None:
        raise MissingEngineError(f"Microsoft {display_name} COM 转换进程意外退出")
    if bool(result.get("ok")):
        if not target.is_file() or target.stat().st_size == 0:
            raise MissingEngineError(
                f"Microsoft Office COM 未生成有效输出文件：{target.name}"
            )
        return
    error = str(result.get("error") or "Microsoft Office COM 转换失败")
    kind = str(result.get("kind") or "engine")
    if kind == "validation":
        raise ValidationError(error)
    if kind == "exists":
        raise FileExistsError(error)
    raise MissingEngineError(error)


def microsoft_pdf_to_docx(
    source: PathLike,
    target: PathLike,
    *,
    timeout: float = 900,
) -> Path:
    """Convert a digital PDF with Microsoft Word's native PDF Reflow engine.

    This is intentionally an exact-target adapter rather than a generic Office
    conversion engine.  The caller requests Word PDF Reflow explicitly; the
    conversion pipeline then validates the generated DOCX without mixing it
    with LayoutLoom's built-in reconstruction result.
    """

    source_path = Path(source).expanduser().resolve()
    target_path = Path(target).expanduser().resolve()
    if not source_path.is_file():
        raise ValidationError(f"文件不存在：{source_path}")
    if source_path.suffix.lower() != ".pdf":
        raise ValidationError("Microsoft Word PDF Reflow 仅支持 PDF 输入")
    if target_path.suffix.lower() != ".docx":
        raise ValidationError("Microsoft Word PDF Reflow 目标必须是 DOCX")
    try:
        timeout_value = float(timeout)
    except (TypeError, ValueError, OverflowError) as exc:
        raise ValidationError("timeout 必须是有限且大于 0 的秒数") from exc
    if not math.isfinite(timeout_value) or timeout_value <= 0:
        raise ValidationError("timeout 必须是有限且大于 0 的秒数")

    status = detect_office_engines().get("microsoft_word")
    if status is None or not status.available:
        reason = status.reason if status is not None else "未检测到 Microsoft Word"
        raise MissingEngineError(reason)
    word_executable = _registered_office_executable("Word.Application")
    if word_executable is None:
        raise MissingEngineError("Microsoft Word COM 注册路径无效")
    target_path.parent.mkdir(parents=True, exist_ok=True)
    if target_path.exists():
        raise FileExistsError(f"目标文件已存在：{target_path}")
    reflow_executable = word_executable.parent / "PDFREFLOW.EXE"
    auxiliaries = (reflow_executable,) if reflow_executable.is_file() else ()
    with _word_pdf_reflow_prompt_suppressed():
        _convert_com_supervised(
            source_path,
            target_path,
            "docx",
            timeout=timeout_value,
            component_override="microsoft_word",
            auxiliary_executables=auxiliaries,
        )
    if not target_path.is_file() or target_path.stat().st_size == 0:
        raise MissingEngineError("Microsoft Word PDF Reflow 未生成有效 DOCX")
    return target_path


def _select_conversion_engine(
    source: Path,
    requested: str,
    office_statuses: Mapping[str, OfficeEngineStatus] | None = None,
    wps_statuses: Mapping[str, Any] | None = None,
) -> str:
    """Select an engine for the source's actual Word/Excel/PowerPoint family."""

    engine_name = requested.lower().strip()
    if engine_name not in {
        "auto",
        "microsoft_office",
        "com",
        "wps",
        "libreoffice",
    }:
        raise ValidationError(
            "engine 必须是 auto、microsoft_office/com、wps 或 libreoffice"
        )
    component = _component_for_source(source)
    statuses = dict(office_statuses or detect_office_engines())
    microsoft = statuses.get(component)
    libreoffice = statuses.get("libreoffice")
    if microsoft is None or libreoffice is None:
        raise MissingEngineError("Office 组件探测结果不完整")

    if engine_name in {"microsoft_office", "com"}:
        if not microsoft.available:
            raise MissingEngineError(microsoft.reason)
        return "microsoft_office"
    if engine_name == "libreoffice":
        if not libreoffice.available:
            raise MissingEngineError(libreoffice.reason)
        return "libreoffice"

    if wps_statuses is None and engine_name in {"auto", "wps"}:
        try:
            from .wps import detect_wps_engines

            wps_statuses = {
                str(key): value for key, value in detect_wps_engines().items()
            }
        except Exception:
            wps_statuses = {}
    wps_kind = _WPS_KIND_FOR_COMPONENT[component]
    wps = (wps_statuses or {}).get(wps_kind)
    if engine_name == "wps":
        if wps is None or not wps.available:
            reason = getattr(wps, "reason", "未检测到对应的 WPS Office 组件")
            raise MissingEngineError(reason)
        return "wps"

    # Auto is deliberately family-specific.  A registered Word component must
    # never make an Excel or PowerPoint source select Microsoft Office.
    if wps is not None and wps.available:
        return "wps"
    if microsoft.available:
        return "microsoft_office"
    if libreoffice.available:
        return "libreoffice"
    wps_reason = getattr(wps, "reason", "未检测到对应的 WPS Office 组件")
    raise MissingEngineError(
        "没有可用于当前文档类型的 Office 转换引擎。"
        f"WPS：{wps_reason}；Microsoft：{microsoft.reason}；"
        f"LibreOffice：{libreoffice.reason}"
    )


def _require_conversion_outputs(
    outputs: Iterable[PathLike], engine_label: str
) -> list[Path]:
    paths = [Path(path) for path in outputs]
    if not paths:
        raise MissingEngineError(f"{engine_label} 未返回任何输出文件")
    invalid = next(
        (path for path in paths if not path.is_file() or path.stat().st_size == 0),
        None,
    )
    if invalid is not None:
        raise MissingEngineError(f"{engine_label} 未生成有效输出文件：{invalid}")
    return paths


def _convert_with_selected_engine(
    source: Path,
    output_dir: PathLike,
    target_format: str,
    *,
    selected: str,
    statuses: Mapping[str, OfficeEngineStatus],
    overwrite: bool,
    timeout: float,
    excel_pdf_layout: str,
    excel_pdf_paper: str,
    excel_pdf_orientation: str,
    excel_pdf_margin: str,
) -> list[Path]:
    if selected == "wps":
        from .wps import convert_with_wps

        outputs = convert_with_wps(
            source,
            output_dir,
            target_format,
            overwrite=overwrite,
            excel_pdf_layout=excel_pdf_layout,
            excel_pdf_paper=excel_pdf_paper,
            excel_pdf_orientation=excel_pdf_orientation,
            excel_pdf_margin=excel_pdf_margin,
        )
        return _require_conversion_outputs(outputs, "WPS Office")

    target = _output_path(
        source, output_dir, suffix=f".{target_format}", overwrite=overwrite
    )
    if selected == "microsoft_office":
        # COM needs the final extension but writes its own output; stage beside target.
        with atomic_output(target) as temporary:
            temporary.unlink(missing_ok=True)
            _convert_com_supervised(
                source,
                temporary,
                target_format,
                timeout=timeout,
                excel_pdf_layout=excel_pdf_layout,
                excel_pdf_paper=excel_pdf_paper,
                excel_pdf_orientation=excel_pdf_orientation,
                excel_pdf_margin=excel_pdf_margin,
            )
            if not temporary.is_file() or temporary.stat().st_size == 0:
                raise MissingEngineError("Microsoft Office COM 未生成有效输出文件")
    elif selected == "libreoffice":
        executable = statuses["libreoffice"].executable
        if executable is None:  # defensive; status invariant
            raise MissingEngineError("LibreOffice 可执行文件路径无效")
        _convert_libreoffice(source, target, target_format, executable, timeout)
    else:  # defensive; all public callers validate the engine first
        raise ValidationError(f"未知 Office 转换引擎：{selected}")
    return _require_conversion_outputs(
        [target],
        {
            "microsoft_office": "Microsoft Office",
            "libreoffice": "LibreOffice",
        }[selected],
    )


def convert_with_office(
    source: PathLike,
    output_dir: PathLike,
    target_format: str = "pdf",
    *,
    engine: str = "auto",
    overwrite: bool = False,
    timeout: float = 180,
    excel_pdf_layout: str = "smart",
    excel_pdf_paper: str = "auto",
    excel_pdf_orientation: str = "auto",
    excel_pdf_margin: str = "auto",
) -> list[Path]:
    """Convert with a real Office renderer; never simulates high fidelity.

    ``engine='auto'`` checks the source document's own component and prefers
    WPS Office, then Microsoft Office, then LibreOffice.  Installing Word alone
    therefore never makes an Excel workbook select Microsoft Office COM.
    """

    from docuforge.runner import check_cancelled

    check_cancelled("任务已取消；已完成的文件会保留")
    source_path = Path(source).expanduser().resolve()
    if not source_path.is_file():
        raise ValidationError(f"文件不存在：{source_path}")
    format_name = target_format.lower().lstrip(".")
    if not re.fullmatch(r"[a-z0-9]+", format_name):
        raise ValidationError(f"无效目标格式：{target_format}")
    try:
        timeout_value = float(timeout)
    except (TypeError, ValueError, OverflowError) as exc:
        raise ValidationError("timeout 必须是有限且大于 0 的秒数") from exc
    if not math.isfinite(timeout_value) or timeout_value <= 0:
        raise ValidationError("timeout 必须是有限且大于 0 的秒数")
    timeout = timeout_value
    engine_name = engine.lower().strip()
    if engine_name not in {
        "auto",
        "microsoft_office",
        "com",
        "wps",
        "libreoffice",
    }:
        raise ValidationError(
            "engine 必须是 auto、microsoft_office/com、wps 或 libreoffice"
        )
    component = _component_for_source(source_path)
    if component == "microsoft_excel" and format_name == "pdf":
        from .excel_pdf_layout import normalize_excel_pdf_options

        (
            excel_pdf_layout,
            excel_pdf_paper,
            excel_pdf_orientation,
            excel_pdf_margin,
        ) = normalize_excel_pdf_options(
            excel_pdf_layout,
            excel_pdf_paper,
            excel_pdf_orientation,
            excel_pdf_margin,
        )
    statuses = detect_office_engines()
    conversion_options = {
        "overwrite": overwrite,
        "timeout": timeout,
        "excel_pdf_layout": excel_pdf_layout,
        "excel_pdf_paper": excel_pdf_paper,
        "excel_pdf_orientation": excel_pdf_orientation,
        "excel_pdf_margin": excel_pdf_margin,
    }
    if engine_name != "auto":
        selected = _select_conversion_engine(source_path, engine_name, statuses)
        return _convert_with_selected_engine(
            source_path,
            output_dir,
            format_name,
            selected=selected,
            statuses=statuses,
            **conversion_options,
        )

    candidates: list[str] = []
    wps_statuses: Mapping[str, Any] = {}
    try:
        from .wps import detect_wps_engines

        wps_statuses = {str(key): value for key, value in detect_wps_engines().items()}
    except Exception:
        wps_statuses = {}
    wps = wps_statuses.get(_WPS_KIND_FOR_COMPONENT[component])
    if wps is not None and bool(getattr(wps, "available", False)):
        candidates.append("wps")
    if statuses[component].available:
        candidates.append("microsoft_office")
    if statuses["libreoffice"].available:
        candidates.append("libreoffice")
    if not candidates:
        # Reuse the detailed family-specific diagnostic.
        _select_conversion_engine(
            source_path, "auto", statuses, wps_statuses=wps_statuses
        )
        raise MissingEngineError("没有可用的 Office 转换引擎")  # pragma: no cover

    failures: list[str] = []
    labels = {
        "microsoft_office": "Microsoft Office",
        "wps": "WPS Office",
        "libreoffice": "LibreOffice",
    }
    for selected in candidates:
        try:
            return _convert_with_selected_engine(
                source_path,
                output_dir,
                format_name,
                selected=selected,
                statuses=statuses,
                **conversion_options,
            )
        except MissingEngineError as exc:
            failures.append(f"{labels[selected]}：{exc}")
    raise MissingEngineError("自动转换失败；" + "；".join(failures))


def office_to_pdf(
    source: PathLike,
    output_dir: PathLike,
    *,
    engine: str = "auto",
    overwrite: bool = False,
    timeout: float = 180,
    excel_pdf_layout: str = "smart",
    excel_pdf_paper: str = "auto",
    excel_pdf_orientation: str = "auto",
    excel_pdf_margin: str = "auto",
) -> list[Path]:
    return convert_with_office(
        source,
        output_dir,
        "pdf",
        engine=engine,
        overwrite=overwrite,
        timeout=timeout,
        excel_pdf_layout=excel_pdf_layout,
        excel_pdf_paper=excel_pdf_paper,
        excel_pdf_orientation=excel_pdf_orientation,
        excel_pdf_margin=excel_pdf_margin,
    )


# Short aliases used by operation registries.
excel_sort = excel_sort_rows
excel_filter = excel_filter_rows
excel_remove_blank = excel_remove_blank_rows_columns
excel_replace = excel_replace_text
excel_delete_formulas = excel_formulas_to_values
ppt_add_text_watermark = ppt_add_watermark
convert_office = convert_with_office


__all__ = [
    name
    for name in globals()
    if (
        name.startswith("word_")
        or name.startswith("excel_")
        or name.startswith("ppt_")
        or name
        in {
            "detect_office_engines",
            "convert_with_office",
            "convert_office",
            "office_to_pdf",
            "microsoft_pdf_to_docx",
            "OfficeEngineStatus",
        }
    )
]
