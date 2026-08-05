from __future__ import annotations

import hashlib
import math
import os
import tempfile
import time
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

try:  # Keep the public catalog importable when optional video CV wheels are absent.
    import cv2  # type: ignore[import-not-found]
    import numpy as np  # type: ignore[import-not-found]
except Exception:  # pragma: no cover - exercised through the capability probe
    cv2 = None
    np = None

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from ..models import MissingEngineError, ValidationError
from ..runner import check_cancelled, report_progress
from ..utils import atomic_output
from .image_enhancement import enhance_bgr, multiframe_fuse

PathLike = str | Path
RGBColor = tuple[int, int, int]
BGRColor = tuple[int, int, int]

_SCAN_INTERVALS = {
    "accurate": 0.5,
    "balanced": 1.0,
    "fast": 2.0,
}
_CHANGE_THRESHOLDS = {
    "conservative": (0.145, 0.115, 17),
    "balanced": (0.105, 0.085, 13),
    "sensitive": (0.075, 0.065, 10),
}
_MAX_SCAN_FRAMES = 6000
_MAX_CANDIDATES_PER_SLIDE = 16
_MAX_REFINED_CANDIDATES_PER_SLIDE = 32
_TRANSITION_REFINE_SECONDS = 2.0
_TRANSITION_REFINE_FPS = 8.0
_TRANSITION_TAIL_SECONDS = 1.25
_STABLE_WINDOW_CANDIDATES = 3
_STABLE_MIN_SECONDS = 0.55
_TERMINAL_STABLE_MIN_SECONDS = 0.32
_TERMINAL_STABLE_MIN_CANDIDATES = 4
_STABLE_DENSE_GAP_SECONDS = 0.30
_STABLE_MAX_GAP_SECONDS = 1.05
_POST_BOUNDARY_SETTLE_SECONDS = 0.45
_REFINED_ANALYSIS_WIDTH = 320
_MIN_SEGMENT_SECONDS = 1.0
_REPORT_SCHEMA = "docuforge.video-slides.v1"
_SHORT_PREVIEW_MAX_SECONDS = 2.25
_FORMAL_OCCURRENCE_MIN_SECONDS = 4.0
_FORMAL_OCCURRENCE_DURATION_RATIO = 2.5


@dataclass(frozen=True)
class _ScanFrame:
    frame_index: int
    timestamp: float
    gray: "np.ndarray"
    packed_edges: "np.ndarray"
    edge_pixels: int
    phash: "np.ndarray"
    sharpness: float


@dataclass
class _SlideSegment:
    first: int
    last: int

    @property
    def length(self) -> int:
        return self.last - self.first + 1


@dataclass(frozen=True)
class _CleanResult:
    image: "np.ndarray"
    annotation_pixels: int
    colour_guided_pixels: int
    watermark_pixels: int
    low_confidence_pixels: int
    residual_annotation_pixels: int
    residual_colour_match_pixels: int
    selected_timestamp: float
    selected_stability_support: int
    transition_candidates_rejected: int
    fusion_input_frames: int = 0
    fusion_registered_frames: int = 0
    fusion_pixels: int = 0
    fusion_rejected_frames: int = 0
    enhancement_engine: str = "未启用"
    enhancement_scale: float = 1.0
    ai_attempted: bool = False
    ai_accepted: bool = False
    enhancement_fallback_blocks: int = 0
    enhancement_total_blocks: int = 0
    enhancement_reason: str = ""


@dataclass(frozen=True)
class _SlideSignature:
    """Compact duplicate-detection data kept after a slide is written to disk."""

    shape: tuple[int, ...]
    phash: "np.ndarray"
    gray: "np.ndarray"
    edges: "np.ndarray"
    dilated_edges: "np.ndarray"
    colored_ink: "np.ndarray"
    annotation_ink: "np.ndarray"


def _require_dependencies() -> None:
    if cv2 is None or np is None:
        raise MissingEngineError(
            "讲解视频提取 PPT 需要 OpenCV 与 NumPy；请重新运行安装脚本补齐依赖"
        )


def _input_video(path: PathLike) -> Path:
    source = Path(path)
    if not source.is_file():
        raise ValidationError(f"视频不存在：{source}")
    if source.stat().st_size <= 0:
        raise ValidationError(f"视频为空：{source.name}")
    return source


def _parse_percent_boxes(
    value: object,
    label: str,
    *,
    multiple: bool,
    allow_empty: bool = True,
) -> list[tuple[float, float, float, float]]:
    text = str(value or "").strip()
    if not text:
        if allow_empty:
            return []
        raise ValidationError(f"请填写{label}，格式为 x,y,宽,高（百分比）")
    pieces = [part.strip() for part in text.split(";") if part.strip()]
    if not multiple and len(pieces) != 1:
        raise ValidationError(f"{label}只能填写一个区域")
    boxes: list[tuple[float, float, float, float]] = []
    for piece in pieces:
        raw = [item.strip() for item in piece.replace("，", ",").split(",")]
        if len(raw) != 4:
            raise ValidationError(f"{label}格式应为 x,y,宽,高（百分比）")
        try:
            x, y, width, height = (float(item) for item in raw)
        except ValueError as exc:
            raise ValidationError(f"{label}必须由四个数字组成") from exc
        if not all(math.isfinite(item) for item in (x, y, width, height)):
            raise ValidationError(f"{label}必须使用有限数字")
        if x < 0 or y < 0 or width <= 0 or height <= 0:
            raise ValidationError(f"{label}的坐标不能为负，宽高必须大于 0")
        if x + width > 100.0001 or y + height > 100.0001:
            raise ValidationError(f"{label}不能超出画面 0–100% 范围")
        boxes.append((x / 100.0, y / 100.0, width / 100.0, height / 100.0))
    return boxes


def _parse_rgb_colors(
    value: object,
    label: str,
    *,
    allow_empty: bool = False,
) -> list[RGBColor]:
    """Parse one or more user-facing RGB colours.

    The UI stores colours as ``#RRGGBB`` but accepting ``R,G,B`` keeps the
    processor usable from scripts and older task records.  Multiple colours
    are separated with an ASCII semicolon so commas remain unambiguous inside
    the RGB form.
    """

    text = str(value or "").strip()
    if not text:
        if allow_empty:
            return []
        raise ValidationError(f"请填写{label}")
    colours: list[RGBColor] = []
    for raw_piece in (piece.strip() for piece in text.split(";")):
        if not raw_piece:
            continue
        if raw_piece.startswith("#"):
            hexadecimal = raw_piece[1:]
            if len(hexadecimal) != 6:
                raise ValidationError(f"{label}应为 #RRGGBB 或 R,G,B")
            try:
                channels = tuple(
                    int(hexadecimal[index : index + 2], 16)
                    for index in (0, 2, 4)
                )
            except ValueError as exc:
                raise ValidationError(f"{label}包含无效十六进制颜色") from exc
        else:
            pieces = [piece.strip() for piece in raw_piece.replace("，", ",").split(",")]
            if len(pieces) != 3:
                raise ValidationError(f"{label}应为 #RRGGBB 或 R,G,B")
            try:
                channels = tuple(int(piece) for piece in pieces)
            except ValueError as exc:
                raise ValidationError(f"{label}的 RGB 通道必须是整数") from exc
        if any(channel < 0 or channel > 255 for channel in channels):
            raise ValidationError(f"{label}的 RGB 通道必须位于 0–255")
        colour = (int(channels[0]), int(channels[1]), int(channels[2]))
        if colour not in colours:
            colours.append(colour)
    if not colours and not allow_empty:
        raise ValidationError(f"请填写{label}")
    return colours


def _rgb_to_bgr(colour: RGBColor) -> BGRColor:
    red, green, blue = colour
    return blue, green, red


def _validate_annotation_colour_options(
    mode: str,
    colours: object,
    tolerance: object,
) -> tuple[str, tuple[BGRColor, ...], int]:
    normalized_mode = str(mode or "auto").strip().lower()
    if normalized_mode not in {"auto", "manual", "off"}:
        raise ValidationError("手写 / 异色标记颜色辅助选项无效")
    try:
        normalized_tolerance = int(tolerance)
    except (TypeError, ValueError) as exc:
        raise ValidationError("颜色容差必须是 0–100 的整数") from exc
    if not 0 <= normalized_tolerance <= 100:
        raise ValidationError("颜色容差必须位于 0–100")
    parsed = _parse_rgb_colors(
        colours,
        "手写 / 异色标记基准色",
        allow_empty=normalized_mode != "manual",
    )
    return (
        normalized_mode,
        tuple(_rgb_to_bgr(colour) for colour in parsed),
        normalized_tolerance,
    )


def _validate_fixed_fill_options(
    mode: str,
    colour: object,
) -> tuple[str, BGRColor]:
    normalized_mode = str(mode or "temporal").strip().lower()
    if normalized_mode not in {"temporal", "background", "color"}:
        raise ValidationError("固定水印填充方式无效")
    parsed = _parse_rgb_colors(
        colour,
        "固定水印指定填充色",
        allow_empty=normalized_mode != "color",
    )
    rgb = parsed[0] if parsed else (255, 255, 255)
    return normalized_mode, _rgb_to_bgr(rgb)


def _pixel_box(
    box: tuple[float, float, float, float], width: int, height: int
) -> tuple[int, int, int, int]:
    x, y, box_width, box_height = box
    left = max(0, min(width - 1, int(round(x * width))))
    top = max(0, min(height - 1, int(round(y * height))))
    right = max(left + 1, min(width, int(round((x + box_width) * width))))
    bottom = max(top + 1, min(height, int(round((y + box_height) * height))))
    return left, top, right, bottom


def _scaled_boxes(
    boxes: Iterable[tuple[float, float, float, float]], width: int, height: int
) -> list[tuple[int, int, int, int]]:
    return [_pixel_box(box, width, height) for box in boxes]


def _watermark_boxes(
    mode: str,
    custom_box: object,
) -> list[tuple[float, float, float, float]]:
    if mode == "off":
        return []
    if mode == "auto":
        return [(0.0, 0.0, 1.0, 0.42)]
    if mode == "top":
        return [(0.0, 0.0, 1.0, 0.20)]
    if mode == "bottom":
        return [(0.0, 0.80, 1.0, 0.20)]
    if mode == "full":
        return [(0.0, 0.0, 1.0, 1.0)]
    if mode == "custom":
        return _parse_percent_boxes(
            custom_box, "水印搜索区域", multiple=False, allow_empty=False
        )
    raise ValidationError("水印搜索范围选项无效")


def _presenter_boxes(
    policy: str,
    custom_box: object,
) -> list[tuple[float, float, float, float]]:
    if policy in {"keep", "auto_crop"}:
        return []
    if policy == "right_bottom":
        return [(0.78, 0.67, 0.22, 0.33)]
    if policy == "custom":
        return _parse_percent_boxes(
            custom_box, "讲师画面区域", multiple=False, allow_empty=False
        )
    raise ValidationError("讲师画面处理选项无效")


def _analysis_valid_mask(
    width: int,
    height: int,
    ignored_boxes: Sequence[tuple[float, float, float, float]],
) -> "np.ndarray":
    valid = np.ones((height, width), dtype=np.uint8)
    for left, top, right, bottom in _scaled_boxes(ignored_boxes, width, height):
        valid[top:bottom, left:right] = 0
    if int(valid.sum()) < max(64, width * height // 8):
        valid[:, :] = 1
    return valid.astype(bool)


def _masked_gray(gray: "np.ndarray", valid: "np.ndarray") -> "np.ndarray":
    result = gray.copy()
    values = gray[valid]
    fill = int(np.median(values)) if values.size else 0
    result[~valid] = fill
    return result


def _perceptual_hash(gray: "np.ndarray") -> "np.ndarray":
    sample = cv2.resize(gray, (32, 32), interpolation=cv2.INTER_AREA)
    transformed = cv2.dct(sample.astype(np.float32))[:8, :8]
    flattened = transformed.reshape(-1)
    median = float(np.median(flattened[1:]))
    return flattened > median


def _sharpness(gray: "np.ndarray", valid: "np.ndarray") -> float:
    laplacian = cv2.Laplacian(gray, cv2.CV_32F)
    values = laplacian[valid]
    return float(values.var()) if values.size else 0.0


def _decoded_frame_index(capture: object, fallback: int) -> int:
    """Return the index of the frame that was just decoded.

    OpenCV reports ``CAP_PROP_POS_FRAMES`` as the position of the *next* frame
    after ``read()``.  Using that value without subtracting one makes the image
    and its timestamp disagree by one frame.  A few backends return zero or a
    stale value, so a tightly bounded sequential fallback remains necessary.
    """

    try:
        position = float(capture.get(cv2.CAP_PROP_POS_FRAMES))
    except (AttributeError, TypeError, ValueError):
        return fallback
    if not math.isfinite(position):
        return fallback
    decoded = int(round(position)) - 1
    # Frame identity must follow the number of successful sequential reads.
    # Some OpenCV/FFmpeg builds occasionally report a position one frame ahead
    # or behind while decoding B-frames.  Accepting that +/-1 drift moves every
    # later scheduled candidate onto the wrong physical image.  The backend
    # property is therefore only an exact consistency check.
    return decoded if decoded == fallback else fallback


def _decoded_timestamp(
    capture: object,
    fallback_index: int,
    fps: float,
    previous: float | None,
) -> float:
    """Return the presentation time of the frame that was just decoded.

    ``CAP_PROP_PTS`` is preferred for FFmpeg-backed variable-frame-rate input;
    OpenCV documents it in the FPS time base.  ``CAP_PROP_POS_MSEC`` is the
    secondary source.  Both properties are known to be zero, stale, or briefly
    non-monotonic with some codecs, so frame ordinal divided by nominal FPS is
    retained as a safe last resort.
    """

    nominal_step = 1.0 / max(float(fps), 1e-6)
    fallback = max(0.0, float(fallback_index) * nominal_step)
    candidates: list[float] = []
    pts_property = getattr(cv2, "CAP_PROP_PTS", None)
    if pts_property is not None:
        try:
            pts = float(capture.get(pts_property))
        except (AttributeError, TypeError, ValueError):
            pts = math.nan
        if math.isfinite(pts) and pts >= 0.0:
            candidates.append(pts * nominal_step)
    try:
        milliseconds = float(capture.get(cv2.CAP_PROP_POS_MSEC))
    except (AttributeError, TypeError, ValueError):
        milliseconds = math.nan
    if math.isfinite(milliseconds) and milliseconds >= 0.0:
        candidates.append(milliseconds / 1000.0)

    if previous is None:
        valid_candidates = [candidate for candidate in candidates if candidate >= 0.0]
        if valid_candidates:
            # Some FFmpeg builds expose CAP_PROP_PTS one FPS tick ahead after
            # read(), while POS_MSEC correctly starts at zero.  Prefer the
            # backend value that best agrees with the first decoded ordinal.
            return min(valid_candidates, key=lambda value: abs(value - fallback))
        return fallback

    # A decoded lecture frame cannot legitimately jump many seconds ahead of
    # the preceding frame.  Reject stale seek values while still allowing low
    # frame-rate/VFR material.
    maximum_step = max(10.0, nominal_step * 20.0)
    monotonic_candidates = [
        candidate
        for candidate in candidates
        if previous + 1e-6 < candidate <= previous + maximum_step
    ]
    if monotonic_candidates:
        expected = previous + nominal_step
        return min(monotonic_candidates, key=lambda value: abs(value - expected))
    if fallback > previous + 1e-6:
        return fallback
    return previous + nominal_step


def _scan_video(
    source: Path,
    *,
    scan_mode: str,
    ignored_boxes: Sequence[tuple[float, float, float, float]],
) -> tuple[list[_ScanFrame], dict[str, float | int]]:
    _require_dependencies()
    if scan_mode not in _SCAN_INTERVALS:
        raise ValidationError("扫描精度选项无效")
    capture = cv2.VideoCapture(str(source))
    if not capture.isOpened():
        raise MissingEngineError(f"OpenCV 无法解码该视频：{source.name}")
    try:
        fps = float(capture.get(cv2.CAP_PROP_FPS) or 0.0)
        frame_count = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
        width = int(capture.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
        height = int(capture.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
        if fps <= 0 or width <= 0 or height <= 0:
            raise ValidationError("无法读取视频帧率或分辨率")
        duration = frame_count / fps if frame_count > 0 else 0.0
        interval = _SCAN_INTERVALS[scan_mode]
        if duration > 0 and duration / interval > _MAX_SCAN_FRAMES:
            interval = duration / _MAX_SCAN_FRAMES
        thumb_width = 160
        thumb_height = max(64, int(round(height * thumb_width / width)))
        valid = _analysis_valid_mask(thumb_width, thumb_height, ignored_boxes)
        scans: list[_ScanFrame] = []
        index = 0
        last_decoded_index = -1
        last_raw_timestamp: float | None = None
        timestamp_origin: float | None = None
        next_sample_time = 0.0
        while True:
            check_cancelled("已取消视频幻灯片扫描")
            ok, frame = capture.read()
            if not ok:
                break
            decoded_index = _decoded_frame_index(capture, index)
            if decoded_index <= last_decoded_index:
                decoded_index = index
            last_decoded_index = decoded_index
            raw_timestamp = _decoded_timestamp(
                capture, index, fps, last_raw_timestamp
            )
            last_raw_timestamp = raw_timestamp
            if timestamp_origin is None:
                timestamp_origin = raw_timestamp
            timestamp = max(0.0, raw_timestamp - timestamp_origin)
            if not scans or timestamp + 1e-6 >= next_sample_time:
                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                gray = cv2.resize(
                    gray, (thumb_width, thumb_height), interpolation=cv2.INTER_AREA
                )
                masked = _masked_gray(gray, valid)
                edges = cv2.Canny(masked, 55, 150) > 0
                packed = np.packbits(edges.reshape(-1))
                scans.append(
                    _ScanFrame(
                        frame_index=decoded_index,
                        timestamp=timestamp,
                        gray=gray,
                        packed_edges=packed,
                        edge_pixels=int(edges.sum()),
                        phash=_perceptual_hash(masked),
                        sharpness=_sharpness(masked, valid),
                    )
                )
                while next_sample_time <= timestamp + 1e-6:
                    next_sample_time += interval
            index += 1
            if frame_count > 0 and index % max(1, int(fps * 4)) == 0:
                report_progress(
                    min(0.32, 0.32 * index / frame_count),
                    f"扫描换页与稳定画面 · {timestamp:.0f}/{duration:.0f} 秒",
                )
        if not scans:
            raise ValidationError("视频中没有可读取画面")
        if len(scans) == 1:
            scans.append(scans[0])
        decoded_duration = (
            max(0.0, last_raw_timestamp - (timestamp_origin or 0.0)) + 1.0 / fps
            if last_raw_timestamp is not None
            else index / fps
        )
        timestamp_gaps = [
            current.timestamp - previous.timestamp
            for previous, current in zip(scans, scans[1:])
            if current.timestamp > previous.timestamp
        ]
        measured_interval = (
            float(np.median(timestamp_gaps)) if timestamp_gaps else interval
        )
        return scans, {
            "fps": fps,
            "frame_count": frame_count or index,
            "width": width,
            "height": height,
            "duration": decoded_duration,
            "sample_interval": measured_interval,
        }
    finally:
        capture.release()


def _rebuild_scan_features(
    scans: Sequence[_ScanFrame],
    ignored_boxes: Sequence[tuple[float, float, float, float]],
) -> list[_ScanFrame]:
    """Recalculate compact scan features after auto-overlay detection.

    The first pass deliberately keeps the raw thumbnail in ``gray`` so an
    automatic presenter decision can be made without decoding the video a
    third time.  Only the compact edge/hash features need rebuilding.
    """

    if not scans:
        return []
    height, width = scans[0].gray.shape
    valid = _analysis_valid_mask(width, height, ignored_boxes)
    rebuilt: list[_ScanFrame] = []
    for item in scans:
        masked = _masked_gray(item.gray, valid)
        edges = cv2.Canny(masked, 55, 150) > 0
        rebuilt.append(
            _ScanFrame(
                frame_index=item.frame_index,
                timestamp=item.timestamp,
                gray=item.gray,
                packed_edges=np.packbits(edges.reshape(-1)),
                edge_pixels=int(edges.sum()),
                phash=_perceptual_hash(masked),
                sharpness=_sharpness(masked, valid),
            )
        )
    return rebuilt


def _detect_auto_presenter(
    scans: Sequence[_ScanFrame],
    box: tuple[float, float, float, float] = (0.78, 0.67, 0.22, 0.33),
) -> bool:
    """Detect sustained local motion typical of a lecturer camera window.

    Page transitions are rejected by comparing the candidate region with the
    rest of the frame.  A static lower-right table or logo therefore remains
    untouched, while a continuously moving talking-head overlay is ignored by
    page segmentation and cleaned during the second pass.
    """

    if len(scans) < 8:
        return False
    height, width = scans[0].gray.shape
    left, top, right, bottom = _pixel_box(box, width, height)
    outside = np.ones((height, width), dtype=bool)
    outside[top:bottom, left:right] = False
    region_area = max(1, (right - left) * (bottom - top))
    outside_area = max(1, int(np.count_nonzero(outside)))
    active = 0
    stable_pairs = 0
    region_levels: list[float] = []
    excess_levels: list[float] = []
    for previous, current in zip(scans, scans[1:]):
        delta = cv2.absdiff(previous.gray, current.gray)
        region_changed = (
            float(np.count_nonzero(delta[top:bottom, left:right] >= 13)) / region_area
        )
        outside_changed = (
            float(np.count_nonzero((delta >= 13) & outside)) / outside_area
        )
        # Whole-page changes are transitions or animations, not evidence for a
        # presenter window.
        if outside_changed > 0.055:
            continue
        stable_pairs += 1
        region_levels.append(region_changed)
        excess = max(0.0, region_changed - 1.7 * outside_changed)
        excess_levels.append(excess)
        if region_changed >= 0.045 and excess >= 0.025:
            active += 1
    if stable_pairs < max(6, len(scans) // 12):
        return False
    active_ratio = active / stable_pairs
    percentile = float(np.percentile(region_levels, 70)) if region_levels else 0.0
    excess_percentile = (
        float(np.percentile(excess_levels, 70)) if excess_levels else 0.0
    )
    return bool(
        active_ratio >= 0.18 and percentile >= 0.055 and excess_percentile >= 0.028
    )


def _edge_difference(first: _ScanFrame, second: _ScanFrame) -> float:
    xor = np.bitwise_xor(first.packed_edges, second.packed_edges)
    changed = int(np.unpackbits(xor).sum())
    total = max(1, first.gray.size)
    return changed / total


def _scan_difference(
    first: _ScanFrame,
    second: _ScanFrame,
    valid: "np.ndarray",
) -> tuple[float, float, float, int]:
    delta = cv2.absdiff(first.gray, second.gray)
    values = delta[valid]
    changed = float(np.mean(values >= 18)) if values.size else 0.0
    mean_delta = float(values.mean() / 255.0) if values.size else 0.0
    edge_delta = _edge_difference(first, second)
    hash_distance = int(np.count_nonzero(first.phash != second.phash))
    return changed, mean_delta, edge_delta, hash_distance


def _segment_scans(
    scans: Sequence[_ScanFrame],
    *,
    ignored_boxes: Sequence[tuple[float, float, float, float]],
    sensitivity: str,
) -> list[_SlideSegment]:
    if sensitivity not in _CHANGE_THRESHOLDS:
        raise ValidationError("换页灵敏度选项无效")
    height, width = scans[0].gray.shape
    valid = _analysis_valid_mask(width, height, ignored_boxes)
    metrics = [
        _scan_difference(scans[index - 1], scans[index], valid)
        for index in range(1, len(scans))
    ]
    changed_values = np.asarray([item[0] for item in metrics], dtype=np.float32)
    edge_values = np.asarray([item[2] for item in metrics], dtype=np.float32)
    median_changed = float(np.median(changed_values)) if changed_values.size else 0.0
    median_edge = float(np.median(edge_values)) if edge_values.size else 0.0
    mad_changed = (
        float(np.median(np.abs(changed_values - median_changed)))
        if changed_values.size
        else 0.0
    )
    mad_edge = (
        float(np.median(np.abs(edge_values - median_edge))) if edge_values.size else 0.0
    )
    base_changed, base_edge, base_hash = _CHANGE_THRESHOLDS[sensitivity]
    changed_threshold = max(base_changed, median_changed + 6.0 * mad_changed)
    edge_threshold = max(base_edge, median_edge + 5.0 * mad_edge)
    sparse_thresholds = {
        "conservative": (0.062, 0.040),
        "balanced": (0.048, 0.030),
        "sensitive": (0.034, 0.022),
    }
    sparse_changed, sparse_edge = sparse_thresholds[sensitivity]
    boundaries: list[int] = []
    last_boundary = 0
    interval = max(0.001, scans[1].timestamp - scans[0].timestamp)
    minimum_gap = max(1, int(round(_MIN_SEGMENT_SECONDS / interval)))
    # Adjacent thumbnails are deliberately sampled rather sparsely.  A slow
    # dissolve can therefore stay below every one-step gate even though the
    # accumulated start/end difference is a genuine page change.  Keep this
    # look-back bounded so the extra work remains linear with a small constant.
    cumulative_lookback = max(2, min(6, int(math.ceil(3.0 / interval))))
    for index, (changed, mean_delta, edge_delta, hash_distance) in enumerate(
        metrics, start=1
    ):
        strong_area = changed >= changed_threshold and (
            edge_delta >= edge_threshold
            or hash_distance >= base_hash
            or mean_delta >= 0.065
        )
        strong_structure = (
            edge_delta >= edge_threshold * 1.45
            and hash_distance >= base_hash + 3
            and changed >= changed_threshold * 0.72
        )
        global_tone_change = (
            mean_delta >= 0.05
            and hash_distance >= base_hash + 3
            and edge_delta >= max(0.012, median_edge + 2.5 * mad_edge)
        )
        # Some course templates replace only a question or a short title.  The
        # changed area can be below the normal 7.5-14.5% gate even though the
        # structural hash changes abruptly.  This extra gate catches such
        # sparse page changes while remaining well above ordinary pen strokes.
        sparse_structure_change = (
            changed >= sparse_changed
            and edge_delta >= sparse_edge
            and mean_delta >= 0.010
            and hash_distance >= base_hash + 3
        )
        is_boundary = (
            strong_area
            or strong_structure
            or global_tone_change
            or sparse_structure_change
        )
        if not is_boundary and index - last_boundary >= minimum_gap:
            earliest = max(last_boundary, index - cumulative_lookback)
            for start in range(index - 2, earliest - 1, -1):
                cumulative = _scan_difference(scans[start], scans[index], valid)
                (
                    cumulative_changed,
                    cumulative_mean,
                    cumulative_edge,
                    cumulative_hash,
                ) = cumulative
                path = metrics[start:index]
                active_steps = sum(
                    1
                    for step_changed, step_mean, step_edge, _step_hash in path
                    if step_changed >= 0.010
                    or step_mean >= 0.003
                    or step_edge >= 0.008
                )
                gradual_change = active_steps >= max(2, len(path) // 2)
                cumulative_structure = (
                    cumulative_changed >= changed_threshold * 0.88
                    and (
                        cumulative_edge >= edge_threshold * 0.82
                        or cumulative_hash >= base_hash + 2
                        or cumulative_mean >= 0.050
                    )
                ) or (
                    cumulative_changed >= sparse_changed
                    and cumulative_edge >= sparse_edge
                    and cumulative_mean >= 0.012
                    and cumulative_hash >= base_hash + 2
                )
                if not (gradual_change and cumulative_structure):
                    continue

                # A build animation commonly adds printed bullets while keeping
                # the previous page intact.  It is not a new slide: defer to the
                # final-state selector and retain only the most complete state.
                first = scans[start]
                final = scans[index]
                progressive_reveal = (
                    final.edge_pixels >= first.edge_pixels * 0.96
                    and _edge_containment(first, final) >= 0.88
                )
                if not progressive_reveal:
                    is_boundary = True
                    break
        if is_boundary and index - last_boundary >= minimum_gap:
            boundaries.append(index)
            last_boundary = index
    starts = [0, *boundaries]
    ends = [item - 1 for item in boundaries] + [len(scans) - 1]
    segments = [_SlideSegment(first, last) for first, last in zip(starts, ends)]
    return _merge_short_and_incremental_segments(segments, scans, valid)


def _edge_containment(first: _ScanFrame, second: _ScanFrame) -> float:
    first_edges = np.unpackbits(first.packed_edges)[: first.gray.size].astype(bool)
    second_edges = np.unpackbits(second.packed_edges)[: second.gray.size].astype(bool)
    if not first_edges.any():
        return 0.0
    shape = first.gray.shape
    first_image = first_edges.reshape(shape).astype(np.uint8)
    second_image = second_edges.reshape(shape).astype(np.uint8)
    second_dilated = cv2.dilate(second_image, np.ones((3, 3), np.uint8)) > 0
    return float(np.count_nonzero(first_image.astype(bool) & second_dilated)) / max(
        1, int(first_image.sum())
    )


def _merge_short_and_incremental_segments(
    segments: list[_SlideSegment],
    scans: Sequence[_ScanFrame],
    valid: "np.ndarray",
) -> list[_SlideSegment]:
    if len(segments) <= 1:
        return segments
    merged: list[_SlideSegment] = []
    for segment in segments:
        if not merged:
            merged.append(segment)
            continue
        previous = merged[-1]
        too_short = segment.length <= 1
        # The scan immediately adjacent to a boundary may already contain the
        # first partially switched frame.  Comparing that contaminated frame
        # with the next segment makes two template-similar pages look like a
        # progressive reveal and silently removes a real page.  Use the latest
        # observation at least one second before the boundary whenever the
        # previous segment is long enough.
        previous_stable_position = previous.last
        stable_before = scans[previous.last].timestamp - 1.0
        for position in range(previous.last, previous.first - 1, -1):
            if scans[position].timestamp <= stable_before:
                previous_stable_position = position
                break
        previous_final = scans[previous_stable_position]
        current_final = scans[segment.last]
        changed, mean_delta, edge_delta, _hash_distance = _scan_difference(
            previous_final, current_final, valid
        )
        previous_contained = _edge_containment(previous_final, current_final)
        edge_growth = current_final.edge_pixels / max(1, previous_final.edge_pixels)
        incremental_reveal = (
            previous_contained >= 0.975
            and edge_growth >= 0.96
            and changed <= 0.24
            and mean_delta <= 0.050
            and edge_delta <= 0.14
        )
        if too_short or incremental_reveal:
            previous.last = segment.last
        else:
            merged.append(segment)
    return merged


def _auto_presentation_crop(
    scans: Sequence[_ScanFrame], width: int, height: int
) -> tuple[int, int, int, int]:
    if width < 320 or height < 180 or len(scans) < 3:
        return 0, 0, width, height
    # Only remove genuine uniform outer letterbox bars.  Interior vertical
    # seams are often part of the slide template (the user sample has one at
    # x≈1360), so treating them as a crop edge can silently discard tables.
    sample_count = min(32, len(scans))
    positions = np.linspace(0, len(scans) - 1, sample_count, dtype=int)
    median_frame = np.median(
        np.stack([scans[int(position)].gray for position in positions], axis=0),
        axis=0,
    )
    row_mean = median_frame.mean(axis=1)
    row_std = median_frame.std(axis=1)
    column_mean = median_frame.mean(axis=0)
    column_std = median_frame.std(axis=0)

    def outer_run(means: "np.ndarray", deviations: "np.ndarray", reverse: bool) -> int:
        indices = range(len(means) - 1, -1, -1) if reverse else range(len(means))
        count = 0
        for index in indices:
            if means[index] <= 9.0 and deviations[index] <= 5.0:
                count += 1
            else:
                break
        return count

    top = outer_run(row_mean, row_std, False)
    bottom = outer_run(row_mean, row_std, True)
    left = outer_run(column_mean, column_std, False)
    right = outer_run(column_mean, column_std, True)
    thumb_height, thumb_width = median_frame.shape
    if top < thumb_height * 0.015:
        top = 0
    if bottom < thumb_height * 0.015:
        bottom = 0
    if left < thumb_width * 0.015:
        left = 0
    if right < thumb_width * 0.015:
        right = 0
    crop_left = int(round(left * width / thumb_width))
    crop_top = int(round(top * height / thumb_height))
    crop_right = width - int(round(right * width / thumb_width))
    crop_bottom = height - int(round(bottom * height / thumb_height))
    if crop_right - crop_left < width * 0.60 or crop_bottom - crop_top < height * 0.60:
        return 0, 0, width, height
    return crop_left, crop_top, crop_right, crop_bottom


def _presentation_crop(
    mode: str,
    custom_box: object,
    scans: Sequence[_ScanFrame],
    width: int,
    height: int,
) -> tuple[int, int, int, int]:
    if mode == "full":
        return 0, 0, width, height
    if mode == "auto":
        return _auto_presentation_crop(scans, width, height)
    if mode == "custom":
        boxes = _parse_percent_boxes(
            custom_box, "课件画面区域", multiple=False, allow_empty=False
        )
        return _pixel_box(boxes[0], width, height)
    raise ValidationError("课件画面范围选项无效")


def _candidate_scan_positions(
    segment: _SlideSegment, scans: Sequence[_ScanFrame]
) -> list[int]:
    positions = list(range(segment.first, segment.last + 1))
    if len(positions) <= _MAX_CANDIDATES_PER_SLIDE:
        return positions
    selected = set(
        int(round(value))
        for value in np.linspace(
            segment.first, segment.last, _MAX_CANDIDATES_PER_SLIDE - 4
        )
    )
    selected.update(positions[:2])
    selected.update(positions[-2:])
    return sorted(selected)


def _refined_candidate_frame_indices(
    segment: _SlideSegment,
    scans: Sequence[_ScanFrame],
    fps: float,
) -> set[int]:
    """Add short, dense sampling after a coarse boundary.

    The first pass deliberately stays sparse for long-video efficiency.  A
    dissolve/wipe can therefore begin between two scan samples.  The second
    pass already decodes every frame sequentially, so retaining a small 8 FPS
    window after each boundary gives the selector enough temporal evidence to
    wait for a stable platform without materially increasing decode cost.
    """

    base_positions = _candidate_scan_positions(segment, scans)
    base_indices = {scans[position].frame_index for position in base_positions}
    indices = set(base_indices)
    start = scans[segment.first].frame_index
    end = scans[segment.last].frame_index
    protected: set[int] = {start, end}
    if fps > 0 and end > start:
        start_time = scans[segment.first].timestamp
        end_time = scans[segment.last].timestamp
        elapsed = end_time - start_time
        local_fps = (end - start) / elapsed if elapsed > 1e-6 else fps
        if not math.isfinite(local_fps) or local_fps <= 0:
            local_fps = fps
        dense_stride = max(1, int(round(local_fps / _TRANSITION_REFINE_FPS)))

        dense_end_target = start_time + _TRANSITION_REFINE_SECONDS
        dense_end = end
        for position in range(segment.first, segment.last + 1):
            if scans[position].timestamp >= dense_end_target:
                dense_end = scans[position].frame_index
                break
        dense_start = set(range(start, dense_end + 1, dense_stride))
        dense_start.add(dense_end)
        # A short tail window gives late progressive pages a confirmed final
        # platform too.  It also avoids choosing the first frame after a return
        # to an earlier slide when the segment ends soon afterwards.
        tail_target = end_time - _TRANSITION_TAIL_SECONDS
        tail_start = start
        for position in range(segment.first, segment.last + 1):
            if scans[position].timestamp <= tail_target:
                tail_start = scans[position].frame_index
            else:
                break
        dense_tail = set(range(tail_start, end + 1, dense_stride))
        dense_tail.update((tail_start, end))
        protected.update(dense_start)
        protected.update(dense_tail)
        indices.update(protected)
    if len(indices) <= _MAX_REFINED_CANDIDATES_PER_SLIDE:
        return indices
    protected_values = sorted(protected)
    if len(protected_values) > _MAX_REFINED_CANDIDATES_PER_SLIDE:
        positions = np.linspace(
            0,
            len(protected_values) - 1,
            _MAX_REFINED_CANDIDATES_PER_SLIDE,
        )
        return {protected_values[int(round(position))] for position in positions}
    remaining = sorted(base_indices - protected)
    slots = max(0, _MAX_REFINED_CANDIDATES_PER_SLIDE - len(protected_values))
    if slots and remaining:
        positions = np.linspace(0, len(remaining) - 1, min(slots, len(remaining)))
        protected_values.extend(
            remaining[int(round(position))] for position in positions
        )
    return set(protected_values[:_MAX_REFINED_CANDIDATES_PER_SLIDE])


def _scan_frame_from_image(
    frame_index: int,
    timestamp: float,
    image: "np.ndarray",
    ignored_boxes: Sequence[tuple[float, float, float, float]],
) -> _ScanFrame:
    height, width = image.shape[:2]
    thumb_width = _REFINED_ANALYSIS_WIDTH
    thumb_height = max(64, int(round(height * thumb_width / max(1, width))))
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    gray = cv2.resize(gray, (thumb_width, thumb_height), interpolation=cv2.INTER_AREA)
    valid = _analysis_valid_mask(thumb_width, thumb_height, ignored_boxes)
    masked = _masked_gray(gray, valid)
    edges = cv2.Canny(masked, 55, 150) > 0
    return _ScanFrame(
        frame_index=frame_index,
        timestamp=timestamp,
        gray=gray,
        packed_edges=np.packbits(edges.reshape(-1)),
        edge_pixels=int(edges.sum()),
        phash=_perceptual_hash(masked),
        sharpness=_sharpness(masked, valid),
    )


def _candidate_stability_support(
    frames: Sequence[tuple[_ScanFrame, "np.ndarray"]],
    excluded_boxes: Sequence[tuple[int, int, int, int]],
) -> tuple["np.ndarray", "np.ndarray"]:
    """Return per-candidate stable-neighbour support and quiet pair flags.

    A partially dissolved page often has *more* edges than either endpoint.
    Stability is therefore a hard temporal gate, not another small quality
    bonus.  The block vote tolerates a moving pointer/watermark or a retained
    lecturer window while rejecting changes that cover a substantial part of
    the slide.
    """

    count = len(frames)
    support = np.zeros(count, dtype=np.int16)
    if count < 2:
        return support, np.zeros(0, dtype=bool)
    height, width = frames[0][1].shape[:2]
    thumb_width = _REFINED_ANALYSIS_WIDTH
    thumb_height = max(96, int(round(height * thumb_width / max(1, width))))
    valid = np.ones((thumb_height, thumb_width), dtype=bool)
    scale_x = thumb_width / max(1, width)
    scale_y = thumb_height / max(1, height)
    for left, top, right, bottom in excluded_boxes:
        scaled_left = max(0, min(thumb_width, int(round(left * scale_x))))
        scaled_top = max(0, min(thumb_height, int(round(top * scale_y))))
        scaled_right = max(0, min(thumb_width, int(round(right * scale_x))))
        scaled_bottom = max(0, min(thumb_height, int(round(bottom * scale_y))))
        valid[scaled_top:scaled_bottom, scaled_left:scaled_right] = False
    if int(np.count_nonzero(valid)) < max(64, valid.size // 8):
        valid[:, :] = True
    grays = [
        cv2.resize(
            cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY),
            (thumb_width, thumb_height),
            interpolation=cv2.INTER_AREA,
        )
        for _scan, frame in frames
    ]
    edges = [cv2.Canny(gray, 48, 138) > 0 for gray in grays]
    block_rows = 6
    block_columns = 10
    focus_kernel = np.ones((3, 3), np.uint8)

    def is_quiet(first_index: int, second_index: int, *, strict: bool) -> bool:
        delta = cv2.absdiff(grays[first_index], grays[second_index])
        values = delta[valid]
        if not values.size:
            return False
        changed = float(np.mean(values >= 18))
        mean_delta = float(values.mean() / 255.0)
        first_values = grays[first_index][valid].astype(np.float32)
        second_values = grays[second_index][valid].astype(np.float32)
        first_values -= float(first_values.mean())
        second_values -= float(second_values.mean())
        denominator = float(
            np.linalg.norm(first_values) * np.linalg.norm(second_values)
        )
        correlation = (
            float(np.dot(first_values, second_values) / denominator)
            if denominator > 1e-6
            else 1.0
        )
        edge_delta = float(np.mean((edges[first_index] != edges[second_index])[valid]))

        # Whole-page averages hide a dissolve that changes only one title or a
        # sparse line of text.  Re-evaluate change on the union of both frames'
        # edge neighbourhoods so printed glyph drift is not divided by empty
        # page area.
        focus = cv2.dilate(
            (edges[first_index] | edges[second_index]).astype(np.uint8),
            focus_kernel,
            iterations=1,
        ).astype(bool)
        focus &= valid
        focus_values = delta[focus]
        foreground_changed = (
            float(np.mean(focus_values >= 12)) if focus_values.size else changed
        )
        foreground_mean = (
            float(focus_values.mean() / 255.0) if focus_values.size else mean_delta
        )

        active_blocks = 0
        usable_blocks = 0
        active_rows: set[int] = set()
        active_columns: set[int] = set()
        for row in range(block_rows):
            top = int(round(row * thumb_height / block_rows))
            bottom = int(round((row + 1) * thumb_height / block_rows))
            for column in range(block_columns):
                left = int(round(column * thumb_width / block_columns))
                right = int(round((column + 1) * thumb_width / block_columns))
                block_valid = valid[top:bottom, left:right]
                if int(np.count_nonzero(block_valid)) < max(4, block_valid.size // 3):
                    continue
                usable_blocks += 1
                block_values = delta[top:bottom, left:right][block_valid]
                if (
                    float(np.mean(block_values >= 16)) >= 0.075
                    and float(block_values.mean() / 255.0) >= 0.018
                ):
                    active_blocks += 1
                    active_rows.add(row)
                    active_columns.add(column)
        active_ratio = active_blocks / max(1, usable_blocks)
        motion_mask = ((delta >= 16) & valid).astype(np.uint8)
        motion_mask = cv2.morphologyEx(
            motion_mask,
            cv2.MORPH_CLOSE,
            np.ones((3, 7), np.uint8),
        )
        component_count, _labels, stats, _centroids = cv2.connectedComponentsWithStats(
            motion_mask, 8
        )
        compact_components = True
        component_area = 0
        for label in range(1, component_count):
            box_width = int(stats[label, cv2.CC_STAT_WIDTH])
            box_height = int(stats[label, cv2.CC_STAT_HEIGHT])
            area = int(stats[label, cv2.CC_STAT_AREA])
            if area < 4:
                continue
            component_area += area
            if box_width > thumb_width * 0.32 or box_height > thumb_height * 0.20:
                compact_components = False
                break
        compact_motion = bool(
            compact_components
            and component_area <= int(np.count_nonzero(valid) * 0.04)
            and active_ratio <= (0.10 if strict else 0.12)
            and len(active_rows) <= 3
        )
        if strict:
            global_quiet = bool(
                changed <= 0.050
                and mean_delta <= 0.016
                and edge_delta <= 0.050
                and active_ratio <= 0.12
            )
            structural_quiet = bool(
                correlation >= 0.982
                and foreground_changed <= 0.14
                and foreground_mean <= 0.025
            )
            local_overlay_only = bool(
                compact_motion
                and changed <= 0.025
                and mean_delta <= 0.010
                and edge_delta <= 0.025
                and active_ratio <= 0.10
            )
            return global_quiet and (structural_quiet or local_overlay_only)
        global_quiet = bool(
            changed <= 0.075
            and mean_delta <= 0.024
            and edge_delta <= 0.070
            and active_ratio <= 0.16
        )
        structural_quiet = bool(
            correlation >= 0.960
            and foreground_changed <= 0.23
            and foreground_mean <= 0.045
        )
        local_overlay_only = bool(
            compact_motion
            and changed <= 0.035
            and mean_delta <= 0.014
            and edge_delta <= 0.038
            and active_ratio <= 0.12
        )
        return global_quiet and (structural_quiet or local_overlay_only)

    timestamps = np.asarray(
        [float(item[0].timestamp) for item in frames], dtype=np.float64
    )
    gaps = np.diff(timestamps)
    positive_gaps = gaps[np.isfinite(gaps) & (gaps > 1e-6)]
    if positive_gaps.size:
        cadence = float(np.percentile(positive_gaps, 30))
        maximum_gap = min(
            _STABLE_MAX_GAP_SECONDS,
            max(_STABLE_DENSE_GAP_SECONDS, cadence * 2.25),
        )
    else:
        maximum_gap = _STABLE_DENSE_GAP_SECONDS
    continuous = (gaps > 1e-6) & (gaps <= maximum_gap + 1e-6)
    quiet_array = np.asarray(
        [
            bool(continuous[index]) and is_quiet(index, index + 1, strict=False)
            for index in range(count - 1)
        ],
        dtype=bool,
    )

    # Require elapsed-time stability, not merely three nearby candidates.
    # This rejects slow dissolves and also prevents sparse samples several
    # seconds apart from being combined into a fictitious stable platform.
    maximum_window = max(1.15, min(2.5, maximum_gap * 2.2))
    for start in range(count - _STABLE_WINDOW_CANDIDATES + 1):
        for stop in range(start + _STABLE_WINDOW_CANDIDATES - 1, count):
            if stop > start and not bool(quiet_array[stop - 1]):
                break
            duration = float(timestamps[stop] - timestamps[start])
            if duration < _STABLE_MIN_SECONDS:
                continue
            if duration > maximum_window:
                break
            if not is_quiet(start, stop, strict=True):
                continue
            support[start : stop + 1] += 1

    # A printed answer can appear during the final half-second.  Admit that
    # short terminal platform only after an earlier full platform proves the
    # segment itself is real; it cannot bootstrap a transition-only fragment.
    if np.any(support) and count >= _TERMINAL_STABLE_MIN_CANDIDATES:
        terminal_start = count - _TERMINAL_STABLE_MIN_CANDIDATES
        terminal_duration = float(timestamps[-1] - timestamps[terminal_start])
        if (
            terminal_duration >= _TERMINAL_STABLE_MIN_SECONDS
            and bool(np.all(quiet_array[terminal_start:]))
            and is_quiet(terminal_start, count - 1, strict=True)
        ):
            support[terminal_start:] += 1
    if np.any(support) and quiet_array.size and bool(quiet_array[-1]):
        support[-2:] += 1
    return support, quiet_array


def _annotation_colour_mask(
    image: "np.ndarray",
    *,
    mode: str,
    colours: Sequence[BGRColor],
    tolerance: int,
    excluded_boxes: Sequence[tuple[int, int, int, int]] = (),
) -> "np.ndarray":
    """Return pixels matching the automatic or user supplied ink colours.

    Manual matching is performed in CIE Lab rather than independently
    thresholding RGB channels.  That makes the tolerance useful for
    anti-aliased stroke edges and video compression colour drift while keeping
    perceptually different printed colours separate.
    """

    height, width = image.shape[:2]
    if mode == "off":
        return np.zeros((height, width), dtype=np.uint8)
    if mode == "auto":
        hsv = cv2.cvtColor(image, cv2.COLOR_BGR2HSV)
        values = image.astype(np.int16)
        cyan_chroma = (
            np.minimum(values[:, :, 0], values[:, :, 1]) - values[:, :, 2]
        )
        mask = (
            (hsv[:, :, 0] >= 74)
            & (hsv[:, :, 0] <= 108)
            & (hsv[:, :, 1] >= 90)
            & (hsv[:, :, 2] >= 68)
            & (cyan_chroma >= 38)
        )
    else:
        if not colours:
            return np.zeros((height, width), dtype=np.uint8)
        lab = cv2.cvtColor(image, cv2.COLOR_BGR2LAB).astype(np.int16)
        references = np.asarray(colours, dtype=np.uint8).reshape(-1, 1, 3)
        reference_lab = cv2.cvtColor(references, cv2.COLOR_BGR2LAB).reshape(-1, 3)
        # OpenCV's 8-bit Lab uses an approximately perceptual scale.  A small
        # floor keeps exact colour-pick mode tolerant of ordinary H.264 noise.
        distance_limit = max(3.0, float(tolerance) * 1.15)
        limit_squared = distance_limit * distance_limit
        mask = np.zeros((height, width), dtype=bool)
        for reference in reference_lab.astype(np.int16):
            delta = lab - reference[None, None, :]
            distance_squared = np.sum(delta.astype(np.int32) ** 2, axis=2)
            mask |= distance_squared <= limit_squared
    result = mask.astype(np.uint8) * 255
    for left, top, right, bottom in excluded_boxes:
        result[top:bottom, left:right] = 0
    return result


def _structural_compatibility(
    first: "np.ndarray",
    second: "np.ndarray",
    *,
    excluded_boxes: Sequence[tuple[int, int, int, int]],
    annotation_mode: str,
    annotation_colours: Sequence[BGRColor],
    annotation_tolerance: int,
) -> bool:
    """Test whether two observations are the same printed page structure.

    This comparison intentionally ignores pixels matching the lecturer's ink
    colour.  It is used only after a real continuous stable platform has been
    found, so an isolated early frame can never bootstrap a slide by itself.
    """

    height, width = first.shape[:2]
    thumb_width = _REFINED_ANALYSIS_WIDTH
    thumb_height = max(96, int(round(height * thumb_width / max(1, width))))
    first_small = cv2.resize(first, (thumb_width, thumb_height), interpolation=cv2.INTER_AREA)
    second_small = cv2.resize(second, (thumb_width, thumb_height), interpolation=cv2.INTER_AREA)
    valid = np.ones((thumb_height, thumb_width), dtype=bool)
    scale_x = thumb_width / max(1, width)
    scale_y = thumb_height / max(1, height)
    scaled_excluded: list[tuple[int, int, int, int]] = []
    for left, top, right, bottom in excluded_boxes:
        box = (
            max(0, min(thumb_width, int(round(left * scale_x)))),
            max(0, min(thumb_height, int(round(top * scale_y)))),
            max(0, min(thumb_width, int(round(right * scale_x)))),
            max(0, min(thumb_height, int(round(bottom * scale_y)))),
        )
        scaled_excluded.append(box)
        valid[box[1] : box[3], box[0] : box[2]] = False
    first_ink = _annotation_colour_mask(
        first_small,
        mode=annotation_mode,
        colours=annotation_colours,
        tolerance=annotation_tolerance,
        excluded_boxes=scaled_excluded,
    )
    second_ink = _annotation_colour_mask(
        second_small,
        mode=annotation_mode,
        colours=annotation_colours,
        tolerance=annotation_tolerance,
        excluded_boxes=scaled_excluded,
    )
    if np.any(first_ink) or np.any(second_ink):
        ink = cv2.dilate(
            ((first_ink > 0) | (second_ink > 0)).astype(np.uint8),
            np.ones((5, 5), np.uint8),
            iterations=1,
        ).astype(bool)
        valid &= ~ink
    if int(np.count_nonzero(valid)) < max(128, valid.size // 5):
        return False
    first_gray = cv2.cvtColor(first_small, cv2.COLOR_BGR2GRAY)
    second_gray = cv2.cvtColor(second_small, cv2.COLOR_BGR2GRAY)
    delta = cv2.absdiff(first_gray, second_gray)
    values = delta[valid]
    changed = float(np.mean(values >= 18))
    mean_delta = float(values.mean() / 255.0)
    first_values = first_gray[valid].astype(np.float32)
    second_values = second_gray[valid].astype(np.float32)
    first_values -= float(first_values.mean())
    second_values -= float(second_values.mean())
    denominator = float(np.linalg.norm(first_values) * np.linalg.norm(second_values))
    correlation = (
        float(np.dot(first_values, second_values) / denominator)
        if denominator > 1e-6
        else 1.0
    )
    first_edges = cv2.Canny(first_gray, 48, 138) > 0
    second_edges = cv2.Canny(second_gray, 48, 138) > 0
    edge_delta = float(np.mean((first_edges != second_edges)[valid]))
    focus = cv2.dilate(
        (first_edges | second_edges).astype(np.uint8),
        np.ones((3, 3), np.uint8),
        iterations=1,
    ).astype(bool)
    focus &= valid
    focus_values = delta[focus]
    foreground_changed = (
        float(np.mean(focus_values >= 13)) if focus_values.size else changed
    )
    sparse_print_addition = bool(
        changed <= 0.040
        and mean_delta <= 0.018
        and edge_delta <= 0.035
        and foreground_changed <= 0.36
    )
    correlated_structure = bool(
        correlation >= 0.955
        and (foreground_changed <= 0.30 or changed <= 0.025)
    )
    return bool(
        changed <= 0.085
        and mean_delta <= 0.030
        and edge_delta <= 0.075
        and (correlated_structure or sparse_print_addition)
    )


def _reinclude_structurally_compatible_positions(
    frames: Sequence[tuple[_ScanFrame, "np.ndarray"]],
    stability_support: "np.ndarray",
    *,
    excluded_boxes: Sequence[tuple[int, int, int, int]],
    annotation_mode: str,
    annotation_colours: Sequence[BGRColor],
    annotation_tolerance: int,
) -> tuple[list[int], "np.ndarray"]:
    """Back stable, structurally matching sparse observations with real support."""

    count = len(frames)
    augmented = np.asarray(stability_support, dtype=np.int16).copy()
    stable_positions = [
        int(position) for position in np.flatnonzero(augmented > 0)
    ]
    if not stable_positions:
        return [], augmented
    stable_set = {int(position) for position in stable_positions}
    # The stable endpoint with the richest structure is the safest comparison
    # anchor.  Two anchors are retained when available to avoid accepting a
    # one-frame codec artefact.
    ranked_anchors = sorted(
        stable_set,
        key=lambda position: frames[position][0].edge_pixels,
        reverse=True,
    )
    anchors = ranked_anchors[: min(3, len(ranked_anchors))]
    settled_at = frames[0][0].timestamp + _POST_BOUNDARY_SETTLE_SECONDS
    accepted = set(stable_set)
    for position, (scan, image) in enumerate(frames):
        if position in stable_set or scan.timestamp < settled_at:
            continue
        matches = 0
        for anchor in anchors:
            if _structural_compatibility(
                image,
                frames[anchor][1],
                excluded_boxes=excluded_boxes,
                annotation_mode=annotation_mode,
                annotation_colours=annotation_colours,
                annotation_tolerance=annotation_tolerance,
            ):
                matches += 1
        required = 1 if len(anchors) == 1 else 2
        if matches >= required:
            accepted.add(position)
            anchor_support = max(int(augmented[anchor]) for anchor in anchors)
            augmented[position] = max(1, anchor_support)
    return sorted(accepted), augmented


def _select_final_candidate(
    frames: Sequence[tuple[_ScanFrame, "np.ndarray"]],
    *,
    excluded_boxes: Sequence[tuple[int, int, int, int]],
    watermark_boxes: Sequence[tuple[int, int, int, int]] = (),
    annotation_mode: str = "auto",
    annotation_colours: Sequence[BGRColor] = (),
    annotation_tolerance: int = 24,
    stability_support_override: "np.ndarray | None" = None,
) -> int:
    if len(frames) <= 1:
        return 0
    analysis_area = max(1, int(frames[0][0].gray.size))
    edges = np.asarray(
        [item[0].edge_pixels / analysis_area for item in frames],
        dtype=np.float64,
    )
    sharpness = np.asarray([item[0].sharpness for item in frames], dtype=np.float64)
    finite_sharpness = sharpness[np.isfinite(sharpness)]
    if finite_sharpness.size >= 2:
        sharp_low = float(np.percentile(finite_sharpness, 10))
        sharp_high = float(np.percentile(finite_sharpness, 90))
        sharp_span = sharp_high - sharp_low
        if sharp_span > max(1e-6, abs(sharp_low) * 1e-6):
            normalized_sharpness = np.clip(
                (sharpness - sharp_low) / sharp_span,
                0.0,
                1.0,
            )
        else:
            normalized_sharpness = np.zeros_like(sharpness)
    else:
        normalized_sharpness = np.zeros_like(sharpness)
    first = frames[0][1]
    height, width = first.shape[:2]
    valid = np.ones((height, width), dtype=bool)
    for left, top, right, bottom in excluded_boxes:
        valid[top:bottom, left:right] = False
    annotation_penalties: list[float] = []
    first_values = first.astype(np.int16)
    for _scan, frame in frames:
        frame_values = frame.astype(np.int16)
        difference = np.max(np.abs(frame_values - first_values), axis=2) >= 24
        colour_ink = _annotation_colour_mask(
            frame,
            mode=annotation_mode,
            colours=annotation_colours,
            tolerance=annotation_tolerance,
            excluded_boxes=excluded_boxes,
        ).astype(bool)
        if annotation_mode == "auto":
            newly_bright = (
                frame_values[:, :, 0] - first_values[:, :, 0] >= 24
            ) | (frame_values[:, :, 1] - first_values[:, :, 1] >= 24)
            colour_ink &= newly_bright
        annotation_penalties.append(
            float(np.count_nonzero(difference & colour_ink & valid))
            / max(1, int(np.count_nonzero(valid)))
        )
    # Printed animation objects add stable structural edges.  Handwriting can
    # also add edges.  Penalising only *new* blue/cyan pixels preserves template
    # graphics that were present from page birth while preferring the latest
    # print-complete frame before the lecturer starts writing.
    scores = (
        edges
        # Laplacian variance is resolution/content dependent and routinely
        # reaches 10,000 on 1080p text.  Multiplying the raw value allowed a
        # sharp handwritten/watermarked terminal frame to overwhelm every
        # cleanliness signal.  Candidate-relative robust normalization keeps
        # sharpness as a small tie-breaker only.
        + 0.006 * normalized_sharpness
        - 2.40 * np.asarray(annotation_penalties)
    )
    if stability_support_override is None:
        stability_support, _quiet_pairs = _candidate_stability_support(
            frames, excluded_boxes
        )
    else:
        stability_support = np.asarray(stability_support_override, dtype=np.int16)
        if stability_support.shape != (len(frames),):
            raise ValueError("stability support override length does not match frames")
    stable_candidates = np.flatnonzero(stability_support > 0)
    candidate_pool = (
        stable_candidates if stable_candidates.size else np.arange(len(frames))
    )
    scores = scores + 0.0015 * stability_support
    best = float(scores[candidate_pool].max())
    # Moving white text adds real edges, so the clean candidate can sit just
    # outside the ordinary tie window even though its stable slide structure
    # is otherwise identical.  A slightly wider window is used only when a
    # watermark search region exists; the temporal motion test below then
    # decides between those structurally close candidates.
    eligibility_margin = 0.005 if watermark_boxes else 0.0035
    eligible = candidate_pool[
        scores[candidate_pool] >= best - eligibility_margin
    ]
    if eligible.size > 1 and watermark_boxes:
        sample_width = min(320, width)
        sample_height = max(1, int(round(height * sample_width / max(1, width))))
        bright_masks: list[np.ndarray] = []
        for _scan, frame in frames:
            sample = cv2.resize(
                frame,
                (sample_width, sample_height),
                interpolation=cv2.INTER_AREA,
            )
            hsv = cv2.cvtColor(sample, cv2.COLOR_BGR2HSV)
            bright_masks.append((hsv[:, :, 2] >= 155) & (hsv[:, :, 1] <= 105))
        bright_stack = np.stack(bright_masks, axis=0)
        bright_support = np.count_nonzero(bright_stack, axis=0)
        low_occupancy = bright_support <= max(
            2, int(math.ceil(len(frames) * 0.45))
        )
        bright_motion = np.zeros((sample_height, sample_width), dtype=np.uint8)
        motion_kernel = np.ones((3, 7), np.uint8)
        for previous_bright, current_bright in zip(
            bright_stack,
            bright_stack[1:],
        ):
            changed = cv2.dilate(
                np.logical_xor(previous_bright, current_bright).astype(np.uint8),
                motion_kernel,
                iterations=1,
            )
            bright_motion += changed
        moving_transient = low_occupancy & (bright_motion >= 2)
        sample_region = np.zeros((sample_height, sample_width), dtype=bool)
        scale_x = sample_width / max(1, width)
        scale_y = sample_height / max(1, height)
        for left, top, right, bottom in watermark_boxes:
            sample_left = max(0, min(sample_width - 1, int(math.floor(left * scale_x))))
            sample_top = max(0, min(sample_height - 1, int(math.floor(top * scale_y))))
            sample_right = max(
                sample_left + 1,
                min(sample_width, int(math.ceil(right * scale_x))),
            )
            sample_bottom = max(
                sample_top + 1,
                min(sample_height, int(math.ceil(bottom * scale_y))),
            )
            sample_region[
                sample_top:sample_bottom,
                sample_left:sample_right,
            ] = True
        suspicious = moving_transient & sample_region
        if np.any(suspicious):
            pollution = np.asarray(
                [
                    int(np.count_nonzero(mask & suspicious))
                    for mask in bright_stack
                ],
                dtype=np.int32,
            )
            eligible_pollution = pollution[eligible]
            minimum_pollution = int(eligible_pollution.min())
            region_area = max(1, int(np.count_nonzero(sample_region)))
            meaningful_gap = max(8, int(round(region_area * 0.0004)))
            if int(eligible_pollution.max()) - minimum_pollution >= meaningful_gap:
                clean_limit = minimum_pollution + max(
                    3, int(round(region_area * 0.00008))
                )
                motion_clean = eligible[eligible_pollution <= clean_limit]
                if motion_clean.size:
                    eligible = motion_clean
        region = np.zeros((height, width), dtype=bool)
        for left, top, right, bottom in watermark_boxes:
            region[top:bottom, left:right] = True
        region_area = int(np.count_nonzero(region))
        if region_area:
            bright_counts: list[int] = []
            for _scan, frame in frames:
                hsv = cv2.cvtColor(frame, cv2.COLOR_BGR2HSV)
                bright_overlay = (hsv[:, :, 2] >= 180) & (hsv[:, :, 1] <= 80) & region
                bright_counts.append(int(np.count_nonzero(bright_overlay)))
            counts = np.asarray(bright_counts, dtype=np.float64)
            median = float(np.median(counts))
            mad = float(np.median(np.abs(counts - median)))
            abnormal_limit = median + max(6.0 * mad, 0.0015 * region_area, 500.0)
            watermark_clean = eligible[counts[eligible] <= abnormal_limit]
            # If every structurally complete candidate is bright, keep the
            # original decision.  This protects a legitimate late white title.
            if watermark_clean.size:
                eligible = watermark_clean
    return int(eligible[-1]) if eligible.size else len(frames) - 1


def _annotation_mask(
    frames: Sequence["np.ndarray"],
    *,
    excluded_boxes: Sequence[tuple[int, int, int, int]],
    annotation_mode: str = "auto",
    annotation_colours: Sequence[BGRColor] = (),
    annotation_tolerance: int = 24,
) -> "np.ndarray":
    height, width = frames[0].shape[:2]
    if len(frames) < 3:
        return np.zeros((height, width), dtype=np.uint8)
    first = frames[0]
    final = frames[-1]
    difference = np.max(
        np.abs(final.astype(np.int16) - first.astype(np.int16)), axis=2
    ).astype(np.uint8)
    raw = (difference >= 28).astype(np.uint8)
    for left, top, right, bottom in excluded_boxes:
        raw[top:bottom, left:right] = 0
    raw = cv2.morphologyEx(raw, cv2.MORPH_OPEN, np.ones((2, 2), np.uint8))
    raw = cv2.morphologyEx(raw, cv2.MORPH_CLOSE, np.ones((3, 3), np.uint8))
    count, labels, stats, _centroids = cv2.connectedComponentsWithStats(raw, 8)
    result = np.zeros_like(raw)
    final_hsv = cv2.cvtColor(final, cv2.COLOR_BGR2HSV)
    page_area = height * width
    for label in range(1, count):
        x, y, box_width, box_height, area = (int(value) for value in stats[label])
        if area < max(8, page_area // 500_000) or area > page_area * 0.08:
            continue
        component = labels[y : y + box_height, x : x + box_width] == label
        density = area / max(1, box_width * box_height)
        areas: list[int] = []
        for frame in frames:
            delta = np.max(
                np.abs(
                    frame[y : y + box_height, x : x + box_width].astype(np.int16)
                    - first[y : y + box_height, x : x + box_width].astype(np.int16)
                ),
                axis=2,
            )
            areas.append(int(np.count_nonzero((delta >= 28) & component)))
        increments = np.diff(np.asarray(areas, dtype=np.int32))
        meaningful = int(
            np.count_nonzero(increments >= max(3, int(round(area * 0.06))))
        )
        monotonic = float(np.count_nonzero(increments >= -max(2, area * 0.03))) / max(
            1, len(increments)
        )
        hsv = final_hsv[y : y + box_height, x : x + box_width]
        saturated_pixels = (hsv[:, :, 1] >= 95) & (hsv[:, :, 2] >= 75) & component
        saturated = float(np.count_nonzero(saturated_pixels)) / max(1, area)
        aspect = max(box_width, box_height) / max(1, min(box_width, box_height))
        elongated = aspect >= 4.0 and min(box_width, box_height) <= max(
            18, int(round(min(width, height) * 0.05))
        )
        # Achromatic additions are commonly late printed answers, subtitles or
        # bullet text.  Treating every thin monotonic shape as handwriting used
        # to erase legitimate white text.  Temporal cleanup is therefore gated
        # by real colour evidence; uncertain black/white marks are preserved.
        thin_growth = (
            meaningful >= 2
            and monotonic >= 0.68
            and density <= 0.58
            and saturated >= 0.18
        )
        colored_stroke = (
            meaningful >= 2 and saturated >= 0.45 and (density <= 0.42 or elongated)
        )
        moving_pointer = (
            area <= page_area * 0.0018 and saturated >= 0.55 and meaningful >= 1
        )
        if thin_growth or colored_stroke or moving_pointer:
            # Restore only the coloured ink and its anti-aliased fringe.  A
            # handwritten stroke can touch a newly printed white answer; using
            # the whole connected change component would erase both.
            ink = cv2.dilate(
                saturated_pixels.astype(np.uint8),
                np.ones((3, 3), np.uint8),
                iterations=1,
            ).astype(bool)
            ink &= component
            result[y : y + box_height, x : x + box_width][ink] = 255
    if annotation_mode == "manual" and annotation_colours:
        final_prior = _annotation_colour_mask(
            final,
            mode="manual",
            colours=annotation_colours,
            tolerance=annotation_tolerance,
            excluded_boxes=excluded_boxes,
        ) > 0
        initial_count = min(4, len(frames))
        initial_support = np.zeros((height, width), dtype=np.uint8)
        for frame in frames[:initial_count]:
            initial_support += (
                _annotation_colour_mask(
                    frame,
                    mode="manual",
                    colours=annotation_colours,
                    tolerance=annotation_tolerance,
                    excluded_boxes=excluded_boxes,
                )
                > 0
            ).astype(np.uint8)
        initially_persistent = initial_support >= max(
            2, int(math.ceil(initial_count * 0.75))
        )
        changed_from_birth = difference >= max(12, annotation_tolerance // 2)
        guided = final_prior & changed_from_birth & (~initially_persistent)
        if np.any(guided):
            result[guided] = 255
    if np.any(result):
        result = cv2.dilate(result, np.ones((3, 3), np.uint8), iterations=1)
    return result


def _low_frequency_fill_region(
    output: "np.ndarray", box: tuple[int, int, int, int]
) -> None:
    """Fill an always-occluded presenter box without hard seams or cloning text.

    Reconstruction is deliberately low-frequency because the real pixels do
    not exist in the video.  The result is visually quiet and all affected
    pixels remain reported as low confidence.
    """

    left, top, right, bottom = box
    height, width = output.shape[:2]
    if right <= left or bottom <= top:
        return
    scale = min(1.0, 360.0 / max(width, height))
    small_width = max(48, int(round(width * scale)))
    small_height = max(32, int(round(height * scale)))
    small = cv2.resize(
        output, (small_width, small_height), interpolation=cv2.INTER_AREA
    )
    small_mask = np.zeros((small_height, small_width), dtype=np.uint8)
    scaled_box = (
        max(0, min(small_width - 1, int(round(left * scale)))),
        max(0, min(small_height - 1, int(round(top * scale)))),
        max(1, min(small_width, int(round(right * scale)))),
        max(1, min(small_height, int(round(bottom * scale)))),
    )
    sl, st, sr, sb = scaled_box
    small_mask[st:sb, sl:sr] = 255
    radius = max(3, int(round(min(small_width, small_height) * 0.035)))
    filled_small = cv2.inpaint(small, small_mask, radius, cv2.INPAINT_NS)
    filled = cv2.resize(filled_small, (width, height), interpolation=cv2.INTER_CUBIC)
    patch = filled[top:bottom, left:right].astype(np.float32)
    original = output[top:bottom, left:right].astype(np.float32)
    box_height = bottom - top
    box_width = right - left
    feather = max(6, min(32, box_height // 5, box_width // 5))
    yy, xx = np.indices((box_height, box_width))
    distances: list[np.ndarray] = []
    if top > 0:
        distances.append(yy + 1)
    if left > 0:
        distances.append(xx + 1)
    if bottom < height:
        distances.append(box_height - yy)
    if right < width:
        distances.append(box_width - xx)
    distance = (
        np.minimum.reduce(distances).astype(np.float32)
        if distances
        else np.full((box_height, box_width), feather, dtype=np.float32)
    )
    alpha = np.clip(distance / max(1.0, float(feather)), 0.0, 1.0)[..., None]
    output[top:bottom, left:right] = np.clip(
        original * (1.0 - alpha) + patch * alpha, 0, 255
    ).astype(np.uint8)


def _temporal_modal_medoid(
    stack: "np.ndarray",
) -> tuple["np.ndarray", "np.ndarray"]:
    """Return an observed colour from the dominant temporal colour mode.

    A per-channel median can synthesize a colour that never existed and is
    easily biased by a translucent moving watermark.  Here colours are first
    clustered into compact RGB cubes, the most occupied cube wins per pixel,
    and the closest *observed* sample becomes the representative medoid.  A
    low-occupancy overlay therefore loses to the repeatedly revealed real
    slide pixel.
    """

    if stack.shape[0] == 1:
        return stack[0].copy(), np.ones(stack.shape[1:3], dtype=np.uint8)
    codes = (
        ((stack[:, :, :, 0].astype(np.uint16) >> 4) << 8)
        | ((stack[:, :, :, 1].astype(np.uint16) >> 4) << 4)
        | (stack[:, :, :, 2].astype(np.uint16) >> 4)
    )
    ordered_codes = np.sort(codes, axis=0)
    current_code = ordered_codes[0].copy()
    current_count = np.ones(current_code.shape, dtype=np.uint8)
    best_code = current_code.copy()
    best_count = current_count.copy()
    for index in range(1, ordered_codes.shape[0]):
        same = ordered_codes[index] == current_code
        current_count = np.where(same, current_count + 1, 1).astype(np.uint8)
        current_code = ordered_codes[index]
        better = current_count > best_count
        best_count[better] = current_count[better]
        best_code[better] = current_code[better]
    blue = ((best_code >> 8) & 15).astype(np.int16) * 16 + 8
    green = ((best_code >> 4) & 15).astype(np.int16) * 16 + 8
    red = (best_code & 15).astype(np.int16) * 16 + 8
    centre = np.stack((blue, green, red), axis=2)
    # Keep the working set bounded for 1080p/4K lectures: three channel-wise
    # int16 passes need far less peak memory than one N×H×W×3 int16 temporary.
    distance = np.zeros(codes.shape, dtype=np.int16)
    for channel in range(3):
        distance += np.abs(
            stack[:, :, :, channel].astype(np.int16)
            - centre[None, :, :, channel]
        )
    distance[codes != best_code[None, ...]] = np.iinfo(np.int16).max
    best = np.argmin(distance, axis=0)
    rows, columns = np.indices(best.shape)
    return stack[best, rows, columns], best_count


def _fill_region_with_colour(
    output: "np.ndarray",
    box: tuple[int, int, int, int],
    colour: BGRColor,
) -> None:
    """Cover a declared fixed region with a feathered user-selected colour."""

    left, top, right, bottom = box
    if right <= left or bottom <= top:
        return
    box_height = bottom - top
    box_width = right - left
    feather = max(2, min(14, box_height // 8, box_width // 8))
    yy, xx = np.indices((box_height, box_width))
    distance = np.minimum.reduce((yy + 1, xx + 1, box_height - yy, box_width - xx))
    # Even the outermost pixel receives a strong replacement.  The feather is
    # only for hiding a hard seam, not for leaving a visible watermark rim.
    alpha = 0.72 + 0.28 * np.clip(distance / max(1.0, float(feather)), 0.0, 1.0)
    original = output[top:bottom, left:right].astype(np.float32)
    fill = np.asarray(colour, dtype=np.float32)[None, None, :]
    output[top:bottom, left:right] = np.clip(
        original * (1.0 - alpha[..., None]) + fill * alpha[..., None],
        0,
        255,
    ).astype(np.uint8)


def _temporal_restore_region(
    output: "np.ndarray",
    frames: Sequence["np.ndarray"],
    box: tuple[int, int, int, int],
    *,
    variation_threshold: int,
    force: bool,
    force_all: bool = False,
    selected_index: int | None = None,
    fallback_fill: str = "background",
    fallback_colour: BGRColor = (255, 255, 255),
) -> tuple[int, int]:
    left, top, right, bottom = box
    if right <= left or bottom <= top or len(frames) < 2:
        return 0, 0
    if force_all:
        stack = np.stack([frame[top:bottom, left:right] for frame in frames], axis=0)
        variation = np.max(
            stack.max(axis=0).astype(np.int16) - stack.min(axis=0).astype(np.int16),
            axis=2,
        )
        motion = (variation >= 13).astype(np.uint8)
        box_height = bottom - top
        box_width = right - left
        kernel_size = max(5, int(round(min(box_height, box_width) * 0.035)))
        if kernel_size % 2 == 0:
            kernel_size += 1
        kernel_size = min(25, kernel_size)
        kernel = np.ones((kernel_size, kernel_size), np.uint8)
        motion = cv2.morphologyEx(motion, cv2.MORPH_CLOSE, kernel)
        motion = cv2.dilate(motion, kernel, iterations=2)
        count, labels, stats, _centroids = cv2.connectedComponentsWithStats(motion, 8)
        filtered = np.zeros_like(motion)
        minimum_component = max(48, int(box_height * box_width * 0.003))
        for label in range(1, count):
            if int(stats[label, cv2.CC_STAT_AREA]) >= minimum_component:
                filtered[labels == label] = 255
        if int(np.count_nonzero(filtered)) < box_height * box_width * 0.015:
            filtered[:, :] = 255
        fill = output.copy()
        _low_frequency_fill_region(fill, box)
        global_mask = np.zeros(output.shape[:2], dtype=np.uint8)
        global_mask[top:bottom, left:right] = filtered
        output[global_mask > 0] = fill[global_mask > 0]
        return 0, int(np.count_nonzero(global_mask))
    if not force and len(frames) < 3:
        # With only two or three observations there is insufficient temporal
        # evidence to distinguish a moving watermark from a late printed
        # title.  Preserving authentic slide content is the safer choice.
        return 0, 0
    restored_pixels = 0
    low_confidence = 0
    unsupported_mask = np.zeros(output.shape[:2], dtype=np.uint8)
    tile_height = 192
    kernel = np.ones((3, 3), np.uint8)
    for tile_top in range(top, bottom, tile_height):
        tile_bottom = min(bottom, tile_top + tile_height)
        stack = np.stack(
            [frame[tile_top:tile_bottom, left:right] for frame in frames], axis=0
        )
        minimum = stack.min(axis=0)
        maximum = stack.max(axis=0)
        variation = np.max(maximum.astype(np.int16) - minimum.astype(np.int16), axis=2)
        median = np.median(stack, axis=0).astype(np.uint8)
        modal_medoid, modal_support = _temporal_modal_medoid(stack)
        target = output[tile_top:tile_bottom, left:right]
        print_protection = np.zeros(target.shape[:2], dtype=np.uint8)
        slow_bright_trajectory = np.zeros(target.shape[:2], dtype=bool)
        slow_trajectory_support = np.zeros(target.shape[:2], dtype=bool)
        if force:
            dynamic = (variation >= variation_threshold).astype(np.uint8)
        else:
            # Only alter pixels that look like a transient bright text overlay
            # in the selected authentic frame.  Using the entire temporal
            # variance mask destroys legitimate animated bullet text.
            selected_delta = np.max(
                np.abs(target.astype(np.int16) - median.astype(np.int16)), axis=2
            )
            hsv = cv2.cvtColor(target, cv2.COLOR_BGR2HSV)
            similarity = (
                np.max(
                    np.abs(stack.astype(np.int16) - target[None, ...].astype(np.int16)),
                    axis=3,
                )
                <= 12
            )
            support = np.count_nonzero(similarity, axis=0)
            low_occupancy = support <= max(
                2, int(math.ceil(len(frames) * 0.35))
            )
            stable_support = support >= max(3, int(math.ceil(len(frames) * 0.25)))
            if selected_index is not None and 0 <= selected_index < len(frames) - 1:
                future = similarity[selected_index:]
                stable_support |= np.mean(future, axis=0) >= 0.80
            anchor_index = (
                selected_index
                if selected_index is not None and 0 <= selected_index < len(frames)
                else len(frames) - 1
            )
            # Late printed answers often fade in only during the final second.
            # Their exact anti-aliased RGB values can fluctuate by more than 12
            # levels, so exact-pixel support alone misclassifies them as a
            # transient bright watermark.  Confirm a neutral bright glyph by
            # its consecutive presence in the recent dense candidates instead.
            consecutive = np.zeros(target.shape[:2], dtype=np.uint8)
            active = np.ones(target.shape[:2], dtype=bool)
            recent_start = max(0, anchor_index - 5)
            for candidate in reversed(stack[recent_start : anchor_index + 1]):
                channel_min = np.min(candidate, axis=2)
                channel_max = np.max(candidate, axis=2)
                neutral_bright = (
                    (channel_max >= 95)
                    & (channel_min >= 58)
                    & (
                        (channel_max.astype(np.int16) - channel_min.astype(np.int16))
                        <= 105
                    )
                )
                active &= neutral_bright
                consecutive += active.astype(np.uint8)
            target_bright = (hsv[:, :, 2] >= 112) & (hsv[:, :, 1] <= 118)
            target_gray = cv2.cvtColor(target, cv2.COLOR_BGR2GRAY)
            glyph_edges = cv2.dilate(
                (cv2.Canny(target_gray, 45, 135) > 0).astype(np.uint8),
                np.ones((3, 3), np.uint8),
                iterations=1,
            ).astype(bool)
            recent_required = 2 if len(frames) <= 3 else 3
            persistent_bright_candidate = (
                target_bright
                & (consecutive >= recent_required)
                & glyph_edges
                & ((len(frames) <= 3) | (~low_occupancy))
            )
            recent_similarity = np.max(
                np.abs(
                    stack[recent_start : anchor_index + 1].astype(np.int16)
                    - target[None, ...].astype(np.int16)
                ),
                axis=3,
            ) <= 18
            recent_consecutive = np.zeros(target.shape[:2], dtype=np.uint8)
            recent_active = np.ones(target.shape[:2], dtype=bool)
            for candidate_similarity in reversed(recent_similarity):
                recent_active &= candidate_similarity
                recent_consecutive += recent_active.astype(np.uint8)
            local_motion_count = np.zeros(target.shape[:2], dtype=np.uint8)
            bright_motion_count = np.zeros(target.shape[:2], dtype=np.uint8)
            recent_candidates = stack[recent_start : anchor_index + 1]
            for previous_candidate, current_candidate in zip(
                recent_candidates,
                recent_candidates[1:],
            ):
                pair_motion = (
                    np.max(
                        np.abs(
                            current_candidate.astype(np.int16)
                            - previous_candidate.astype(np.int16)
                        ),
                        axis=2,
                    )
                    >= max(24, variation_threshold + 6)
                ).astype(np.uint8)
                # A scrolling glyph may leave one pixel white for several
                # samples even though the surrounding word keeps translating.
                # Count motion in a small neighbourhood so those pixels are
                # not mistaken for stationary late-printed text.
                pair_motion = cv2.dilate(
                    pair_motion,
                    np.ones((5, 9), np.uint8),
                    iterations=1,
                )
                local_motion_count += pair_motion
                previous_min = np.min(previous_candidate, axis=2)
                previous_max = np.max(previous_candidate, axis=2)
                current_min = np.min(current_candidate, axis=2)
                current_max = np.max(current_candidate, axis=2)
                previous_bright = (
                    (previous_max >= 95)
                    & (previous_min >= 58)
                    & (
                        previous_max.astype(np.int16)
                        - previous_min.astype(np.int16)
                        <= 105
                    )
                )
                current_bright = (
                    (current_max >= 95)
                    & (current_min >= 58)
                    & (
                        current_max.astype(np.int16)
                        - current_min.astype(np.int16)
                        <= 105
                    )
                )
                bright_motion = cv2.dilate(
                    np.logical_xor(previous_bright, current_bright).astype(np.uint8),
                    np.ones((5, 9), np.uint8),
                    iterations=1,
                )
                bright_motion_count += bright_motion
            # A long watermark can translate so slowly that the same white
            # glyph occupies a pixel in most recent candidates.  Recent-frame
            # persistence then looks exactly like printed text and the modal
            # medoid can even reinforce the watermark.  Track the complete
            # same-page history instead: a stationary title normally creates
            # one appearance boundary, while a scrolling word produces
            # repeated neighbouring XOR motion along its trajectory.
            bright_history: list[np.ndarray] = []
            for candidate in stack:
                channel_min = np.min(candidate, axis=2)
                channel_max = np.max(candidate, axis=2)
                bright_history.append(
                    (
                        (channel_max >= 95)
                        & (channel_min >= 58)
                        & (
                            channel_max.astype(np.int16)
                            - channel_min.astype(np.int16)
                            <= 105
                        )
                    )
                )
            trajectory_motion_count = np.zeros(target.shape[:2], dtype=np.uint16)
            trajectory_kernel = np.ones((5, 17), np.uint8)
            for previous_bright, current_bright in zip(
                bright_history,
                bright_history[1:],
            ):
                trajectory_motion = cv2.dilate(
                    np.logical_xor(previous_bright, current_bright).astype(np.uint8),
                    trajectory_kernel,
                    iterations=1,
                )
                trajectory_motion_count += trajectory_motion.astype(np.uint16)
            slow_trajectory_support = (
                (trajectory_motion_count >= 2)
                & (variation >= variation_threshold)
            )
            slow_bright_trajectory = target_bright & slow_trajectory_support
            # Full-history trajectory evidence is intentionally restricted to
            # the narrow top advertising strip.  Across ordinary slide body
            # text, line-by-line reveals can produce the same repeated XOR
            # count as a slow ticker; treating them globally damages authentic
            # questions and answers.  Outside this strip the established
            # conservative occupancy/modal logic remains in control.
            if top <= int(round(output.shape[0] * 0.05)):
                ticker_bottom = int(round(output.shape[0] * 0.10))
                slow_bright_trajectory &= (
                    np.arange(tile_top, tile_bottom)[:, None] < ticker_bottom
                )
            else:
                slow_bright_trajectory[:, :] = False
            # A very slow scrolling white watermark can keep the same glyph
            # pixel bright for three or more dense candidates.  Persistence
            # alone therefore cannot prove that it is late printed content;
            # require the surrounding neighbourhood to be stationary too.
            stable_bright_print = persistent_bright_candidate & (
                bright_motion_count <= 1
            ) & (~slow_bright_trajectory)
            modal_difference = np.max(
                np.abs(target.astype(np.int16) - modal_medoid.astype(np.int16)),
                axis=2,
            )
            exact_recent_required = 3 if len(frames) <= 6 else 4
            stable_recent_print = (
                (recent_consecutive >= exact_recent_required)
                & glyph_edges
                & (modal_difference >= max(12, variation_threshold - 2))
                & (local_motion_count <= 1)
            )
            stable_support |= stable_recent_print
            if np.any(stable_bright_print | stable_recent_print):
                print_protection = cv2.dilate(
                    (stable_bright_print | stable_recent_print).astype(np.uint8),
                    np.ones((5, 5), np.uint8),
                    iterations=1,
                )
            median_hsv = cv2.cvtColor(median, cv2.COLOR_BGR2HSV)
            modal_dominant = modal_support >= max(2, int(math.ceil(len(frames) * 0.28)))
            bright_overlay = (
                (hsv[:, :, 2] >= 118)
                & (hsv[:, :, 1] <= 105)
                & (
                    hsv[:, :, 2].astype(np.int16) - median_hsv[:, :, 2].astype(np.int16)
                    >= 10
                )
            )
            overlay_like = (
                (selected_delta >= variation_threshold)
                & (~stable_support)
                & (bright_overlay | (low_occupancy & modal_dominant))
            )
            dynamic = (
                ((variation >= variation_threshold) & overlay_like)
                | slow_bright_trajectory
            ).astype(np.uint8)
        if np.any(dynamic):
            dynamic = cv2.morphologyEx(dynamic, cv2.MORPH_CLOSE, kernel)
            dynamic = cv2.dilate(dynamic, np.ones((5, 5), np.uint8), iterations=2)
            if np.any(print_protection):
                protected = print_protection > 0
                if not force:
                    override_seed = slow_bright_trajectory.astype(np.uint8)
                    # Repeated motion may also come from text being revealed
                    # line by line.  Only a narrow top advertising band is
                    # allowed to override confirmed persistent print; elsewhere
                    # the late/stable print protection remains authoritative.
                    if top <= int(round(output.shape[0] * 0.05)):
                        ticker_bottom = int(round(output.shape[0] * 0.10))
                        override_seed = override_seed.copy()
                        override_seed[
                            np.arange(tile_top, tile_bottom) >= ticker_bottom,
                            :,
                        ] = 0
                    else:
                        override_seed[:, :] = 0
                    trajectory_override = cv2.dilate(
                        override_seed,
                        np.ones((5, 5), np.uint8),
                        iterations=1,
                    ).astype(bool)
                    protected &= ~trajectory_override
                dynamic[protected] = 0
        authentic = modal_medoid
        mask = dynamic.astype(bool)
        target[mask] = authentic[mask]
        if not force and np.any(slow_bright_trajectory):
            # Restore a trajectory pixel from one real candidate frame, never
            # from a per-channel minimum that could synthesize a colour which
            # did not exist in the video.  The lowest-V observation is the
            # strongest evidence for the background beneath a neutral bright
            # overlay.
            brightness = np.max(stack, axis=3)
            darkest_index = np.argmin(brightness, axis=0)
            rows, columns = np.indices(darkest_index.shape)
            darkest_real_pixel = stack[darkest_index, rows, columns]
            restore_support = slow_bright_trajectory.copy()
            # Modal restoration can introduce white fragments from other
            # positions of a very long ticker.  Broaden trajectory restoration
            # only inside the narrow top advertising band; doing this across
            # the whole watermark-search rectangle can mistake progressively
            # revealed exercise text for a moving overlay.
            if top <= int(round(output.shape[0] * 0.05)):
                ticker_bottom = int(round(output.shape[0] * 0.10))
                global_rows = np.arange(tile_top, tile_bottom)[:, None]
                top_ticker_band = global_rows < ticker_bottom
                restore_support |= slow_trajectory_support & top_ticker_band
            trajectory_restore = mask & cv2.dilate(
                restore_support.astype(np.uint8),
                np.ones((5, 5), np.uint8),
                iterations=1,
            ).astype(bool)
            target[trajectory_restore] = darkest_real_pixel[trajectory_restore]
            # A very long ticker may cover a coordinate throughout the whole
            # lifetime of a short slide, so even the darkest authentic sample
            # remains white.  Once repeated trajectory motion has proved that
            # those pixels belong to a moving overlay, model only the still-
            # bright residue from its immediate local background.  Stationary
            # printed titles never enter ``slow_bright_trajectory`` and remain
            # protected.
            restored_hsv = cv2.cvtColor(target, cv2.COLOR_BGR2HSV)
            unrevealed_residue = (
                trajectory_restore
                & (restored_hsv[:, :, 2] >= 118)
                & (restored_hsv[:, :, 1] <= 118)
            ).astype(np.uint8) * 255
            if np.any(unrevealed_residue):
                # This is a watermark-only mask, not the generic pointer mask:
                # do not restore bright horizontal/vertical lines from its
                # boundary because neighbouring ticker glyphs can be mistaken
                # for authentic rules.  Narrow Telea inpainting follows the
                # surrounding slide texture and removes anti-aliased rims.
                unrevealed_residue = cv2.dilate(
                    unrevealed_residue,
                    np.ones((3, 3), np.uint8),
                    iterations=1,
                )
                inpainted = cv2.inpaint(
                    target,
                    unrevealed_residue,
                    3,
                    cv2.INPAINT_TELEA,
                )
                target[unrevealed_residue > 0] = inpainted[
                    unrevealed_residue > 0
                ]
        restored_pixels += int(np.count_nonzero(mask))
        if force:
            unsupported = (~mask).astype(np.uint8) * 255
            unsupported_mask[tile_top:tile_bottom, left:right] = unsupported
            low_confidence += int(np.count_nonzero(unsupported))
    if force and np.any(unsupported_mask):
        # A user-declared fixed region may never reveal its real background.
        # Model it from the *surrounding full slide* (or an explicit colour)
        # and report every such pixel as low-confidence instead of pretending
        # it was recovered.
        filled = output.copy()
        if fallback_fill == "color":
            _fill_region_with_colour(filled, box, fallback_colour)
        else:
            _low_frequency_fill_region(filled, box)
        output[unsupported_mask > 0] = filled[unsupported_mask > 0]
    return restored_pixels, low_confidence


def _fill_mask_with_local_background(image: "np.ndarray", mask: "np.ndarray") -> None:
    points = cv2.findNonZero((mask > 0).astype(np.uint8))
    if points is None:
        return
    x, y, width, height = cv2.boundingRect(points)
    mask_pixels = mask > 0
    hsv = cv2.cvtColor(image, cv2.COLOR_BGR2HSV)
    bright_neutral = (hsv[:, :, 2] >= 170) & (hsv[:, :, 1] <= 90)
    horizontal_lines: list[tuple[int, "np.ndarray"]] = []
    vertical_lines: list[tuple[int, "np.ndarray"]] = []
    sample_width = 8
    sample_gap = 4
    image_height, image_width = mask.shape
    if (
        x >= sample_width + sample_gap
        and x + width + sample_gap + sample_width <= image_width
    ):
        for row in range(y, min(image_height, y + height)):
            left_band = bright_neutral[
                row, x - sample_gap - sample_width : x - sample_gap
            ]
            right_band = bright_neutral[
                row,
                x + width + sample_gap : x + width + sample_gap + sample_width,
            ]
            if float(left_band.mean()) >= 0.50 and float(right_band.mean()) >= 0.50:
                samples = np.concatenate(
                    [
                        image[
                            row,
                            x - sample_gap - sample_width : x - sample_gap,
                        ],
                        image[
                            row,
                            x
                            + width
                            + sample_gap : x
                            + width
                            + sample_gap
                            + sample_width,
                        ],
                    ],
                    axis=0,
                )
                horizontal_lines.append(
                    (row, np.median(samples, axis=0).astype(np.uint8))
                )
    if (
        y >= sample_width + sample_gap
        and y + height + sample_gap + sample_width <= image_height
    ):
        for column in range(x, min(image_width, x + width)):
            top_band = bright_neutral[
                y - sample_gap - sample_width : y - sample_gap, column
            ]
            bottom_band = bright_neutral[
                y + height + sample_gap : y + height + sample_gap + sample_width,
                column,
            ]
            if float(top_band.mean()) >= 0.50 and float(bottom_band.mean()) >= 0.50:
                samples = np.concatenate(
                    [
                        image[
                            y - sample_gap - sample_width : y - sample_gap,
                            column,
                        ],
                        image[
                            y
                            + height
                            + sample_gap : y
                            + height
                            + sample_gap
                            + sample_width,
                            column,
                        ],
                    ],
                    axis=0,
                )
                vertical_lines.append(
                    (column, np.median(samples, axis=0).astype(np.uint8))
                )
    radius = max(6, max(width, height))
    kernel_size = min(81, radius * 2 + 1)
    if kernel_size % 2 == 0:
        kernel_size += 1
    expanded = (
        cv2.dilate(
            (mask > 0).astype(np.uint8),
            np.ones((kernel_size, kernel_size), np.uint8),
            iterations=1,
        )
        > 0
    )
    ring = expanded & (~mask_pixels)
    if not np.any(ring):
        return
    ring_values = hsv[:, :, 2][ring]
    value_limit = float(np.percentile(ring_values, 65))
    background = ring & (hsv[:, :, 2] <= value_limit)
    if int(np.count_nonzero(background)) < 24:
        background = ring
    colour = np.median(image[background], axis=0).astype(np.uint8)
    image[mask_pixels] = colour
    boundary = mask_pixels & (
        cv2.erode(mask_pixels.astype(np.uint8), np.ones((3, 3), np.uint8)) == 0
    )
    if np.any(boundary):
        softened = cv2.GaussianBlur(image, (5, 5), 0)
        image[boundary] = softened[boundary]
    for row, line_colour in horizontal_lines:
        row_mask = mask_pixels[row]
        image[row, row_mask] = line_colour
    for column, line_colour in vertical_lines:
        column_mask = mask_pixels[:, column]
        image[column_mask, column] = line_colour


def _restore_transient_colored_marks(
    output: "np.ndarray",
    frames: Sequence["np.ndarray"],
    *,
    selected_index: int,
    excluded_boxes: Sequence[tuple[int, int, int, int]],
) -> int:
    if len(frames) < 2:
        return 0
    height, width = output.shape[:2]
    target_hsv = cv2.cvtColor(output, cv2.COLOR_BGR2HSV)
    # The lecture player uses a compact olive/yellow circular cursor.  Limit
    # automatic colour cleanup to that distinctive geometry; legal cyan or
    # yellow printed text must never be erased merely because of its hue.
    yellow = (
        (target_hsv[:, :, 0] >= 25)
        & (target_hsv[:, :, 0] <= 45)
        & (target_hsv[:, :, 1] >= 100)
        & (target_hsv[:, :, 2] >= 90)
    ).astype(np.uint8)
    for left, top, right, bottom in excluded_boxes:
        yellow[top:bottom, left:right] = 0
    yellow = cv2.morphologyEx(yellow, cv2.MORPH_CLOSE, np.ones((3, 3), np.uint8))
    count, labels, stats, _centroids = cv2.connectedComponentsWithStats(yellow, 8)
    pointer_mask = np.zeros((height, width), dtype=np.uint8)
    fallback_mask = np.zeros((height, width), dtype=np.uint8)
    page_area = height * width
    max_size = max(16, int(round(min(width, height) * 0.038)))
    anchor_index = (
        selected_index if 0 <= selected_index < len(frames) else len(frames) - 1
    )
    for label in range(1, count):
        x, y, box_width, box_height, area = (int(value) for value in stats[label])
        if area < max(45, page_area // 80_000) or area > page_area * 0.0012:
            continue
        if max(box_width, box_height) > max_size or min(box_width, box_height) < 6:
            continue
        aspect = max(box_width, box_height) / max(1, min(box_width, box_height))
        density = area / max(1, box_width * box_height)
        if aspect > 1.65 or density < 0.22:
            continue
        component = (labels[y : y + box_height, x : x + box_width] == label).astype(
            np.uint8
        )
        contours, _hierarchy = cv2.findContours(
            component * 255, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
        )
        perimeter = max(
            (cv2.arcLength(contour, True) for contour in contours), default=0.0
        )
        circularity = 4.0 * math.pi * area / max(1.0, perimeter * perimeter)
        seed_hsv = target_hsv[y : y + box_height, x : x + box_width][
            component.astype(bool)
        ]
        mean_hue = float(seed_hsv[:, 0].mean()) if seed_hsv.size else 0.0
        mean_value = float(seed_hsv[:, 2].mean()) if seed_hsv.size else 255.0
        olive_cursor = mean_hue >= 30.0 and mean_value <= 205.0
        round_seed = density >= 0.50 and circularity >= 0.42
        fragment_seed = circularity >= 0.20
        if not olive_cursor or not (round_seed or fragment_seed):
            continue
        padding = max(3, int(round(0.25 * max(box_width, box_height))))
        center = (x + box_width // 2, y + box_height // 2)
        axes = (
            max(3, int(math.ceil(box_width / 2.0 + padding))),
            max(3, int(math.ceil(box_height / 2.0 + padding))),
        )
        full_component = np.zeros((height, width), dtype=np.uint8)
        cv2.ellipse(full_component, center, axes, 0, 0, 360, 255, -1)
        for left, top, right, bottom in excluded_boxes:
            full_component[top:bottom, left:right] = 0
        component_pixels = full_component > 0
        if not np.any(component_pixels):
            continue
        cyan_context = (
            (target_hsv[:, :, 0] >= 70)
            & (target_hsv[:, :, 0] <= 112)
            & (target_hsv[:, :, 1] >= 80)
            & (target_hsv[:, :, 2] >= 55)
            & component_pixels
        )
        target_cyan_count = int(np.count_nonzero(cyan_context))
        multicolour_fragment = fragment_seed and target_cyan_count >= max(
            12, int(round(area * 0.12))
        )
        if not round_seed and not multicolour_fragment:
            continue

        left = max(0, center[0] - axes[0])
        top = max(0, center[1] - axes[1])
        right = min(width, center[0] + axes[0] + 1)
        bottom = min(height, center[1] + axes[1] + 1)
        local_component = component_pixels[top:bottom, left:right]
        yellow_presence: list[int] = []
        cyan_presence: list[int] = []
        for frame in frames:
            patch_hsv = cv2.cvtColor(frame[top:bottom, left:right], cv2.COLOR_BGR2HSV)
            patch_yellow = (
                (patch_hsv[:, :, 0] >= 25)
                & (patch_hsv[:, :, 0] <= 45)
                & (patch_hsv[:, :, 1] >= 100)
                & (patch_hsv[:, :, 2] >= 90)
            )
            yellow_presence.append(
                int(np.count_nonzero(patch_yellow & local_component))
            )
            patch_cyan = (
                (patch_hsv[:, :, 0] >= 70)
                & (patch_hsv[:, :, 0] <= 112)
                & (patch_hsv[:, :, 1] >= 80)
                & (patch_hsv[:, :, 2] >= 55)
            )
            cyan_presence.append(int(np.count_nonzero(patch_cyan & local_component)))
        # A lecturer may park the olive cursor for the entire lifetime of one
        # slide, so persistence alone cannot make it legitimate slide content.
        # Bright printed yellow icons are protected earlier by the olive/dim
        # colour gate; only the player's distinctive seed reaches this point.

        clean_threshold = max(3, int(round(area * 0.15)))
        clean_cyan_threshold = max(6, int(round(target_cyan_count * 0.15)))
        source_index = next(
            (
                index
                for index in sorted(
                    range(len(frames)),
                    key=lambda item: (abs(item - anchor_index), item),
                )
                if index != anchor_index
                and yellow_presence[index] <= clean_threshold
                and (
                    not multicolour_fragment
                    or cyan_presence[index] <= clean_cyan_threshold
                )
            ),
            None,
        )
        if source_index is None and multicolour_fragment:
            # Every observation still contains part of the coloured mark.  A
            # blind inpaint can destroy covered print, so retain the smaller
            # residue instead of replacing real text with a synthetic blob.
            continue
        pointer_mask[component_pixels] = 255
        if source_index is None:
            fallback_mask[component_pixels] = 255
        else:
            output[component_pixels] = frames[source_index][component_pixels]
    if not np.any(pointer_mask):
        return 0
    for left, top, right, bottom in excluded_boxes:
        pointer_mask[top:bottom, left:right] = 0
        fallback_mask[top:bottom, left:right] = 0
    restored = int(np.count_nonzero(pointer_mask))
    if np.any(fallback_mask):
        _fill_mask_with_local_background(output, fallback_mask)
    return restored


def _axis_line_coverage(component: "np.ndarray") -> float:
    height, width = component.shape
    horizontal_length = max(5, min(width, int(round(width * 0.22))))
    vertical_length = max(5, min(height, int(round(height * 0.22))))
    horizontal = cv2.morphologyEx(
        component,
        cv2.MORPH_OPEN,
        np.ones((1, horizontal_length), np.uint8),
    )
    vertical = cv2.morphologyEx(
        component,
        cv2.MORPH_OPEN,
        np.ones((vertical_length, 1), np.uint8),
    )
    lines = (horizontal > 0) | (vertical > 0)
    return float(np.count_nonzero(lines)) / max(1, int(np.count_nonzero(component)))


def _photo_like_patch(
    image: "np.ndarray", x: int, y: int, width: int, height: int
) -> bool:
    image_height, image_width = image.shape[:2]
    margin = max(4, int(round(min(image_width, image_height) * 0.008)))
    left = max(0, x - margin)
    top = max(0, y - margin)
    right = min(image_width, x + width + margin)
    bottom = min(image_height, y + height + margin)
    patch = image[top:bottom, left:right]
    if patch.size == 0:
        return False
    sample = cv2.resize(patch, (40, 40), interpolation=cv2.INTER_AREA)
    quantized = (sample // 32).reshape(-1, 3)
    colours = int(np.unique(quantized, axis=0).shape[0])
    colour_std = float(np.mean(sample.reshape(-1, 3).std(axis=0)))
    gray = cv2.cvtColor(sample, cv2.COLOR_BGR2GRAY)
    edge_ratio = float(np.count_nonzero(cv2.Canny(gray, 45, 135))) / gray.size
    return bool(colours >= 24 and colour_std >= 24.0 and edge_ratio >= 0.045)


def _colored_handwriting_mask(
    image: "np.ndarray",
    *,
    excluded_boxes: Sequence[tuple[int, int, int, int]],
    annotation_mode: str = "auto",
    annotation_colours: Sequence[BGRColor] = (),
    annotation_tolerance: int = 24,
) -> "np.ndarray":
    """Find high-confidence cyan/blue pen strokes without erasing print.

    Printed text is protected when glyphs form an aligned row/column, when a
    component is a clean axis-aligned diagram line, or when it lies in a
    photo-like region.  The remaining thick, isolated, irregular components
    match circles, underlines, arrows, stars and handwritten notes.
    """

    height, width = image.shape[:2]
    if annotation_mode == "off":
        return np.zeros((height, width), dtype=np.uint8)
    cyan = (
        _annotation_colour_mask(
            image,
            mode=annotation_mode,
            colours=annotation_colours,
            tolerance=annotation_tolerance,
            excluded_boxes=excluded_boxes,
        )
        > 0
    ).astype(np.uint8)
    cyan = cv2.morphologyEx(cyan, cv2.MORPH_CLOSE, np.ones((2, 2), np.uint8))
    if annotation_mode == "manual":
        # A picked colour is an explicit user prior, so retain its complete
        # anti-aliased stroke rather than requiring the cyan-specific geometry
        # heuristics below.  Persistent first-platform print is protected by
        # _persistent_content_protection before this mask is applied.
        result = cv2.dilate(cyan, np.ones((3, 3), np.uint8), iterations=1) * 255
        for left, top, right, bottom in excluded_boxes:
            result[top:bottom, left:right] = 0
        return result
    count, labels, stats, centroids = cv2.connectedComponentsWithStats(cyan, 8)
    page_area = height * width
    compact: list[int] = []
    protected: set[int] = set()
    component_info: dict[int, tuple[int, int, int, int, int, float]] = {}
    for label in range(1, count):
        x, y, box_width, box_height, area = (int(value) for value in stats[label])
        if area < max(8, page_area // 500_000):
            continue
        density = area / max(1, box_width * box_height)
        component_info[label] = (x, y, box_width, box_height, area, density)
        aspect = max(box_width, box_height) / max(1, min(box_width, box_height))
        if (
            6 <= box_width <= max(72, int(width * 0.06))
            and 6 <= box_height <= max(72, int(height * 0.08))
            and aspect <= 2.4
            and 0.10 <= density <= 0.88
        ):
            compact.append(label)
    # Rows of similarly sized glyphs and columns of repeated template icons are
    # overwhelmingly likely to be printed content rather than freehand ink.
    for label in compact:
        x, y, box_width, box_height, _area, _density = component_info[label]
        center_x, center_y = centroids[label]
        row_matches = 0
        column_matches = 0
        for other in compact:
            if other == label:
                continue
            ox, oy, ow, oh, _oa, _od = component_info[other]
            other_x, other_y = centroids[other]
            height_ratio = max(box_height, oh) / max(1, min(box_height, oh))
            width_ratio = max(box_width, ow) / max(1, min(box_width, ow))
            if (
                height_ratio <= 1.75
                and abs(center_y - other_y) <= max(5.0, 0.32 * max(box_height, oh))
                and abs(center_x - other_x) >= min(box_width, ow) * 0.75
            ):
                row_matches += 1
            if (
                width_ratio <= 1.75
                and abs(center_x - other_x) <= max(5.0, 0.32 * max(box_width, ow))
                and abs(center_y - other_y) >= min(box_height, oh) * 0.75
            ):
                column_matches += 1
        if row_matches >= 2 or column_matches >= 1:
            protected.add(label)
    result = np.zeros((height, width), dtype=np.uint8)
    for label, (x, y, box_width, box_height, area, density) in component_info.items():
        if label in protected or area > page_area * 0.045:
            continue
        component = (labels[y : y + box_height, x : x + box_width] == label).astype(
            np.uint8
        )
        distance = cv2.distanceTransform(component, cv2.DIST_L2, 3)
        radii = distance[component > 0]
        stroke_radius = float(np.percentile(radii, 70)) if radii.size else 0.0
        axis_coverage = _axis_line_coverage(component)
        aspect = max(box_width, box_height) / max(1, min(box_width, box_height))
        clean_diagram_line = (axis_coverage >= 0.62 and stroke_radius <= 2.15) or (
            axis_coverage >= 0.86 and density >= 0.62
        )
        if clean_diagram_line:
            continue
        if area >= max(1200, int(page_area * 0.003)) and _photo_like_patch(
            image, x, y, box_width, box_height
        ):
            continue
        long_irregular_stroke = (
            aspect >= 3.0
            and min(box_width, box_height) <= max(55, int(min(width, height) * 0.06))
            and axis_coverage < 0.62
        )
        thick_scribble = (
            area >= max(70, page_area // 30_000)
            and density <= 0.58
            and stroke_radius >= 1.35
        )
        isolated_large_ink = (
            area >= max(420, page_area // 5_000)
            and stroke_radius >= 1.85
            and density <= 0.66
        )
        if long_irregular_stroke or thick_scribble or isolated_large_ink:
            result[y : y + box_height, x : x + box_width][component > 0] = 255
    if np.any(result):
        result = cv2.dilate(result, np.ones((5, 5), np.uint8), iterations=1)
    for left, top, right, bottom in excluded_boxes:
        result[top:bottom, left:right] = 0
    return result


def _persistent_content_protection(
    final: "np.ndarray", frames: Sequence["np.ndarray"]
) -> "np.ndarray":
    """Protect coloured print that remains at the same coordinates over time.

    Cyan is common for both lecturer ink and legitimate course-template text.
    Geometry alone cannot always separate them.  When a suspected coloured
    component is already present, pixel-consistent, in most stable candidates,
    it is safer to keep it as printed content than to erase it destructively.
    """

    height, width = final.shape[:2]
    if len(frames) < 3:
        return np.zeros((height, width), dtype=np.uint8)
    final_values = final.astype(np.int16)
    support = np.zeros((height, width), dtype=np.uint16)
    initial_support = np.zeros((height, width), dtype=np.uint8)
    usable = 0
    initial_count = min(4, len(frames))
    for index, frame in enumerate(frames):
        if frame.shape != final.shape:
            continue
        difference = np.max(
            np.abs(frame.astype(np.int16) - final_values),
            axis=2,
        )
        similar = difference <= 24
        support += similar.astype(np.uint16)
        if index < initial_count:
            initial_support += similar.astype(np.uint8)
        usable += 1
    if usable < 3:
        return np.zeros((height, width), dtype=np.uint8)
    required = max(2, int(math.ceil(usable * 0.70)))
    initial_required = max(2, int(math.ceil(initial_count * 0.75)))
    temporal = (support >= required) & (initial_support >= initial_required)
    hsv = cv2.cvtColor(final, cv2.COLOR_BGR2HSV)
    gray = cv2.cvtColor(final, cv2.COLOR_BGR2GRAY)
    glyph_edges = cv2.dilate(
        (cv2.Canny(gray, 45, 135) > 0).astype(np.uint8),
        np.ones((3, 3), np.uint8),
        iterations=1,
    ).astype(bool)
    coloured_foreground = (hsv[:, :, 1] >= 48) & (hsv[:, :, 2] >= 45)
    neutral_glyph = glyph_edges & (
        ((hsv[:, :, 2] >= 110) & (hsv[:, :, 1] <= 105))
        | (hsv[:, :, 2] <= 100)
    )
    protected = (temporal & (coloured_foreground | neutral_glyph)).astype(np.uint8) * 255
    return cv2.dilate(protected, np.ones((5, 5), np.uint8), iterations=1)


def _restore_annotation_mask(
    output: "np.ndarray",
    frames: Sequence["np.ndarray"],
    mask: "np.ndarray",
    *,
    annotation_mode: str = "auto",
    annotation_colours: Sequence[BGRColor] = (),
    annotation_tolerance: int = 24,
) -> int:
    if not np.any(mask):
        return 0
    height, width = output.shape[:2]
    restored = 0
    tile_height = 160
    for top in range(0, height, tile_height):
        bottom = min(height, top + tile_height)
        local_mask = mask[top:bottom] > 0
        if not np.any(local_mask):
            continue
        stack = np.stack([frame[top:bottom] for frame in frames], axis=0)
        median = np.median(stack, axis=0).astype(np.uint8)
        distance = np.sum(
            np.abs(stack.astype(np.int16) - median[None, ...].astype(np.int16)),
            axis=3,
        )
        best = np.argmin(distance, axis=0)
        rows, columns = np.indices(best.shape)
        authentic = stack[best, rows, columns]
        ink_stack = np.stack(
            [
                _annotation_colour_mask(
                    frame,
                    mode=annotation_mode,
                    colours=annotation_colours,
                    tolerance=annotation_tolerance,
                )
                > 0
                for frame in stack
            ],
            axis=0,
        )
        non_ink = ~ink_stack
        has_non_ink = np.any(non_ink, axis=0)
        first_non_ink = np.argmax(non_ink, axis=0)
        clean_candidate = stack[first_non_ink, rows, columns]
        authentic[has_non_ink] = clean_candidate[has_non_ink]
        target = output[top:bottom]
        target[local_mask] = authentic[local_mask]
        restored += int(np.count_nonzero(local_mask))
    # If every observed frame contains the same handwritten stroke, real
    # pixels do not exist.  Remove only the still-cyan residue with a narrow
    # local inpaint; all non-cyan printed pixels remain untouched.
    residue = (
        (mask > 0)
        & (
            _annotation_colour_mask(
                output,
                mode=annotation_mode,
                colours=annotation_colours,
                tolerance=annotation_tolerance,
            )
            > 0
        )
    ).astype(np.uint8) * 255
    if np.any(residue):
        inpainted = cv2.inpaint(output, residue, 2, cv2.INPAINT_TELEA)
        output[residue > 0] = inpainted[residue > 0]
    return restored


def _clean_slide(
    candidates: Sequence[tuple[_ScanFrame, "np.ndarray"]],
    *,
    watermark_boxes: Sequence[tuple[float, float, float, float]],
    fixed_boxes: Sequence[tuple[float, float, float, float]],
    presenter_boxes: Sequence[tuple[float, float, float, float]],
    watermark_hint: str,
    annotation_mode: str = "auto",
    annotation_colours: Sequence[BGRColor] = (),
    annotation_tolerance: int = 24,
    fixed_fill: str = "temporal",
    fixed_fill_colour: BGRColor = (255, 255, 255),
    enhancement_mode: str = "off",
) -> _CleanResult | None:
    ordered = sorted(candidates, key=lambda item: item[0].frame_index)
    if not ordered:
        return None
    height, width = ordered[0][1].shape[:2]
    watermark_pixels = 0
    low_confidence = 0
    annotation_excluded = _scaled_boxes([*fixed_boxes, *presenter_boxes], width, height)
    watermark_regions = _scaled_boxes(watermark_boxes, width, height)
    compatibility_excluded = list(annotation_excluded)
    page_area = max(1, width * height)
    compatibility_excluded.extend(
        box
        for box in watermark_regions
        if (box[2] - box[0]) * (box[3] - box[1]) <= page_area * 0.48
    )
    # Watermark *search* rectangles can cover the title band or even the full
    # page.  They are not proof that every pixel there is watermark, so using
    # them as stability exclusions hides real slide changes.  Only explicit
    # fixed/presenter regions are excluded; compact moving overlays are handled
    # by the block vote inside the stability detector.
    stability_support, quiet_pairs = _candidate_stability_support(
        ordered, annotation_excluded
    )
    stable_positions = [int(position) for position in np.flatnonzero(stability_support > 0)]
    if stable_positions:
        settled_at = ordered[0][0].timestamp + _POST_BOUNDARY_SETTLE_SECONDS
        settled_positions = [
            position
            for position in stable_positions
            if ordered[position][0].timestamp >= settled_at
        ]
        if settled_positions:
            stable_positions = settled_positions
            settled_set = set(stable_positions)
            stability_support = stability_support.copy()
            for position in range(len(ordered)):
                if position not in settled_set:
                    stability_support[position] = 0
    # Three or more decoded candidates without any stable platform describe a
    # transition/animation fragment, not a safe slide.  Such fragments must be
    # merged away by the timeline rather than exported as half-switched pages.
    if len(ordered) >= _STABLE_WINDOW_CANDIDATES and not stable_positions:
        return None
    if stable_positions:
        history_positions, augmented_support = _reinclude_structurally_compatible_positions(
            ordered,
            stability_support,
            excluded_boxes=compatibility_excluded,
            annotation_mode=annotation_mode,
            annotation_colours=annotation_colours,
            annotation_tolerance=annotation_tolerance,
        )
    else:
        history_positions = list(range(len(ordered)))
        augmented_support = stability_support.copy()
    history_ordered = [ordered[position] for position in history_positions]
    transition_rejected = len(ordered) - len(history_ordered)
    history_images = [item[1] for item in history_ordered]
    final_index = _select_final_candidate(
        history_ordered,
        excluded_boxes=annotation_excluded,
        watermark_boxes=watermark_regions,
        annotation_mode=annotation_mode,
        annotation_colours=annotation_colours,
        annotation_tolerance=annotation_tolerance,
        stability_support_override=augmented_support[history_positions],
    )
    final = history_images[final_index].copy()
    authentic_selected = final.copy()
    selected_position = history_positions[final_index]

    # Every observation below has either direct continuous stability support or
    # is structurally matched to such a platform after the colour prior is
    # ignored.  Keeping the whole same-page history is what lets a clean frame
    # from 10–25 s repair a dirty terminal frame at 96 s.
    restoration_frames = history_images
    restoration_index = final_index
    annotation_history = [
        ordered[position][1]
        for position in history_positions
        if position <= selected_position
    ]
    if not annotation_history:
        annotation_history = [final]
    temporal_annotation = _annotation_mask(
        annotation_history,
        excluded_boxes=annotation_excluded,
        annotation_mode=annotation_mode,
        annotation_colours=annotation_colours,
        annotation_tolerance=annotation_tolerance,
    )
    static_annotation = _colored_handwriting_mask(
        final,
        excluded_boxes=annotation_excluded,
        annotation_mode=annotation_mode,
        annotation_colours=annotation_colours,
        annotation_tolerance=annotation_tolerance,
    )
    persistent_print = _persistent_content_protection(final, history_images)
    static_annotation[persistent_print > 0] = 0
    annotation = cv2.bitwise_or(temporal_annotation, static_annotation)
    colour_prior_before = _annotation_colour_mask(
        final,
        mode=annotation_mode,
        colours=annotation_colours,
        tolerance=annotation_tolerance,
        excluded_boxes=annotation_excluded,
    )
    colour_guided_pixels = int(
        np.count_nonzero((annotation > 0) & (colour_prior_before > 0))
    )
    annotation_pixels = _restore_annotation_mask(
        final,
        restoration_frames,
        annotation,
        annotation_mode=annotation_mode,
        annotation_colours=annotation_colours,
        annotation_tolerance=annotation_tolerance,
    )
    threshold = 16 if watermark_hint.strip() else 21
    for box in _scaled_boxes(watermark_boxes, width, height):
        restored, uncertain = _temporal_restore_region(
            final,
            restoration_frames,
            box,
            variation_threshold=threshold,
            force=False,
            selected_index=restoration_index,
        )
        watermark_pixels += restored
        low_confidence += uncertain
    for box in _scaled_boxes(fixed_boxes, width, height):
        if fixed_fill == "temporal":
            restored, uncertain = _temporal_restore_region(
                final,
                restoration_frames,
                box,
                variation_threshold=14,
                force=True,
                fallback_fill="background",
                fallback_colour=fixed_fill_colour,
            )
            watermark_pixels += restored
            low_confidence += uncertain
        else:
            left, top, right, bottom = box
            if fixed_fill == "color":
                _fill_region_with_colour(final, box, fixed_fill_colour)
            else:
                _low_frequency_fill_region(final, box)
            affected = max(0, right - left) * max(0, bottom - top)
            watermark_pixels += affected
            low_confidence += affected
    for box in _scaled_boxes(presenter_boxes, width, height):
        restored, uncertain = _temporal_restore_region(
            final,
            restoration_frames,
            box,
            variation_threshold=18,
            force=True,
            force_all=True,
        )
        watermark_pixels += restored
        low_confidence += uncertain
    annotation_pixels += _restore_transient_colored_marks(
        final,
        restoration_frames,
        selected_index=restoration_index,
        excluded_boxes=_scaled_boxes([*fixed_boxes, *presenter_boxes], width, height),
    )
    cleanup_protected = np.max(
        np.abs(final.astype(np.int16) - authentic_selected.astype(np.int16)), axis=2
    ) > 5
    requested_enhancement = str(enhancement_mode or "auto").strip().lower()
    normalized_enhancement = {
        "multiframe": "compatible",
        "high_fidelity": "gpu_ai",
    }.get(requested_enhancement, requested_enhancement)
    if normalized_enhancement not in {"off", "compatible", "gpu_ai", "auto"}:
        raise ValidationError("视频幻灯片清晰增强模式无效")
    if normalized_enhancement != "off":
        fusion = multiframe_fuse(
            authentic_selected,
            restoration_frames,
            repaired_reference=final,
            excluded_boxes=annotation_excluded,
        )
        final = fusion.image
    else:
        fusion = None
    residual_annotation = _colored_handwriting_mask(
        final,
        excluded_boxes=_scaled_boxes([*fixed_boxes, *presenter_boxes], width, height),
        annotation_mode=annotation_mode,
        annotation_colours=annotation_colours,
        annotation_tolerance=annotation_tolerance,
    )
    residual_annotation[persistent_print > 0] = 0
    residual_colour = _annotation_colour_mask(
        final,
        mode=annotation_mode,
        colours=annotation_colours,
        tolerance=annotation_tolerance,
        excluded_boxes=_scaled_boxes([*fixed_boxes, *presenter_boxes], width, height),
    )
    residual_colour[persistent_print > 0] = 0
    if normalized_enhancement in {"auto", "compatible", "gpu_ai"}:
        enhanced = enhance_bgr(
            final,
            mode=(
                "auto"
                if normalized_enhancement == "auto"
                else (
                    "preprocess"
                    if normalized_enhancement == "compatible"
                    else "high_fidelity"
                )
            ),
            content_type="document",
            scale=2,
            max_dimension=4096,
            tile_size=256,
            protected_mask=(
                cleanup_protected
                | (residual_annotation > 0)
                | (residual_colour > 0)
            ),
        )
        final = enhanced.image
    else:
        enhanced = None
    return _CleanResult(
        image=final,
        annotation_pixels=annotation_pixels,
        colour_guided_pixels=colour_guided_pixels,
        watermark_pixels=watermark_pixels,
        low_confidence_pixels=low_confidence,
        residual_annotation_pixels=int(np.count_nonzero(residual_annotation)),
        residual_colour_match_pixels=int(np.count_nonzero(residual_colour)),
        selected_timestamp=float(ordered[selected_position][0].timestamp),
        selected_stability_support=int(augmented_support[selected_position]),
        transition_candidates_rejected=transition_rejected,
        fusion_input_frames=fusion.input_frames if fusion is not None else 0,
        fusion_registered_frames=(
            fusion.registered_frames if fusion is not None else 0
        ),
        fusion_pixels=fusion.fused_pixels if fusion is not None else 0,
        fusion_rejected_frames=(
            fusion.rejected_frames if fusion is not None else 0
        ),
        enhancement_engine=(
            enhanced.engine
            if enhanced is not None
            else ("多帧配准融合" if fusion is not None else "未启用")
        ),
        enhancement_scale=enhanced.scale if enhanced is not None else 1.0,
        ai_attempted=enhanced.ai_attempted if enhanced is not None else False,
        ai_accepted=enhanced.ai_accepted if enhanced is not None else False,
        enhancement_fallback_blocks=(
            enhanced.fallback_blocks if enhanced is not None else 0
        ),
        enhancement_total_blocks=(
            enhanced.total_blocks if enhanced is not None else 0
        ),
        enhancement_reason=enhanced.reason if enhanced is not None else "",
    )


def _image_hash(image: "np.ndarray") -> "np.ndarray":
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    return _perceptual_hash(gray)


def _slide_signature(image: "np.ndarray") -> _SlideSignature:
    sample = cv2.resize(image, (160, 90), interpolation=cv2.INTER_AREA)
    gray = cv2.cvtColor(sample, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 55, 150) > 0
    dilated_edges = cv2.dilate(edges.astype(np.uint8), np.ones((3, 3), np.uint8)) > 0
    hsv = cv2.cvtColor(sample, cv2.COLOR_BGR2HSV)
    values = sample.astype(np.int16)
    cyan_chroma = np.minimum(values[:, :, 0], values[:, :, 1]) - values[:, :, 2]
    colored_ink = (
        (hsv[:, :, 0] >= 70)
        & (hsv[:, :, 0] <= 112)
        & (hsv[:, :, 1] >= 75)
        & (hsv[:, :, 2] >= 55)
        & (cyan_chroma >= 35)
    )
    colored_ink = cv2.dilate(
        colored_ink.astype(np.uint8), np.ones((3, 3), np.uint8)
    ).astype(bool)
    annotation = _colored_handwriting_mask(image, excluded_boxes=[])
    annotation_ink = cv2.resize(annotation, (160, 90), interpolation=cv2.INTER_NEAREST)
    annotation_ink = cv2.dilate(
        (annotation_ink > 0).astype(np.uint8), np.ones((3, 3), np.uint8)
    ).astype(bool)
    return _SlideSignature(
        shape=tuple(int(value) for value in image.shape),
        phash=_image_hash(image),
        gray=gray,
        edges=edges,
        dilated_edges=dilated_edges,
        colored_ink=np.packbits(colored_ink.reshape(-1)),
        annotation_ink=np.packbits(annotation_ink.reshape(-1)),
    )


def _signatures_near_duplicate(first: _SlideSignature, second: _SlideSignature) -> bool:
    if first.shape != second.shape:
        return False
    hash_distance = int(np.count_nonzero(first.phash != second.phash))
    if hash_distance > 12:
        return False
    correlation = cv2.matchTemplate(first.gray, second.gray, cv2.TM_CCOEFF_NORMED)[0, 0]
    mean_difference = float(
        np.mean(np.abs(first.gray.astype(np.int16) - second.gray.astype(np.int16)))
        / 255.0
    )
    first_containment = float(
        np.count_nonzero(first.edges & second.dilated_edges)
    ) / max(1, int(first.edges.sum()))
    second_containment = float(
        np.count_nonzero(second.edges & first.dilated_edges)
    ) / max(1, int(second.edges.sum()))
    edge_iou = float(np.count_nonzero(first.edges & second.edges)) / max(
        1, int(np.count_nonzero(first.edges | second.edges))
    )
    strict_match = bool(
        hash_distance <= 10
        and correlation >= 0.938
        and mean_difference <= 0.060
        and min(first_containment, second_containment) >= 0.78
        and edge_iou >= 0.55
    )
    if strict_match:
        return True
    # A returned-to slide can contain substantial cyan handwriting even after
    # conservative cleanup.  Permit a one-sided "clean base + coloured ink"
    # match only when the unmatched structure is predominantly coloured and a
    # high-confidence handwriting mask exists on exactly one side.  Legitimate
    # blue printed text on both pages therefore remains a hard non-duplicate.
    first_colored_ink = (
        np.unpackbits(first.colored_ink)[: first.edges.size]
        .reshape(first.edges.shape)
        .astype(bool)
    )
    second_colored_ink = (
        np.unpackbits(second.colored_ink)[: second.edges.size]
        .reshape(second.edges.shape)
        .astype(bool)
    )
    first_annotation_ink = (
        np.unpackbits(first.annotation_ink)[: first.edges.size]
        .reshape(first.edges.shape)
        .astype(bool)
    )
    second_annotation_ink = (
        np.unpackbits(second.annotation_ink)[: second.edges.size]
        .reshape(second.edges.shape)
        .astype(bool)
    )
    first_extra = first.edges & (~second.dilated_edges)
    second_extra = second.edges & (~first.dilated_edges)
    unmatched = first_extra | second_extra
    unmatched_pixels = int(np.count_nonzero(unmatched))
    colored_unmatched = int(
        np.count_nonzero(
            (first_extra & first_colored_ink) | (second_extra & second_colored_ink)
        )
    )
    colored_coverage = colored_unmatched / max(1, unmatched_pixels)
    annotation_ratios = (
        float(np.mean(first_annotation_ink)),
        float(np.mean(second_annotation_ink)),
    )
    asymmetric_annotation = (
        min(annotation_ratios) <= 0.001 and max(annotation_ratios) >= 0.004
    )
    return bool(
        correlation >= 0.82
        and mean_difference <= 0.015
        and max(first_containment, second_containment) >= 0.98
        and edge_iou >= 0.50
        and colored_coverage >= 0.45
        and asymmetric_annotation
    )


def _near_duplicate(first: "np.ndarray", second: "np.ndarray") -> bool:
    return _signatures_near_duplicate(_slide_signature(first), _slide_signature(second))


def _content_quality(
    image: "np.ndarray",
    *,
    residual_annotation_pixels: int = 0,
    restored_annotation_pixels: int = 0,
    restored_watermark_pixels: int = 0,
    low_confidence_pixels: int = 0,
) -> float:
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    area = max(1, gray.size)
    edges = float(np.count_nonzero(cv2.Canny(gray, 55, 150))) / area
    sharpness = float(cv2.Laplacian(gray, cv2.CV_32F).var())
    residual_ratio = residual_annotation_pixels / area
    restored_annotation_ratio = restored_annotation_pixels / area
    restored_watermark_ratio = restored_watermark_pixels / area
    low_confidence_ratio = low_confidence_pixels / area
    return (
        edges
        + 0.00002 * sharpness
        - 0.90 * residual_ratio
        - 0.45 * restored_annotation_ratio
        - 0.36 * restored_watermark_ratio
        - 0.08 * low_confidence_ratio
    )


def _enhancement_record(clean: _CleanResult) -> dict[str, object]:
    """Return serializable per-page fusion and enhancement diagnostics."""

    return {
        "fusion_input_frames": clean.fusion_input_frames,
        "fusion_registered_frames": clean.fusion_registered_frames,
        "fusion_pixels": clean.fusion_pixels,
        "fusion_rejected_frames": clean.fusion_rejected_frames,
        "enhancement_engine": clean.enhancement_engine,
        "enhancement_scale": round(clean.enhancement_scale, 3),
        "ai_attempted": clean.ai_attempted,
        "ai_accepted": clean.ai_accepted,
        "enhancement_fallback_blocks": clean.enhancement_fallback_blocks,
        "enhancement_total_blocks": clean.enhancement_total_blocks,
        "enhancement_reason": clean.enhancement_reason,
    }


def _read_cv_image(path: Path) -> "np.ndarray | None":
    try:
        payload = np.fromfile(str(path), dtype=np.uint8)
    except OSError:
        return None
    if payload.size == 0:
        return None
    return cv2.imdecode(payload, cv2.IMREAD_COLOR)


def _write_png(path: Path, image: "np.ndarray") -> None:
    encoded, payload = cv2.imencode(".png", image, [cv2.IMWRITE_PNG_COMPRESSION, 2])
    if not encoded:
        raise ValidationError(f"无法编码幻灯片图片：{path.name}")
    with atomic_output(path) as temporary:
        temporary.write_bytes(payload.tobytes())


def _write_pptx(images: Sequence[Path], target: Path) -> None:
    if not images:
        raise ValidationError("没有可写入 PPT 的幻灯片画面")
    first = _read_cv_image(images[0])
    if first is None:
        raise ValidationError("首张幻灯片图片无法读取")
    height, width = first.shape[:2]
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = int(
        round(presentation.slide_width * height / max(1, width))
    )
    blank = presentation.slide_layouts[6]
    for path in images:
        check_cancelled("已取消 PPT 生成")
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
    if len(presentation.slides) > len(images):
        slide_id = presentation.slides._sldIdLst[0]
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[0]
    with atomic_output(target) as temporary:
        presentation.save(temporary)


def _format_timestamp(seconds: object) -> str:
    total_milliseconds = max(0, int(round(float(seconds or 0.0) * 1000)))
    total_seconds, milliseconds = divmod(total_milliseconds, 1000)
    hours, remainder = divmod(total_seconds, 3600)
    minutes, whole_seconds = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{whole_seconds:02d}.{milliseconds:03d}"


def _format_time_range(occurrence: object) -> str:
    if not isinstance(occurrence, dict):
        return "-"
    return (
        f"{_format_timestamp(occurrence.get('start_seconds', 0.0))}–"
        f"{_format_timestamp(occurrence.get('end_seconds', 0.0))}"
    )


def _ordered_occurrences(record: dict[str, object]) -> list[dict[str, object]]:
    occurrences = record.get("occurrences")
    if not isinstance(occurrences, list):
        return []
    typed = [item for item in occurrences if isinstance(item, dict)]
    return sorted(
        typed,
        key=lambda item: (
            float(item.get("start_seconds", 0.0)),
            float(item.get("end_seconds", 0.0)),
        ),
    )


def _effective_ordering_occurrence(
    record: dict[str, object],
) -> tuple[dict[str, object] | None, str]:
    """Choose the first *formal* occurrence without deleting short pages.

    A lecturer may accidentally reveal a later page for one or two seconds and
    then return to the real sequence.  Keeping that interval in the audit trail
    is important, but using it as the PPT ordering anchor swaps the two pages.
    Only a duplicated page whose later occurrence is both clearly longer and
    at least four seconds is deferred.  A unique short page remains untouched.
    """

    occurrences = _ordered_occurrences(record)
    if not occurrences:
        return None, "missing"
    first = occurrences[0]
    first_duration = max(
        0.0,
        float(first.get("end_seconds", 0.0))
        - float(first.get("start_seconds", 0.0)),
    )
    if first_duration <= _SHORT_PREVIEW_MAX_SECONDS and len(occurrences) > 1:
        formal_minimum = max(
            _FORMAL_OCCURRENCE_MIN_SECONDS,
            first_duration * _FORMAL_OCCURRENCE_DURATION_RATIO,
        )
        for occurrence in occurrences[1:]:
            duration = max(
                0.0,
                float(occurrence.get("end_seconds", 0.0))
                - float(occurrence.get("start_seconds", 0.0)),
            )
            if duration >= formal_minimum:
                return occurrence, "short_preview_deferred"
    return first, "first_occurrence"


def _reorder_slide_outputs_by_effective_timeline(
    image_paths: Sequence[Path],
    slide_records: Sequence[dict[str, object]],
    signatures: Sequence[_SlideSignature],
    scores: Sequence[float],
) -> tuple[
    list[Path],
    list[dict[str, object]],
    list[_SlideSignature],
    list[float],
]:
    """Synchronously reorder PNGs, records and validation signatures."""

    count = len(image_paths)
    if not (
        count == len(slide_records) == len(signatures) == len(scores)
    ):
        raise ValidationError("幻灯片顺序二重检查失败：页面关联清单长度不一致")
    bundles: list[
        tuple[
            int,
            float,
            float,
            Path,
            dict[str, object],
            _SlideSignature,
            float,
        ]
    ] = []
    for original_index, (path, record, signature, score) in enumerate(
        zip(image_paths, slide_records, signatures, scores)
    ):
        occurrences = _ordered_occurrences(record)
        if not occurrences:
            raise ValidationError("幻灯片顺序二重检查失败：缺少出现时间")
        record["occurrences"] = occurrences
        first = occurrences[0]
        record["start_seconds"] = float(first.get("start_seconds", 0.0))
        record["end_seconds"] = float(first.get("end_seconds", 0.0))
        anchor, reason = _effective_ordering_occurrence(record)
        if anchor is None:
            raise ValidationError("幻灯片顺序二重检查失败：排序锚点无效")
        ordering_start = float(anchor.get("start_seconds", 0.0))
        ordering_end = float(anchor.get("end_seconds", ordering_start))
        record["ordering_start_seconds"] = ordering_start
        record["ordering_end_seconds"] = ordering_end
        record["ordering_reason"] = reason
        bundles.append(
            (
                original_index,
                ordering_start,
                float(first.get("start_seconds", 0.0)),
                path,
                record,
                signature,
                float(score),
            )
        )
    bundles.sort(key=lambda item: (item[1], item[2], item[0]))
    reordered_paths = [item[3] for item in bundles]
    reordered_records = [item[4] for item in bundles]
    reordered_signatures = [item[5] for item in bundles]
    reordered_scores = [item[6] for item in bundles]

    if [item[0] for item in bundles] != list(range(count)):
        parents = {path.parent for path in reordered_paths}
        if len(parents) != 1:
            raise ValidationError("幻灯片顺序二重检查失败：PNG 不在同一暂存目录")
        temporary_paths: list[Path] = []
        for index, path in enumerate(reordered_paths):
            temporary = path.with_name(
                f".df-order-{os.getpid()}-{index:04d}{path.suffix}"
            )
            os.replace(path, temporary)
            temporary_paths.append(temporary)
        final_paths: list[Path] = []
        parent = reordered_paths[0].parent
        for slide_number, temporary in enumerate(temporary_paths, start=1):
            final_path = parent / f"幻灯片_{slide_number:04d}.png"
            os.replace(temporary, final_path)
            final_paths.append(final_path)
        reordered_paths = final_paths

    for slide_number, record in enumerate(reordered_records, start=1):
        record["slide"] = slide_number
    return (
        reordered_paths,
        reordered_records,
        reordered_signatures,
        reordered_scores,
    )


def _validate_slide_timeline(
    slide_records: Sequence[dict[str, object]], video_duration: float
) -> None:
    """Verify the unique-page timeline before it is exposed to users."""

    previous_ordering_start = -1.0
    duration = max(0.0, float(video_duration))
    for expected_slide, record in enumerate(slide_records, start=1):
        if int(record.get("slide", 0)) != expected_slide:
            raise ValidationError("幻灯片顺序二重检查失败：页码不连续")
        occurrences = record.get("occurrences")
        if not isinstance(occurrences, list) or not occurrences:
            raise ValidationError("幻灯片顺序二重检查失败：缺少出现时间")
        normalized: list[tuple[float, float]] = []
        for occurrence in occurrences:
            if not isinstance(occurrence, dict):
                raise ValidationError("幻灯片顺序二重检查失败：时间记录无效")
            start = float(occurrence.get("start_seconds", -1.0))
            end = float(occurrence.get("end_seconds", -1.0))
            if start < 0 or end < start or (duration and end > duration + 0.05):
                raise ValidationError("幻灯片顺序二重检查失败：时间范围无效")
            normalized.append((start, end))
        if normalized != sorted(normalized):
            raise ValidationError("幻灯片顺序二重检查失败：重复页时间未排序")
        first_start, first_end = normalized[0]
        ordering_occurrence, ordering_reason = _effective_ordering_occurrence(record)
        if ordering_occurrence is None:
            raise ValidationError("幻灯片顺序二重检查失败：排序锚点无效")
        ordering_start = float(ordering_occurrence.get("start_seconds", -1.0))
        ordering_end = float(ordering_occurrence.get("end_seconds", -1.0))
        if ordering_start + 1e-6 < previous_ordering_start:
            raise ValidationError("幻灯片顺序二重检查失败：PPT 未按有效出现时间排序")
        previous_ordering_start = ordering_start
        if "ordering_start_seconds" in record and abs(
            float(record.get("ordering_start_seconds", -1.0)) - ordering_start
        ) > 0.001:
            raise ValidationError("幻灯片顺序二重检查失败：有效开始时间不一致")
        if "ordering_end_seconds" in record and abs(
            float(record.get("ordering_end_seconds", -1.0)) - ordering_end
        ) > 0.001:
            raise ValidationError("幻灯片顺序二重检查失败：有效结束时间不一致")
        if "ordering_reason" in record and record.get("ordering_reason") != ordering_reason:
            raise ValidationError("幻灯片顺序二重检查失败：有效排序原因不一致")
        if abs(float(record.get("start_seconds", -1.0)) - first_start) > 0.001:
            raise ValidationError("幻灯片顺序二重检查失败：首次开始时间不一致")
        if abs(float(record.get("end_seconds", -1.0)) - first_end) > 0.001:
            raise ValidationError("幻灯片顺序二重检查失败：首次结束时间不一致")
        if int(record.get("merged_duplicate_segments", -1)) != len(normalized) - 1:
            raise ValidationError("幻灯片顺序二重检查失败：重复页计数不一致")
        selected = record.get("selected_occurrence")
        if not isinstance(selected, dict) or selected not in occurrences:
            raise ValidationError("幻灯片顺序二重检查失败：实际采用区间无效")
        selected_time = float(record.get("selected_timestamp_seconds", -1.0))
        selected_start = float(selected.get("start_seconds", -1.0))
        selected_end = float(selected.get("end_seconds", -1.0))
        if not selected_start - 0.05 <= selected_time <= selected_end + 0.05:
            raise ValidationError("幻灯片顺序二重检查失败：实际采用帧时间无效")
        if "selected_stability_support" in record:
            support = int(record.get("selected_stability_support", 0))
            rejected = int(record.get("transition_candidates_rejected", 0))
            candidates = int(record.get("candidate_frames", 0))
            if candidates >= _STABLE_WINDOW_CANDIDATES and support <= 0:
                raise ValidationError(
                    "幻灯片稳定性二重检查失败：采用帧没有连续稳定支持"
                )
            if rejected < 0 or rejected > candidates:
                raise ValidationError("幻灯片稳定性二重检查失败：转场候选计数无效")


def _configure_docx_style(document: object, style_name: str, size: float) -> None:
    from docx.oxml.ns import qn
    from docx.shared import Pt

    style = document.styles[style_name]
    style.font.name = "Microsoft YaHei"
    style.font.size = Pt(size)
    rpr = style.element.get_or_add_rPr()
    fonts = rpr.get_or_add_rFonts()
    for attribute, font_name in (
        ("w:ascii", "Microsoft YaHei"),
        ("w:hAnsi", "Microsoft YaHei"),
        ("w:eastAsia", "微软雅黑"),
    ):
        fonts.set(qn(attribute), font_name)


def _set_docx_cell(
    cell: object,
    value: object,
    *,
    bold: bool = False,
    size: float = 8.0,
    center: bool = False,
) -> None:
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt

    cell.text = ""
    paragraph = cell.paragraphs[0]
    paragraph.paragraph_format.space_after = Pt(0)
    if center:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(str(value))
    run.bold = bold
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    fonts = run._element.get_or_add_rPr().get_or_add_rFonts()
    fonts.set(qn("w:ascii"), "Microsoft YaHei")
    fonts.set(qn("w:hAnsi"), "Microsoft YaHei")
    fonts.set(qn("w:eastAsia"), "微软雅黑")
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    if bold:
        properties = cell._tc.get_or_add_tcPr()
        shading = OxmlElement("w:shd")
        shading.set(qn("w:fill"), "D9EAF7")
        properties.append(shading)


def _write_report_docx(target: Path, payload: dict[str, object]) -> None:
    try:
        from docx import Document
        from docx.enum.section import WD_ORIENT
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Mm, Pt
    except ImportError as exc:  # pragma: no cover - dependency probe covers packaging
        raise MissingEngineError("视频提取 Word 报告需要 python-docx") from exc

    document = Document()
    section = document.sections[0]
    section.orientation = WD_ORIENT.LANDSCAPE
    section.page_width = Mm(297)
    section.page_height = Mm(210)
    section.top_margin = Mm(13)
    section.bottom_margin = Mm(13)
    section.left_margin = Mm(13)
    section.right_margin = Mm(13)
    for style_name, size in (
        ("Normal", 9.0),
        ("Title", 18.0),
        ("Heading 1", 13.0),
        ("Heading 2", 11.0),
    ):
        _configure_docx_style(document, style_name, size)
    properties = document.core_properties
    properties.title = "讲解视频幻灯片提取报告"
    properties.subject = "页织工坊视频生成完整 PPT 二重检查报告"
    properties.keywords = _REPORT_SCHEMA

    title = document.add_heading("讲解视频幻灯片提取报告", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = document.add_paragraph(Path(str(payload.get("source", ""))).name)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.paragraph_format.space_after = Pt(8)

    metadata = payload.get("metadata")
    metadata = metadata if isinstance(metadata, dict) else {}
    slides = payload.get("slides")
    slides = slides if isinstance(slides, list) else []
    total_duplicates = sum(
        int(item.get("merged_duplicate_segments", 0))
        for item in slides
        if isinstance(item, dict)
    )
    document.add_heading("任务概况", level=1)
    summary_rows = [
        ("结构版本", _REPORT_SCHEMA),
        ("源视频", str(payload.get("source", ""))),
        ("视频时长", _format_timestamp(metadata.get("duration", 0.0))),
        (
            "原始分辨率",
            f"{int(metadata.get('width', 0))} × {int(metadata.get('height', 0))}",
        ),
        ("帧率", f"{float(metadata.get('fps', 0.0)):.3f} FPS"),
        ("扫描间隔", f"{float(metadata.get('sample_interval', 0.0)):.3f} 秒"),
        ("扫描帧数", str(int(payload.get("scan_frames", 0)))),
        ("识别片段", str(int(payload.get("detected_segments", 0)))),
        (
            "跳过未稳定转场片段",
            str(int(payload.get("skipped_unstable_segments", 0))),
        ),
        ("最终 PPT 页数", str(int(payload.get("output_slides", 0)))),
        ("合并重复片段", str(total_duplicates)),
        (
            "页面顺序",
            "按有效首次出现时间排列；正式讲解前的极短误翻不抢占页序",
        ),
    ]
    summary = document.add_table(rows=1, cols=2)
    summary.style = "Table Grid"
    summary.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_docx_cell(summary.rows[0].cells[0], "项目", bold=True, center=True)
    _set_docx_cell(summary.rows[0].cells[1], "结果", bold=True, center=True)
    for label, value in summary_rows:
        cells = summary.add_row().cells
        _set_docx_cell(cells[0], label, bold=True)
        _set_docx_cell(cells[1], value)

    document.add_heading("处理设置与二检", level=1)
    crop = payload.get("presentation_crop_pixels")
    crop_text = (
        "，".join(str(int(value)) for value in crop) if isinstance(crop, list) else "-"
    )
    settings_rows = [
        ("课件裁剪区域（x,y,宽,高）", crop_text),
        ("移动水印搜索", str(payload.get("watermark_search", "-"))),
        (
            "水印文字提示",
            str(payload.get("watermark_text_hint", "")).strip() or "未提供",
        ),
        ("讲师画面策略", str(payload.get("presenter_policy", "-"))),
        ("颜色辅助模式", str(payload.get("annotation_color_mode", "auto"))),
        (
            "颜色辅助基准色",
            str(payload.get("annotation_colors", "")).strip() or "未指定",
        ),
        (
            "颜色辅助容差",
            str(int(payload.get("annotation_color_tolerance", 24))),
        ),
        ("固定水印填充方式", str(payload.get("fixed_watermark_fill", "temporal"))),
        (
            "固定水印填充色",
            str(payload.get("fixed_watermark_fill_color", "#FFFFFF")),
        ),
        (
            "清晰增强链路",
            {
                "auto": "自动检测：有 Vulkan GPU 时使用 AI，无显卡时使用兼容增强",
                "compatible": "同页多帧配准融合 → CPU/OpenCV 高保真预处理（无独显可用）",
                "gpu_ai": "同页多帧配准融合 → 高保真预处理 → Real-ESRGAN 2× → 结构二检与回退",
                "high_fidelity": "同页多帧配准融合 → 高保真预处理 → Real-ESRGAN 2× → 结构二检与回退（旧任务兼容）",
                "multiframe": "同页多帧配准融合 → CPU/OpenCV 高保真预处理（旧任务兼容）",
                "off": "关闭，保持旧版单帧清理结果",
            }.get(str(payload.get("enhancement_mode", "auto")), "-"),
        ),
        (
            "自动讲师区",
            (
                "已检测并裁除"
                if payload.get("auto_presenter_cropped")
                else (
                    "已检测" if payload.get("auto_presenter_detected") else "未检测到"
                )
            ),
        ),
        ("PNG 二检", "通过：文件可解码，尺寸与页面签名一致"),
        ("PPT 二检", "通过：页数、单图结构及页面图片顺序一致"),
        ("报告二检", "通过后与 PPT、可选 PNG 一次性整体提交"),
    ]
    settings = document.add_table(rows=1, cols=2)
    settings.style = "Table Grid"
    settings.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_docx_cell(settings.rows[0].cells[0], "检查项", bold=True, center=True)
    _set_docx_cell(settings.rows[0].cells[1], "说明", bold=True, center=True)
    for label, value in settings_rows:
        cells = settings.add_row().cells
        _set_docx_cell(cells[0], label, bold=True)
        _set_docx_cell(cells[1], value)

    document.add_heading("页面明细", level=1)
    headers = (
        "PPT 页码",
        "首次出现区间",
        "实际采用帧",
        "出现次数",
        "全部出现区间",
        "候选帧",
        "清理像素\n批注 / 颜色引导 / 水印",
        "风险像素\n疑似笔迹 / 颜色匹配 / 低可信",
        "稳定支持\n/ 剔除转场",
        "多帧融合 / 清晰增强",
        "状态",
    )
    details = document.add_table(rows=1, cols=len(headers))
    details.style = "Table Grid"
    details.alignment = WD_TABLE_ALIGNMENT.CENTER
    details.autofit = True
    for cell, header in zip(details.rows[0].cells, headers):
        _set_docx_cell(cell, header, bold=True, size=7.5, center=True)
    for item in slides:
        if not isinstance(item, dict):
            continue
        occurrences = item.get("occurrences")
        occurrences = occurrences if isinstance(occurrences, list) else []
        all_ranges = "\n".join(_format_time_range(value) for value in occurrences)
        residual = int(item.get("residual_annotation_pixels", 0))
        residual_colour = int(item.get("residual_colour_match_pixels", 0))
        uncertain = int(item.get("low_confidence_pixels", 0))
        status = (
            "需抽查：存在低可信区域"
            if uncertain
            else (
                "建议抽查：仍有颜色匹配或疑似笔迹"
                if residual or residual_colour
                else "二检通过"
            )
        )
        if item.get("ordering_reason") == "short_preview_deferred":
            status = f"{status}；短暂预览已延后排序"
        enhancement_status = (
            "AI通过"
            if item.get("ai_accepted")
            else (
                "AI安全回退"
                if item.get("ai_attempted")
                else "传统增强 / 未调用AI"
            )
        )
        values = (
            int(item.get("slide", 0)),
            _format_time_range(occurrences[0] if occurrences else None),
            _format_timestamp(item.get("selected_timestamp_seconds", 0.0)),
            len(occurrences),
            all_ranges or "-",
            int(item.get("candidate_frames", 0)),
            (
                f"{int(item.get('annotation_pixels_restored', 0))} / "
                f"{int(item.get('colour_guided_pixels_restored', 0))} / "
                f"{int(item.get('dynamic_watermark_pixels_restored', 0))}"
            ),
            f"{residual} / {residual_colour} / {uncertain}",
            (
                f"{int(item.get('selected_stability_support', 0))} / "
                f"{int(item.get('transition_candidates_rejected', 0))}"
            ),
            (
                f"{int(item.get('fusion_registered_frames', 0))} / "
                f"{int(item.get('fusion_input_frames', 0))} 帧；"
                f"{int(item.get('fusion_pixels', 0))} 像素\n"
                f"{float(item.get('enhancement_scale', 1.0)):.2f}×；"
                f"{enhancement_status}；"
                f"回退块 {int(item.get('enhancement_fallback_blocks', 0))}"
            ),
            status,
        )
        row = details.add_row()
        for column, (cell, value) in enumerate(zip(row.cells, values)):
            _set_docx_cell(
                cell,
                value,
                size=7.2,
                center=column in {0, 1, 2, 3, 5, 6, 7, 8, 9, 10},
            )

    document.add_heading("风险提示与使用说明", level=1)
    warnings = payload.get("warnings")
    warnings = warnings if isinstance(warnings, list) else []
    if not warnings:
        warnings = ["本次未发现需要单独提示的低可信区域，仍建议抽查关键公式和表格。"]
    for warning in warnings:
        document.add_paragraph(str(warning), style="List Bullet")
    document.add_paragraph(
        "页面顺序说明：同页逐渐出现的内容只保留印刷内容最完整的最终状态；"
        "全局相同页面会合并。若某页第一次只短暂闪现、随后存在明显更长的正式讲解区间，"
        "程序保留短暂闪现记录，但使用正式区间参与页序；只出现一次的短页不会被此规则删除。"
        "后续更清晰画面可以替换该页图像，PPT、PNG 文件名和报告页码会同步重排。"
    )
    document.add_paragraph(
        "本报告不嵌入高清页面预览，以控制体积并提高 WPS 与手机端打开速度；"
        "最终视觉结果请直接查看同目录中的完整 PPTX。"
    )
    footer = section.footer.paragraphs[0]
    footer.text = "页织工坊 · 讲解视频幻灯片提取报告"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    with atomic_output(target) as temporary:
        document.save(temporary)


def _validate_staged_images(
    paths: Sequence[Path], signatures: Sequence[_SlideSignature]
) -> None:
    if not paths or len(paths) != len(signatures):
        raise ValidationError("幻灯片图片二重检查失败：页面清单不完整")
    for path, signature in zip(paths, signatures):
        if not path.is_file() or path.stat().st_size < 256:
            raise ValidationError(f"幻灯片图片二重检查失败：{path.name}")
        image = _read_cv_image(path)
        if image is None or tuple(image.shape) != signature.shape:
            raise ValidationError(f"幻灯片图片二重检查失败：{path.name}")


def _validate_pptx(
    path: Path,
    expected_slides: int,
    expected_images: Sequence[Path] | None = None,
) -> None:
    if not path.is_file() or path.stat().st_size < 1000:
        raise ValidationError("生成的 PPT 文件为空或不完整")
    presentation = Presentation(path)
    if presentation.slide_width <= 0 or presentation.slide_height <= 0:
        raise ValidationError("PPT 二重检查失败：页面尺寸无效")
    if len(presentation.slides) != expected_slides:
        raise ValidationError(
            f"PPT 二重检查失败：预期 {expected_slides} 页，实际 {len(presentation.slides)} 页"
        )
    if expected_images is not None and len(expected_images) != expected_slides:
        raise ValidationError("PPT 二重检查失败：有序图片清单不完整")
    for index, slide in enumerate(presentation.slides):
        pictures = [
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        if len(pictures) != 1:
            raise ValidationError("PPT 二重检查发现空白页或图片结构异常")
        picture = pictures[0]
        if picture.width <= 0 or picture.height <= 0:
            raise ValidationError("PPT 二重检查发现无效图片")
        if (
            picture.left != 0
            or picture.top != 0
            or picture.width != presentation.slide_width
            or picture.height != presentation.slide_height
        ):
            raise ValidationError("PPT 二重检查发现页面图片未完整铺满画布")
        if expected_images is not None:
            expected_digest = hashlib.sha256(
                expected_images[index].read_bytes()
            ).digest()
            embedded_digest = hashlib.sha256(picture.image.blob).digest()
            if expected_digest != embedded_digest:
                raise ValidationError(
                    f"PPT 二重检查失败：第 {index + 1} 页图片与时间排序清单不一致"
                )


def _validate_report_docx(
    path: Path,
    expected_slides: int,
    payload: dict[str, object] | None = None,
) -> None:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - dependency probe covers packaging
        raise MissingEngineError("视频提取 Word 报告需要 python-docx") from exc
    if not path.is_file() or path.stat().st_size < 1500:
        raise ValidationError("提取报告二重检查失败：Word 文件为空或不完整")
    try:
        with zipfile.ZipFile(path) as archive:
            if archive.testzip() is not None:
                raise ValidationError("提取报告二重检查失败：Word 压缩包校验失败")
            required = {"[Content_Types].xml", "word/document.xml"}
            if not required.issubset(archive.namelist()):
                raise ValidationError("提取报告二重检查失败：Word 核心结构缺失")
        document = Document(path)
    except ValidationError:
        raise
    except (OSError, zipfile.BadZipFile, ValueError) as exc:
        raise ValidationError("提取报告二重检查失败：Word 无法重新打开") from exc
    if document.core_properties.keywords != _REPORT_SCHEMA:
        raise ValidationError("提取报告二重检查失败：结构版本无效")
    if not any(
        "讲解视频幻灯片提取报告" in paragraph.text for paragraph in document.paragraphs
    ):
        raise ValidationError("提取报告二重检查失败：报告标题缺失")
    detail_tables = [
        table
        for table in document.tables
        if table.rows and table.rows[0].cells[0].text.strip() == "PPT 页码"
    ]
    if len(detail_tables) != 1:
        raise ValidationError("提取报告二重检查失败：页面明细表缺失")
    rows = detail_tables[0].rows[1:]
    if len(rows) != expected_slides:
        raise ValidationError("提取报告二重检查失败：页面记录数与 PPT 不一致")
    try:
        page_numbers = [int(row.cells[0].text.strip()) for row in rows]
    except (ValueError, IndexError) as exc:
        raise ValidationError("提取报告二重检查失败：页码无法读取") from exc
    if page_numbers != list(range(1, expected_slides + 1)):
        raise ValidationError("提取报告二重检查失败：页码不连续")
    if payload is None:
        return
    if payload.get("schema") != _REPORT_SCHEMA:
        raise ValidationError("提取报告二重检查失败：内存结构版本无效")
    if int(payload.get("output_slides", -1)) != expected_slides:
        raise ValidationError("提取报告二重检查失败：页数与 PPT 不一致")
    slides = payload.get("slides")
    if not isinstance(slides, list) or len(slides) != expected_slides:
        raise ValidationError("提取报告二重检查失败：页面记录不完整")
    metadata = payload.get("metadata")
    metadata = metadata if isinstance(metadata, dict) else {}
    records = [item for item in slides if isinstance(item, dict)]
    if len(records) != expected_slides:
        raise ValidationError("提取报告二重检查失败：页面记录类型无效")
    _validate_slide_timeline(records, float(metadata.get("duration", 0.0)))
    for row, record in zip(rows, records):
        occurrences = record.get("occurrences")
        occurrences = occurrences if isinstance(occurrences, list) else []
        expected_values = (
            _format_time_range(occurrences[0] if occurrences else None),
            _format_timestamp(record.get("selected_timestamp_seconds", 0.0)),
            str(len(occurrences)),
            "\n".join(_format_time_range(value) for value in occurrences),
        )
        actual_values = tuple(row.cells[index].text.strip() for index in (1, 2, 3, 4))
        if actual_values != expected_values:
            raise ValidationError("提取报告二重检查失败：时间线内容与内存记录不一致")


def _publish_staged_directory(staging_dir: Path, target_dir: Path) -> None:
    """Publish a fully validated task payload without exposing half-results."""

    target_dir.parent.mkdir(parents=True, exist_ok=True)
    if target_dir.exists():
        raise ValidationError(f"输出文件夹已存在，请改用新的输出位置：{target_dir}")
    retry_delays = (0.0, 0.05, 0.10, 0.20, 0.40, 0.80, 1.60)
    last_error: OSError | None = None
    for delay in retry_delays:
        if delay:
            time.sleep(delay)
        try:
            os.replace(staging_dir, target_dir)
            return
        except OSError as exc:
            last_error = exc
            # Windows search indexing, antivirus scanning, thumbnail previews,
            # or a just-closed image viewer can briefly retain a handle inside
            # the fully validated staging directory. Retrying the atomic
            # directory rename preserves the all-or-nothing publish guarantee.
            if os.name != "nt" or getattr(exc, "winerror", None) not in {
                5,   # ERROR_ACCESS_DENIED
                32,  # ERROR_SHARING_VIOLATION
                33,  # ERROR_LOCK_VIOLATION
            }:
                break
            if target_dir.exists():
                break
    raise ValidationError(f"无法提交已完成的幻灯片结果：{target_dir}") from last_error


def _crop_normalized_boxes(
    boxes: Sequence[tuple[float, float, float, float]],
    crop: tuple[int, int, int, int],
    full_width: int,
    full_height: int,
) -> list[tuple[float, float, float, float]]:
    crop_left, crop_top, crop_right, crop_bottom = crop
    crop_width = crop_right - crop_left
    crop_height = crop_bottom - crop_top
    result: list[tuple[float, float, float, float]] = []
    for box in boxes:
        left, top, right, bottom = _pixel_box(box, full_width, full_height)
        left = max(left, crop_left)
        top = max(top, crop_top)
        right = min(right, crop_right)
        bottom = min(bottom, crop_bottom)
        if right <= left or bottom <= top:
            continue
        result.append(
            (
                (left - crop_left) / crop_width,
                (top - crop_top) / crop_height,
                (right - left) / crop_width,
                (bottom - top) / crop_height,
            )
        )
    return result


def extract_slides_to_pptx(
    input_video: PathLike,
    output_dir: PathLike,
    *,
    scan_mode: str = "accurate",
    change_sensitivity: str = "balanced",
    crop_mode: str = "auto",
    crop_rect: object = "",
    watermark_search: str = "auto",
    watermark_rect: object = "",
    watermark_text_hint: str = "",
    annotation_color_mode: str = "auto",
    annotation_colors: object = "#00AEEF",
    annotation_color_tolerance: object = 24,
    fixed_watermark_regions: object = "",
    fixed_watermark_fill: str = "temporal",
    fixed_watermark_fill_color: object = "#FFFFFF",
    presenter_policy: str = "auto_crop",
    presenter_rect: object = "",
    enhancement_mode: str = "auto",
    keep_images: bool = True,
    keep_report: bool = True,
) -> list[Path]:
    """Extract final, printed slide states from a lecture video.

    The output is intentionally visual (one PNG per slide and one image-only
    PPTX).  It does not OCR or invent text hidden for the full lifetime of a
    slide.  User-declared fixed regions are inpainted only as a low-confidence
    fallback and are recorded in the user-readable Word report.
    """

    _require_dependencies()
    source = _input_video(input_video)
    target_dir = Path(output_dir)
    if target_dir.exists():
        if not target_dir.is_dir():
            raise ValidationError(f"输出位置不是文件夹：{target_dir}")
        raise ValidationError(f"输出文件夹已存在，请改用新的输出位置：{target_dir}")
    target_dir.parent.mkdir(parents=True, exist_ok=True)
    annotation_mode, annotation_bgr_colours, annotation_tolerance = (
        _validate_annotation_colour_options(
            annotation_color_mode,
            annotation_colors,
            annotation_color_tolerance,
        )
    )
    fixed_fill_mode, fixed_fill_bgr = _validate_fixed_fill_options(
        fixed_watermark_fill,
        fixed_watermark_fill_color,
    )
    watermark = _watermark_boxes(watermark_search, watermark_rect)
    fixed = _parse_percent_boxes(
        fixed_watermark_regions,
        "固定水印区域",
        multiple=True,
        allow_empty=True,
    )
    presenter = _presenter_boxes(presenter_policy, presenter_rect)
    # 换页检测必须继续观察正文/标题区域。自动水印搜索范围较大，若将其
    # 整块忽略，会把仅在顶部标题或目录文字不同的页面错误合并。讲师窗口
    # 则应在第一遍扫描时就屏蔽，避免人物运动抬高换页阈值。
    initial_ignored_for_scan = [*fixed, *presenter]
    report_progress(0.01, "读取视频并建立高精度索引")
    scans, metadata = _scan_video(
        source,
        scan_mode=scan_mode,
        ignored_boxes=initial_ignored_for_scan,
    )
    auto_presenter_detected = False
    if presenter_policy == "auto_crop":
        auto_presenter_detected = _detect_auto_presenter(scans)
        presenter = [(0.78, 0.67, 0.22, 0.33)] if auto_presenter_detected else []
    ignored_for_scan = [*fixed, *presenter]
    if ignored_for_scan != initial_ignored_for_scan:
        scans = _rebuild_scan_features(scans, ignored_for_scan)
    check_cancelled("已取消视频幻灯片提取")
    report_progress(0.34, "识别换页、动画最终状态与重复画面")
    segments = _segment_scans(
        scans,
        ignored_boxes=ignored_for_scan,
        sensitivity=change_sensitivity,
    )
    full_width = int(metadata["width"])
    full_height = int(metadata["height"])
    crop = _presentation_crop(crop_mode, crop_rect, scans, full_width, full_height)
    auto_presenter_cropped = False
    if presenter_policy == "auto_crop" and auto_presenter_detected and presenter:
        presenter_left, presenter_top, presenter_right, presenter_bottom = _pixel_box(
            presenter[0], full_width, full_height
        )
        crop_left, crop_top, crop_right, crop_bottom = crop
        # The automatic preset is intentionally a crop, not synthetic content
        # generation.  When the lecturer window is anchored to the lower-right
        # edge, removing its side strip keeps every retained pixel authentic
        # and avoids visible inpaint/clone artefacts.
        if (
            presenter_right >= full_width - 2
            and presenter_bottom >= full_height - 2
            and presenter_left - crop_left >= (crop_right - crop_left) * 0.60
        ):
            crop = (
                crop_left,
                crop_top,
                min(crop_right, presenter_left),
                crop_bottom,
            )
            presenter = []
            auto_presenter_cropped = True
    cropped_watermark = _crop_normalized_boxes(watermark, crop, full_width, full_height)
    cropped_fixed = _crop_normalized_boxes(fixed, crop, full_width, full_height)
    cropped_presenter = _crop_normalized_boxes(presenter, crop, full_width, full_height)
    refined_fps = float(metadata["fps"])
    schedules: dict[int, int] = {}
    expected_by_segment: dict[int, set[int]] = {}
    for segment_index, segment in enumerate(segments):
        frame_indices = _refined_candidate_frame_indices(segment, scans, refined_fps)
        expected_by_segment[segment_index] = frame_indices
        for frame_index in frame_indices:
            schedules[frame_index] = segment_index
    capture = cv2.VideoCapture(str(source))
    if not capture.isOpened():
        raise MissingEngineError(f"OpenCV 无法再次解码该视频：{source.name}")
    staging_context = tempfile.TemporaryDirectory(
        prefix=".dfvs-", dir=target_dir.parent
    )
    staging_dir = Path(staging_context.name)
    image_dir = staging_dir / "高清幻灯片"
    image_dir.mkdir(parents=True, exist_ok=True)
    candidates: dict[int, list[tuple[_ScanFrame, np.ndarray]]] = {}
    image_paths: list[Path] = []
    slide_records: list[dict[str, object]] = []
    accepted_signatures: list[_SlideSignature] = []
    accepted_scores: list[float] = []
    completed_segments: set[int] = set()
    skipped_unstable_segments = 0
    frame_count = max(1, int(metadata["frame_count"]))
    current_index = 0
    last_decoded_index = -1
    last_raw_timestamp: float | None = None
    timestamp_origin: float | None = None
    try:
        while True:
            check_cancelled("已取消高清候选帧提取")
            ok, frame = capture.read()
            if not ok:
                break
            decoded_index = _decoded_frame_index(capture, current_index)
            if decoded_index <= last_decoded_index:
                decoded_index = current_index
            last_decoded_index = decoded_index
            raw_timestamp = _decoded_timestamp(
                capture, current_index, refined_fps, last_raw_timestamp
            )
            last_raw_timestamp = raw_timestamp
            if timestamp_origin is None:
                timestamp_origin = raw_timestamp
            decoded_timestamp = max(0.0, raw_timestamp - timestamp_origin)
            segment_index = schedules.get(decoded_index)
            if segment_index is not None:
                left, top, right, bottom = crop
                cropped = frame[top:bottom, left:right].copy()
                candidate_scan = _scan_frame_from_image(
                    decoded_index,
                    decoded_timestamp,
                    cropped,
                    [*cropped_fixed, *cropped_presenter],
                )
                candidates.setdefault(segment_index, []).append(
                    (candidate_scan, cropped)
                )
                if len(candidates[segment_index]) == len(
                    expected_by_segment[segment_index]
                ):
                    clean = _clean_slide(
                        candidates.pop(segment_index),
                        watermark_boxes=cropped_watermark,
                        fixed_boxes=cropped_fixed,
                        presenter_boxes=cropped_presenter,
                        watermark_hint=str(watermark_text_hint or ""),
                        annotation_mode=annotation_mode,
                        annotation_colours=annotation_bgr_colours,
                        annotation_tolerance=annotation_tolerance,
                        fixed_fill=fixed_fill_mode,
                        fixed_fill_colour=fixed_fill_bgr,
                        enhancement_mode=enhancement_mode,
                    )
                    completed_segments.add(segment_index)
                    if clean is None:
                        skipped_unstable_segments += 1
                        current_index += 1
                        continue
                    clean_signature = _slide_signature(clean.image)
                    duplicate_index = next(
                        (
                            index
                            for index, existing in enumerate(accepted_signatures)
                            if _signatures_near_duplicate(existing, clean_signature)
                        ),
                        None,
                    )
                    segment = segments[segment_index]
                    occurrence = {
                        "start_seconds": round(scans[segment.first].timestamp, 3),
                        "end_seconds": round(scans[segment.last].timestamp, 3),
                    }
                    if duplicate_index is not None:
                        record = slide_records[duplicate_index]
                        occurrences = record.setdefault("occurrences", [])
                        if isinstance(occurrences, list):
                            occurrences.append(occurrence)
                        record["merged_duplicate_segments"] = (
                            int(record.get("merged_duplicate_segments", 0)) + 1
                        )
                        new_score = _content_quality(
                            clean.image,
                            residual_annotation_pixels=(
                                clean.residual_annotation_pixels
                            ),
                            restored_annotation_pixels=clean.annotation_pixels,
                            restored_watermark_pixels=clean.watermark_pixels,
                            low_confidence_pixels=clean.low_confidence_pixels,
                        )
                        old_score = accepted_scores[duplicate_index]
                        if new_score - old_score > 0.01 * max(
                            abs(old_score), 0.001
                        ) and clean.low_confidence_pixels <= int(
                            record["low_confidence_pixels"]
                        ):
                            image_path = image_paths[duplicate_index]
                            _write_png(image_path, clean.image)
                            accepted_signatures[duplicate_index] = clean_signature
                            accepted_scores[duplicate_index] = new_score
                            record["candidate_frames"] = len(
                                expected_by_segment[segment_index]
                            )
                            record["annotation_pixels_restored"] = (
                                clean.annotation_pixels
                            )
                            record["colour_guided_pixels_restored"] = (
                                clean.colour_guided_pixels
                            )
                            record["dynamic_watermark_pixels_restored"] = (
                                clean.watermark_pixels
                            )
                            record["low_confidence_pixels"] = (
                                clean.low_confidence_pixels
                            )
                            record["residual_annotation_pixels"] = (
                                clean.residual_annotation_pixels
                            )
                            record["residual_colour_match_pixels"] = (
                                clean.residual_colour_match_pixels
                            )
                            record["selected_occurrence"] = occurrence
                            record["selected_timestamp_seconds"] = round(
                                clean.selected_timestamp, 3
                            )
                            record["selected_stability_support"] = (
                                clean.selected_stability_support
                            )
                            record["transition_candidates_rejected"] = (
                                clean.transition_candidates_rejected
                            )
                            record.update(_enhancement_record(clean))
                    else:
                        slide_number = len(image_paths) + 1
                        image_path = image_dir / f"幻灯片_{slide_number:04d}.png"
                        _write_png(image_path, clean.image)
                        verified = _read_cv_image(image_path)
                        if verified is None or verified.shape != clean.image.shape:
                            raise ValidationError(
                                f"幻灯片图片二重检查失败：{image_path.name}"
                            )
                        image_paths.append(image_path)
                        slide_records.append(
                            {
                                "slide": slide_number,
                                "start_seconds": round(
                                    scans[segment.first].timestamp, 3
                                ),
                                "end_seconds": round(scans[segment.last].timestamp, 3),
                                "candidate_frames": len(
                                    expected_by_segment[segment_index]
                                ),
                                "annotation_pixels_restored": clean.annotation_pixels,
                                "colour_guided_pixels_restored": (
                                    clean.colour_guided_pixels
                                ),
                                "residual_annotation_pixels": (
                                    clean.residual_annotation_pixels
                                ),
                                "residual_colour_match_pixels": (
                                    clean.residual_colour_match_pixels
                                ),
                                "dynamic_watermark_pixels_restored": clean.watermark_pixels,
                                "low_confidence_pixels": clean.low_confidence_pixels,
                                "occurrences": [occurrence],
                                "merged_duplicate_segments": 0,
                                "selected_occurrence": occurrence,
                                "selected_timestamp_seconds": round(
                                    clean.selected_timestamp, 3
                                ),
                                "selected_stability_support": (
                                    clean.selected_stability_support
                                ),
                                "transition_candidates_rejected": (
                                    clean.transition_candidates_rejected
                                ),
                                **_enhancement_record(clean),
                            }
                        )
                        accepted_signatures.append(clean_signature)
                        accepted_scores.append(
                            _content_quality(
                                clean.image,
                                residual_annotation_pixels=(
                                    clean.residual_annotation_pixels
                                ),
                                restored_annotation_pixels=(clean.annotation_pixels),
                                restored_watermark_pixels=clean.watermark_pixels,
                                low_confidence_pixels=clean.low_confidence_pixels,
                            )
                        )
            current_index += 1
            if current_index % max(1, int(float(metadata["fps"]) * 4)) == 0:
                report_progress(
                    0.36 + 0.48 * min(1.0, current_index / frame_count),
                    "提取高清真帧并去除手写批注",
                )
    finally:
        capture.release()
    expected_segments = set(range(len(segments)))
    if completed_segments != expected_segments or candidates:
        missing = len(expected_segments - completed_segments)
        raise ValidationError(
            f"视频第二遍解码不完整，仍有 {missing} 个页面片段未完成；未提交半成品"
        )
    if not image_paths:
        raise ValidationError("未识别到可输出的稳定幻灯片")
    (
        image_paths,
        slide_records,
        accepted_signatures,
        accepted_scores,
    ) = _reorder_slide_outputs_by_effective_timeline(
        image_paths,
        slide_records,
        accepted_signatures,
        accepted_scores,
    )
    _validate_slide_timeline(slide_records, float(metadata["duration"]))
    _validate_staged_images(image_paths, accepted_signatures)
    check_cancelled("已取消 PPT 写入")
    report_progress(0.86, "生成一页一图的高清不可编辑 PPT")
    pptx_path = staging_dir / f"{source.stem}_高清幻灯片.pptx"
    _write_pptx(image_paths, pptx_path)
    report_progress(0.94, "执行页数、空白页与图片完整性二重检查")
    _validate_pptx(pptx_path, len(image_paths), image_paths)
    crop_left, crop_top, crop_right, crop_bottom = crop
    warnings: list[str] = []
    if skipped_unstable_segments:
        warnings.append(
            f"已跳过 {skipped_unstable_segments} 个未形成连续稳定平台的短暂转场片段，"
            "以避免把半切换画面写入 PPT。"
        )
    low_confidence_total = sum(
        int(item["low_confidence_pixels"]) for item in slide_records
    )
    if low_confidence_total:
        warnings.append(
            "部分用户指定固定区域或讲师区域在同页期间始终被遮挡，已保守修补并标为低可信；无法保证恢复原始正文。"
        )
    if fixed:
        if fixed_fill_mode == "temporal":
            warnings.append(
                "固定水印区域采用“同页全历史真实像素优先、周边背景建模兜底”；建议抽查报告中的低可信像素。"
            )
        elif fixed_fill_mode == "background":
            warnings.append(
                "固定水印区域直接使用周边背景建模覆盖；区域内若原本存在正文，无法恢复被遮挡的真实像素。"
            )
        else:
            warnings.append(
                "固定水印区域使用用户指定纯色羽化覆盖；区域内若原本存在正文，无法恢复被遮挡的真实像素。"
            )
    residual_annotation_total = sum(
        int(item["residual_annotation_pixels"]) for item in slide_records
    )
    if residual_annotation_total:
        warnings.append(
            "仍检测到少量高置信彩色笔迹；若批注从该页首个候选帧起就始终存在，程序会优先保护被覆盖的印刷文字，不能保证强行抹除。"
        )
    residual_colour_total = sum(
        int(item.get("residual_colour_match_pixels", 0)) for item in slide_records
    )
    if annotation_mode == "manual" and residual_colour_total:
        warnings.append(
            f"二检仍发现 {residual_colour_total} 个像素接近用户指定标记色；其中可能包含受保护的同色印刷内容，建议按报告逐页抽查。"
        )
    if watermark_search in {"auto", "top", "bottom"}:
        warnings.append(
            "移动水印只在所选搜索范围内清理；若广告会移动到页面中部或底部，请改用“全画面”或自定义区域。"
        )
    normalized_report_enhancement = {
        "high_fidelity": "gpu_ai",
        "multiframe": "compatible",
    }.get(str(enhancement_mode), str(enhancement_mode))
    if normalized_report_enhancement in {"auto", "gpu_ai"}:
        ai_accepted_pages = sum(bool(item.get("ai_accepted")) for item in slide_records)
        ai_attempted_pages = sum(bool(item.get("ai_attempted")) for item in slide_records)
        fallback_blocks = sum(
            int(item.get("enhancement_fallback_blocks", 0)) for item in slide_records
        )
        if ai_attempted_pages == 0 and normalized_report_enhancement == "gpu_ai":
            warnings.append(
                "本次明确选择了 GPU AI 增强，但未检测到可用 Vulkan 显卡或引擎；"
                "已自动改用无独显兼容增强，PPT 提取未中断。"
            )
        elif ai_attempted_pages and ai_accepted_pages < len(slide_records):
            warnings.append(
                f"清晰增强二检：{ai_accepted_pages}/{len(slide_records)} 页采用 Real-ESRGAN，"
                f"其余页面因引擎不可用、原分辨率已足够或结构检查未通过而安全回退；"
                f"共有 {fallback_blocks} 个高风险分块恢复为真实像素放大结果。"
            )
        elif fallback_blocks:
            warnings.append(
                f"全部页面已通过 Real-ESRGAN 全局二检；其中 {fallback_blocks} 个高风险分块"
                "自动恢复为真实像素放大结果，以保护文字、数字、公式和表格线。"
            )
    deferred_previews = sum(
        item.get("ordering_reason") == "short_preview_deferred"
        for item in slide_records
    )
    if deferred_previews:
        warnings.append(
            f"检测到 {deferred_previews} 个页面曾在正式讲解前短暂误翻；"
            "已保留全部出现时间，并按后续稳定讲解区间重排 PPT、PNG 与报告。"
        )
    report_payload = {
        "schema": _REPORT_SCHEMA,
        "source": str(source.resolve()),
        "mode": "final_printed_state_only",
        "handwriting": "best_effort_temporal_and_high_confidence_cleanup",
        "metadata": metadata,
        "scan_frames": len(scans),
        "detected_segments": len(segments),
        "skipped_unstable_segments": skipped_unstable_segments,
        "output_slides": len(image_paths),
        "presentation_crop_pixels": [
            crop_left,
            crop_top,
            crop_right - crop_left,
            crop_bottom - crop_top,
        ],
        "watermark_search": watermark_search,
        "watermark_text_hint": str(watermark_text_hint or ""),
        "annotation_color_mode": annotation_mode,
        "annotation_colors": str(annotation_colors or ""),
        "annotation_color_tolerance": annotation_tolerance,
        "fixed_watermark_fill": fixed_fill_mode,
        "fixed_watermark_fill_color": str(fixed_watermark_fill_color or ""),
        "enhancement_mode": enhancement_mode,
        "presenter_policy": presenter_policy,
        "auto_presenter_detected": auto_presenter_detected,
        "auto_presenter_cropped": auto_presenter_cropped,
        "slides": slide_records,
        "warnings": warnings,
    }
    report_path = staging_dir / f"{source.stem}_提取报告.docx"
    if keep_report:
        report_progress(0.965, "生成可阅读的 Word 提取与二检报告")
        _write_report_docx(report_path, report_payload)
        _validate_report_docx(report_path, len(image_paths), report_payload)
    staged_outputs: list[Path] = [pptx_path]
    if keep_images:
        staged_outputs.extend(image_paths)
    else:
        for path in image_paths:
            path.unlink(missing_ok=True)
        try:
            image_dir.rmdir()
        except OSError:
            pass
    if keep_report:
        staged_outputs.append(report_path)
    relative_outputs = [path.relative_to(staging_dir) for path in staged_outputs]
    check_cancelled("任务已取消；未提交当前视频的半成品")
    report_progress(0.99, "全部二检通过，提交完整结果")
    _publish_staged_directory(staging_dir, target_dir)
    outputs = [target_dir / relative for relative in relative_outputs]
    staging_context.cleanup()
    report_progress(1.0, f"完成：输出 {len(image_paths)} 页高清幻灯片")
    return outputs


__all__ = [
    "extract_slides_to_pptx",
]
