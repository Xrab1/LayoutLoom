from __future__ import annotations

import csv
import json
import subprocess
import sys
import time
import traceback
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Iterable

from PIL import Image, ImageDraw
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from openpyxl import Workbook, load_workbook
from openpyxl.drawing.image import Image as ExcelImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pypdf import PdfReader
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen.canvas import Canvas

from docuforge.processors.office import convert_with_office
from docuforge.processors.video import detect_video_engine
from docuforge.processors.video_slide_repair import make_plan
from docuforge.registry import get_operations
from docuforge.runner import TaskRunner


ROOT = Path(__file__).resolve().parent
RUN_ROOT = ROOT / "validation_output" / (
    "full_feature_audit_" + datetime.now().strftime("%Y%m%d_%H%M%S")
)
FIXTURES = RUN_ROOT / "fixtures"
OUTPUTS = RUN_ROOT / "outputs"
OPERATIONS = {item.id: item for item in get_operations()}
RESULTS: list[dict[str, Any]] = []
COVERED: set[str] = set()


def add_hyperlink(paragraph: Any, text: str, url: str) -> None:
    relationship = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship)
    run = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")
    colour = OxmlElement("w:color")
    colour.set(qn("w:val"), "0563C1")
    properties.append(colour)
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    properties.append(underline)
    run.append(properties)
    node = OxmlElement("w:t")
    node.text = text
    run.append(node)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def make_fixtures() -> dict[str, Path]:
    FIXTURES.mkdir(parents=True, exist_ok=False)
    OUTPUTS.mkdir(parents=True, exist_ok=False)

    image_a = FIXTURES / "source_a.png"
    image_b = FIXTURES / "source_b.jpg"
    watermark = FIXTURES / "watermark.png"
    overlay = FIXTURES / "overlay.png"

    canvas_image = Image.new("RGBA", (640, 360), "white")
    draw = ImageDraw.Draw(canvas_image)
    draw.rectangle((25, 25, 615, 335), outline="#2457A7", width=5)
    draw.rectangle((55, 92, 585, 280), fill="#EAF1FC", outline="#7A9FD4", width=3)
    draw.text((70, 45), "LayoutLoom image fixture 2026", fill="#111111")
    draw.text((80, 135), "Editable text / table / image audit", fill="#202020")
    draw.line((80, 190, 560, 190), fill="#4E72A4", width=3)
    draw.ellipse((470, 210, 555, 295), fill="#F59E0B")
    canvas_image.save(image_a)
    exif = Image.Exif()
    exif[0x010E] = "LayoutLoom audit EXIF"
    canvas_image.convert("RGB").save(image_b, quality=93, exif=exif)
    mark = Image.new("RGBA", (180, 70), (0, 0, 0, 0))
    mark_draw = ImageDraw.Draw(mark)
    mark_draw.rounded_rectangle((2, 2, 177, 67), radius=15, fill=(220, 30, 55, 210))
    mark_draw.text((22, 25), "AUDIT", fill="white")
    mark.save(watermark)
    layer = Image.new("RGBA", (220, 120), (20, 150, 100, 160))
    ImageDraw.Draw(layer).text((45, 48), "OVERLAY", fill="white")
    layer.save(overlay)

    pdf_path = FIXTURES / "source.pdf"
    pdf = Canvas(str(pdf_path), pagesize=A4)
    width, height = A4
    pdf.setFont("Helvetica-Bold", 22)
    pdf.drawString(60, height - 72, "LayoutLoom PDF audit - page 1")
    pdf.setFont("Helvetica", 12)
    for index in range(18):
        pdf.drawString(60, height - 110 - index * 22, f"Line {index + 1}: searchable PDF text and layout.")
    pdf.showPage()
    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(60, height - 72, "Page 2 - table")
    for row in range(6):
        y = height - 130 - row * 45
        pdf.line(60, y, 500, y)
    for column in range(5):
        x = 60 + column * 110
        pdf.line(x, height - 130, x, height - 355)
    pdf.setFont("Helvetica", 11)
    for row in range(5):
        for column in range(4):
            pdf.drawString(70 + column * 110, height - 158 - row * 45, f"R{row + 1} C{column + 1}")
    pdf.showPage()
    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(60, height - 72, "Page 3 - embedded raster image")
    pdf.drawImage(str(image_b), 70, 250, width=460, height=259, preserveAspectRatio=True)
    pdf.save()

    one_page_pdf = FIXTURES / "insert.pdf"
    insert_canvas = Canvas(str(one_page_pdf), pagesize=A4)
    insert_canvas.setFont("Helvetica-Bold", 24)
    insert_canvas.drawString(70, height - 100, "Inserted appendix page")
    insert_canvas.save()

    docx_path = FIXTURES / "source.docx"
    document = Document()
    heading = document.add_heading("LayoutLoom Word audit", level=1)
    heading.runs[0].font.name = "Arial"
    document.add_paragraph("旧文字：This paragraph contains English 123 and 中文。")
    document.add_paragraph("   ")
    paragraph = document.add_paragraph("Link: ")
    add_hyperlink(paragraph, "LayoutLoom", "https://example.com")
    table = document.add_table(rows=3, cols=3)
    for row in range(3):
        for column in range(3):
            table.cell(row, column).text = f"R{row + 1}C{column + 1}"
    document.add_picture(str(image_a), width=Inches(4.5))
    document.add_paragraph("Page two content").add_run().add_break(WD_BREAK.PAGE)
    document.add_paragraph("Second page body")
    document.save(docx_path)

    fixed_docx = FIXTURES / "legacy_fixed_coordinate.docx"
    fixed_document = Document()
    for page_number in range(1, 3):
        if page_number > 1:
            fixed_document.add_section(WD_SECTION.NEW_PAGE)
        for text, x, y, text_width, size, bold in (
            (f"Page {page_number} heading", 1500, 1000, 5000, 16.0, True),
            (f"Page {page_number} editable paragraph begins with normal text", 1450, 1450, 7200, 10.5, False),
            ("and continues on another positioned visual line for copying.", 1450, 1740, 7200, 10.5, False),
            ("A second paragraph remains editable after the region rebuild.", 1450, 2260, 7200, 10.5, False),
        ):
            paragraph = fixed_document.add_paragraph()
            run = paragraph.add_run(text)
            run.font.size = Pt(size)
            run.bold = bold
            properties = paragraph._p.get_or_add_pPr()
            frame = OxmlElement("w:framePr")
            for name, value in (
                ("w:hAnchor", "page"),
                ("w:vAnchor", "page"),
                ("w:x", str(x)),
                ("w:y", str(y)),
                ("w:w", str(text_width)),
                ("w:h", str(max(200, round(size * 24)))),
            ):
                frame.set(qn(name), value)
            properties.insert(0, frame)
            fit = OxmlElement("w:fitText")
            fit.set(qn("w:val"), str(text_width))
            run._r.get_or_add_rPr().append(fit)
    fixed_document.save(fixed_docx)

    xlsx_path = FIXTURES / "source.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", "Score", "Tags", "Text"])
    sheet.append(["B", 20, "red,blue", "旧值"])
    sheet.append(["A", 10, "green,yellow", "normal"])
    sheet.append(["B", 20, "red,blue", "旧值"])
    sheet.append(["C", 30, "white,black", "normal"])
    sheet.add_image(ExcelImage(str(image_a)), "F2")
    workbook.create_sheet("Second").append(["Other", 1])
    workbook.save(xlsx_path)
    workbook.close()

    xlsx_plain = FIXTURES / "source_plain.xlsx"
    plain_workbook = Workbook()
    plain_sheet = plain_workbook.active
    plain_sheet.title = "Data"
    plain_sheet.append(["Name", "Score", "Tags", "Text"])
    plain_sheet.append(["B", 20, "red,blue", "旧值"])
    plain_sheet.append(["A", 10, "green,yellow", "normal"])
    plain_sheet.append(["B", 20, "red,blue", "旧值"])
    plain_sheet.append(["C", 30, "white,black", "normal"])
    plain_workbook.save(xlsx_plain)
    plain_workbook.close()

    pptx_path = FIXTURES / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.5), Inches(0.8))
    title.text_frame.paragraphs[0].text = "LayoutLoom PPT audit"
    title.text_frame.paragraphs[0].runs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    slide.shapes.add_picture(str(image_a), Inches(1), Inches(1.5), width=Inches(7.5))
    presentation.slides.add_slide(presentation.slide_layouts[6]).shapes.add_textbox(
        Inches(1), Inches(1), Inches(6), Inches(1)
    ).text_frame.text = "Second slide"
    presentation.save(pptx_path)

    legacy_dir = FIXTURES / "legacy"
    legacy_dir.mkdir()
    legacy_doc = convert_with_office(docx_path, legacy_dir, "doc", engine="wps")[0]
    legacy_xls = convert_with_office(xlsx_path, legacy_dir, "xls", engine="wps")[0]

    ffmpeg_status = detect_video_engine()
    if not ffmpeg_status.available or ffmpeg_status.executable is None:
        raise RuntimeError(ffmpeg_status.reason)
    audio_video = FIXTURES / "audio_video.mp4"
    subprocess.run(
        [
            str(ffmpeg_status.executable),
            "-hide_banner",
            "-loglevel",
            "error",
            "-f",
            "lavfi",
            "-i",
            "testsrc=size=640x360:rate=10",
            "-f",
            "lavfi",
            "-i",
            "sine=frequency=880:sample_rate=44100",
            "-t",
            "4",
            "-c:v",
            "libx264",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            "-shortest",
            "-y",
            str(audio_video),
        ],
        check=True,
        timeout=90,
    )
    return {
        "image_a": image_a,
        "image_b": image_b,
        "watermark": watermark,
        "overlay": overlay,
        "pdf": pdf_path,
        "insert_pdf": one_page_pdf,
        "docx": docx_path,
        "fixed_docx": fixed_docx,
        "xlsx": xlsx_path,
        "xlsx_plain": xlsx_plain,
        "pptx": pptx_path,
        "doc": legacy_doc,
        "xls": legacy_xls,
        "audio_video": audio_video,
    }


def validate_generic(path: Path) -> None:
    if not path.is_file() or path.stat().st_size <= 0:
        raise AssertionError(f"输出为空：{path}")
    suffix = path.suffix.casefold()
    if suffix == ".pdf":
        reader = PdfReader(path)
        try:
            if not reader.is_encrypted and len(reader.pages) < 1:
                raise AssertionError("PDF 没有页面")
        finally:
            reader.close()
    elif suffix in {".docx", ".docm"}:
        Document(path)
    elif suffix in {".xlsx", ".xlsm"}:
        workbook = load_workbook(path, read_only=True, data_only=False)
        if not workbook.sheetnames:
            raise AssertionError("Excel 没有工作表")
        workbook.close()
    elif suffix == ".pptx":
        if len(Presentation(path).slides) < 1:
            raise AssertionError("PPT 没有幻灯片")
    elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff", ".gif", ".ico"}:
        with Image.open(path) as image:
            image.load()
            if image.width < 1 or image.height < 1:
                raise AssertionError("图片尺寸无效")
    elif suffix == ".json":
        json.loads(path.read_text(encoding="utf-8-sig"))
    elif suffix == ".csv":
        with path.open("r", encoding="utf-8-sig", newline="") as stream:
            if not list(csv.reader(stream)):
                raise AssertionError("CSV 没有数据")
    elif suffix in {".mp4", ".mkv", ".mov", ".wav"}:
        status = detect_video_engine()
        ffprobe = status.ffprobe_executable
        if ffprobe is None:
            raise AssertionError("缺少 ffprobe，无法验证视频/音频")
        checked = subprocess.run(
            [
                str(ffprobe),
                "-v",
                "error",
                "-show_entries",
                "format=duration:stream=codec_type,codec_name",
                "-of",
                "json",
                str(path),
            ],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=30,
            check=True,
        )
        payload = json.loads(checked.stdout)
        if not payload.get("streams"):
            raise AssertionError("媒体文件没有可读流")


def run_case(
    case_name: str,
    operation_id: str,
    inputs: Iterable[Path],
    params: dict[str, Any] | None = None,
    validator: Callable[[list[Path]], None] | None = None,
) -> list[Path]:
    operation = OPERATIONS[operation_id]
    case_dir = OUTPUTS / case_name.replace("/", "_").replace("[", "_").replace("]", "")
    case_dir.mkdir(parents=True, exist_ok=False)
    started = time.perf_counter()
    print(f"RUN  {case_name}: {operation.name}", flush=True)
    try:
        result = TaskRunner().run(operation, list(inputs), case_dir, params or {})
        for output in result.outputs:
            validate_generic(output)
        if validator is not None:
            validator(result.outputs)
        elapsed = round(time.perf_counter() - started, 3)
        RESULTS.append(
            {
                "case": case_name,
                "operation": operation_id,
                "status": "PASS",
                "seconds": elapsed,
                "outputs": [str(path) for path in result.outputs],
                "warnings": list(result.warnings),
                "details": dict(result.details),
            }
        )
        COVERED.add(operation_id)
        print(f"PASS {case_name} ({elapsed:.3f}s, {len(result.outputs)} outputs)", flush=True)
        return result.outputs
    except BaseException as exc:
        elapsed = round(time.perf_counter() - started, 3)
        RESULTS.append(
            {
                "case": case_name,
                "operation": operation_id,
                "status": "FAIL",
                "seconds": elapsed,
                "outputs": [],
                "error": f"{type(exc).__name__}: {exc}",
                "traceback": traceback.format_exc(),
            }
        )
        COVERED.add(operation_id)
        print(f"FAIL {case_name}: {type(exc).__name__}: {exc}", flush=True)
        return []


def expect_failure(
    case_name: str,
    operation_id: str,
    inputs: Iterable[Path],
    params: dict[str, Any],
    expected_text: str,
) -> None:
    operation = OPERATIONS[operation_id]
    case_dir = OUTPUTS / case_name
    case_dir.mkdir(parents=True, exist_ok=False)
    started = time.perf_counter()
    try:
        TaskRunner().run(operation, list(inputs), case_dir, params)
    except BaseException as exc:
        message = str(exc)
        status = "PASS" if expected_text in message else "FAIL"
        RESULTS.append(
            {
                "case": case_name,
                "operation": operation_id,
                "status": status,
                "seconds": round(time.perf_counter() - started, 3),
                "expected_failure": expected_text,
                "error": f"{type(exc).__name__}: {message}",
            }
        )
        print(f"{status} {case_name}: {message}", flush=True)
        return
    RESULTS.append(
        {
            "case": case_name,
            "operation": operation_id,
            "status": "FAIL",
            "seconds": round(time.perf_counter() - started, 3),
            "error": "预期失败，但任务意外成功",
        }
    )


def first_with_suffix(outputs: Iterable[Path], suffix: str) -> Path:
    return next(path for path in outputs if path.suffix.casefold() == suffix.casefold())


def audit(fixtures: dict[str, Path]) -> None:
    pdf_hybrid = run_case(
        "pdf_to_word_hybrid",
        "pdf.to_word",
        [fixtures["pdf"]],
        {"mode": "hybrid", "column_layout": "single", "low_quality_policy": "keep", "dpi": 180},
    )
    run_case(
        "pdf_to_word_editable",
        "pdf.to_word",
        [fixtures["pdf"]],
        {"mode": "editable", "column_layout": "single", "low_quality_policy": "keep", "dpi": 180},
    )
    run_case(
        "pdf_to_word_visual",
        "pdf.to_word",
        [fixtures["pdf"]],
        {"mode": "visual", "column_layout": "auto", "low_quality_policy": "keep", "dpi": 180},
    )
    run_case("pdf_to_ppt", "pdf.to_ppt", [fixtures["pdf"]], {"dpi": 150})
    run_case("pdf_to_images", "pdf.to_images", [fixtures["pdf"]], {"format": "jpg", "dpi": 150})
    run_case(
        "pdf_extract_images_all",
        "pdf.extract_images",
        [fixtures["pdf"]],
        {"mode": "all", "dpi": 180, "format": "png", "write_manifest": True},
    )
    run_case("pdf_to_text", "pdf.to_text", [fixtures["pdf"]], {"layout": True})
    image_pdf = run_case(
        "images_to_pdf",
        "image.to_pdf",
        [fixtures["image_a"], fixtures["image_b"]],
        {"filename": "image_bundle", "dpi": 96, "background": "white"},
    )
    run_case("pdf_merge", "pdf.merge", [fixtures["pdf"], fixtures["insert_pdf"]], {"filename": "merged"})
    run_case("pdf_split", "pdf.split", [fixtures["pdf"]], {"ranges": "1,2-3"})
    run_case("pdf_extract_pages", "pdf.extract_pages", [fixtures["pdf"]], {"pages": "2-3"})
    run_case("pdf_delete_pages", "pdf.delete_pages", [fixtures["pdf"]], {"pages": "2"})
    run_case("pdf_insert_pages", "pdf.insert_pages", [fixtures["pdf"], fixtures["insert_pdf"]], {"position": 2})
    run_case("pdf_rotate", "pdf.rotate", [fixtures["pdf"]], {"pages": "1,3", "angle": "90"})
    run_case("pdf_compress_lossless", "pdf.compress", [fixtures["pdf"]], {"level": "9"})
    run_case(
        "pdf_compress_lossy_smart",
        "pdf.compress_lossy",
        [fixtures["pdf"]],
        {"strategy": "smart", "dpi": "150", "jpeg_quality": "82", "color_mode": "color"},
    )
    run_case(
        "pdf_compress_lossy_raster",
        "pdf.compress_lossy",
        [fixtures["pdf"]],
        {"strategy": "raster", "dpi": "150", "jpeg_quality": "82", "color_mode": "grayscale"},
    )
    encrypted = run_case(
        "pdf_encrypt",
        "pdf.encrypt",
        [fixtures["pdf"]],
        {"user_password": "Audit-2026", "owner_password": "Owner-2026", "algorithm": "AES-128"},
        lambda outputs: (
            None
            if PdfReader(outputs[0]).is_encrypted
            else (_ for _ in ()).throw(AssertionError("输出 PDF 未加密"))
        ),
    )
    if encrypted:
        run_case("pdf_decrypt", "pdf.decrypt", encrypted, {"password": "Audit-2026"})
    run_case(
        "pdf_watermark_9",
        "pdf.watermark",
        [fixtures["pdf"]],
        {"text": "AUDIT", "count": 9, "opacity": 0.25, "angle": 35, "scale": 0.2, "font_size": 24},
    )
    run_case(
        "pdf_header_footer",
        "pdf.header_footer",
        [fixtures["pdf"]],
        {"header": "LayoutLoom audit", "footer": "verified", "page_numbers": True, "page_format": "{page}/{total}"},
    )

    run_case("word_to_pdf_wps", "word.to_pdf", [fixtures["docx"]], {"engine": "wps"})
    run_case(
        "word_legacy_compatibility",
        "word.full_compatibility",
        [fixtures["fixed_docx"]],
        {"verification_engine": "none"},
    )
    run_case(
        "excel_to_pdf_wps",
        "excel.to_pdf",
        [fixtures["xlsx"]],
        {"engine": "wps", "excel_pdf_layout": "smart", "excel_pdf_paper": "auto", "excel_pdf_orientation": "auto", "excel_pdf_margin": "auto"},
    )
    run_case("ppt_to_pdf_wps", "ppt.to_pdf", [fixtures["pptx"]], {"engine": "wps"})
    run_case("word_to_text", "word.to_text", [fixtures["docx"]], {"include_tables": True})
    run_case("excel_to_csv", "excel.to_csv", [fixtures["xlsx"]], {"delimiter": ","})
    run_case("excel_to_json", "excel.to_json", [fixtures["xlsx"]], {"header": True})
    run_case("excel_to_txt", "excel.to_txt", [fixtures["xlsx"]], {"delimiter": "\t"})
    run_case("legacy_doc_to_docx", "legacy.doc_to_docx", [fixtures["doc"]], {"engine": "wps"})
    run_case("legacy_xls_to_xlsx", "legacy.xls_to_xlsx", [fixtures["xls"]], {"engine": "wps"})
    run_case("ppt_to_images_wps", "ppt.to_images", [fixtures["pptx"]], {"renderer": "wps", "format": "png", "width": 1280, "dpi": 150})

    run_case("word_replace", "word.replace", [fixtures["docx"]], {"replacements": '{"旧文字":"新文字"}', "case_sensitive": True})
    run_case("word_remove_blank_lines", "word.remove_blank_lines", [fixtures["docx"]])
    run_case("word_remove_images", "word.remove_images", [fixtures["docx"]])
    run_case("word_typography", "word.typography", [fixtures["docx"]], {"font_name": "微软雅黑", "font_size": 11, "line_spacing": 1.25})
    run_case("word_headers_footers", "word.headers_footers", [fixtures["docx"]], {"header": "AUDIT HEADER", "footer": "AUDIT FOOTER", "replace": True})
    run_case("word_extract_images", "word.extract_images", [fixtures["docx"]])
    run_case("word_remove_hyperlinks", "word.remove_hyperlinks", [fixtures["docx"]])

    run_case("excel_sort", "excel.sort", [fixtures["xlsx"]], {"column": "B", "sheet": "Data", "header": True, "reverse": False})
    run_case("excel_filter", "excel.filter", [fixtures["xlsx"]], {"column": "B", "operator": "greater_equal", "value": "20", "sheet": "Data", "header": True})
    run_case("excel_deduplicate", "excel.deduplicate", [fixtures["xlsx_plain"]], {"columns": "A,B,C,D", "sheet": "Data", "header": True})
    run_case("excel_replace", "excel.replace", [fixtures["xlsx"]], {"replacements": '{"旧值":"新值"}', "sheets": "Data", "case_sensitive": True, "exact": False})
    run_case("excel_split_column", "excel.split_column", [fixtures["xlsx_plain"]], {"column": "C", "delimiter": ",", "sheet": "Data", "header": True, "maxsplit": -1})
    run_case("excel_conditional_format", "excel.conditional_format", [fixtures["xlsx"]], {"range": "B2:B10", "rule": "cell", "operator": "greaterThan", "threshold": "15", "color": "FFF2CC", "sheet": "Data"})
    run_case("excel_extract_images", "excel.extract_images", [fixtures["xlsx"]])

    run_case("ppt_replace_fonts", "ppt.replace_fonts", [fixtures["pptx"]], {"replacements": '{"Arial":"微软雅黑"}', "default_font": ""})
    run_case("ppt_watermark", "ppt.watermark", [fixtures["pptx"]], {"text": "AUDIT", "font_size": 28, "color": "B7B7B7", "rotation": -25, "opacity": 0.3})
    run_case("ppt_extract_media", "ppt.extract_media", [fixtures["pptx"]], {"images": True, "audio": True, "video": True})
    run_case("ppt_compress_images", "ppt.compress_images", [fixtures["pptx"]], {"quality": 75, "max_dimension": 1200})

    run_case("image_convert", "image.convert", [fixtures["image_a"]], {"format": "jpg", "quality": 90, "background": "white"})
    run_case("image_resize", "image.resize", [fixtures["image_a"]], {"width": 320, "height": 180, "keep_aspect": False})
    run_case("image_scale", "image.scale", [fixtures["image_a"]], {"percent": 62.5})
    run_case("image_crop", "image.crop", [fixtures["image_a"]], {"left": 40, "top": 30, "right": 500, "bottom": 310})
    run_case("image_rotate", "image.rotate", [fixtures["image_a"]], {"angle": 33, "expand": True, "background": "white"})
    run_case("image_flip", "image.flip", [fixtures["image_a"]], {"direction": "vertical"})
    run_case("image_compress", "image.compress", [fixtures["image_b"]], {"quality": 70, "target_kb": 20, "min_quality": 35, "allow_resize": True})
    run_case("image_remove_exif", "image.remove_exif", [fixtures["image_b"]])
    run_case("image_enhance_compatible", "image.enhance", [fixtures["image_a"]], {"mode": "compatible", "content_type": "document", "scale": "2", "max_dimension": 1600, "output_format": "png"})
    run_case("image_enhance_gpu", "image.enhance", [fixtures["image_a"]], {"mode": "gpu_ai", "content_type": "document", "scale": "2", "max_dimension": 1600, "output_format": "png"})
    run_case("image_adjust", "image.adjust", [fixtures["image_a"]], {"brightness": 1.15, "contrast": 1.1, "saturation": 0.8})
    run_case("image_filter", "image.filter", [fixtures["image_a"]], {"filter": "sepia", "intensity": 0.75})
    run_case("image_text_watermark", "image.text_watermark", [fixtures["image_a"]], {"text": "AUDIT", "position": "center", "font_size": 42, "color": "#FF0000", "opacity": 0.65, "margin": 20})
    run_case("image_image_watermark", "image.image_watermark", [fixtures["image_a"]], {"watermark_path": fixtures["watermark"], "position": "bottom-right", "opacity": 0.7, "scale": 0.25, "margin": 12})
    run_case("image_border", "image.border", [fixtures["image_a"]], {"width": 18, "color": "#22AA66"})
    run_case("image_mosaic", "image.mosaic", [fixtures["image_a"]], {"left": 70, "top": 110, "right": 360, "bottom": 210, "block_size": 14})
    run_case("image_stitch", "image.stitch", [fixtures["image_a"], fixtures["image_b"]], {"filename": "stitched", "format": "png", "direction": "horizontal", "spacing": 12, "background": "white", "alignment": "center"})
    run_case("image_overlay", "image.overlay", [fixtures["image_a"]], {"overlay_path": fixtures["overlay"], "position": "center", "opacity": 0.7, "scale": 0.5, "margin": 0})
    run_case("image_rename", "image.rename", [fixtures["image_a"], fixtures["image_b"]], {"pattern": "audit_{index:03d}", "start": 5, "move": False})

    slide_video = run_case(
        "images_to_video",
        "image.to_video",
        [fixtures["image_a"], fixtures["image_b"], fixtures["watermark"]],
        {"filename": "slides", "slide_duration": 4, "fps": 5, "resolution": "720p", "transition": "none", "transition_duration": 0, "background": "black", "quality": 28},
    )
    extracted_deck: list[Path] = []
    if slide_video:
        extracted_outputs = run_case(
            "video_extract_slides",
            "video.extract_slides_ppt",
            slide_video,
            {"scan_mode": "fast", "change_sensitivity": "sensitive", "crop_mode": "full", "watermark_search": "off", "annotation_color_mode": "off", "presenter_policy": "keep", "enhancement_mode": "off", "keep_images": True, "keep_report": True},
        )
        extracted_deck = [path for path in extracted_outputs if path.suffix.casefold() == ".pptx"]
        if extracted_deck:
            plan = make_plan(
                [
                    {
                        "kind": "replace_page_frame",
                        "page": 1,
                        "timestamp": 0.2,
                        "regions": [],
                        "method": "temporal",
                        "colour": "#FFFFFF",
                    }
                ]
            )
            run_case(
                "video_repair_slides",
                "video.repair_slides_ppt",
                [extracted_deck[0], slide_video[0]],
                {"repair_plan": plan},
            )
    if image_pdf:
        run_case(
            "pdf_to_video",
            "pdf.to_video",
            image_pdf,
            {"dpi": 150, "slide_duration": 0.7, "fps": 5, "resolution": "720p", "transition": "none", "transition_duration": 0, "background": "black", "quality": 28},
        )
    if slide_video:
        run_case("video_transcode", "video.transcode", slide_video, {"format": "mkv", "video_codec": "h264", "audio_codec": "none", "quality": 28, "resolution": "original", "target_fps": "", "audio_bitrate": 128})
        run_case("video_compress", "video.compress", slide_video, {"format": "mp4", "quality": 30, "resolution": "720p", "target_fps": "5", "audio_bitrate": 96})
        run_case("video_trim", "video.trim", slide_video, {"format": "mp4", "start": "0.5", "end": "2.5", "duration": "", "quality": 28})
    run_case("video_extract_audio", "video.extract_audio", [fixtures["audio_video"]], {"sample_rate": "22050", "channels": "1"})

    for operation_id in ("word.to_pdf", "excel.to_pdf", "ppt.to_pdf"):
        expect_failure(
            "explicit_microsoft_unavailable_" + operation_id.replace(".", "_"),
            operation_id,
            [fixtures[{"word.to_pdf": "docx", "excel.to_pdf": "xlsx", "ppt.to_pdf": "pptx"}[operation_id]]],
            {"engine": "microsoft_office"},
            "Microsoft",
        )


def write_report() -> None:
    passed = sum(item["status"] == "PASS" for item in RESULTS)
    failed = sum(item["status"] == "FAIL" for item in RESULTS)
    all_ids = set(OPERATIONS)
    missing = sorted(all_ids - COVERED)
    payload = {
        "run_root": str(RUN_ROOT),
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "operation_count": len(OPERATIONS),
        "covered_operation_count": len(COVERED),
        "missing_operations": missing,
        "case_count": len(RESULTS),
        "passed": passed,
        "failed": failed,
        "results": RESULTS,
    }
    (RUN_ROOT / "audit_report.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    lines = [
        "# LayoutLoom 全功能实际执行审计",
        "",
        f"- 功能总数：{len(OPERATIONS)}",
        f"- 已覆盖功能：{len(COVERED)}",
        f"- 实际测试场景：{len(RESULTS)}",
        f"- 通过：{passed}",
        f"- 失败：{failed}",
        f"- 未覆盖：{', '.join(missing) if missing else '无'}",
        "",
        "| 状态 | 场景 | 功能 ID | 秒数 | 说明 |",
        "|---|---|---|---:|---|",
    ]
    for item in RESULTS:
        note = item.get("error") or f"输出 {len(item.get('outputs', []))} 个"
        note = str(note).replace("|", "\\|").replace("\n", " ")
        lines.append(
            f"| {item['status']} | {item['case']} | {item['operation']} | {item['seconds']} | {note} |"
        )
    (RUN_ROOT / "audit_report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(json.dumps({key: payload[key] for key in ("run_root", "operation_count", "covered_operation_count", "case_count", "passed", "failed", "missing_operations")}, ensure_ascii=False, indent=2))


def main() -> int:
    try:
        fixtures = make_fixtures()
        audit(fixtures)
    except BaseException as exc:
        RESULTS.append(
            {
                "case": "audit_setup_or_fatal",
                "operation": "-",
                "status": "FAIL",
                "seconds": 0,
                "error": f"{type(exc).__name__}: {exc}",
                "traceback": traceback.format_exc(),
            }
        )
        print(traceback.format_exc(), file=sys.stderr, flush=True)
    finally:
        RUN_ROOT.mkdir(parents=True, exist_ok=True)
        write_report()
    return 0 if not any(item["status"] == "FAIL" for item in RESULTS) else 1


if __name__ == "__main__":
    import multiprocessing

    multiprocessing.freeze_support()
    raise SystemExit(main())
